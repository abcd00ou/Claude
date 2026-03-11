"""
lecture/agents/slide_builder.py — Professional Analyst Edition v3
디자인 철학: McKinsey/BCG 인스파이어 — 절제된 색상, 명확한 계층, 데이터 중심
"""

import io
import json
import os
import sys
import pathlib
import textwrap
from datetime import datetime

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

sys.path.insert(0, str(pathlib.Path(__file__).parent.parent))
from config import (
    CHAPTERS, AI_TIMELINE, WORKFLOW_CHANGES,
    SLIDE_DESIGN, LECTURE_TITLE, LECTURE_SUBTITLE,
    LECTURE_AUTHOR, LECTURE_DATE, OUTPUT_DIR, DATA_DIR,
    INDUSTRY_STATS, JOB_STATS, TRIE_FORMULA,
    TOOL_MATRIX, THIRTY_DAY_PLAN, CAUTION_ITEMS,
    TOKEN_MEMORY, CONTEXT_EVOLUTION,
    GPU_DATA, TRANSFORMER_DATA, TRAINING_INFERENCE_DATA,
    SCALING_LAWS_DATA, VRAM_KVCACHE_DATA, MULTIMODAL_DATA,
    AI_FRONTIER_2026, AI_PREPARATION_2026, PHASE_DATA,
)

# ─────────────────────────────────────────────────────────────────
# DESIGN SYSTEM — Professional Analyst Theme
# ─────────────────────────────────────────────────────────────────
D = {
    # Backgrounds
    "bg_white":   "#FFFFFF",
    "bg_light":   "#F8FAFC",   # near-white for alternating rows
    "bg_dark":    "#0A1628",   # deep navy for title/divider slides
    "bg_card":    "#FFFFFF",
    "bg_subtle":  "#F1F5F9",   # subtle grey for secondary panels

    # Brand colors
    "navy":       "#0F2B4C",   # primary dark blue
    "blue":       "#1D4ED8",   # accent blue
    "blue_light": "#EFF6FF",   # very light blue for backgrounds
    "red":        "#C0392B",   # data highlight red
    "green":      "#0D6B4F",   # positive indicator
    "amber":      "#B45309",   # warning / caution

    # Text hierarchy
    "text_h1":    "#0F172A",   # headings
    "text_h2":    "#1E293B",
    "text_body":  "#334155",
    "text_muted": "#64748B",
    "text_light": "#94A3B8",
    "text_white": "#FFFFFF",

    # Borders & dividers
    "border":     "#CBD5E1",
    "border_light": "#E2E8F0",
    "divider":    "#F1F5F9",

    # Chapter accent colors (more muted, professional)
    "ch1": "#1D4ED8",   # blue
    "ch2": "#0D6B4F",   # green
    "ch3": "#6D28D9",   # purple
    "ch4": "#B91C1C",   # red
    "ch5": "#0E7490",   # teal
}

CHAPTER_COLORS = [D["ch1"], D["ch2"], D["ch3"], D["ch4"], D["ch5"]]

# Slide dimensions
W = SLIDE_DESIGN["slide_width_emu"]
H = SLIDE_DESIGN["slide_height_emu"]

# Font
FONT = SLIDE_DESIGN.get("fallback_font", "Apple SD Gothic Neo")

# PIL font
KOREAN_FONT_CANDIDATES = [
    "/System/Library/Fonts/AppleSDGothicNeo.ttc",
    "/System/Library/Fonts/Supplemental/AppleSDGothicNeo.ttc",
    "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
]
_KFP = next((p for p in KOREAN_FONT_CANDIDATES if os.path.exists(p)), None)


# ─────────────────────────────────────────────────────────────────
# LOW-LEVEL HELPERS
# ─────────────────────────────────────────────────────────────────
def _rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _rect(slide, x, y, w, h, fill: str):
    sh = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(w), Emu(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = _rgb(fill)
    sh.line.fill.background()
    return sh


def _txt(slide, x, y, w, h, text: str, pt: float,
         bold=False, color="#0F172A", align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(pt)
    r.font.bold = bold
    r.font.color.rgb = _rgb(color)
    r.font.name = FONT
    return tb


def _txt_lines(slide, x, y, w, h, lines: list, pt: float,
               bold=False, color="#0F172A"):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(pt)
        r.font.bold = bold
        r.font.color.rgb = _rgb(color)
        r.font.name = FONT
    return tb


def _accent_bar(slide, color: str, height_frac=0.006):
    """Top accent bar — thin line."""
    _rect(slide, 0, 0, W, int(H * height_frac), color)


def _slide_label(slide, text: str):
    """Bottom-right slide label / source."""
    _txt(slide, int(W * 0.72), int(H * 0.935), int(W * 0.26), int(H * 0.05),
         text, 8, color=D["text_light"], align=PP_ALIGN.RIGHT)


def _card(slide, x, y, w, h, accent_color: str = None, bg: str = None):
    """White card with optional left accent bar."""
    bg = bg or D["bg_card"]
    _rect(slide, x, y, w, h, bg)
    if accent_color:
        _rect(slide, x, y, int(W * 0.005), h, accent_color)


def _shadow_card(slide, x, y, w, h, accent_color: str = None):
    """Card with subtle drop shadow effect."""
    _rect(slide, x + int(W * 0.003), y + int(H * 0.005), w, h, D["border_light"])
    _card(slide, x, y, w, h, accent_color)


# ─────────────────────────────────────────────────────────────────
# PIL INFRASTRUCTURE
# ─────────────────────────────────────────────────────────────────
def _kfont(size: int, bold=False):
    if not PIL_AVAILABLE:
        return None
    if _KFP:
        try:
            return ImageFont.truetype(_KFP, size, index=1 if bold else 0)
        except Exception:
            pass
    try:
        return ImageFont.load_default(size=size)
    except Exception:
        return ImageFont.load_default()


def _pil_stream(img) -> io.BytesIO:
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def _draw_center(draw, text, cx, cy, font, fill):
    try:
        bb = draw.textbbox((0, 0), text, font=font)
        draw.text((cx - (bb[2]-bb[0])//2, cy - (bb[3]-bb[1])//2), text, font=font, fill=fill)
    except Exception:
        draw.text((cx, cy), text, font=font, fill=fill)


def _rrect(draw, x0, y0, x1, y1, r, fill, outline=None, lw=1):
    try:
        draw.rounded_rectangle([x0, y0, x1, y1], radius=r, fill=fill, outline=outline, width=lw)
    except Exception:
        draw.rectangle([x0, y0, x1, y1], fill=fill, outline=outline)


def _pil_img(slide, stream: io.BytesIO, x, y, w, h):
    slide.shapes.add_picture(stream, Emu(x), Emu(y), Emu(w), Emu(h))


# ─────────────────────────────────────────────────────────────────
# PIL IMAGE GENERATORS — Professional Style
# ─────────────────────────────────────────────────────────────────

def pil_stat_card(value: str, label: str, sublabel: str, color_hex: str,
                  w=860, h=460) -> io.BytesIO | None:
    """
    Professional stat card: clean white/dark split design.
    Left 60%: large number on dark navy.
    Right 40%: label + sublabel on white.
    """
    if not PIL_AVAILABLE:
        return None

    cv = tuple(int(color_hex.lstrip("#")[i:i+2], 16) for i in (0, 2, 4))
    img = Image.new("RGB", (w, h), (255, 255, 255))
    draw = ImageDraw.Draw(img)

    # Left dark panel
    lw = int(w * 0.60)
    draw.rectangle([0, 0, lw, h], fill=(10, 22, 40))

    # Accent bar on left panel (top)
    draw.rectangle([0, 0, lw, 6], fill=cv)

    # Big number
    fn_big = _kfont(88, bold=True)
    _draw_center(draw, value, lw // 2, h // 2 - 20, fn_big, (255, 255, 255))

    # Thin divider
    draw.rectangle([int(lw * 0.12), h // 2 + 50, int(lw * 0.88), h // 2 + 54],
                   fill=(*cv, 180))

    # Value footnote
    fn_sm = _kfont(18)
    _draw_center(draw, sublabel, lw // 2, h // 2 + 80, fn_sm, (148, 163, 184))

    # Right panel
    draw.rectangle([lw, 0, w, h], fill=(255, 255, 255))
    draw.rectangle([lw, 0, lw + 4, h], fill=cv)  # separator line

    # Label text (wrapped)
    fn_lbl = _kfont(22, bold=True)
    fn_sub = _kfont(15)
    lines = textwrap.wrap(label, width=14)
    y_start = h // 2 - (len(lines) * 32) // 2 - 20
    for i, line in enumerate(lines):
        bb = draw.textbbox((0, 0), line, font=fn_lbl)
        x_pos = lw + (w - lw) // 2 - (bb[2] - bb[0]) // 2
        draw.text((x_pos, y_start + i * 34), line, font=fn_lbl, fill=(15, 23, 42))

    return _pil_stream(img)


def pil_flow(steps: list, color_hex: str, w=1360, h=480) -> io.BytesIO | None:
    """
    Clean professional flow diagram: horizontal or 3+3 layout.
    White background, clean boxes, minimal decoration.
    """
    if not PIL_AVAILABLE:
        return None

    cv = tuple(int(color_hex.lstrip("#")[i:i+2], 16) for i in (0, 2, 4))
    img = Image.new("RGB", (w, h), (248, 250, 252))
    draw = ImageDraw.Draw(img)

    n = len(steps)
    row1 = steps[:3]
    row2 = steps[3:]

    bw, bh = 380, 160
    gap = 28
    total_w1 = len(row1) * bw + (len(row1) - 1) * gap
    ox1 = (w - total_w1) // 2
    ry1 = 30
    ry2 = 270

    def draw_box(bx, by, idx, text):
        # Shadow
        draw.rectangle([bx + 3, by + 3, bx + bw + 3, by + bh + 3], fill=(220, 225, 235))
        # Box
        _rrect(draw, bx, by, bx + bw, by + bh, 8, fill=(255, 255, 255),
               outline=(203, 213, 225), lw=2)
        # Top accent
        _rrect(draw, bx, by, bx + bw, by + 6, 4, fill=cv)
        draw.rectangle([bx + 40, by, bx + bw, by + 6], fill=cv)
        # Number badge
        _rrect(draw, bx + 12, by + 14, bx + 44, by + 46, 6, fill=cv)
        fn_num = _kfont(18, bold=True)
        _draw_center(draw, str(idx), bx + 28, by + 30, fn_num, (255, 255, 255))
        # Text
        fn_txt = _kfont(15)
        wrapped = textwrap.wrap(text, width=21)
        y0 = by + 22 + (bh - 44 - len(wrapped) * 24) // 2
        for li, line in enumerate(wrapped[:3]):
            draw.text((bx + 56, y0 + li * 24), line, font=fn_txt, fill=(30, 41, 59))

    def draw_arrow_h(x, y):
        ax = x + bw + 4
        ay = y + bh // 2
        draw.polygon([(ax, ay - 9), (ax + 24, ay), (ax, ay + 9)], fill=cv)
        draw.line([(x + bw, ay), (ax, ay)], fill=cv, width=2)

    def draw_arrow_d(x, y_from, y_to):
        # Down-right corner: from end of row1 going down then left
        # "⌐" shape connector
        ax = x + bw // 2
        draw.line([(ax, y_from + bh), (ax, y_from + bh + 28)], fill=cv, width=2)
        draw.polygon([(ax - 9, y_from + bh + 8), (ax, y_from + bh + 32),
                       (ax + 9, y_from + bh + 8)], fill=cv)

    # Row 1
    for i, step in enumerate(row1):
        bx = ox1 + i * (bw + gap)
        draw_box(bx, ry1, i + 1, step)
        if i < len(row1) - 1:
            draw_arrow_h(bx, ry1)

    # Row 2 (reverse order, right to left)
    if row2:
        # Down arrow from last box of row1
        last_bx1 = ox1 + (len(row1) - 1) * (bw + gap)
        draw_arrow_d(last_bx1, ry1, ry2)

        total_w2 = len(row2) * bw + (len(row2) - 1) * gap
        ox2 = ox1 + (len(row1) - 1) * (bw + gap)
        for i, step in enumerate(row2):
            ri = len(row2) - 1 - i
            bx = ox2 - ri * (bw + gap)
            draw_box(bx, ry2, len(row1) + i + 1, step)
            if i < len(row2) - 1:
                # Left-pointing arrow
                px = bx - gap
                py = ry2 + bh // 2
                draw.line([(bx, py), (bx - gap + 4, py)], fill=cv, width=2)
                draw.polygon([(bx - gap + 4, py - 9), (bx - gap - 22, py),
                               (bx - gap + 4, py + 9)], fill=cv)

    return _pil_stream(img)


def pil_trie(formula_data: dict, w=1300, h=440) -> io.BytesIO | None:
    """Professional TRIE formula — clean 4-column grid."""
    if not PIL_AVAILABLE:
        return None

    img = Image.new("RGB", (w, h), (248, 250, 252))
    draw = ImageDraw.Draw(img)

    elements = formula_data.get("elements", [])
    cw = (w - 40) // 4
    colors = [e["color"] for e in elements]

    for i, el in enumerate(elements):
        x = 10 + i * (cw + 6)
        cv = tuple(int(el["color"].lstrip("#")[j:j+2], 16) for j in (0, 2, 4))

        # Card
        _rrect(draw, x, 0, x + cw, h - 10, 10, fill=(255, 255, 255),
               outline=(203, 213, 225), lw=1)

        # Header
        _rrect(draw, x, 0, x + cw, 80, 10, fill=cv)
        draw.rectangle([x, 40, x + cw, 80], fill=cv)  # flatten bottom

        # Letter
        fn_big = _kfont(56, bold=True)
        _draw_center(draw, el["letter"], x + cw // 2, 40, fn_big, (255, 255, 255))

        # Word
        fn_word = _kfont(17, bold=True)
        fn_kor = _kfont(13)
        draw.text((x + 14, 88), el["word"], font=fn_word, fill=cv)
        draw.text((x + 14, 112), f'({el["korean"]})', font=fn_kor, fill=(100, 116, 139))

        # Desc
        fn_desc = _kfont(14)
        wrapped_desc = textwrap.wrap(el["description"], 18)
        for li, ln in enumerate(wrapped_desc):
            draw.text((x + 14, 140 + li * 22), ln, font=fn_desc, fill=(51, 65, 85))

        # Divider
        draw.line([(x + 14, 200), (x + cw - 14, 200)], fill=(226, 232, 240), width=1)

        # Good example
        fn_ex = _kfont(12)
        draw.text((x + 14, 210), "✅ 올바른 예:", font=fn_ex, fill=(13, 107, 79))
        ex_text = el["example"].strip('"')
        wrapped_ex = textwrap.wrap(ex_text, 22)
        for li, ln in enumerate(wrapped_ex[:2]):
            draw.text((x + 14, 230 + li * 20), f'"{ln}"', font=fn_ex, fill=(15, 118, 110))

        # Bad example
        draw.text((x + 14, 285), "❌ 잘못된 예:", font=fn_ex, fill=(153, 27, 27))
        draw.text((x + 14, 305), el["bad_example"], font=fn_ex, fill=(185, 28, 28))

        # Separator between cards
        if i < 3:
            draw.line([(x + cw + 3, 20), (x + cw + 3, h - 30)],
                      fill=(226, 232, 240), width=1)

    return _pil_stream(img)


def pil_week_plan(plan_data: dict, w=1320, h=460) -> io.BytesIO | None:
    """Professional 4-week plan — timeline card layout."""
    if not PIL_AVAILABLE:
        return None

    img = Image.new("RGB", (w, h), (255, 255, 255))
    draw = ImageDraw.Draw(img)

    weeks = ["week1", "week2", "week3", "week4"]
    cw = (w - 60) // 4
    pad = 12

    for i, wk in enumerate(weeks):
        d = plan_data.get(wk, {})
        x = 10 + i * (cw + 12)
        c = d.get("color", "#1565C0")
        cv = tuple(int(c.lstrip("#")[j:j+2], 16) for j in (0, 2, 4))

        # Shadow
        draw.rectangle([x + 3, 3, x + cw + 3, h - 3], fill=(226, 232, 240))
        # Card
        draw.rectangle([x, 0, x + cw, h - 6], fill=(255, 255, 255))
        # Top bar
        draw.rectangle([x, 0, x + cw, 6], fill=cv)
        # Thin outline
        draw.rectangle([x, 0, x + cw, h - 6], fill=None, outline=(203, 213, 225))

        # Week number
        fn_wn = _kfont(13, bold=True)
        draw.text((x + pad, 14), f"WEEK {i+1}", font=fn_wn, fill=cv)

        # Theme
        fn_theme = _kfont(18, bold=True)
        theme = d.get("theme", "")
        emoji = d.get("emoji", "")
        draw.text((x + pad, 36), f"{emoji} {theme}", font=fn_theme, fill=(15, 23, 42))

        # Divider
        draw.line([(x + pad, 72), (x + cw - pad, 72)], fill=(226, 232, 240), width=1)

        # Tasks
        fn_task = _kfont(13)
        tasks = d.get("tasks", [])
        for ti, task in enumerate(tasks[:3]):
            ty = 84 + ti * 118
            # Circle number
            _rrect(draw, x + pad, ty + 2, x + pad + 22, ty + 24, 11, fill=cv)
            fn_n = _kfont(12, bold=True)
            _draw_center(draw, str(ti + 1), x + pad + 11, ty + 13, fn_n, (255, 255, 255))
            # Task text
            wrapped = textwrap.wrap(task, width=22)
            for li, ln in enumerate(wrapped[:3]):
                draw.text((x + pad + 28, ty + 2 + li * 20), ln,
                          font=fn_task, fill=(51, 65, 85))

        # Arrow to next
        if i < 3:
            ax = x + cw + 2
            ay = h // 2 - 20
            draw.polygon([(ax, ay), (ax + 10, ay + 8), (ax, ay + 16)], fill=cv)

    return _pil_stream(img)


def pil_context_bar(data: list, w: int = 2400, h: int = 900) -> io.BytesIO:
    """로그 스케일 수평 막대 그래프 — 컨텍스트 윈도우 진화 시각화."""
    import math
    img = Image.new("RGB", (w, h), "#0A1628")
    draw = ImageDraw.Draw(img)
    fn = _kfont(22)
    fn_sm = _kfont(18)
    fn_bold = _kfont(28, bold=True)

    pad_l, pad_r, pad_t, pad_b = 240, 60, 60, 60
    bar_h = 48
    gap = (h - pad_t - pad_b - len(data) * bar_h) // (len(data) + 1)
    max_tokens = max(d["tokens"] for d in data)

    for i, d in enumerate(data):
        y = pad_t + gap + i * (bar_h + gap)
        # log scale bar width
        log_ratio = math.log10(d["tokens"] + 1) / math.log10(max_tokens + 1)
        bar_w = int((w - pad_l - pad_r) * log_ratio)

        # bar background (dim)
        _rrect(draw, pad_l, y, pad_l + (w - pad_l - pad_r), y + bar_h, 6,
               fill=(255, 255, 255, 0), outline=(255, 255, 255, 20))
        draw.rectangle([(pad_l, y), (pad_l + (w - pad_l - pad_r), y + bar_h)],
                       fill="#1E293B")

        # colored bar
        try:
            r2, g2, b2 = int(d["color"][1:3], 16), int(d["color"][3:5], 16), int(d["color"][5:7], 16)
        except Exception:
            r2, g2, b2 = 29, 78, 216
        draw.rectangle([(pad_l, y), (pad_l + bar_w, y + bar_h)], fill=(r2, g2, b2))

        # model name (left)
        _draw_center(draw, d["model"], pad_l // 2, y + bar_h // 2, fn, "#E2E8F0")

        # year badge
        _draw_center(draw, str(d["year"]), pad_l - 60, y + bar_h // 2, fn_sm, "#64748B")

        # token count label
        t = d["tokens"]
        label = f"{t:,}" if t < 100_000 else (f"{t//1000:,}K" if t < 1_000_000 else f"{t//1_000_000:.1f}M")
        tx = pad_l + bar_w + 10
        try:
            draw.text((tx, y + bar_h // 2), label, font=fn_bold, fill=(r2, g2, b2), anchor="lm")
        except TypeError:
            draw.text((tx, y + bar_h // 2 - 14), label, font=fn_bold, fill=(r2, g2, b2))

    # title
    _draw_center(draw, "컨텍스트 윈도우 진화 (로그 스케일)", w // 2, 28, _kfont(26, bold=True),
                 "#60A5FA")

    return _pil_stream(img)


def pil_cost_bar(data: list, w: int = 2000, h: int = 700) -> io.BytesIO:
    """토큰 비용 감소 막대 그래프."""
    img = Image.new("RGB", (w, h), "#0A1628")
    draw = ImageDraw.Draw(img)
    fn = _kfont(22)
    fn_bold = _kfont(30, bold=True)

    n = len(data)
    bar_w = int((w - 160) / n * 0.55)
    group_w = (w - 160) // n
    pad_t, pad_b = 80, 100
    max_rel = max(d["relative"] for d in data)
    chart_h = h - pad_t - pad_b

    for i, d in enumerate(data):
        cx = 80 + i * group_w + group_w // 2
        bh = int(chart_h * d["relative"] / max_rel)
        by = pad_t + (chart_h - bh)

        # bar
        rel = d["relative"] / 100
        r2 = int(29 + (185 - 29) * rel)
        g2 = int(78 + (28 - 78) * rel)
        b2 = int(216 + (28 - 216) * rel)
        color = (max(0, min(255, r2)), max(0, min(255, g2)), max(0, min(255, b2)))

        draw.rectangle([(cx - bar_w // 2, by), (cx + bar_w // 2, pad_t + chart_h)], fill=color)

        # value label (top of bar)
        _draw_center(draw, f"${d['cost_per_1m_input']}", cx, by - 20, fn_bold,
                     (r2, g2, b2))

        # x-axis label (multiline)
        lines = d["period"].split("\n")
        for j, line in enumerate(lines):
            _draw_center(draw, line, cx, pad_t + chart_h + 20 + j * 28, fn, "#94A3B8")

    # title
    _draw_center(draw, "입력 토큰 100만개당 API 비용 (USD)", w // 2, 20,
                 _kfont(24, bold=True), "#60A5FA")

    return _pil_stream(img)


# ─────────────────────────────────────────────────────────────────
# SLIDE BUILDERS — Professional Edition
# ─────────────────────────────────────────────────────────────────

def build_title_slide(prs: Presentation) -> None:
    """Clean dark title slide — editorial style."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_dark"])

    # Accent line (left vertical)
    _rect(slide, int(W * 0.05), int(H * 0.20), int(W * 0.004), int(H * 0.58), D["blue"])

    # Year tag
    _txt(slide, int(W * 0.08), int(H * 0.20), int(W * 0.40), int(H * 0.07),
         f"INSIGHT REPORT  {LECTURE_DATE}", 11, color=D["blue"])

    # Main title
    _txt(slide, int(W * 0.08), int(H * 0.30), int(W * 0.78), int(H * 0.28),
         LECTURE_TITLE, 46, bold=True, color=D["text_white"])

    # Subtitle
    _txt(slide, int(W * 0.08), int(H * 0.59), int(W * 0.75), int(H * 0.10),
         LECTURE_SUBTITLE, 19, color="#94A3B8")

    # Thin divider
    _rect(slide, int(W * 0.08), int(H * 0.72), int(W * 0.60), int(H * 0.002), "#334155")

    # Author / date
    _txt(slide, int(W * 0.08), int(H * 0.75), int(W * 0.50), int(H * 0.06),
         f"{LECTURE_AUTHOR}", 13, color="#64748B")

    # Target audience tag (right side)
    _rect(slide, int(W * 0.74), int(H * 0.45), int(W * 0.22), int(H * 0.18), "#0F2B4C")
    _txt(slide, int(W * 0.74), int(H * 0.46), int(W * 0.22), int(H * 0.06),
         "대상", 10, color=D["blue"], align=PP_ALIGN.CENTER)
    _txt(slide, int(W * 0.74), int(H * 0.52), int(W * 0.22), int(H * 0.11),
         "B2C 영업\n마케팅\n상품기획", 12, color="#CBD5E1", align=PP_ALIGN.CENTER)


def build_toc_slide(prs: Presentation) -> None:
    """Clean table of contents — list format."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_light"])
    _accent_bar(slide, D["blue"])

    _txt(slide, int(W * 0.06), int(H * 0.06), int(W * 0.30), int(H * 0.10),
         "CONTENTS", 11, bold=True, color=D["blue"])
    _txt(slide, int(W * 0.06), int(H * 0.13), int(W * 0.60), int(H * 0.10),
         "강의 목차", 30, bold=True, color=D["text_h1"])
    _rect(slide, int(W * 0.06), int(H * 0.24), int(W * 0.88), int(H * 0.001), D["border"])

    for i, ch in enumerate(CHAPTERS):
        y = int(H * 0.27) + i * int(H * 0.13)
        col = CHAPTER_COLORS[i]
        # Hover line on hover
        _rect(slide, int(W * 0.06), y, int(W * 0.88), int(H * 0.12), D["bg_white"])
        # Number
        _txt(slide, int(W * 0.07), y + int(H * 0.015), int(W * 0.06), int(H * 0.09),
             f"0{i+1}", 26, bold=True, color=col)
        # Divider
        _rect(slide, int(W * 0.14), y + int(H * 0.025), int(W * 0.001), int(H * 0.07), D["border"])
        # Title
        _txt(slide, int(W * 0.16), y + int(H * 0.018), int(W * 0.46), int(H * 0.05),
             ch["title"], 16, bold=True, color=D["text_h1"])
        # Subtitle
        _txt(slide, int(W * 0.16), y + int(H * 0.065), int(W * 0.46), int(H * 0.04),
             ch["subtitle"], 11, color=D["text_muted"])
        # Thin bottom border
        _rect(slide, int(W * 0.06), y + int(H * 0.12), int(W * 0.88), int(H * 0.001), D["border_light"])


def build_chapter_divider(prs: Presentation, chapter: dict) -> None:
    """Minimal dark chapter divider."""
    ch_idx = chapter["id"] - 1
    col = CHAPTER_COLORS[ch_idx]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_dark"])
    _accent_bar(slide, col)

    # Large number (watermark style)
    _txt(slide, int(W * 0.68), int(H * 0.10), int(W * 0.30), int(H * 0.70),
         f"0{chapter['id']}", 140, bold=True, color="#0F2B4C")

    # Chapter label
    _txt(slide, int(W * 0.07), int(H * 0.18), int(W * 0.40), int(H * 0.08),
         f"CHAPTER  0{chapter['id']}", 11, bold=True, color=col)

    # Title
    _txt(slide, int(W * 0.07), int(H * 0.30), int(W * 0.68), int(H * 0.30),
         chapter["title"], 40, bold=True, color=D["text_white"])

    # Thin divider
    _rect(slide, int(W * 0.07), int(H * 0.62), int(W * 0.35), int(H * 0.002), col)

    # Subtitle
    _txt(slide, int(W * 0.07), int(H * 0.65), int(W * 0.60), int(H * 0.10),
         chapter["subtitle"], 17, color="#64748B")


def build_content_slide(prs: Presentation, title: str, bullets: list,
                         chapter_color: str) -> None:
    """Clean content slide — numbered bullet cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_light"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.80), int(H * 0.10),
         title, 25, bold=True, color=D["text_h1"])
    _rect(slide, int(W * 0.06), int(H * 0.15), int(W * 0.88), int(H * 0.001), D["border"])

    if not bullets:
        return

    y_start = int(H * 0.18)
    avail_h = int(H * 0.76)
    card_h = min(avail_h // len(bullets), int(H * 0.165)) - int(H * 0.008)

    for i, bullet in enumerate(bullets):
        y = y_start + i * (card_h + int(H * 0.008))
        # Shadow
        _rect(slide, int(W * 0.07) + int(W * 0.003), y + int(H * 0.004),
              int(W * 0.87), card_h, D["border_light"])
        # Card
        _rect(slide, int(W * 0.07), y, int(W * 0.87), card_h, D["bg_white"])
        # Left accent
        _rect(slide, int(W * 0.07), y, int(W * 0.005), card_h, chapter_color)
        # Number
        _txt(slide, int(W * 0.075), y + int(card_h * 0.20),
             int(W * 0.038), int(card_h * 0.60),
             f"{i+1:02d}", 14, bold=True, color=chapter_color, align=PP_ALIGN.CENTER)
        # Content
        _txt(slide, int(W * 0.12), y + int(card_h * 0.15),
             int(W * 0.79), int(card_h * 0.70),
             bullet, 13, color=D["text_body"])


def build_timeline_slide(prs: Presentation) -> None:
    """AI evolution timeline — clean horizontal flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_dark"])
    _accent_bar(slide, D["blue"])

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.80), int(H * 0.10),
         "AI 진화 타임라인 — 2022 → 2025", 25, bold=True, color=D["text_white"])
    _txt(slide, int(W * 0.06), int(H * 0.14), int(W * 0.60), int(H * 0.06),
         "질문-답변 → 이미지 생성 → 에이전트 → CLI 자율 실행", 13, color="#64748B")

    items = AI_TIMELINE
    n = len(items)
    x_step = int(W * 0.90 / n)
    x_start = int(W * 0.05)
    line_y = int(H * 0.52)

    # Timeline line
    _rect(slide, x_start, line_y, int(W * 0.90), int(H * 0.003), "#1E3A5F")

    phase_colors = [D["ch1"], D["ch1"], D["ch2"], D["ch2"],
                    D["ch3"], D["ch3"], D["ch4"], D["ch4"]]

    for i, ev in enumerate(items):
        cx = x_start + i * x_step + x_step // 2
        col = phase_colors[i]

        # Dot
        dr = int(W * 0.016)
        _rect(slide, cx - dr, line_y - dr, dr * 2, dr * 2, col)
        _rect(slide, cx - int(W * 0.007), line_y - int(W * 0.007),
              int(W * 0.014), int(W * 0.014), D["bg_dark"])
        _rect(slide, cx - int(W * 0.009), line_y - int(W * 0.009),
              int(W * 0.018), int(W * 0.018), col)

        tbox_w = int(W * 0.11)
        if i % 2 == 0:  # above
            _txt(slide, cx - tbox_w // 2, int(H * 0.16),
                 tbox_w, int(H * 0.09), f"{ev['year']}", 11, bold=True,
                 color=col, align=PP_ALIGN.CENTER)
            _txt(slide, cx - tbox_w // 2, int(H * 0.25),
                 tbox_w, int(H * 0.23), ev["event"], 10, color="#CBD5E1",
                 align=PP_ALIGN.CENTER)
            # Vertical connector
            _rect(slide, cx - int(W * 0.001), int(H * 0.44),
                  int(W * 0.002), int(H * 0.08), "#1E3A5F")
        else:  # below
            _rect(slide, cx - int(W * 0.001), line_y + int(H * 0.02),
                  int(W * 0.002), int(H * 0.08), "#1E3A5F")
            _txt(slide, cx - tbox_w // 2, int(H * 0.62),
                 tbox_w, int(H * 0.23), ev["event"], 10, color="#CBD5E1",
                 align=PP_ALIGN.CENTER)
            _txt(slide, cx - tbox_w // 2, int(H * 0.84),
                 tbox_w, int(H * 0.08), f"{ev['year']}", 11, bold=True,
                 color=col, align=PP_ALIGN.CENTER)

    # Phase labels
    phases = [("Phase 1\nQ&A", D["ch1"], 0.08), ("Phase 2\n이미지", D["ch2"], 0.31),
              ("Phase 3\n에이전트", D["ch3"], 0.56), ("Phase 4\nCLI", D["ch4"], 0.83)]
    for label, c, xf in phases:
        _txt(slide, int(W * xf), int(H * 0.92), int(W * 0.20), int(H * 0.07),
             label, 9, color=c, align=PP_ALIGN.CENTER)


def build_before_after_slide(prs: Presentation, job: str, data: dict,
                              chapter_color: str) -> None:
    """Professional Before/After comparison — two-panel table layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_light"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.80), int(H * 0.09),
         f"{job} 업무 — AI 도입 전후 비교", 25, bold=True, color=D["text_h1"])

    # Column headers
    half = int(W * 0.42)
    x_before = int(W * 0.05)
    x_after  = int(W * 0.53)
    y_header = int(H * 0.16)
    col_h    = int(H * 0.075)

    _rect(slide, x_before, y_header, half, col_h, "#FEE2E2")
    _txt(slide, x_before, y_header, half, col_h,
         "BEFORE  — AI 도입 전", 14, bold=True, color="#991B1B", align=PP_ALIGN.CENTER)

    _rect(slide, x_after, y_header, half, col_h, "#D1FAE5")
    _txt(slide, x_after, y_header, half, col_h,
         "AFTER  — AI 도입 후", 14, bold=True, color="#065F46", align=PP_ALIGN.CENTER)

    # Rows
    before_items = data.get("before", [])[:4]
    after_items  = data.get("after", [])[:4]
    n_rows = max(len(before_items), len(after_items))

    row_h = int(H * 0.14)
    for i in range(n_rows):
        y = int(H * 0.25) + i * (row_h + int(H * 0.01))
        bg_b = "#FFF1F2" if i % 2 == 0 else "#FFF5F5"
        bg_a = "#ECFDF5" if i % 2 == 0 else "#F0FDF4"

        if i < len(before_items):
            _rect(slide, x_before, y, half, row_h, bg_b)
            _rect(slide, x_before, y, int(W * 0.004), row_h, "#FCA5A5")
            _txt(slide, x_before + int(W * 0.01), y + int(row_h * 0.15),
                 int(half - int(W * 0.02)), int(row_h * 0.70),
                 before_items[i], 12, color="#7F1D1D")

        if i < len(after_items):
            _rect(slide, x_after, y, half, row_h, bg_a)
            _rect(slide, x_after, y, int(W * 0.004), row_h, "#6EE7B7")
            _txt(slide, x_after + int(W * 0.01), y + int(row_h * 0.15),
                 int(half - int(W * 0.02)), int(row_h * 0.70),
                 after_items[i], 12, color="#064E3B")

    # Center divider + arrow
    _rect(slide, int(W * 0.48), int(H * 0.16), int(W * 0.04), int(H * 0.76), D["bg_light"])
    _txt(slide, int(W * 0.48), int(H * 0.48), int(W * 0.04), int(H * 0.08),
         "→", 22, bold=True, color=chapter_color, align=PP_ALIGN.CENTER)


def build_agent_example_slide(prs: Presentation, agents: list, title="AI 에이전트 도입 사례") -> None:
    """Professional agent showcase — 3-column card layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_dark"])
    _accent_bar(slide, D["ch4"])

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.80), int(H * 0.09),
         title, 25, bold=True, color=D["text_white"])
    _txt(slide, int(W * 0.06), int(H * 0.13), int(W * 0.60), int(H * 0.05),
         "실무에서 즉시 적용 가능한 자동화 에이전트", 13, color="#64748B")

    cols = min(len(agents), 3)
    cw = int(W * 0.28)
    spacing = int((W * 0.94 - cols * cw) / (cols + 1))
    y_top = int(H * 0.21)
    ch = int(H * 0.68)
    col_list = [D["ch1"], D["ch2"], D["ch3"], D["ch4"], D["ch5"]]

    for i, ag in enumerate(agents[:cols]):
        x = spacing + i * (cw + spacing)
        col = col_list[i % len(col_list)]

        # Card shadow
        _rect(slide, x + int(W * 0.003), y_top + int(H * 0.005), cw, ch, "#0F2B4C")
        # Card
        _rect(slide, x, y_top, cw, ch, "#0F2040")
        # Top accent
        _rect(slide, x, y_top, cw, int(H * 0.008), col)

        name = ag.get("name", "")
        job = ag.get("job_role", "")
        desc = ag.get("description", "")
        tools = ag.get("tools_used", [])
        output = ag.get("output", "")
        difficulty = ag.get("difficulty", "중급")

        # Job tag
        _txt(slide, x + int(cw * 0.05), y_top + int(ch * 0.03),
             int(cw * 0.90), int(ch * 0.08), job, 10, color=col)

        # Name
        _txt(slide, x + int(cw * 0.05), y_top + int(ch * 0.11),
             int(cw * 0.90), int(ch * 0.16), name, 13, bold=True, color="#F1F5F9")

        # Divider
        _rect(slide, x + int(cw * 0.05), y_top + int(ch * 0.28),
              int(cw * 0.90), int(H * 0.001), "#1E3A5F")

        # Description
        _txt(slide, x + int(cw * 0.05), y_top + int(ch * 0.30),
             int(cw * 0.90), int(ch * 0.25),
             desc[:90] + "..." if len(desc) > 90 else desc,
             11, color="#94A3B8")

        # Tools
        tools_str = "  ·  ".join(tools[:3])
        _txt(slide, x + int(cw * 0.05), y_top + int(ch * 0.57),
             int(cw * 0.90), int(ch * 0.10),
             f"도구: {tools_str}", 9, color=col)

        # Output
        _txt(slide, x + int(cw * 0.05), y_top + int(ch * 0.68),
             int(cw * 0.90), int(ch * 0.12),
             f"산출물: {output[:50]}", 10, color="#CBD5E1")

        # Difficulty badge
        dc = {"초급": "#059669", "중급": "#D97706", "고급": "#DC2626"}.get(difficulty, "#64748B")
        _rect(slide, x + int(cw * 0.62), y_top + int(ch * 0.84),
              int(cw * 0.34), int(ch * 0.10), dc)
        _txt(slide, x + int(cw * 0.62), y_top + int(ch * 0.84),
             int(cw * 0.34), int(ch * 0.10),
             difficulty, 9, bold=True, color="#FFFFFF", align=PP_ALIGN.CENTER)


def build_closing_slide(prs: Presentation) -> None:
    """Professional closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_dark"])
    _accent_bar(slide, D["blue"])

    # Left decorative bar
    _rect(slide, int(W * 0.05), int(H * 0.18), int(W * 0.004), int(H * 0.62), D["blue"])

    _txt(slide, int(W * 0.08), int(H * 0.20), int(W * 0.70), int(H * 0.10),
         "KEY TAKEAWAY", 11, bold=True, color=D["blue"])

    _txt(slide, int(W * 0.08), int(H * 0.29), int(W * 0.80), int(H * 0.18),
         "AI는 대체가 아닌 증폭입니다", 38, bold=True, color=D["text_white"])

    _txt(slide, int(W * 0.08), int(H * 0.50), int(W * 0.75), int(H * 0.10),
         "반복 업무는 AI에게, 창의적 판단은 당신에게.", 20, color="#94A3B8")

    _rect(slide, int(W * 0.08), int(H * 0.62), int(W * 0.50), int(H * 0.001), "#1E3A5F")

    msgs = [
        "01  반복 업무 자동화로 고부가 업무에 집중",
        "02  작은 자동화부터 시작해 점진적으로 확장",
        "03  AI와 협업하는 역량이 2026년의 핵심 경쟁력",
    ]
    for i, m in enumerate(msgs):
        _txt(slide, int(W * 0.08), int(H * 0.65) + i * int(H * 0.09),
             int(W * 0.75), int(H * 0.08), m, 14, color="#CBD5E1")

    # Q&A
    _txt(slide, int(W * 0.06), int(H * 0.92), int(W * 0.88), int(H * 0.06),
         "Q & A", 16, bold=True, color=D["blue"], align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────
# NEW SLIDE BUILDERS — Professional Style
# ─────────────────────────────────────────────────────────────────

def build_ai_definition_slide(prs, chapter_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_light"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.80), int(H * 0.09),
         "AI란 무엇인가? — 3가지 비유로 이해하기", 25, bold=True, color=D["text_h1"])
    _rect(slide, int(W * 0.06), int(H * 0.14), int(W * 0.88), int(H * 0.001), D["border"])

    analogies = [
        ("🚗", "검색 → 대화",
         "Google에게 '제안서 써줘'라고 해도 결과 없음.\nAI는 해당 요청을 이해하고 실제로 작성해 줌.\n→ 도구가 아닌 협업 파트너로의 전환.",
         D["ch1"]),
        ("✍️", "타이핑 → 말하기",
         "스마트폰 등장 전 우리는 모두 버튼을 눌렀다.\nAI는 자연어로 명령하는 시대를 열었다.\n→ 누구나 전문 도구를 쓸 수 있는 시대.",
         D["ch2"]),
        ("🤝", "도구 → 동료",
         "엑셀은 시킨 것만 한다. 창의성 없음.\nAI 에이전트는 목표를 주면 스스로 계획 실행.\n→ '경쟁사 분석해줘' → 크롤링·분석·정리 자동.",
         D["ch3"]),
    ]
    cw = int(W * 0.28)
    spacing = int(W * 0.045)
    ch = int(H * 0.70)
    y0 = int(H * 0.18)

    for i, (icon, title, body, col) in enumerate(analogies):
        x = spacing + i * (cw + spacing)
        _rect(slide, x + int(W * 0.003), y0 + int(H * 0.005), cw, ch, D["border_light"])
        _rect(slide, x, y0, cw, ch, D["bg_white"])
        _rect(slide, x, y0, cw, int(H * 0.006), col)
        # Icon + title
        _txt(slide, x + int(cw * 0.05), y0 + int(ch * 0.04),
             int(cw * 0.90), int(ch * 0.12), icon, 30, color=D["text_h1"])
        _txt(slide, x + int(cw * 0.05), y0 + int(ch * 0.20),
             int(cw * 0.90), int(ch * 0.12), title, 16, bold=True, color=col)
        _rect(slide, x + int(cw * 0.05), y0 + int(ch * 0.34),
              int(cw * 0.90), int(H * 0.001), D["border_light"])
        _txt(slide, x + int(cw * 0.05), y0 + int(ch * 0.37),
             int(cw * 0.90), int(ch * 0.58), body, 12, color=D["text_body"])


def build_stats_slide(prs, stat_key: str, chapter_color: str):
    stat = INDUSTRY_STATS.get(stat_key, {})
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_light"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.80), int(H * 0.09),
         stat.get("label", ""), 25, bold=True, color=D["text_h1"])

    stream = pil_stat_card(
        value=stat.get("value", ""),
        label=stat.get("label", ""),
        sublabel=stat.get("timeframe", ""),
        color_hex=stat.get("color", chapter_color),
    )
    if stream:
        _pil_img(slide, stream,
                 int(W * 0.05), int(H * 0.17), int(W * 0.46), int(H * 0.70))

    _txt(slide, int(W * 0.57), int(H * 0.20), int(W * 0.38), int(H * 0.08),
         "왜 중요한가?", 16, bold=True, color=chapter_color)
    _rect(slide, int(W * 0.57), int(H * 0.29), int(W * 0.36), int(H * 0.001), D["border"])

    for i, b in enumerate(stat.get("bullets", [])[:3]):
        y = int(H * 0.32) + i * int(H * 0.165)
        _rect(slide, int(W * 0.57), y, int(W * 0.005), int(H * 0.12), chapter_color)
        _txt(slide, int(W * 0.59), y + int(H * 0.01),
             int(W * 0.36), int(H * 0.13), b, 13, color=D["text_body"])

    _txt(slide, int(W * 0.57), int(H * 0.88), int(W * 0.38), int(H * 0.06),
         f"출처: {stat.get('source', '')}", 10, color=D["text_light"])


def build_economic_impact_slide(prs, chapter_color: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_light"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.85), int(H * 0.09),
         "AI가 만들어내는 경제적 가치 — 글로벌 데이터", 25, bold=True, color=D["text_h1"])
    _rect(slide, int(W * 0.06), int(H * 0.14), int(W * 0.88), int(H * 0.001), D["border"])

    for i, key in enumerate(["mckinsey_value", "microsoft_time_saved"]):
        stat = INDUSTRY_STATS.get(key, {})
        x = int(W * 0.05) + i * int(W * 0.48)
        stream = pil_stat_card(
            value=stat.get("value", ""),
            label=stat.get("label", ""),
            sublabel=stat.get("timeframe", ""),
            color_hex=stat.get("color", chapter_color),
            w=780, h=400,
        )
        if stream:
            _pil_img(slide, stream, x, int(H * 0.17), int(W * 0.44), int(H * 0.58))
        _txt(slide, x, int(H * 0.78), int(W * 0.44), int(H * 0.07),
             f"출처: {stat.get('source', '')}", 10, color=D["text_light"],
             align=PP_ALIGN.CENTER)


def build_job_stats_slide(prs, job: str, chapter_color: str):
    data = JOB_STATS.get(job, {})
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_light"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.85), int(H * 0.09),
         f"{job} — AI 도입 효과 수치", 25, bold=True, color=D["text_h1"])

    stream = pil_stat_card(
        value=data.get("headline_stat", ""),
        label=data.get("headline_label", ""),
        sublabel=f"출처: {data.get('headline_source', '')}",
        color_hex=chapter_color,
        w=700, h=360,
    )
    if stream:
        _pil_img(slide, stream, int(W * 0.05), int(H * 0.17), int(W * 0.40), int(H * 0.52))

    positions = [
        (int(W * 0.50), int(H * 0.17)),
        (int(W * 0.76), int(H * 0.17)),
        (int(W * 0.50), int(H * 0.50)),
        (int(W * 0.76), int(H * 0.50)),
    ]
    scw = int(W * 0.23)
    sch = int(H * 0.28)
    for i, stat in enumerate(data.get("stats", [])[:4]):
        x, y = positions[i]
        _rect(slide, x + int(W * 0.003), y + int(H * 0.005), scw, sch, D["border_light"])
        _rect(slide, x, y, scw, sch, D["bg_white"])
        _rect(slide, x, y, scw, int(H * 0.005), chapter_color)
        _txt(slide, x + int(scw * 0.06), y + int(sch * 0.12),
             int(scw * 0.88), int(sch * 0.38),
             stat.get("value", ""), 22, bold=True, color=chapter_color)
        _txt(slide, x + int(scw * 0.06), y + int(sch * 0.52),
             int(scw * 0.88), int(sch * 0.28), stat.get("metric", ""), 11, color=D["text_body"])
        _txt(slide, x + int(scw * 0.06), y + int(sch * 0.80),
             int(scw * 0.88), int(sch * 0.15), stat.get("source", ""), 9, color=D["text_light"])

    # Agent hint bar
    ae = data.get("agent_examples", [{}])[0]
    _rect(slide, int(W * 0.05), int(H * 0.73), int(W * 0.90), int(H * 0.16),
          D["blue_light"])
    _rect(slide, int(W * 0.05), int(H * 0.73), int(W * 0.005), int(H * 0.16), D["blue"])
    _txt(slide, int(W * 0.07), int(H * 0.75), int(W * 0.85), int(H * 0.12),
         f"에이전트 도입 예시: {ae.get('name', '')} — {ae.get('description', '')[:65]}",
         12, color=D["navy"])


def build_job_agent_workflow_slide(prs, job: str, chapter_color: str):
    data = JOB_STATS.get(job, {})
    steps = data.get("agent_workflow_steps", [])
    ae = data.get("agent_example", {})

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_light"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.85), int(H * 0.09),
         f"{job} 에이전트 자동화 워크플로우 — 6단계", 25, bold=True, color=D["text_h1"])

    stream = pil_flow(steps=steps, color_hex=chapter_color)
    if stream:
        _pil_img(slide, stream, int(W * 0.04), int(H * 0.16), int(W * 0.92), int(H * 0.55))

    # Meta row
    meta = [
        ("트리거", ae.get("trigger", "")),
        ("사용 도구", "  ·  ".join(ae.get("tools", [])[:3])),
        ("산출물", ae.get("output", "")),
        ("절약 시간", ae.get("time_saved", "")),
    ]
    mw = int(W * 0.22)
    y_meta = int(H * 0.75)
    mh = int(H * 0.18)
    for i, (k, v) in enumerate(meta):
        x = int(W * 0.04) + i * (mw + int(W * 0.01))
        _rect(slide, x + int(W * 0.003), y_meta + int(H * 0.005), mw, mh, D["border_light"])
        _rect(slide, x, y_meta, mw, mh, D["bg_white"])
        _rect(slide, x, y_meta, mw, int(H * 0.005), chapter_color)
        _txt(slide, x + int(mw * 0.06), y_meta + int(mh * 0.08),
             int(mw * 0.88), int(mh * 0.32), k, 10, bold=True, color=chapter_color)
        _txt(slide, x + int(mw * 0.06), y_meta + int(mh * 0.44),
             int(mw * 0.88), int(mh * 0.48), v or "—", 10, color=D["text_body"])


def build_job_agent_example_slide(prs, job: str, chapter_color: str):
    data = JOB_STATS.get(job, {})
    agents = data.get("agent_examples", [])
    build_agent_example_slide(prs, agents, title=f"{job} 에이전트 도입 사례")


def build_common_patterns_slide(prs, chapter_color: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_light"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.85), int(H * 0.09),
         "3개 직무 공통 AI 자동화 패턴", 25, bold=True, color=D["text_h1"])
    _txt(slide, int(W * 0.06), int(H * 0.13), int(W * 0.70), int(H * 0.05),
         "영업·마케팅·상품기획 모두 동일한 3가지 패턴으로 생산성이 향상됩니다",
         13, color=D["text_muted"])

    patterns = [
        ("반복 업무 자동화", D["ch1"],
         ["정기 리포트 자동 생성·발송", "데이터 수집 자동 크롤링", "문서·제안서 초안 자동 생성"]),
        ("데이터 분석 가속화", D["ch2"],
         ["대용량 리뷰·데이터 즉시 분석", "트렌드 3~6개월 조기 감지", "경쟁사 인텔리전스 자동화"]),
        ("콘텐츠 대량 생산", D["ch3"],
         ["기획서/제안서 초안: 수 시간→수 분", "SNS 7일치 콘텐츠 30분 생성", "A/B 변형 50개 즉시 생성"]),
    ]
    cw = int(W * 0.28)
    spacing = int(W * 0.045)
    ch = int(H * 0.68)
    y0 = int(H * 0.22)

    for i, (title, col, items) in enumerate(patterns):
        x = spacing + i * (cw + spacing)
        _rect(slide, x + int(W * 0.003), y0 + int(H * 0.005), cw, ch, D["border_light"])
        _rect(slide, x, y0, cw, ch, D["bg_white"])
        _rect(slide, x, y0, cw, int(H * 0.006), col)
        _txt(slide, x + int(cw * 0.05), y0 + int(ch * 0.06),
             int(cw * 0.90), int(ch * 0.18), title, 16, bold=True, color=col)
        _rect(slide, x + int(cw * 0.05), y0 + int(ch * 0.25),
              int(cw * 0.90), int(H * 0.001), D["border_light"])
        for j, item in enumerate(items):
            y = y0 + int(ch * 0.28) + j * int(ch * 0.22)
            _rect(slide, x + int(cw * 0.05), y + int(ch * 0.04),
                  int(W * 0.008), int(ch * 0.12), col)
            _txt(slide, x + int(cw * 0.10), y,
                 int(cw * 0.86), int(ch * 0.20), item, 12, color=D["text_body"])


def build_agent_definition_slide(prs, chapter_color: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_light"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.80), int(H * 0.09),
         "AI 에이전트란 무엇인가?", 25, bold=True, color=D["text_h1"])

    # Comparison box
    half = int(W * 0.42)
    _rect(slide, int(W * 0.05), int(H * 0.17), half, int(H * 0.22), "#FEE2E2")
    _rect(slide, int(W * 0.05), int(H * 0.17), half, int(H * 0.008), "#DC2626")
    _txt(slide, int(W * 0.07), int(H * 0.19), int(half - int(W * 0.04)), int(H * 0.06),
         "일반 AI (ChatGPT)", 13, bold=True, color="#991B1B")
    _txt(slide, int(W * 0.07), int(H * 0.26), int(half - int(W * 0.04)), int(H * 0.12),
         "질문에 답하는 기계\n한 번에 하나의 답변만 생성\n외부 도구 사용 불가", 12, color="#7F1D1D")

    _rect(slide, int(W * 0.53), int(H * 0.17), half, int(H * 0.22), "#D1FAE5")
    _rect(slide, int(W * 0.53), int(H * 0.17), half, int(H * 0.008), "#059669")
    _txt(slide, int(W * 0.55), int(H * 0.19), int(half - int(W * 0.04)), int(H * 0.06),
         "AI 에이전트 (Claude Code 등)", 13, bold=True, color="#065F46")
    _txt(slide, int(W * 0.55), int(H * 0.26), int(half - int(W * 0.04)), int(H * 0.12),
         "목표를 주면 스스로 계획하고 실행\n다단계 작업을 자율 처리\n외부 도구·인터넷·코드 사용 가능", 12, color="#064E3B")

    _txt(slide, int(W * 0.47), int(H * 0.22), int(W * 0.06), int(H * 0.10),
         "→", 22, bold=True, color=chapter_color, align=PP_ALIGN.CENTER)

    # 3 components
    comps = [
        ("Tool Use", "🔧", "인터넷 검색, 파일 읽기, 코드 실행, API 호출 등 외부 도구를 직접 사용"),
        ("Memory", "🧠", "이전 대화와 작업을 기억. 문맥을 유지하며 일관된 다단계 작업 수행"),
        ("Orchestration", "🎯", "큰 목표를 작은 단계로 자동 분해하고 순서대로 실행하는 계획 능력"),
    ]
    cw = int(W * 0.28)
    spacing = int(W * 0.045)
    ch = int(H * 0.26)
    y0 = int(H * 0.64)

    for i, (name, icon, desc) in enumerate(comps):
        x = spacing + i * (cw + spacing)
        _rect(slide, x + int(W * 0.003), y0 + int(H * 0.005), cw, ch, D["border_light"])
        _rect(slide, x, y0, cw, ch, D["bg_dark"])
        _rect(slide, x, y0, cw, int(H * 0.006), CHAPTER_COLORS[i])
        _txt(slide, x + int(cw * 0.08), y0 + int(ch * 0.08),
             int(cw * 0.20), int(ch * 0.35), icon, 22, color=D["text_white"])
        _txt(slide, x + int(cw * 0.32), y0 + int(ch * 0.08),
             int(cw * 0.64), int(ch * 0.30), name, 14, bold=True, color=CHAPTER_COLORS[i])
        _txt(slide, x + int(cw * 0.06), y0 + int(ch * 0.45),
             int(cw * 0.88), int(ch * 0.50), desc, 11, color="#94A3B8")


def build_agent_architecture_slide(prs, chapter_color: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_light"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.85), int(H * 0.09),
         "에이전트 실행 메커니즘 — 6단계 루프", 25, bold=True, color=D["text_h1"])

    arch_steps = [
        "목표 입력 (자연어 명령)", "계획 수립 (단계 분해)",
        "도구 선택·실행 (검색/코드/API)",
        "결과 관찰 (성공/실패 판단)", "반복 또는 수정 (자율 판단)",
        "최종 결과 출력 (리포트/파일)",
    ]
    stream = pil_flow(steps=arch_steps, color_hex=chapter_color)
    if stream:
        _pil_img(slide, stream, int(W * 0.04), int(H * 0.17), int(W * 0.92), int(H * 0.56))

    _rect(slide, int(W * 0.05), int(H * 0.78), int(W * 0.90), int(H * 0.16), D["blue_light"])
    _rect(slide, int(W * 0.05), int(H * 0.78), int(W * 0.005), int(H * 0.16), D["blue"])
    _txt(slide, int(W * 0.07), int(H * 0.80), int(W * 0.86), int(H * 0.12),
         '실제 사례: "경쟁사 가격 분석 리포트 만들어줘" →  에이전트가 URL 확인 → 크롤링 → 정리 → 비교 분석 → PDF 생성 → 이메일 발송까지 자동 실행',
         12, color=D["navy"])


def build_no_code_agents_slide(prs, chapter_color: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_light"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.85), int(H * 0.09),
         "코딩 없이 에이전트 만들기 — 노코드 3대 플랫폼", 25, bold=True, color=D["text_h1"])

    tools = [
        {"name": "Zapier", "icon": "⚡",
         "desc": "7,000개 이상 앱 연결\n드래그&드롭 자동화\n가장 많이 쓰이는 노코드",
         "best": "영업 알림·CRM 자동화",
         "ex": "Gmail 수신 → Claude 요약 → Notion 저장",
         "cost": "무료 / 월$19.99~", "col": "#FF4F00"},
        {"name": "Make (Integromat)", "icon": "🔗",
         "desc": "복잡한 워크플로우 시각화\nZapier 대비 저렴한 가격\n데이터 가공 유연성 높음",
         "best": "마케팅 콘텐츠 파이프라인",
         "ex": "신제품 등록 → SNS 카피 생성 → Buffer 예약",
         "cost": "무료 / 월$9~", "col": "#6D00CC"},
        {"name": "Claude / ChatGPT 직접 활용", "icon": "🤖",
         "desc": "대화형 프롬프트만으로\n별도 설정 없이 즉시 사용\n가장 낮은 진입 장벽",
         "best": "상품기획·문서 작성",
         "ex": "경쟁사 URL 붙여넣기 → 분석 + 기획서 초안 즉시",
         "cost": "무료 / 월$20~", "col": D["ch1"]},
    ]
    cw = int(W * 0.28)
    spacing = int(W * 0.045)
    ch = int(H * 0.74)
    y0 = int(H * 0.18)

    for i, t in enumerate(tools):
        x = spacing + i * (cw + spacing)
        cv = t["col"]
        _rect(slide, x + int(W * 0.003), y0 + int(H * 0.005), cw, ch, D["border_light"])
        _rect(slide, x, y0, cw, ch, D["bg_white"])
        _rect(slide, x, y0, cw, int(H * 0.006), cv)

        _txt(slide, x + int(cw * 0.05), y0 + int(ch * 0.04),
             int(cw * 0.90), int(ch * 0.10),
             f"{t['icon']}  {t['name']}", 16, bold=True, color=cv)
        _txt(slide, x + int(cw * 0.05), y0 + int(ch * 0.16),
             int(cw * 0.90), int(ch * 0.22), t["desc"], 11, color=D["text_body"])
        _rect(slide, x + int(cw * 0.05), y0 + int(ch * 0.39),
              int(cw * 0.90), int(H * 0.001), D["border_light"])
        _txt(slide, x + int(cw * 0.05), y0 + int(ch * 0.41),
             int(cw * 0.90), int(ch * 0.08),
             f"최적 용도: {t['best']}", 11, bold=True, color=cv)
        _txt(slide, x + int(cw * 0.05), y0 + int(ch * 0.52),
             int(cw * 0.90), int(ch * 0.20), t["ex"], 11, color=D["text_muted"])
        _rect(slide, x + int(cw * 0.05), y0 + int(ch * 0.76),
              int(cw * 0.90), int(ch * 0.12), D["bg_subtle"])
        _txt(slide, x + int(cw * 0.05), y0 + int(ch * 0.77),
             int(cw * 0.90), int(ch * 0.10),
             f"비용: {t['cost']}", 11, bold=True, color=D["text_h1"], align=PP_ALIGN.CENTER)


def build_first_agent_guide_slide(prs, chapter_color: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_light"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.85), int(H * 0.09),
         "지금 당장 첫 번째 에이전트 만들기 — 5단계", 25, bold=True, color=D["text_h1"])

    steps = [
        ("반복 업무 선택", "매주 하는 작업 중 '자동화하면 좋겠다' 싶은 것 1개 선택", D["ch1"]),
        ("입력·출력 정의", "'무엇을 넣으면 → 무엇이 나와야 하는가' 명확히 정의", D["ch2"]),
        ("필요 도구 결정", "검색? 이메일? CRM? 파일 읽기? 필요한 도구 목록 메모", D["ch3"]),
        ("AI에게 설계 요청", '"내 업무는 ~이야. 이걸 자동화하는 에이전트 만들어줘" 그대로 입력', D["ch4"]),
        ("테스트 → 개선 반복", "처음엔 70점이어도 OK — 피드백으로 점진적 개선", D["ch5"]),
    ]
    rh = int(H * 0.14)
    rw = int(W * 0.87)
    y0 = int(H * 0.18)

    for i, (title, desc, col) in enumerate(steps):
        y = y0 + i * (rh + int(H * 0.01))
        _rect(slide, int(W * 0.065) + int(W * 0.003), y + int(H * 0.005), rw, rh, D["border_light"])
        _rect(slide, int(W * 0.065), y, rw, rh, D["bg_white"])
        _rect(slide, int(W * 0.065), y, int(W * 0.005), rh, col)
        # Step number
        _rect(slide, int(W * 0.09), y + int(rh * 0.2),
              int(W * 0.05), int(rh * 0.6), col)
        _txt(slide, int(W * 0.09), y + int(rh * 0.2),
             int(W * 0.05), int(rh * 0.6),
             f"0{i+1}", 16, bold=True, color=D["text_white"], align=PP_ALIGN.CENTER)
        # Title
        _txt(slide, int(W * 0.15), y + int(rh * 0.10),
             int(W * 0.28), int(rh * 0.45), title, 14, bold=True, color=D["text_h1"])
        # Desc
        _txt(slide, int(W * 0.15), y + int(rh * 0.57),
             int(W * 0.72), int(rh * 0.36), desc, 11, color=D["text_muted"])


def build_tool_matrix_slide(prs, chapter_color: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_light"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.85), int(H * 0.09),
         "직무별 AI 도구 추천 매트릭스", 25, bold=True, color=D["text_h1"])

    jobs = ["영업", "마케팅", "상품기획"]
    job_cols = [D["ch1"], D["ch2"], D["ch3"]]
    cw = int(W * 0.28)
    spacing = int(W * 0.045)
    rh = int(H * 0.155)
    y_header = int(H * 0.17)

    for i, (job, col) in enumerate(zip(jobs, job_cols)):
        x = spacing + i * (cw + spacing)
        _rect(slide, x, y_header, cw, int(H * 0.07), col)
        _txt(slide, x, y_header, cw, int(H * 0.07),
             job, 15, bold=True, color=D["text_white"], align=PP_ALIGN.CENTER)

        for j, t in enumerate(TOOL_MATRIX.get(job, [])[:4]):
            y = y_header + int(H * 0.075) + j * (rh + int(H * 0.005))
            bg = D["bg_white"] if j % 2 == 0 else D["bg_light"]
            _rect(slide, x, y, cw, rh, bg)
            _rect(slide, x, y, int(W * 0.005), rh, col)
            _txt(slide, x + int(cw * 0.05), y + int(rh * 0.06),
                 int(cw * 0.92), int(rh * 0.38),
                 t["tool"], 12, bold=True, color=D["text_h1"])
            _txt(slide, x + int(cw * 0.05), y + int(rh * 0.44),
                 int(cw * 0.92), int(rh * 0.28),
                 t["use_case"], 10, color=D["text_body"])
            dc = {"쉬움": "#059669", "보통": "#D97706"}.get(t["difficulty"], "#64748B")
            _rect(slide, x + int(cw * 0.62), y + int(rh * 0.74),
                  int(cw * 0.35), int(rh * 0.22), dc)
            _txt(slide, x + int(cw * 0.62), y + int(rh * 0.74),
                 int(cw * 0.35), int(rh * 0.22),
                 t["difficulty"], 9, bold=True, color=D["text_white"],
                 align=PP_ALIGN.CENTER)


def build_trie_formula_slide(prs, chapter_color: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_light"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.85), int(H * 0.09),
         "TRIE 프롬프트 공식 — 4요소로 AI 품질 10배 향상", 25, bold=True, color=D["text_h1"])

    stream = pil_trie(TRIE_FORMULA)
    if stream:
        _pil_img(slide, stream, int(W * 0.04), int(H * 0.17), int(W * 0.92), int(H * 0.73))

    _txt(slide, int(W * 0.06), int(H * 0.93), int(W * 0.88), int(H * 0.05),
         "이 4가지 공식은 ChatGPT, Claude, Gemini 등 모든 AI 서비스에 동일하게 적용됩니다",
         11, color=D["text_muted"], align=PP_ALIGN.CENTER)


def build_trie_example_slide(prs, chapter_color: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_light"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.85), int(H * 0.09),
         "나쁜 프롬프트 vs TRIE 공식 적용 — 실전 비교", 25, bold=True, color=D["text_h1"])

    examples = [
        {
            "job": "영업", "bad": '"제안서 써줘"',
            "good": (
                "[역할] 당신은 IT 솔루션 영업 전문가입니다.\n"
                "[과제] 삼성전자 물류팀을 위한 SaaS 솔루션 제안서 초안을 작성해주세요.\n"
                "[입력] 예산: 연 2억, 도입 목적: 물류 효율화, 담당자: 구매팀 김과장\n"
                "[기대] 5페이지 분량, 핵심 가치 3가지, ROI 수치 포함, 한국어"
            ),
            "gain": "응답 품질 5배 향상",
        },
        {
            "job": "마케팅", "bad": '"SNS 글 써줘"',
            "good": (
                "[역할] 당신은 MZ세대 디지털 마케터입니다.\n"
                "[과제] 무선이어폰 출시 인스타그램 피드 카피 3개 버전을 작성해주세요.\n"
                "[입력] 제품명: SoundX Pro, 타깃: 20~30대 직장인, 강점: 35시간 배터리\n"
                "[기대] 각 150자 이내, 해시태그 5개, 구매 유도 CTA 포함"
            ),
            "gain": "전환율 +18% 기대",
        },
    ]

    rh = int(H * 0.38)
    for i, ex in enumerate(examples):
        y = int(H * 0.17) + i * (rh + int(H * 0.02))

        # Bad
        _rect(slide, int(W * 0.03), y, int(W * 0.26), rh, "#FFF1F2")
        _rect(slide, int(W * 0.03), y, int(W * 0.26), int(H * 0.006), "#DC2626")
        _txt(slide, int(W * 0.04), y + int(rh * 0.06),
             int(W * 0.24), int(rh * 0.16),
             f"❌  [{ex['job']}] 잘못된 예", 12, bold=True, color="#991B1B")
        _txt(slide, int(W * 0.04), y + int(rh * 0.26),
             int(W * 0.24), int(rh * 0.60),
             ex["bad"], 22, bold=True, color="#B91C1C")

        # Arrow
        _txt(slide, int(W * 0.29), y + int(rh * 0.38),
             int(W * 0.06), int(rh * 0.24), "→", 22,
             bold=True, color=chapter_color, align=PP_ALIGN.CENTER)

        # Good
        _rect(slide, int(W * 0.36), y, int(W * 0.61), rh, "#F0FDF4")
        _rect(slide, int(W * 0.36), y, int(W * 0.61), int(H * 0.006), "#059669")
        _txt(slide, int(W * 0.37), y + int(rh * 0.06),
             int(W * 0.59), int(rh * 0.14),
             f"✅  TRIE 공식 적용 — {ex['gain']}", 12, bold=True, color="#065F46")
        _txt(slide, int(W * 0.37), y + int(rh * 0.23),
             int(W * 0.59), int(rh * 0.70),
             ex["good"], 11, color="#14532D")


def build_thirty_day_plan_slide(prs, chapter_color: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_light"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.85), int(H * 0.09),
         "30일 AI 온보딩 로드맵", 25, bold=True, color=D["text_h1"])

    stream = pil_week_plan(THIRTY_DAY_PLAN)
    if stream:
        _pil_img(slide, stream, int(W * 0.04), int(H * 0.17), int(W * 0.92), int(H * 0.71))

    _txt(slide, int(W * 0.06), int(H * 0.91), int(W * 0.88), int(H * 0.06),
         "처음부터 완벽하지 않아도 됩니다. 매주 1개씩 실행하면 30일 후 완전히 다른 업무 환경이 만들어집니다.",
         11, color=D["text_muted"], align=PP_ALIGN.CENTER)


def build_caution_slide(prs, chapter_color: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_light"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.85), int(H * 0.09),
         "AI 도입 시 반드시 알아야 할 4가지 주의사항", 25, bold=True, color=D["text_h1"])
    _rect(slide, int(W * 0.06), int(H * 0.14), int(W * 0.88), int(H * 0.001), D["border"])

    cw = int(W * 0.44)
    ch = int(H * 0.36)
    positions = [
        (int(W * 0.04), int(H * 0.17)),
        (int(W * 0.53), int(H * 0.17)),
        (int(W * 0.04), int(H * 0.58)),
        (int(W * 0.53), int(H * 0.58)),
    ]

    for i, caution in enumerate(CAUTION_ITEMS[:4]):
        x, y = positions[i]
        c = caution["color"]
        _rect(slide, x + int(W * 0.003), y + int(H * 0.005), cw, ch, D["border_light"])
        _rect(slide, x, y, cw, ch, D["bg_white"])
        _rect(slide, x, y, cw, int(H * 0.005), c)
        # Icon + title
        _txt(slide, x + int(cw * 0.04), y + int(ch * 0.08),
             int(cw * 0.15), int(ch * 0.20), caution["icon"], 20, color=D["text_h1"])
        _txt(slide, x + int(cw * 0.22), y + int(ch * 0.09),
             int(cw * 0.74), int(ch * 0.18), caution["title"], 14, bold=True, color=c)
        _rect(slide, x + int(cw * 0.04), y + int(ch * 0.30),
              int(cw * 0.92), int(H * 0.001), D["border_light"])
        # Risk
        _txt(slide, x + int(cw * 0.04), y + int(ch * 0.33),
             int(cw * 0.18), int(ch * 0.10), "위험", 9, bold=True, color="#991B1B")
        _txt(slide, x + int(cw * 0.04), y + int(ch * 0.43),
             int(cw * 0.92), int(ch * 0.22), caution["risk"], 11, color="#7F1D1D")
        # Action
        _txt(slide, x + int(cw * 0.04), y + int(ch * 0.65),
             int(cw * 0.18), int(ch * 0.10), "대응", 9, bold=True, color="#065F46")
        _txt(slide, x + int(cw * 0.04), y + int(ch * 0.75),
             int(cw * 0.92), int(ch * 0.22), caution["action"], 11, color="#064E3B")


def build_demo_detail_slide(prs, title: str, content: list, chapter_color: str,
                             prompt: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_dark"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.85), int(H * 0.09),
         title, 22, bold=True, color=D["text_white"])
    _rect(slide, int(W * 0.06), int(H * 0.14), int(W * 0.88), int(H * 0.001), "#1E3A5F")

    for i, item in enumerate(content[:4]):
        y = int(H * 0.17) + i * int(H * 0.13)
        _rect(slide, int(W * 0.06), y, int(W * 0.88), int(H * 0.11), "#0F2040")
        _rect(slide, int(W * 0.06), y, int(W * 0.005), int(H * 0.11), chapter_color)
        _txt(slide, int(W * 0.08), y + int(H * 0.02),
             int(W * 0.84), int(H * 0.09), item, 13, color="#CBD5E1")

    if prompt:
        _rect(slide, int(W * 0.06), int(H * 0.72), int(W * 0.88), int(H * 0.22), "#050D1A")
        _rect(slide, int(W * 0.06), int(H * 0.72), int(W * 0.005), int(H * 0.22), D["blue"])
        _txt(slide, int(W * 0.08), int(H * 0.73), int(W * 0.20), int(H * 0.05),
             "사용 프롬프트:", 11, bold=True, color=D["blue"])
        _txt(slide, int(W * 0.08), int(H * 0.79), int(W * 0.84), int(H * 0.14),
             prompt, 11, color="#94A3B8")


# ─────────────────────────────────────────────────────────────────
# CHAPTER AGENT DATA
# ─────────────────────────────────────────────────────────────────
CH1_AGENTS = [
    {"name": "배치 이메일 분류 에이전트", "job_role": "영업 / 마케팅",
     "description": "수백 건의 고객 이메일을 한 번에 입력 → 우선순위·카테고리 자동 분류",
     "tools_used": ["Claude API", "Python", "Gmail"], "output": "분류된 이메일 리스트", "difficulty": "초급"},
    {"name": "웹 검색 뉴스 리포터", "job_role": "상품기획",
     "description": "키워드 설정 → 매일 오전 관련 뉴스 자동 수집·요약 메일 발송",
     "tools_used": ["Claude API", "Perplexity", "Gmail SMTP"], "output": "일간 업계 뉴스 요약", "difficulty": "초급"},
    {"name": "데이터 파이프라인 에이전트", "job_role": "분석",
     "description": "CSV 수신 → 자동 분석 → 차트 생성 → 대시보드 업데이트까지 자율 실행",
     "tools_used": ["Claude Code", "pandas", "matplotlib", "Slack"], "output": "자동 분석 리포트", "difficulty": "고급"},
]
CH3_AGENTS = [
    {"name": "랜딩 페이지 배치 생성기", "job_role": "마케팅",
     "description": "제품 목록 CSV 업로드 → 각 제품별 HTML 랜딩 페이지 일괄 생성",
     "tools_used": ["Claude API", "Python", "Jinja2"], "output": "제품별 HTML 랜딩 페이지", "difficulty": "초급"},
    {"name": "A/B 카피 자동 테스트", "job_role": "마케팅",
     "description": "광고 소재 업로드 → 50개 변형 생성 → Meta/Google Ads API 자동 제출",
     "tools_used": ["Claude API", "Meta Ads API", "Google Ads API"], "output": "50개 카피 + 자동 게재", "difficulty": "중급"},
    {"name": "콘텐츠 품질 검수 에이전트", "job_role": "마케팅",
     "description": "AI 생성 콘텐츠 → 브랜드 가이드라인 체크 → 승인/거절/수정 제안",
     "tools_used": ["Claude API", "Custom Rules Engine"], "output": "품질 검수 리포트", "difficulty": "중급"},
]
CH5_AGENTS = [
    {"name": "학습 진도 트래커", "job_role": "전 직무",
     "description": "매일 AI 활용 일지 작성 → Notion 자동 기록 → 주간 진도 리포트",
     "tools_used": ["Claude API", "Notion API", "Gmail"], "output": "주간 AI 활용 리포트", "difficulty": "초급"},
    {"name": "프롬프트 품질 개선 에이전트", "job_role": "전 직무",
     "description": "작성한 프롬프트 → TRIE 공식 준수 체크 → 자동 개선안 3개 제안",
     "tools_used": ["Claude API"], "output": "개선된 프롬프트 3개", "difficulty": "초급"},
    {"name": "AI 도입 ROI 측정 에이전트", "job_role": "전 직무",
     "description": "업무 시간 로그 기록 → AI 절감 시간 계산 → 월간 효과 리포트",
     "tools_used": ["Claude API", "Google Sheets API", "Gmail"], "output": "월간 ROI 리포트", "difficulty": "중급"},
]


# ─────────────────────────────────────────────────────────────────
# PIL CHART — GPU Hardware Timeline
# ─────────────────────────────────────────────────────────────────

def pil_hardware_bar(data: list, w: int = 2400, h: int = 800) -> "io.BytesIO | None":
    """GPU tflops 진화 수평 로그 막대 차트."""
    if not PIL_AVAILABLE:
        return None
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (w, h), (10, 22, 48))
    draw = ImageDraw.Draw(img)
    font_sm, font_md, font_lg = _kfont(12), _kfont(15), _kfont(18)

    max_val = max(d["tflops"] for d in data)
    import math
    log_max = math.log10(max_val + 1)
    bar_area_w = int(w * 0.55)
    bar_area_x = int(w * 0.22)
    bar_h = int((h - 80) / len(data)) - 14
    pad_top = 50

    # Title
    _draw_center(draw, "GPU Computing Power 진화 (TFLOPS)", w // 2, 25, font_lg, (150, 200, 255))

    colors = ["#94A3B8", "#64748B", "#1D4ED8", "#1D4ED8", "#0D9488", "#6D28D9"]
    for i, d in enumerate(data):
        y = pad_top + i * (bar_h + 14)
        log_val = math.log10(d["tflops"] + 1)
        bar_w = int(bar_area_w * log_val / log_max)

        color = colors[i % len(colors)]
        r2, g2, b2 = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)

        # Year + Model label
        label = f"{d['year']}  {d['model']}"
        try:
            draw.text((10, y + bar_h // 2 - 8), label, font=font_md, fill=(200, 220, 255))
        except Exception:
            draw.text((10, y + bar_h // 2 - 8), label, fill=(200, 220, 255))

        # Bar
        draw.rectangle([bar_area_x, y, bar_area_x + bar_w, y + bar_h], fill=(r2, g2, b2))

        # TFLOPS value
        val_str = f"{d['tflops']:,.0f} TFLOPS"
        try:
            draw.text((bar_area_x + bar_w + 10, y + bar_h // 2 - 8), val_str, font=font_sm, fill=(200, 200, 200))
        except Exception:
            draw.text((bar_area_x + bar_w + 10, y + bar_h // 2 - 8), val_str, fill=(200, 200, 200))

        # Event text (right side)
        ev_x = int(w * 0.78)
        try:
            draw.text((ev_x, y + bar_h // 2 - 8), d["event"], font=font_sm, fill=(150, 200, 150))
        except Exception:
            draw.text((ev_x, y + bar_h // 2 - 8), d["event"], fill=(150, 200, 150))

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def pil_cpu_gpu_compare(w: int = 1600, h: int = 600) -> "io.BytesIO | None":
    """CPU vs GPU 핵심 비교 다이어그램."""
    if not PIL_AVAILABLE:
        return None
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (w, h), (248, 250, 252))
    draw = ImageDraw.Draw(img)
    font_sm, font_md, font_lg = _kfont(12), _kfont(16), _kfont(22)

    half = w // 2
    # CPU panel
    draw.rectangle([20, 20, half - 20, h - 20], fill=(15, 43, 76))
    _draw_center(draw, "CPU", half // 2, 55, font_lg, (255, 255, 255))
    _draw_center(draw, "직렬(순차) 처리", half // 2, 95, font_md, (150, 180, 255))

    # Draw 4 big cores
    core_colors = [(29, 78, 216)] * 4
    for i in range(4):
        cx = int(half * 0.2 + i * half * 0.18)
        cy = 160
        draw.rectangle([cx - 35, cy, cx + 35, cy + 80], fill=(29, 78, 216))
        _draw_center(draw, "CORE", cx, cy + 35, font_sm, (255, 255, 255))

    _draw_center(draw, "천재 수학자 4명", half // 2, 280, font_md, (200, 220, 255))
    _draw_center(draw, "복잡한 계산을", half // 2, 310, font_sm, (150, 170, 200))
    _draw_center(draw, "순서대로 처리", half // 2, 335, font_sm, (150, 170, 200))

    # Small CPU core count
    _draw_center(draw, "8~32개 코어", half // 2, h - 70, font_md, (100, 200, 100))
    _draw_center(draw, "복잡한 로직, 조건 분기에 강함", half // 2, h - 40, font_sm, (150, 170, 200))

    # GPU panel
    draw.rectangle([half + 20, 20, w - 20, h - 20], fill=(10, 43, 10))
    _draw_center(draw, "GPU", half + (w - half) // 2, 55, font_lg, (255, 255, 255))
    _draw_center(draw, "병렬(동시) 처리", half + (w - half) // 2, 95, font_md, (100, 255, 150))

    # Draw many small cores
    core_w, core_h_px = 28, 18
    cols_g, rows_g = 12, 6
    start_x = half + 30
    start_y = 130
    for row in range(rows_g):
        for col in range(cols_g):
            cx = start_x + col * (core_w + 4)
            cy = start_y + row * (core_h_px + 4)
            draw.rectangle([cx, cy, cx + core_w, cy + core_h_px], fill=(13, 107, 79))

    _draw_center(draw, "단순 계산원 수만 명", half + (w - half) // 2, 280, font_md, (150, 255, 180))
    _draw_center(draw, "행렬 곱셈을", half + (w - half) // 2, 310, font_sm, (150, 200, 150))
    _draw_center(draw, "동시에 처리", half + (w - half) // 2, 335, font_sm, (150, 200, 150))

    _draw_center(draw, "수천~수만 개 코어", half + (w - half) // 2, h - 70, font_md, (100, 255, 100))
    _draw_center(draw, "AI 학습·추론 — 행렬 연산에 압도적", half + (w - half) // 2, h - 40, font_sm, (150, 200, 150))

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def pil_training_vs_inference(w: int = 1800, h: int = 560) -> "io.BytesIO | None":
    """학습 vs 추론 비교 다이어그램."""
    if not PIL_AVAILABLE:
        return None
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (w, h), (248, 250, 252))
    draw = ImageDraw.Draw(img)
    font_sm, font_md, font_lg = _kfont(13), _kfont(17), _kfont(22)

    panels = [
        {"title": "사전학습", "subtitle": "Pre-training", "color": (10, 22, 48),
         "text_color": (200, 220, 255), "items": ["수조 개 문서", "수천 GPU", "수개월", "수억~수십억 달러"],
         "icon": "📚", "analogy": "대학 4년 공부"},
        {"title": "파인튜닝", "subtitle": "Fine-tuning", "color": (13, 43, 10),
         "text_color": (180, 255, 200), "items": ["전문 데이터", "수~수십 GPU", "수일~수주", "$1K~$100K"],
         "icon": "🎓", "analogy": "사내 교육 3개월"},
        {"title": "추론", "subtitle": "Inference", "color": (60, 10, 10),
         "text_color": (255, 200, 200), "items": ["가중치 파일 로드", "GPU 1~8개", "즉시 응답", "토큰당 $0.0001"],
         "icon": "⚡", "analogy": "실제 업무 수행"},
    ]

    panel_w = (w - 80) // 3
    for i, p in enumerate(panels):
        x = 20 + i * (panel_w + 20)
        draw.rectangle([x, 20, x + panel_w, h - 20], fill=p["color"])
        # Header
        _draw_center(draw, p["icon"] + "  " + p["title"], x + panel_w // 2, 55, font_lg, (255, 255, 255))
        _draw_center(draw, p["subtitle"], x + panel_w // 2, 90, font_md, p["text_color"])
        # Divider
        draw.rectangle([x + 20, 115, x + panel_w - 20, 117], fill=p["text_color"])
        # Items
        for j, item in enumerate(p["items"]):
            try:
                draw.text((x + 25, 135 + j * 38), "▸  " + item, font=font_md, fill=(220, 220, 220))
            except Exception:
                draw.text((x + 25, 135 + j * 38), "    " + item, fill=(220, 220, 220))
        # Analogy
        _draw_center(draw, p["analogy"], x + panel_w // 2, h - 55, font_sm, p["text_color"])

        # Arrow between panels
        if i < len(panels) - 1:
            ax = x + panel_w + 10
            ay = h // 2
            draw.polygon([(ax, ay - 15), (ax, ay + 15), (ax + 18, ay)], fill=(200, 200, 200))

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def pil_scaling_levels(data: list, w: int = 1800, h: int = 520) -> "io.BytesIO | None":
    """파라미터 규모별 능력 계단 차트."""
    if not PIL_AVAILABLE:
        return None
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (w, h), (10, 22, 48))
    draw = ImageDraw.Draw(img)
    font_sm, font_md, font_lg = _kfont(13), _kfont(16), _kfont(20)

    n = len(data)
    bar_w = (w - 80) // n
    max_h = h - 100

    for i, d in enumerate(data):
        x = 40 + i * bar_w
        bar_height = int(max_h * (i + 1) / n)
        y = h - 50 - bar_height
        color = d.get("color", "#1D4ED8")
        r2, g2, b2 = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)

        draw.rectangle([x + 10, y, x + bar_w - 10, h - 50], fill=(r2, g2, b2, 200))

        # Level badge
        _draw_center(draw, d["level"], x + bar_w // 2, y + 15, font_lg, (255, 255, 255))
        _draw_center(draw, d["params"], x + bar_w // 2, y + 45, font_md, (255, 255, 200))

        # Capability desc — wrap
        words = d["desc"].split(" ")
        lines = []
        cur = ""
        for word in words:
            if len(cur) + len(word) < 12:
                cur += word + " "
            else:
                if cur:
                    lines.append(cur.strip())
                cur = word + " "
        if cur:
            lines.append(cur.strip())

        for j, line in enumerate(lines[:4]):
            _draw_center(draw, line, x + bar_w // 2, y + 75 + j * 22, font_sm, (200, 220, 255))

        # X-axis label
        _draw_center(draw, d["params"], x + bar_w // 2, h - 35, font_sm, (180, 180, 180))

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def pil_multimodal_tokens(data: list, w: int = 1600, h: int = 500) -> "io.BytesIO | None":
    """멀티모달 토큰화 방식 4열 카드."""
    if not PIL_AVAILABLE:
        return None
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (w, h), (248, 250, 252))
    draw = ImageDraw.Draw(img)
    font_sm, font_md, font_lg = _kfont(13), _kfont(17), _kfont(28)

    card_w = (w - 80) // len(data)
    for i, d in enumerate(data):
        x = 20 + i * (card_w + 10)
        color = d.get("color", "#1D4ED8")
        r2, g2, b2 = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
        draw.rectangle([x, 20, x + card_w - 10, h - 20], fill=(r2, g2, b2))
        _draw_center(draw, d["icon"], x + card_w // 2, 60, font_lg, (255, 255, 255))
        _draw_center(draw, d["name"], x + card_w // 2, 105, font_md, (255, 255, 255))
        draw.rectangle([x + 15, 130, x + card_w - 25, 132], fill=(255, 255, 255, 128))
        for j, (key, field) in enumerate([("방식:", "token_method"), ("토큰수:", "token_count"), ("모델:", "examples")]):
            try:
                draw.text((x + 15, 145 + j * 60), key, font=font_sm, fill=(200, 240, 200))
                draw.text((x + 15, 165 + j * 60), d[field], font=font_sm, fill=(240, 240, 240))
            except Exception:
                draw.text((x + 15, 145 + j * 60), key, fill=(200, 240, 200))
                draw.text((x + 15, 165 + j * 60), d[field], fill=(240, 240, 240))

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


# ─────────────────────────────────────────────────────────────────
# PHASE VISUAL SLIDES (CH1 강화 버전)
# ─────────────────────────────────────────────────────────────────

def pil_phase_impact_bar(impacts: list, color_hex: str, w: int = 1400, h: int = 380) -> "io.BytesIO | None":
    """Phase impact before/after 수평 비교 차트."""
    if not PIL_AVAILABLE:
        return None
    from PIL import Image, ImageDraw
    img = Image.new("RGB", (w, h), (10, 22, 48))
    draw = ImageDraw.Draw(img)
    font_sm = _kfont(13)
    font_md = _kfont(17)

    r2, g2, b2 = int(color_hex[1:3], 16), int(color_hex[3:5], 16), int(color_hex[5:7], 16)
    row_h = (h - 40) // len(impacts)

    for i, imp in enumerate(impacts):
        y = 20 + i * row_h
        cx = w // 2

        # Task label
        try:
            draw.text((10, y + row_h // 2 - 10), imp["task"], font=font_md, fill=(200, 220, 255))
        except Exception:
            draw.text((10, y + row_h // 2 - 10), imp["task"], fill=(200, 220, 255))

        # Before box
        bx = cx - 10
        draw.rectangle([cx - 180, y + 8, bx, y + row_h - 8], fill=(80, 10, 10))
        _draw_center(draw, "전: " + imp["before"], cx - 90, y + row_h // 2, font_md, (255, 120, 120))

        # After box
        ax = cx + 10
        draw.rectangle([ax, y + 8, ax + 200, y + row_h - 8], fill=(10, 60, 10))
        _draw_center(draw, "후: " + imp["after"], ax + 100, y + row_h // 2, font_md, (100, 255, 150))

        # Arrow
        _draw_center(draw, "→", cx, y + row_h // 2, font_md, (r2, g2, b2))

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def build_phase_visual_slide(prs: Presentation, phase_idx: int, phases: list) -> None:
    """Phase 진화 단계 — 시각적 카드 레이아웃 (build_content_slide 대체)."""
    if phase_idx >= len(phases):
        return
    pd = phases[phase_idx]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    color = pd["color"]
    bg = pd.get("bg_color", "#0A1628")
    _rect(slide, 0, 0, W, H, bg)
    _accent_bar(slide, color)

    # Phase badge (left big)
    _rect(slide, int(W * 0.04), int(H * 0.06), int(W * 0.22), int(H * 0.88), "#0A1628")
    _rect(slide, int(W * 0.04), int(H * 0.06), int(W * 0.22), int(H * 0.007), color)
    _txt(slide, int(W * 0.04), int(H * 0.10), int(W * 0.22), int(H * 0.14),
         pd["icon"], 48, color=color, align=PP_ALIGN.CENTER)
    _txt(slide, int(W * 0.04), int(H * 0.24), int(W * 0.22), int(H * 0.10),
         pd["phase"], 20, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
    _txt(slide, int(W * 0.04), int(H * 0.33), int(W * 0.22), int(H * 0.08),
         pd["title"], 16, color=color, bold=True, align=PP_ALIGN.CENTER)
    _txt(slide, int(W * 0.04), int(H * 0.43), int(W * 0.22), int(H * 0.06),
         pd["period"], 13, color="#64748B", align=PP_ALIGN.CENTER)
    _rect(slide, int(W * 0.06), int(H * 0.51), int(W * 0.18), int(H * 0.001), "#1E3A5F")
    _txt(slide, int(W * 0.04), int(H * 0.53), int(W * 0.22), int(H * 0.18),
         pd["models"], 11, color="#94A3B8", align=PP_ALIGN.CENTER)

    # Headline stat (center top)
    _rect(slide, int(W * 0.28), int(H * 0.06), int(W * 0.34), int(H * 0.22), "#111827")
    _txt(slide, int(W * 0.28), int(H * 0.08), int(W * 0.34), int(H * 0.12),
         pd["headline_stat"], 44, color=color, bold=True, align=PP_ALIGN.CENTER)
    _txt(slide, int(W * 0.29), int(H * 0.19), int(W * 0.32), int(H * 0.07),
         pd["headline_label"], 12, color="#94A3B8", align=PP_ALIGN.CENTER)

    # Impact table
    _txt(slide, int(W * 0.28), int(H * 0.30), int(W * 0.34), int(H * 0.06),
         "업무별 변화", 13, color=color, bold=True)
    for i, imp in enumerate(pd.get("impacts", [])[:3]):
        y = int(H * (0.37 + i * 0.17))
        _rect(slide, int(W * 0.28), y, int(W * 0.34), int(H * 0.15), "#0A1628")
        _txt(slide, int(W * 0.29), y + int(H * 0.02), int(W * 0.32), int(H * 0.05),
             imp["task"], 12, color="#CBD5E1")
        _rect(slide, int(W * 0.29), y + int(H * 0.07), int(W * 0.07), int(H * 0.06), "#2D1010")
        _txt(slide, int(W * 0.30), y + int(H * 0.075), int(W * 0.06), int(H * 0.05),
             "전: " + imp["before"], 11, color="#F87171")
        _rect(slide, int(W * 0.37), y + int(H * 0.07), int(W * 0.07), int(H * 0.06), "#102D10")
        _txt(slide, int(W * 0.38), y + int(H * 0.075), int(W * 0.06), int(H * 0.05),
             "후: " + imp["after"], 11, color="#34D399")
        _txt(slide, int(W * 0.45), y + int(H * 0.075), int(W * 0.05), int(H * 0.06),
             "→", 14, color=color, bold=True)

    # B2C Examples (right panel)
    _rect(slide, int(W * 0.65), int(H * 0.06), int(W * 0.31), int(H * 0.88), "#0A1628")
    _rect(slide, int(W * 0.65), int(H * 0.06), int(W * 0.31), int(H * 0.004), color)
    _txt(slide, int(W * 0.66), int(H * 0.09), int(W * 0.29), int(H * 0.06),
         "B2C 실무 적용 사례", 14, color=color, bold=True)
    for i, ex in enumerate(pd.get("b2c_examples", [])[:2]):
        y = int(H * (0.17 + i * 0.37))
        _rect(slide, int(W * 0.66), y, int(W * 0.28), int(H * 0.33), "#111827")
        _rect(slide, int(W * 0.66), y, int(W * 0.006), int(H * 0.33), color)
        _txt(slide, int(W * 0.67), y + int(H * 0.03), int(W * 0.26), int(H * 0.06),
             ["영업 사례", "마케팅·기획 사례"][i], 11, color=color, bold=True)
        _txt(slide, int(W * 0.67), y + int(H * 0.10), int(W * 0.26), int(H * 0.20),
             ex, 12, color="#CBD5E1")

    # Who disrupted
    _rect(slide, int(W * 0.28), int(H * 0.88), int(W * 0.34), int(H * 0.08), "#1E1010")
    _txt(slide, int(W * 0.30), int(H * 0.895), int(W * 0.30), int(H * 0.055),
         "⚠️  영향받는 직무: " + pd.get("who_disrupted", ""), 12, color="#F87171")


# ─────────────────────────────────────────────────────────────────
# AI 2026 최전선 + 준비 전략 슬라이드
# ─────────────────────────────────────────────────────────────────

def build_ai_frontier_2026_slide(prs: Presentation, chapter_color: str = "#0F2B4C") -> None:
    """AI 2026 최전선 — MCP 포함 6대 트렌드, 2컬럼 카드 레이아웃."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_dark"])
    _accent_bar(slide, "#60A5FA")

    af = AI_FRONTIER_2026
    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.60), int(H * 0.06),
         "APPENDIX — AI 2026 최전선", 11, color="#60A5FA", bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.09), int(W * 0.75), int(H * 0.09),
         af["headline"], 24, color=D["text_white"], bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.17), int(W * 0.65), int(H * 0.06),
         af["subheading"], 14, color="#F87171")

    trends = af["trends"]
    # 2-column × 3-row grid
    cw = int(W * 0.44)
    ch = int(H * 0.21)
    positions = [
        (int(W * 0.03), int(H * 0.24)), (int(W * 0.52), int(H * 0.24)),
        (int(W * 0.03), int(H * 0.48)), (int(W * 0.52), int(H * 0.48)),
        (int(W * 0.03), int(H * 0.72)), (int(W * 0.52), int(H * 0.72)),
    ]
    for (x, y), tr in zip(positions, trends):
        c = tr["color"]
        r2, g2, b2 = int(c[1:3], 16), int(c[3:5], 16), int(c[5:7], 16)
        _rect(slide, x, y, cw, ch, "#0A1628")
        _rect(slide, x, y, cw, 4, c)

        # Left: icon + title + stat
        _txt(slide, x + 12, y + int(H * 0.02), int(W * 0.05), int(H * 0.08),
             tr["icon"], 22, color="#FFFFFF")
        _txt(slide, x + 12 + int(W * 0.05), y + int(H * 0.02), cw - int(W * 0.06) - 130, int(H * 0.06),
             tr["title"], 14, color="#FFFFFF", bold=True)
        _txt(slide, x + 12 + int(W * 0.05), y + int(H * 0.07), cw - int(W * 0.06) - 130, int(H * 0.05),
             tr["subtitle"], 10, color="#94A3B8")

        # Stat badge (right)
        _rect(slide, x + cw - 115, y + int(H * 0.02), 110, int(H * 0.09), c)
        _txt(slide, x + cw - 115, y + int(H * 0.02), 110, int(H * 0.055),
             tr["stat"], 16, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        _txt(slide, x + cw - 115, y + int(H * 0.065), 110, int(H * 0.04),
             tr["stat_label"][:12], 9, color="#FFFFFF", align=PP_ALIGN.CENTER)

        # Desc + impact
        _rect(slide, x, y + int(H * 0.12), cw, int(H * 0.001), "#1E3A5F")
        _txt(slide, x + 12, y + int(H * 0.13), cw - 24, int(H * 0.055),
             tr["impact"], 12, color=f"#{r2:02X}{g2:02X}{b2:02X}", bold=True)

    # Key message bar
    _rect(slide, int(W * 0.03), int(H * 0.94), int(W * 0.94), int(H * 0.05), "#152032")
    _txt(slide, int(W * 0.05), int(H * 0.952), int(W * 0.90), int(H * 0.035),
         f"💡  {af['urgency_note']}", 13, color="#FBBF24")


def build_ai_preparation_slide(prs: Presentation, chapter_color: str = "#0F2B4C") -> None:
    """우리가 준비해야 할 것 — 6개 영역 실행 가이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_white"])
    _accent_bar(slide, chapter_color)

    ap = AI_PREPARATION_2026
    _txt(slide, int(W * 0.06), int(H * 0.04), int(W * 0.60), int(H * 0.06),
         "APPENDIX — 실전 준비 전략", 11, color=chapter_color, bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.09), int(W * 0.80), int(H * 0.09),
         ap["headline"], 26, color=D["text_h1"], bold=True)

    # Subheading highlight
    _rect(slide, int(W * 0.06), int(H * 0.18), int(W * 0.88), int(H * 0.07), "#FEF3C7")
    _rect(slide, int(W * 0.06), int(H * 0.18), int(W * 0.008), int(H * 0.07), "#F59E0B")
    _txt(slide, int(W * 0.08), int(H * 0.195), int(W * 0.84), int(H * 0.05),
         ap["subheading"], 14, color="#92400E", bold=True)

    # 6 areas in 3-col × 2-row
    areas = ap["areas"]
    cw = int(W * 0.30)
    ch = int(H * 0.31)
    positions = [
        (int(W * 0.03), int(H * 0.27)), (int(W * 0.35), int(H * 0.27)), (int(W * 0.67), int(H * 0.27)),
        (int(W * 0.03), int(H * 0.60)), (int(W * 0.35), int(H * 0.60)), (int(W * 0.67), int(H * 0.60)),
    ]
    for (x, y), area in zip(positions, areas):
        c = area["color"]
        r2, g2, b2 = int(c[1:3], 16), int(c[3:5], 16), int(c[5:7], 16)
        _shadow_card(slide, x, y, cw, ch)
        _rect(slide, x, y, cw, 5, c)

        # Header row
        _rect(slide, x, y, cw, int(H * 0.06), f"#{r2:02X}{g2:02X}{b2:02X}")
        _txt(slide, x + 10, y + int(H * 0.01), int(W * 0.05), int(H * 0.045),
             area["icon"], 18, color="#FFFFFF")
        _txt(slide, x + int(W * 0.045), y + int(H * 0.01), cw - int(W * 0.05) - 60, int(H * 0.04),
             area["category"], 14, color="#FFFFFF", bold=True)

        # Urgency badge
        urgency_colors = {"즉시": "#DC2626", "1개월 내": "#D97706", "3개월 내": "#059669", "지금": "#DC2626"}
        uc = urgency_colors.get(area["urgency"], "#64748B")
        _rect(slide, x + cw - 65, y + int(H * 0.01), 60, int(H * 0.04), uc)
        _txt(slide, x + cw - 65, y + int(H * 0.01), 60, int(H * 0.04),
             area["urgency"], 11, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)

        # Actions (first 3)
        for j, action in enumerate(area["actions"][:3]):
            ay = y + int(H * 0.07) + j * int(H * 0.074)
            _rect(slide, x + 8, ay + int(H * 0.01), 4, int(H * 0.05), c)
            _txt(slide, x + 18, ay, cw - 28, int(H * 0.07),
                 action, 11, color=D["text_body"])

        # KPI badge
        _rect(slide, x + 8, y + ch - int(H * 0.055), cw - 16, int(H * 0.045), "#F8FAFC")
        _txt(slide, x + 10, y + ch - int(H * 0.055), cw - 20, int(H * 0.04),
             "📊  " + area["kpi"], 11, color=c, bold=True)

    # Action matrix (bottom bar)
    _rect(slide, int(W * 0.03), int(H * 0.93), int(W * 0.94), int(H * 0.065), "#EFF6FF")
    matrix = ap["action_matrix"]
    items = list(matrix.items())
    seg_w = int(W * 0.94) // len(items)
    for i, (period, actions) in enumerate(items):
        x = int(W * 0.03) + i * seg_w
        _txt(slide, x + 6, int(H * 0.934), seg_w - 12, int(H * 0.02),
             period, 10, color=chapter_color, bold=True)
        _txt(slide, x + 6, int(H * 0.955), seg_w - 12, int(H * 0.03),
             " · ".join(actions[:2]), 10, color=D["text_body"])


# ─────────────────────────────────────────────────────────────────
# TOKEN / MEMORY CHAPTER SLIDES
# ─────────────────────────────────────────────────────────────────

_TOKEN_COLOR = "#0F2B4C"   # chapter accent — deep navy-blue

def build_token_chapter_divider(prs: Presentation) -> None:
    """Appendix 챕터 구분 슬라이드 — AI 기술 심화."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_dark"])
    # Decorative accent blocks
    _rect(slide, int(W * 0.07), int(H * 0.42), int(W * 0.06), 6, "#60A5FA")
    _rect(slide, int(W * 0.14), int(H * 0.42), int(W * 0.04), 6, "#34D399")
    _rect(slide, int(W * 0.19), int(H * 0.42), int(W * 0.02), 6, "#F87171")

    _txt(slide, int(W * 0.07), int(H * 0.18), int(W * 0.60), int(H * 0.10),
         "APPENDIX — AI 기술 심화", 13, color=D["text_muted"], bold=True)
    _txt(slide, int(W * 0.07), int(H * 0.27), int(W * 0.75), int(H * 0.16),
         "AI는 어떻게 작동하는가", 50, color=D["text_white"], bold=True)
    _txt(slide, int(W * 0.07), int(H * 0.46), int(W * 0.80), int(H * 0.09),
         "GPU · Transformer · 토큰 · 메모리 · 스케일링 · 멀티모달", 22, color="#60A5FA")
    _txt(slide, int(W * 0.07), int(H * 0.58), int(W * 0.78), int(H * 0.09),
         "비전문가도 이해하는 AI 구동 원리 — 비용·속도·한계를 결정하는 핵심 개념", 15,
         color=D["text_muted"])

    # Slide count indicator
    topics = ["GPU·하드웨어", "Transformer", "학습 vs 추론", "스케일링 법칙",
              "토큰·컨텍스트", "메모리", "멀티모달", "2026 최전선"]
    for i, t in enumerate(topics):
        x = int(W * (0.07 + i * 0.117))
        _rect(slide, x, int(H * 0.72), int(W * 0.10), int(H * 0.08), "#152032")
        _txt(slide, x + 8, int(H * 0.745), int(W * 0.09), int(H * 0.05),
             t, 11, color="#60A5FA")


def build_gpu_intro_slide(prs: Presentation, chapter_color: str = _TOKEN_COLOR) -> None:
    """GPU 기초 — CPU vs GPU 비교, AI에 왜 GPU가 필요한가."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_dark"])
    _accent_bar(slide, "#60A5FA")

    _txt(slide, int(W * 0.06), int(H * 0.07), int(W * 0.60), int(H * 0.07),
         "APPENDIX — GPU 기초", 11, color="#60A5FA", bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.13), int(W * 0.85), int(H * 0.10),
         "왜 AI는 GPU가 필요한가?", 28, color=D["text_white"], bold=True)

    cg = GPU_DATA["cpu_vs_gpu"]

    # PIL comparison diagram
    if PIL_AVAILABLE:
        stream = pil_cpu_gpu_compare(w=1600, h=520)
        _pil_img(slide, stream, int(W * 0.05), int(H * 0.25), int(W * 0.58), int(H * 0.47))

    # Right side — key insight cards
    _rect(slide, int(W * 0.66), int(H * 0.25), int(W * 0.30), int(H * 0.20), "#0A1628")
    _txt(slide, int(W * 0.68), int(H * 0.27), int(W * 0.10), int(H * 0.07),
         "⚡", 24, color="#60A5FA")
    _txt(slide, int(W * 0.69), int(H * 0.27), int(W * 0.25), int(H * 0.06),
         "속도 차이", 13, color="#60A5FA", bold=True)
    _txt(slide, int(W * 0.68), int(H * 0.34), int(W * 0.27), int(H * 0.08),
         "GPU가 CPU 대비\n100~1,000배 빠름", 14, color=D["text_white"])

    _rect(slide, int(W * 0.66), int(H * 0.48), int(W * 0.30), int(H * 0.20), "#0A1628")
    _txt(slide, int(W * 0.68), int(H * 0.50), int(W * 0.10), int(H * 0.07),
         "🎯", 24, color="#34D399")
    _txt(slide, int(W * 0.69), int(H * 0.50), int(W * 0.25), int(H * 0.06),
         "핵심 이유", 13, color="#34D399", bold=True)
    _txt(slide, int(W * 0.68), int(H * 0.57), int(W * 0.27), int(H * 0.08),
         "행렬 곱셈을 수천 개 코어가\n동시에 처리", 14, color=D["text_white"])

    _rect(slide, int(W * 0.66), int(H * 0.71), int(W * 0.30), int(H * 0.17), "#0A1628")
    _txt(slide, int(W * 0.68), int(H * 0.73), int(W * 0.10), int(H * 0.07),
         "💬", 24, color="#F87171")
    _txt(slide, int(W * 0.69), int(H * 0.73), int(W * 0.25), int(H * 0.06),
         "비유로 이해", 13, color="#F87171", bold=True)
    _txt(slide, int(W * 0.68), int(H * 0.79), int(W * 0.27), int(H * 0.07),
         "ChatGPT 답변 1번\n≈ A4 수만 장 행렬 곱셈", 13, color=D["text_white"])


def build_ai_hardware_timeline_slide(prs: Presentation, chapter_color: str = _TOKEN_COLOR) -> None:
    """AI 하드웨어 진화 — GPU tflops 로그 스케일 타임라인."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_dark"])
    _accent_bar(slide, "#60A5FA")

    _txt(slide, int(W * 0.06), int(H * 0.07), int(W * 0.60), int(H * 0.07),
         "APPENDIX — AI 하드웨어 진화", 11, color="#60A5FA", bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.13), int(W * 0.85), int(H * 0.10),
         "GPU 성능: 2012~2025, 5,600배 향상", 26, color=D["text_white"], bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.22), int(W * 0.85), int(H * 0.06),
         "GTX 580 (1.6 TFLOPS) → NVIDIA B200 (9,000 TFLOPS) — 로그 스케일", 15,
         color=D["text_muted"])

    if PIL_AVAILABLE:
        stream = pil_hardware_bar(GPU_DATA["gpu_evolution"], w=2160, h=680)
        _pil_img(slide, stream, int(W * 0.04), int(H * 0.30), int(W * 0.92), int(H * 0.58))

    # NVIDIA market info bar
    nm = GPU_DATA["nvidia_market"]
    _rect(slide, int(W * 0.04), int(H * 0.89), int(W * 0.92), int(H * 0.08), "#0A1628")
    _txt(slide, int(W * 0.06), int(H * 0.905), int(W * 0.88), int(H * 0.055),
         f"NVIDIA 시장 점유율 {nm['share']}  |  H100 개당 ${GPU_DATA['training_compute']['gpt4']['cost']} — {nm['moat']}",
         13, color="#60A5FA")


def build_transformer_architecture_slide(prs: Presentation, chapter_color: str = _TOKEN_COLOR) -> None:
    """Transformer 구조 — Self-Attention, FFN, 레이어 스택."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_white"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.07), int(W * 0.60), int(H * 0.07),
         "APPENDIX — Transformer 구조", 11, color=chapter_color, bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.13), int(W * 0.85), int(H * 0.10),
         "Transformer: 모든 현대 AI의 기반 구조", 26, color=D["text_h1"], bold=True)

    td = TRANSFORMER_DATA
    _rect(slide, int(W * 0.06), int(H * 0.23), int(W * 0.88), int(H * 0.08), D["bg_subtle"])
    _txt(slide, int(W * 0.08), int(H * 0.245), int(W * 0.84), int(H * 0.055),
         f"📌  {td['overview']}   |   핵심: {td['key_idea']}", 14, color=D["text_h2"])

    # 4 component cards
    comps = td["components"]
    cw = int(W * 0.20)
    gap = int(W * 0.02)
    total_w = 4 * cw + 3 * gap
    start_x = (W - total_w) // 2

    for i, comp in enumerate(comps):
        x = start_x + i * (cw + gap)
        y = int(H * 0.35)
        ch = int(H * 0.38)
        color = comp["color"]
        r2, g2, b2 = int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
        _rect(slide, x, y, cw, ch, f"#{r2:02X}{g2:02X}{b2:02X}")
        _rect(slide, x, y, cw, 5, "#FFFFFF")
        _txt(slide, x, y + int(H * 0.03), cw, int(H * 0.10),
             comp["icon"], 30, color="#FFFFFF", align=PP_ALIGN.CENTER)
        _txt(slide, x + 8, y + int(H * 0.13), cw - 16, int(H * 0.07),
             comp["name"], 13, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        _txt(slide, x + 10, y + int(H * 0.21), cw - 20, int(H * 0.14),
             comp["desc"], 11, color="#E2E8F0")

        # Arrow between
        if i < len(comps) - 1:
            ax = x + cw + gap // 2
            ay = y + ch // 2
            _rect(slide, ax - 2, ay - 6, 12, 12, "#94A3B8")

    # Token generation note
    _rect(slide, int(W * 0.06), int(H * 0.77), int(W * 0.88), int(H * 0.10), "#E0F2FE")
    _txt(slide, int(W * 0.08), int(H * 0.79), int(W * 0.84), int(H * 0.07),
         f"🔄  {td['token_generation']}", 14, color="#0369A1")

    # Scale table
    _txt(slide, int(W * 0.06), int(H * 0.88), int(W * 0.88), int(H * 0.06),
         "규모 비교:  " + "   |   ".join(
             f"{r['model']} {r['params']} ({r['year']})" for r in td["scale_table"]
         ), 12, color=D["text_muted"])


def build_training_vs_inference_slide(prs: Presentation, chapter_color: str = _TOKEN_COLOR) -> None:
    """학습(Training) vs 추론(Inference) vs 파인튜닝."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_dark"])
    _accent_bar(slide, "#60A5FA")

    _txt(slide, int(W * 0.06), int(H * 0.07), int(W * 0.60), int(H * 0.07),
         "APPENDIX — 학습 vs 추론", 11, color="#60A5FA", bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.13), int(W * 0.85), int(H * 0.10),
         "AI 만들기 vs AI 쓰기 — 비용이 1,000배 다르다", 26, color=D["text_white"], bold=True)

    if PIL_AVAILABLE:
        stream = pil_training_vs_inference(w=2160, h=560)
        _pil_img(slide, stream, int(W * 0.04), int(H * 0.25), int(W * 0.92), int(H * 0.46))

    # RAG vs Fine-tune comparison
    rvf = TRAINING_INFERENCE_DATA["rag_vs_fine"]
    _rect(slide, int(W * 0.06), int(H * 0.73), int(W * 0.43), int(H * 0.18), "#0A1628")
    _txt(slide, int(W * 0.08), int(H * 0.75), int(W * 0.40), int(H * 0.06),
         "🔍  RAG (검색 증강 생성)", 14, color="#60A5FA", bold=True)
    _txt(slide, int(W * 0.08), int(H * 0.81), int(W * 0.39), int(H * 0.07),
         f"{rvf['rag']['desc']}  |  적합: {rvf['when_rag']}", 12, color=D["text_muted"])

    _rect(slide, int(W * 0.52), int(H * 0.73), int(W * 0.43), int(H * 0.18), "#0A1628")
    _txt(slide, int(W * 0.54), int(H * 0.75), int(W * 0.40), int(H * 0.06),
         "🧠  파인튜닝 (Fine-tuning)", 14, color="#34D399", bold=True)
    _txt(slide, int(W * 0.54), int(H * 0.81), int(W * 0.39), int(H * 0.07),
         f"{rvf['fine']['desc']}  |  적합: {rvf['when_fine']}", 12, color=D["text_muted"])


def build_scaling_laws_slide(prs: Presentation, chapter_color: str = _TOKEN_COLOR) -> None:
    """스케일링 법칙 — Chinchilla, 파라미터 규모별 능력, 효율 혁명."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_white"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.07), int(W * 0.60), int(H * 0.07),
         "APPENDIX — 스케일링 법칙", 11, color=chapter_color, bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.13), int(W * 0.85), int(H * 0.10),
         "크면 좋은가? — 스케일링 법칙의 진실", 26, color=D["text_h1"], bold=True)

    sl = SCALING_LAWS_DATA

    # Chinchilla law highlight
    _rect(slide, int(W * 0.06), int(H * 0.24), int(W * 0.88), int(H * 0.12), D["bg_subtle"])
    _rect(slide, int(W * 0.06), int(H * 0.24), int(W * 0.008), int(H * 0.12), chapter_color)
    _txt(slide, int(W * 0.08), int(H * 0.25), int(W * 0.84), int(H * 0.05),
         f"📏  Chinchilla 법칙 ({sl['chinchilla']['source']})", 13, color=chapter_color, bold=True)
    _txt(slide, int(W * 0.08), int(H * 0.30), int(W * 0.84), int(H * 0.05),
         f"{sl['chinchilla']['law']}   →   예: {sl['chinchilla']['example']}", 14, color=D["text_h2"])

    # Capability levels PIL chart
    if PIL_AVAILABLE:
        stream = pil_scaling_levels(sl["capability_levels"], w=1400, h=440)
        _pil_img(slide, stream, int(W * 0.06), int(H * 0.38), int(W * 0.55), int(H * 0.38))

    # Efficiency revolution (right side)
    _txt(slide, int(W * 0.64), int(H * 0.38), int(W * 0.32), int(H * 0.07),
         "효율 혁명 — 같은 성능, 더 작은 모델", 13, color=chapter_color, bold=True)
    for i, ev in enumerate(sl["efficiency_revolution"]):
        y = int(H * 0.46 + i * 0.12)
        _shadow_card(slide, int(W * 0.63), y, int(W * 0.32), int(H * 0.10))
        _txt(slide, int(W * 0.65), y + int(H * 0.015), int(W * 0.10), int(H * 0.04),
             ev["year"], 12, color=chapter_color, bold=True)
        _txt(slide, int(W * 0.65), y + int(H * 0.045), int(W * 0.29), int(H * 0.04),
             f"{ev['model']}  {ev['params']} — {ev['note']}", 12, color=D["text_body"])

    # Trend note
    _rect(slide, int(W * 0.06), int(H * 0.78), int(W * 0.88), int(H * 0.09), "#E0F2FE")
    _txt(slide, int(W * 0.08), int(H * 0.80), int(W * 0.84), int(H * 0.06),
         f"📈  {sl['trend']}   |   {sl['emergent']}", 13, color="#0369A1")


def build_vram_kvcache_slide(prs: Presentation, chapter_color: str = _TOKEN_COLOR) -> None:
    """VRAM과 KV 캐시 — GPU 메모리의 핵심 개념."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_dark"])
    _accent_bar(slide, "#60A5FA")

    _txt(slide, int(W * 0.06), int(H * 0.07), int(W * 0.60), int(H * 0.07),
         "APPENDIX — VRAM & KV 캐시", 11, color="#60A5FA", bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.13), int(W * 0.85), int(H * 0.10),
         "GPU 메모리 — AI 성능의 병목", 26, color=D["text_white"], bold=True)

    vk = VRAM_KVCACHE_DATA

    # VRAM overview
    _rect(slide, int(W * 0.06), int(H * 0.25), int(W * 0.41), int(H * 0.52), "#0A1628")
    _txt(slide, int(W * 0.08), int(H * 0.27), int(W * 0.38), int(H * 0.07),
         "💾  VRAM이란?", 16, color="#60A5FA", bold=True)
    _txt(slide, int(W * 0.08), int(H * 0.34), int(W * 0.38), int(H * 0.10),
         vk["vram_overview"]["what"], 13, color=D["text_white"])
    _txt(slide, int(W * 0.08), int(H * 0.42), int(W * 0.38), int(H * 0.06),
         "비유: " + vk["vram_overview"]["analogy"], 13, color=D["text_muted"])

    # VRAM requirements table
    _txt(slide, int(W * 0.08), int(H * 0.50), int(W * 0.38), int(H * 0.06),
         "모델별 VRAM 필요량", 12, color="#60A5FA", bold=True)
    for i, req in enumerate(vk["vram_requirements"]):
        y = int(H * 0.56 + i * 0.065)
        _rect(slide, int(W * 0.08), y, int(W * 0.38), int(H * 0.055), "#152032")
        _txt(slide, int(W * 0.09), y + int(H * 0.01), int(W * 0.22), int(H * 0.04),
             req["model"], 12, color=D["text_white"])
        _txt(slide, int(W * 0.26), y + int(H * 0.01), int(W * 0.09), int(H * 0.04),
             req["vram"], 12, color="#34D399", bold=True)
        _txt(slide, int(W * 0.32), y + int(H * 0.01), int(W * 0.13), int(H * 0.04),
             req["device"], 11, color=D["text_muted"])

    # KV Cache explanation
    kvc = vk["kv_cache"]
    _rect(slide, int(W * 0.51), int(H * 0.25), int(W * 0.44), int(H * 0.52), "#0A1628")
    _txt(slide, int(W * 0.53), int(H * 0.27), int(W * 0.41), int(H * 0.07),
         "⚡  KV 캐시 (Key-Value Cache)", 16, color="#34D399", bold=True)
    _txt(slide, int(W * 0.53), int(H * 0.34), int(W * 0.41), int(H * 0.08),
         kvc["what"], 13, color=D["text_white"])
    _txt(slide, int(W * 0.53), int(H * 0.42), int(W * 0.41), int(H * 0.07),
         "문제: " + kvc["problem"], 12, color="#F87171")
    _txt(slide, int(W * 0.53), int(H * 0.49), int(W * 0.41), int(H * 0.07),
         "해결: " + kvc["solution"], 12, color="#34D399")
    _txt(slide, int(W * 0.53), int(H * 0.56), int(W * 0.41), int(H * 0.06),
         "효과: " + kvc["benefit"], 13, color="#60A5FA", bold=True)
    _txt(slide, int(W * 0.53), int(H * 0.63), int(W * 0.41), int(H * 0.06),
         "비유: " + kvc["analogy"], 12, color=D["text_muted"])
    _txt(slide, int(W * 0.53), int(H * 0.70), int(W * 0.41), int(H * 0.06),
         "주의: " + kvc["tradeoff"], 11, color="#FBBF24")

    # Memory bandwidth note
    mb = vk["memory_bandwidth"]
    _rect(slide, int(W * 0.06), int(H * 0.80), int(W * 0.88), int(H * 0.10), "#152032")
    _txt(slide, int(W * 0.08), int(H * 0.82), int(W * 0.84), int(H * 0.07),
         f"🚀  메모리 대역폭: {mb['h100']}  |  병목: {mb['bottleneck']}  |  최적화: {mb['solution']}",
         13, color="#60A5FA")


def build_multimodal_slide(prs: Presentation, chapter_color: str = _TOKEN_COLOR) -> None:
    """멀티모달 AI — 텍스트·이미지·음성·동영상 토큰화."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_white"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.07), int(W * 0.60), int(H * 0.07),
         "APPENDIX — 멀티모달 AI", 11, color=chapter_color, bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.13), int(W * 0.85), int(H * 0.10),
         "AI는 이미지·음성·영상을 어떻게 이해하는가", 26, color=D["text_h1"], bold=True)

    md = MULTIMODAL_DATA

    # Unified space concept
    _rect(slide, int(W * 0.06), int(H * 0.24), int(W * 0.88), int(H * 0.08), D["bg_subtle"])
    _txt(slide, int(W * 0.08), int(H * 0.255), int(W * 0.84), int(H * 0.055),
         f"🌐  {md['unified_space']}", 14, color=D["text_h2"])

    # PIL modality cards
    if PIL_AVAILABLE:
        stream = pil_multimodal_tokens(md["modalities"], w=1800, h=460)
        _pil_img(slide, stream, int(W * 0.05), int(H * 0.34), int(W * 0.90), int(H * 0.38))

    # Vision pipeline (simplified)
    _txt(slide, int(W * 0.06), int(H * 0.74), int(W * 0.88), int(H * 0.06),
         "이미지 처리 흐름:  " + "  →  ".join(md["vision_pipeline"]), 13, color=D["text_muted"])

    # Business value
    _txt(slide, int(W * 0.06), int(H * 0.80), int(W * 0.60), int(H * 0.06),
         "실무 활용 사례", 13, color=chapter_color, bold=True)
    for i, bv in enumerate(md["business_value"]):
        col = i % 2
        row = i // 2
        x = int(W * (0.06 + col * 0.45))
        y = int(H * (0.86 + row * 0.07))
        _txt(slide, x, y, int(W * 0.43), int(H * 0.06),
             "✅  " + bv, 13, color=D["text_body"])


def build_what_is_token_slide(prs: Presentation, chapter_color: str = _TOKEN_COLOR) -> None:
    """토큰이란 무엇인가 — 정의·비유·단위 변환."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_white"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.07), int(W * 0.60), int(H * 0.07),
         "APPENDIX — 토큰의 이해", 11, color=chapter_color, bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.13), int(W * 0.75), int(H * 0.10),
         "토큰(Token)이란 무엇인가?", 28, color=D["text_h1"], bold=True)

    tm = TOKEN_MEMORY["what_is_token"]

    _card(slide, int(W * 0.06), int(H * 0.27), int(W * 0.56), int(H * 0.20), chapter_color)
    _txt(slide, int(W * 0.08), int(H * 0.29), int(W * 0.20), int(H * 0.06),
         "정의", 11, color=chapter_color, bold=True)
    _txt(slide, int(W * 0.08), int(H * 0.35), int(W * 0.50), int(H * 0.10),
         tm["definition"], 16, color=D["text_h2"])

    _card(slide, int(W * 0.06), int(H * 0.50), int(W * 0.56), int(H * 0.17), chapter_color)
    _txt(slide, int(W * 0.08), int(H * 0.52), int(W * 0.20), int(H * 0.06),
         "단위 변환", 11, color=chapter_color, bold=True)
    _txt(slide, int(W * 0.08), int(H * 0.57), int(W * 0.50), int(H * 0.09),
         tm["conversion"], 16, color=D["text_h2"])

    _card(slide, int(W * 0.06), int(H * 0.70), int(W * 0.56), int(H * 0.18), chapter_color)
    _txt(slide, int(W * 0.08), int(H * 0.72), int(W * 0.20), int(H * 0.06),
         "왜 중요한가?", 11, color=chapter_color, bold=True)
    _txt(slide, int(W * 0.08), int(H * 0.77), int(W * 0.50), int(H * 0.09),
         tm["why_matters"], 14, color=D["text_body"])

    _txt(slide, int(W * 0.66), int(H * 0.25), int(W * 0.28), int(H * 0.07),
         "실제 예시", 13, color=chapter_color, bold=True)
    for i, ex in enumerate(tm["examples"]):
        _card(slide, int(W * 0.65), int(H * 0.33 + i * 0.19), int(W * 0.29), int(H * 0.15), chapter_color)
        _txt(slide, int(W * 0.67), int(H * 0.36 + i * 0.19), int(W * 0.26), int(H * 0.09),
             ex, 13, color=D["text_body"])


def build_context_evolution_slide(prs: Presentation, chapter_color: str = _TOKEN_COLOR) -> None:
    """컨텍스트 윈도우 진화 — PIL 로그 스케일 막대 차트."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_dark"])
    _accent_bar(slide, "#60A5FA")

    _txt(slide, int(W * 0.06), int(H * 0.07), int(W * 0.60), int(H * 0.07),
         "APPENDIX — 컨텍스트 윈도우 진화", 11, color="#60A5FA", bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.13), int(W * 0.80), int(H * 0.10),
         "컨텍스트 윈도우: 2019 → 2025, 2,000배 확장", 26, color=D["text_white"], bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.22), int(W * 0.80), int(H * 0.07),
         "GPT-2 1,024 토큰(A4 1.5장) → Gemini 2M 토큰(소설 30권)", 15, color=D["text_muted"])

    if PIL_AVAILABLE:
        stream = pil_context_bar(CONTEXT_EVOLUTION, w=2160, h=760)
        _pil_img(slide, stream, int(W * 0.04), int(H * 0.31), int(W * 0.92), int(H * 0.62))


def build_memory_types_slide(prs: Presentation, chapter_color: str = _TOKEN_COLOR) -> None:
    """AI 메모리 4가지 유형 — 카드 그리드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_white"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.07), int(W * 0.60), int(H * 0.07),
         "APPENDIX — AI 메모리 유형", 11, color=chapter_color, bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.13), int(W * 0.85), int(H * 0.10),
         "AI는 어떻게 '기억'하는가 — 4가지 메모리 유형", 26, color=D["text_h1"], bold=True)

    types = TOKEN_MEMORY["memory_types"]
    cols = 4
    cw = int(W * 0.20)
    gap = int(W * 0.025)
    total_w = cols * cw + (cols - 1) * gap
    start_x = (W - total_w) // 2

    for i, mt in enumerate(types):
        x = start_x + i * (cw + gap)
        y = int(H * 0.27)
        ch = int(H * 0.62)
        try:
            r2 = int(mt["color"][1:3], 16)
            g2 = int(mt["color"][3:5], 16)
            b2 = int(mt["color"][5:7], 16)
        except Exception:
            r2, g2, b2 = 15, 43, 76
        _rect(slide, x, y, cw, ch, f"#{r2:02X}{g2:02X}{b2:02X}")
        _rect(slide, x, y, cw, 6, "#FFFFFF")

        _txt(slide, x, y + int(H * 0.03), cw, int(H * 0.10),
             mt["icon"], 32, color="#FFFFFF", align=PP_ALIGN.CENTER)
        _txt(slide, x, y + int(H * 0.13), cw, int(H * 0.08),
             mt["name"], 14, color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        _txt(slide, x + 12, y + int(H * 0.22), cw - 24, int(H * 0.12),
             mt["description"], 12, color=D["text_light"])
        _txt(slide, x + 12, y + int(H * 0.36), cw - 24, int(H * 0.06),
             "한계:", 10, color="#FFFFFF", bold=True)
        _txt(slide, x + 12, y + int(H * 0.41), cw - 24, int(H * 0.10),
             mt["limit"], 11, color=D["text_light"])
        _txt(slide, x + 12, y + int(H * 0.52), cw - 24, int(H * 0.06),
             "비유:", 10, color="#FFFFFF", bold=True)
        _txt(slide, x + 12, y + int(H * 0.57), cw - 24, int(H * 0.09),
             mt["analogy"], 11, color=D["text_light"])


def build_token_cost_trend_slide(prs: Presentation, chapter_color: str = _TOKEN_COLOR) -> None:
    """토큰 API 비용 98% 감소 추세 — PIL 막대 차트."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_dark"])
    _accent_bar(slide, "#60A5FA")

    _txt(slide, int(W * 0.06), int(H * 0.07), int(W * 0.60), int(H * 0.07),
         "APPENDIX — 토큰 비용 트렌드", 11, color="#60A5FA", bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.13), int(W * 0.80), int(H * 0.10),
         "AI API 비용: 2년 만에 98% 감소", 28, color=D["text_white"], bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.22), int(W * 0.85), int(H * 0.07),
         "GPT-4 출시 당시 $30/1M토큰 → 2025년 Claude 3.7 $0.6/1M토큰", 15, color=D["text_muted"])

    if PIL_AVAILABLE:
        stream = pil_cost_bar(TOKEN_MEMORY["cost_trend"], w=1800, h=620)
        _pil_img(slide, stream, int(W * 0.08), int(H * 0.30), int(W * 0.84), int(H * 0.53))

    _rect(slide, int(W * 0.06), int(H * 0.85), int(W * 0.88), int(H * 0.10), D["navy"])
    _txt(slide, int(W * 0.08), int(H * 0.86), int(W * 0.84), int(H * 0.08),
         "💡  비용 장벽 소멸 → 소규모 팀도 GPT-4급 AI를 월 수천원으로 활용 가능", 15, color="#60A5FA")


def build_context_engineering_slide(prs: Presentation, chapter_color: str = _TOKEN_COLOR) -> None:
    """컨텍스트 엔지니어링 — 지금 가장 중요한 기술."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_white"])
    _accent_bar(slide, chapter_color)

    ce = TOKEN_MEMORY["context_engineering"]
    _txt(slide, int(W * 0.06), int(H * 0.07), int(W * 0.60), int(H * 0.07),
         "APPENDIX — 컨텍스트 엔지니어링", 11, color=chapter_color, bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.13), int(W * 0.80), int(H * 0.10),
         "컨텍스트 엔지니어링", 30, color=D["text_h1"], bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.23), int(W * 0.85), int(H * 0.08),
         ce["definition"], 16, color=D["text_muted"])

    _rect(slide, int(W * 0.06), int(H * 0.33), int(W * 0.88), int(H * 0.10), D["bg_subtle"])
    _txt(slide, int(W * 0.08), int(H * 0.34), int(W * 0.84), int(H * 0.08),
         f"핵심: {ce['why_matters']}", 16, color=D["text_h2"], bold=True)

    for i, t in enumerate(ce["techniques"]):
        col = i % 2
        row = i // 2
        x = int(W * (0.06 + col * 0.46))
        y = int(H * (0.47 + row * 0.22))
        cw = int(W * 0.42)
        ch = int(H * 0.19)
        _shadow_card(slide, x, y, cw, ch)
        _rect(slide, x, y, 5, ch, chapter_color)
        _txt(slide, x + 18, y + int(H * 0.02), cw - 24, int(H * 0.07),
             t["icon"] + "  " + t["name"], 15, color=D["text_h1"], bold=True)
        _txt(slide, x + 18, y + int(H * 0.09), cw - 30, int(H * 0.08),
             t["desc"], 13, color=D["text_body"])


def build_new_architectures_slide(prs: Presentation, chapter_color: str = _TOKEN_COLOR) -> None:
    """새로운 AI 아키텍처 — KV Cache, Mamba, Flash Attention, MoE."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_dark"])
    _accent_bar(slide, "#60A5FA")

    _txt(slide, int(W * 0.06), int(H * 0.07), int(W * 0.60), int(H * 0.07),
         "APPENDIX — 새로운 아키텍처", 11, color="#60A5FA", bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.13), int(W * 0.85), int(H * 0.10),
         "Transformer를 넘어서 — 2024~2026 핵심 기술", 26, color=D["text_white"], bold=True)

    archs = TOKEN_MEMORY["new_architectures"]
    cw = int(W * 0.43)
    ch = int(H * 0.30)
    positions = [
        (int(W * 0.05), int(H * 0.27)),
        (int(W * 0.52), int(H * 0.27)),
        (int(W * 0.05), int(H * 0.62)),
        (int(W * 0.52), int(H * 0.62)),
    ]
    for (x, y), arch in zip(positions, archs):
        try:
            r2 = int(arch["color"][1:3], 16)
            g2 = int(arch["color"][3:5], 16)
            b2 = int(arch["color"][5:7], 16)
        except Exception:
            r2, g2, b2 = 29, 78, 216
        _rect(slide, x, y, cw, ch, "#152032")
        _rect(slide, x, y, cw, 4, arch["color"])
        _txt(slide, x + 16, y + int(H * 0.03), cw - 24, int(H * 0.08),
             arch["icon"] + "  " + arch["name"], 16, color="#FFFFFF", bold=True)
        _txt(slide, x + 16, y + int(H * 0.11), cw - 32, int(H * 0.10),
             arch["desc"], 12, color=D["text_muted"])
        _txt(slide, x + 16, y + int(H * 0.20), cw - 24, int(H * 0.07),
             "임팩트: " + arch["impact"], 12, color=f"#{r2:02X}{g2:02X}{b2:02X}", bold=True)


def build_token_summary_slide(prs: Presentation, chapter_color: str = _TOKEN_COLOR) -> None:
    """토큰·컨텍스트·메모리 핵심 요약."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, D["bg_white"])
    _accent_bar(slide, chapter_color)

    _txt(slide, int(W * 0.06), int(H * 0.07), int(W * 0.60), int(H * 0.07),
         "APPENDIX — 핵심 요약", 11, color=chapter_color, bold=True)
    _txt(slide, int(W * 0.06), int(H * 0.13), int(W * 0.80), int(H * 0.10),
         "실무자가 알아야 할 3가지", 30, color=D["text_h1"], bold=True)

    points = [
        ("💡", "컨텍스트 = AI의 작업 공간",
         "컨텍스트가 클수록 더 많은 문서·대화·코드를 한 번에 처리. Gemini 1.5는 책 30권 분량을 한 번에 읽는다."),
        ("💰", "비용은 이미 무시할 수준",
         "2023년 대비 98% 감소. GPT-4급 AI를 하루 1,000번 사용해도 월 3천원 미만. 도구 선택이 아닌 활용 방법이 핵심."),
        ("🏗", "컨텍스트 엔지니어링이 새 핵심 역량",
         "무엇을 어떻게 AI에게 제공하느냐가 결과를 결정. RAG·청킹·요약 전략이 프롬프팅만큼 중요해졌다."),
    ]
    for i, (icon, title, desc) in enumerate(points):
        y = int(H * (0.30 + i * 0.22))
        _shadow_card(slide, int(W * 0.06), y, int(W * 0.88), int(H * 0.19))
        _rect(slide, int(W * 0.06), y, int(W * 0.008), int(H * 0.19), chapter_color)
        _txt(slide, int(W * 0.08), y + int(H * 0.03), int(W * 0.06), int(H * 0.13),
             icon, 28, color=chapter_color)
        _txt(slide, int(W * 0.15), y + int(H * 0.03), int(W * 0.70), int(H * 0.07),
             title, 17, color=D["text_h1"], bold=True)
        _txt(slide, int(W * 0.15), y + int(H * 0.10), int(W * 0.76), int(H * 0.09),
             desc, 13, color=D["text_body"])


# ─────────────────────────────────────────────────────────────────
# MAIN BUILDER
# ─────────────────────────────────────────────────────────────────
def build_presentation(research_data: dict | None = None) -> str:
    prs = Presentation()
    prs.slide_width = Emu(W)
    prs.slide_height = Emu(H)

    total = 0
    def _s(name=""):
        nonlocal total; total += 1
        print(f"  [{total:02d}] {name}")

    print("📊 슬라이드 생성 시작 (Professional Edition)...")

    build_title_slide(prs);                                           _s("타이틀")
    build_toc_slide(prs);                                             _s("목차")

    # CH1
    build_chapter_divider(prs, CHAPTERS[0]);                          _s("CH1 구분")
    build_ai_definition_slide(prs, CHAPTERS[0]["color"]);             _s("AI 정의 — 비유 3가지")
    build_stats_slide(prs, "chatgpt_adoption", CHAPTERS[0]["color"]); _s("ChatGPT 1억 명 통계")
    build_timeline_slide(prs);                                        _s("AI 타임라인")
    build_phase_visual_slide(prs, 0, PHASE_DATA);                     _s("Phase 1 — 텍스트 Q&A")
    build_phase_visual_slide(prs, 1, PHASE_DATA);                     _s("Phase 2 — 멀티모달")
    build_phase_visual_slide(prs, 2, PHASE_DATA);                     _s("Phase 3 — AI 에이전트")
    build_phase_visual_slide(prs, 3, PHASE_DATA);                     _s("Phase 4 — CLI 에이전트")
    build_economic_impact_slide(prs, CHAPTERS[0]["color"]);           _s("McKinsey 경제 임팩트")
    build_agent_example_slide(prs, CH1_AGENTS, "CH1 — AI 진화 단계별 에이전트 사례"); _s("CH1 에이전트")

    # CH2
    build_chapter_divider(prs, CHAPTERS[1]);                          _s("CH2 구분")
    for job in ["영업", "마케팅", "상품기획"]:
        build_before_after_slide(prs, job, WORKFLOW_CHANGES[job], CHAPTERS[1]["color"]); _s(f"{job} Before/After")
        build_job_stats_slide(prs, job, CHAPTERS[1]["color"]);        _s(f"{job} 통계")
        build_job_agent_workflow_slide(prs, job, CHAPTERS[1]["color"]); _s(f"{job} 에이전트 플로우")
        build_job_agent_example_slide(prs, job, CHAPTERS[1]["color"]); _s(f"{job} 에이전트 예시")
    build_common_patterns_slide(prs, CHAPTERS[1]["color"]);           _s("공통 자동화 패턴")

    # CH3
    build_chapter_divider(prs, CHAPTERS[2]);                          _s("CH3 구분")
    build_content_slide(prs, "실전 데모 목록", CHAPTERS[2]["slides"], CHAPTERS[2]["color"]); _s("데모 목록")
    build_demo_detail_slide(prs, "데모 1: 제품 소개 랜딩 페이지 30초 만들기", [
        "✅ 프롬프트 하나로 완성된 HTML 랜딩 페이지 즉시 생성",
        "✅ 히어로 · 특징 3가지 · Before-After · CTA 자동 포함",
        "✅ '색상을 파란색으로 바꿔줘' → 자연어로 즉시 수정",
        "✅ 결과: 기존 2~3일 → 30초 (360배 빠름, 비용 $0)",
    ], CHAPTERS[2]["color"],
    prompt='"[제품명] 소개 랜딩 페이지를 HTML로 만들어줘. 히어로 섹션, 특징 3가지, Before/After 비교, CTA 포함. 다크 배경에 오렌지 포인트 컬러, 모바일 반응형."'); _s("데모1")
    build_demo_detail_slide(prs, "데모 2: AI 이미지 생성 — 프롬프트 품질 비교", [
        "❌ 나쁜 프롬프트: 'SSD 사진' → 평범한 스톡 이미지 수준",
        "✅ 좋은 프롬프트: 'Modern NVMe SSD floating, electric blue particles, 8K, studio lighting, premium tech aesthetic' → 광고 수준",
        "✅ 핵심: 스타일·조명·분위기 키워드가 결과를 10배 바꾼다",
        "✅ 도구별: DALL-E 3 (한국어 OK), Midjourney (최고품질), Firefly (상업용 안전)",
    ], CHAPTERS[2]["color"]);                                         _s("데모2")
    build_agent_example_slide(prs, CH3_AGENTS, "CH3 — 데모 자동화 에이전트 사례"); _s("CH3 에이전트")

    # CH4
    build_chapter_divider(prs, CHAPTERS[3]);                          _s("CH4 구분")
    build_agent_definition_slide(prs, CHAPTERS[3]["color"]);          _s("에이전트 정의")
    build_agent_architecture_slide(prs, CHAPTERS[3]["color"]);        _s("에이전트 아키텍처")
    build_no_code_agents_slide(prs, CHAPTERS[3]["color"]);            _s("노코드 도구")
    build_agent_example_slide(prs, [
        JOB_STATS["영업"]["agent_examples"][0],
        JOB_STATS["마케팅"]["agent_examples"][0],
        JOB_STATS["상품기획"]["agent_examples"][0],
    ], "CH4 — 직무별 에이전트 예시");                                  _s("에이전트 예시")
    build_first_agent_guide_slide(prs, CHAPTERS[3]["color"]);         _s("첫 에이전트 가이드")

    # CH5
    build_chapter_divider(prs, CHAPTERS[4]);                          _s("CH5 구분")
    build_tool_matrix_slide(prs, CHAPTERS[4]["color"]);               _s("도구 매트릭스")
    build_trie_formula_slide(prs, CHAPTERS[4]["color"]);              _s("TRIE 공식")
    build_trie_example_slide(prs, CHAPTERS[4]["color"]);              _s("TRIE 예시 비교")
    build_thirty_day_plan_slide(prs, CHAPTERS[4]["color"]);           _s("30일 플랜")
    build_caution_slide(prs, CHAPTERS[4]["color"]);                   _s("주의사항")
    build_agent_example_slide(prs, CH5_AGENTS, "CH5 — 로드맵 실행 지원 에이전트"); _s("CH5 에이전트")

    build_closing_slide(prs);                                         _s("마무리 Q&A")

    # APPENDIX: AI 기술 심화 (16슬라이드)
    build_token_chapter_divider(prs);                                 _s("APPENDIX 구분 — AI 기술 심화")
    # Part 1: GPU·하드웨어
    build_gpu_intro_slide(prs);                                       _s("GPU 기초 — CPU vs GPU")
    build_ai_hardware_timeline_slide(prs);                            _s("AI 하드웨어 진화")
    # Part 2: AI 모델 구조
    build_transformer_architecture_slide(prs);                        _s("Transformer 구조")
    build_training_vs_inference_slide(prs);                           _s("학습 vs 추론")
    build_scaling_laws_slide(prs);                                    _s("스케일링 법칙")
    # Part 3: 토큰·컨텍스트·메모리
    build_what_is_token_slide(prs);                                   _s("토큰이란?")
    build_context_evolution_slide(prs);                               _s("컨텍스트 윈도우 진화")
    build_vram_kvcache_slide(prs);                                    _s("VRAM & KV 캐시")
    build_memory_types_slide(prs);                                    _s("메모리 4가지 유형")
    # Part 4: 비용·기술 트렌드
    build_token_cost_trend_slide(prs);                                _s("토큰 비용 98% 감소")
    build_context_engineering_slide(prs);                             _s("컨텍스트 엔지니어링")
    build_new_architectures_slide(prs);                               _s("새로운 아키텍처")
    build_multimodal_slide(prs);                                      _s("멀티모달 AI")
    build_ai_frontier_2026_slide(prs);                                _s("AI 2026 최전선 — 6대 트렌드")
    build_ai_preparation_slide(prs);                                  _s("우리가 준비해야 할 것")
    # Summary
    build_token_summary_slide(prs);                                   _s("AI 기술 핵심 요약")

    ts = datetime.now().strftime("%Y%m%d_%H%M")
    out = OUTPUT_DIR / f"AI_Workflow_Lecture_{ts}.pptx"
    prs.save(str(out))
    print(f"\n✅ PPT 저장 완료 ({total}슬라이드) → {out}")
    return str(out)


if __name__ == "__main__":
    cp = DATA_DIR / "research_cache.json"
    rd = json.load(open(cp, encoding="utf-8")) if cp.exists() else None
    build_presentation(rd)
