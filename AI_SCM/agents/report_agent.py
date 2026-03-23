"""
Report / Visualization Agent
- Generates HTML dashboard with dark theme
- Generates PPT presentation
- Sections: Value Chain, Token->HW Demand, Bottleneck, CapEx, Signals, Forecast
"""

import json
import os
import sys
import datetime

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
import config

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("  [ReportAgent] WARNING: python-pptx not available, skipping PPT generation")


# ============================================================
# COLOR SCHEME
# ============================================================
COLORS = {
    "bg_dark": "#0a1628",
    "bg_card": "#0f1f3d",
    "bg_card2": "#152848",
    "accent_blue": "#1e88e5",
    "accent_cyan": "#00bcd4",
    "accent_green": "#00e676",
    "text_primary": "#e0e6ef",
    "text_secondary": "#8fa3c0",
    "critical": "#f44336",
    "high": "#ff9800",
    "medium": "#ffeb3b",
    "low": "#4caf50",
    "border": "#1e3a5f",
    "table_header": "#0d47a1",
}

SEVERITY_COLORS = {
    "critical": COLORS["critical"],
    "high": COLORS["high"],
    "medium": COLORS["medium"],
    "low": COLORS["low"],
}

SIGNAL_COLORS = {
    "Strong Buy": "#00e676",
    "Buy": "#4caf50",
    "Hold": "#ffeb3b",
    "Watch": "#ff9800",
    "Sell": "#f44336",
}


# ============================================================
# HTML GENERATION
# ============================================================

def make_ascii_bar(value, max_value, width=30, fill_char="█", empty_char="░"):
    """Create ASCII progress bar."""
    if max_value == 0:
        return empty_char * width
    filled = int((value / max_value) * width)
    return fill_char * filled + empty_char * (width - filled)


def generate_html_report(data_results, mapping_results, modeling_results,
                          bottleneck_results, strategy_results, date_str):
    """Generate full HTML dashboard."""

    # ---- Value Chain section ----
    value_chain_html = _build_value_chain_section(bottleneck_results)

    # ---- Token -> Hardware demand table ----
    demand_table_html = _build_demand_table(modeling_results)

    # ---- Bottleneck analysis cards ----
    bottleneck_html = _build_bottleneck_section(bottleneck_results)

    # ---- CapEx trends ----
    capex_html = _build_capex_section(data_results)

    # ---- Investment signals table ----
    signals_html = _build_signals_section(strategy_results)

    # ---- Scenario forecast ----
    forecast_html = _build_forecast_section(modeling_results)

    # ---- Investment Network Graph (D3.js) ----
    network_graph_html = _build_network_graph_section(mapping_results)

    html = f"""<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>AI Supply Chain Intelligence Dashboard - {date_str}</title>
<style>
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{
    background: {COLORS['bg_dark']};
    color: {COLORS['text_primary']};
    font-family: 'Segoe UI', -apple-system, BlinkMacSystemFont, sans-serif;
    font-size: 14px;
    line-height: 1.6;
  }}
  .header {{
    background: linear-gradient(135deg, #0d1b35 0%, #1a2f5a 50%, #0d1b35 100%);
    border-bottom: 2px solid {COLORS['accent_blue']};
    padding: 24px 40px;
    display: flex;
    align-items: center;
    justify-content: space-between;
  }}
  .header h1 {{
    font-size: 24px;
    font-weight: 700;
    color: #ffffff;
    letter-spacing: 0.5px;
  }}
  .header .subtitle {{
    color: {COLORS['accent_cyan']};
    font-size: 13px;
    margin-top: 4px;
  }}
  .header .date-badge {{
    background: {COLORS['accent_blue']};
    color: white;
    padding: 6px 16px;
    border-radius: 20px;
    font-size: 12px;
    font-weight: 600;
  }}
  .container {{
    max-width: 1600px;
    margin: 0 auto;
    padding: 24px 40px;
  }}
  .section {{
    margin-bottom: 36px;
  }}
  .section-title {{
    font-size: 18px;
    font-weight: 700;
    color: {COLORS['accent_cyan']};
    border-left: 4px solid {COLORS['accent_blue']};
    padding-left: 12px;
    margin-bottom: 16px;
    letter-spacing: 0.3px;
  }}
  .card {{
    background: {COLORS['bg_card']};
    border: 1px solid {COLORS['border']};
    border-radius: 8px;
    padding: 20px;
    margin-bottom: 12px;
  }}
  .card-grid {{
    display: grid;
    grid-template-columns: repeat(auto-fill, minmax(280px, 1fr));
    gap: 16px;
  }}
  .metric-card {{
    background: {COLORS['bg_card']};
    border: 1px solid {COLORS['border']};
    border-radius: 8px;
    padding: 16px 20px;
  }}
  .metric-card .label {{
    color: {COLORS['text_secondary']};
    font-size: 12px;
    text-transform: uppercase;
    letter-spacing: 0.8px;
    margin-bottom: 8px;
  }}
  .metric-card .value {{
    font-size: 26px;
    font-weight: 700;
    color: {COLORS['accent_cyan']};
  }}
  .metric-card .sub {{
    color: {COLORS['text_secondary']};
    font-size: 12px;
    margin-top: 4px;
  }}
  table {{
    width: 100%;
    border-collapse: collapse;
  }}
  th {{
    background: {COLORS['table_header']};
    color: white;
    padding: 10px 14px;
    text-align: left;
    font-size: 12px;
    text-transform: uppercase;
    letter-spacing: 0.8px;
  }}
  td {{
    padding: 10px 14px;
    border-bottom: 1px solid {COLORS['border']};
    font-size: 13px;
  }}
  tr:hover td {{
    background: rgba(30, 136, 229, 0.08);
  }}
  .badge {{
    display: inline-block;
    padding: 3px 10px;
    border-radius: 12px;
    font-size: 11px;
    font-weight: 700;
    text-transform: uppercase;
    letter-spacing: 0.5px;
  }}
  .badge-critical {{ background: #f44336; color: white; }}
  .badge-high {{ background: #ff9800; color: white; }}
  .badge-medium {{ background: #ffeb3b; color: #000; }}
  .badge-low {{ background: #4caf50; color: white; }}
  .badge-sb {{ background: #00e676; color: #000; }}
  .badge-buy {{ background: #4caf50; color: white; }}
  .badge-hold {{ background: #ffeb3b; color: #000; }}
  .badge-watch {{ background: #ff9800; color: white; }}
  .chain-grid {{
    display: grid;
    grid-template-columns: repeat(auto-fill, minmax(200px, 1fr));
    gap: 10px;
  }}
  .chain-card {{
    border-radius: 6px;
    padding: 12px 16px;
    border-left: 4px solid;
  }}
  .chain-card .layer-name {{
    font-size: 12px;
    font-weight: 700;
    text-transform: uppercase;
    letter-spacing: 0.5px;
    margin-bottom: 4px;
  }}
  .chain-card .company-list {{
    font-size: 11px;
    opacity: 0.8;
  }}
  .progress-bar-wrap {{
    display: flex;
    align-items: center;
    gap: 8px;
  }}
  .progress-bar-bg {{
    flex: 1;
    background: rgba(255,255,255,0.1);
    border-radius: 4px;
    height: 8px;
    overflow: hidden;
  }}
  .progress-bar-fill {{
    height: 100%;
    border-radius: 4px;
    transition: width 0.3s;
  }}
  .progress-label {{
    font-size: 12px;
    font-weight: 600;
    min-width: 40px;
  }}
  .ascii-bar {{
    font-family: monospace;
    font-size: 13px;
    color: {COLORS['accent_cyan']};
  }}
  .capex-row {{
    display: flex;
    align-items: center;
    padding: 8px 0;
    border-bottom: 1px solid {COLORS['border']};
    gap: 12px;
  }}
  .capex-company {{
    width: 100px;
    font-weight: 600;
  }}
  .footer {{
    background: {COLORS['bg_card']};
    border-top: 1px solid {COLORS['border']};
    padding: 16px 40px;
    text-align: center;
    color: {COLORS['text_secondary']};
    font-size: 12px;
  }}
  .scenario-tabs {{
    display: flex;
    gap: 8px;
    margin-bottom: 16px;
  }}
  .scenario-tab {{
    padding: 6px 16px;
    border-radius: 4px;
    font-size: 12px;
    font-weight: 600;
    cursor: pointer;
  }}
  .scenario-bear {{ background: rgba(244,67,54,0.2); border: 1px solid #f44336; color: #f44336; }}
  .scenario-base {{ background: rgba(30,136,229,0.2); border: 1px solid #1e88e5; color: #1e88e5; }}
  .scenario-bull {{ background: rgba(0,230,118,0.2); border: 1px solid #00e676; color: #00e676; }}
  .phase-grid {{
    display: grid;
    grid-template-columns: repeat(4, 1fr);
    gap: 12px;
    margin-top: 12px;
  }}
  .phase-card {{
    padding: 16px;
    border-radius: 8px;
    border: 1px solid {COLORS['border']};
    text-align: center;
  }}
  .phase-card.active {{
    border-color: {COLORS['accent_cyan']};
    background: rgba(0,188,212,0.08);
  }}
  .phase-number {{
    font-size: 28px;
    font-weight: 800;
    margin-bottom: 8px;
  }}
  .risk-card {{
    display: flex;
    gap: 12px;
    padding: 12px;
    border-radius: 6px;
    margin-bottom: 8px;
    border: 1px solid {COLORS['border']};
    background: {COLORS['bg_card2']};
  }}
</style>
</head>
<body>

<div class="header">
  <div>
    <h1>AI Supply Chain Intelligence Dashboard</h1>
    <div class="subtitle">Token Demand → Hardware Demand Quantification | Bottleneck Detection | Investment Signals</div>
  </div>
  <div class="date-badge">Updated: {date_str}</div>
</div>

<div class="container">

  <!-- SECTION 1: VALUE CHAIN -->
  <div class="section">
    <div class="section-title">Section 1 &nbsp;|&nbsp; AI Supply Chain Value Chain</div>
    {value_chain_html}
  </div>

  <!-- SECTION 2: TOKEN -> HARDWARE DEMAND -->
  <div class="section">
    <div class="section-title">Section 2 &nbsp;|&nbsp; Token Demand → Hardware Demand</div>
    {demand_table_html}
  </div>

  <!-- SECTION 3: BOTTLENECK ANALYSIS -->
  <div class="section">
    <div class="section-title">Section 3 &nbsp;|&nbsp; Bottleneck Analysis</div>
    {bottleneck_html}
  </div>

  <!-- SECTION 4: HYPERSCALER CAPEX -->
  <div class="section">
    <div class="section-title">Section 4 &nbsp;|&nbsp; Hyperscaler CapEx Trends</div>
    {capex_html}
  </div>

  <!-- SECTION 5: INVESTMENT SIGNALS -->
  <div class="section">
    <div class="section-title">Section 5 &nbsp;|&nbsp; Investment Signals</div>
    {signals_html}
  </div>

  <!-- SECTION 6: SCENARIO FORECAST -->
  <div class="section">
    <div class="section-title">섹션 6 &nbsp;|&nbsp; 2025-2027 수요 예측 시나리오</div>
    {forecast_html}
  </div>

  <!-- SECTION 7: INVESTMENT NETWORK GRAPH -->
  <div class="section">
    <div class="section-title">섹션 7 &nbsp;|&nbsp; AI 공급망 투자 관계 네트워크</div>
    {network_graph_html}
  </div>

</div>

<div class="footer">
  AI Supply Chain Intelligence System v1.0 &nbsp;|&nbsp;
  데이터: Seed data + 애널리스트 추정 + 공개 공시 &nbsp;|&nbsp;
  생성일: {date_str} &nbsp;|&nbsp;
  본 자료는 정보 제공 목적이며 투자 권유가 아닙니다
</div>

</body>
</html>"""

    return html


def _build_value_chain_section(bottleneck_results):
    """Build value chain visualization HTML."""
    if not bottleneck_results:
        return "<p>Bottleneck data unavailable</p>"

    all_scores = bottleneck_results.get("all_scores", {})

    # Layer order for display
    layer_order = [
        ("end_customer", "End Customer"),
        ("ai_service", "AI Service"),
        ("cloud_dc", "Cloud / DC"),
        ("gpu_server", "GPU Server"),
        ("asic", "ASIC"),
        ("hbm", "HBM"),
        ("dram", "DRAM"),
        ("ssd_storage", "SSD Storage"),
        ("networking", "Networking"),
        ("packaging", "Packaging"),
        ("power", "Power"),
        ("foundry", "Foundry"),
        ("edge_ai", "Edge AI"),
        ("sovereign_ai", "Sovereign AI"),
    ]

    # Map component keys to layer keys
    comp_to_layer = {
        "HBM": "hbm",
        "CoWoS": "packaging",
        "GPU": "gpu_server",
        "Power_DC": "power",
        "Networking": "networking",
        "DRAM": "dram",
        "SSD": "ssd_storage",
        "CPU": "cpu",
        "Foundry_Advanced": "foundry",
        "ASIC": "asic",
        "Edge_AI": "edge_ai",
    }

    # Get layer utilization from bottleneck scores
    layer_util = {}
    for comp, data in all_scores.items():
        layer = comp_to_layer.get(comp)
        if layer:
            layer_util[layer] = data

    html = '<div class="chain-grid">'
    for layer_key, layer_label in layer_order:
        # Get bottleneck risk from config
        layer_info = config.VALUE_CHAIN.get(layer_key, {})
        risk = layer_info.get("bottleneck_risk", "low")

        # Check if we have utilization data
        util_data = layer_util.get(layer_key, {})
        util = util_data.get("utilization", None)
        severity = util_data.get("severity", risk)

        color = SEVERITY_COLORS.get(severity, SEVERITY_COLORS["low"])

        companies = layer_info.get("companies", [])
        companies_str = ", ".join(companies[:3])
        if len(companies) > 3:
            companies_str += f" +{len(companies)-3}"

        util_str = f"{util:.0%}" if util is not None else risk.upper()

        html += f"""
        <div class="chain-card" style="background: rgba(10,22,40,0.9); border-left-color: {color};">
          <div class="layer-name" style="color: {color};">{layer_label}</div>
          <div style="margin: 4px 0;">
            <span class="badge badge-{severity}">{util_str}</span>
          </div>
          <div class="company-list" style="color: {COLORS['text_secondary']};">{companies_str}</div>
        </div>"""

    html += "</div>"

    # Legend
    html += f"""
    <div style="display: flex; gap: 16px; margin-top: 12px; font-size: 12px;">
      <span style="color: {SEVERITY_COLORS['critical']};">■ Critical (&gt;85%)</span>
      <span style="color: {SEVERITY_COLORS['high']};">■ High (&gt;70%)</span>
      <span style="color: {SEVERITY_COLORS['medium']};">■ Medium (&gt;50%)</span>
      <span style="color: {SEVERITY_COLORS['low']};">■ Low (&lt;50%)</span>
    </div>"""

    return html


def _build_demand_table(modeling_results):
    """Build token -> hardware demand conversion table."""
    if not modeling_results:
        return "<p>Modeling data unavailable</p>"

    snap = modeling_results.get("current_snapshot_2025", {})
    service_breakdown = modeling_results.get("service_breakdown_2024", {})

    # Summary metrics
    metrics_html = '<div class="card-grid">'
    metrics = [
        ("Total Tokens/Day", snap.get("total_tokens_per_day_fmt", "N/A"), "2025 Base estimate"),
        ("GPU Required", snap.get("gpu_count_fmt", "N/A"), "H100-equivalent"),
        ("HBM Demand", snap.get("hbm_demand_fmt", "N/A"), "Total installed"),
        ("Power Required", snap.get("power_demand_fmt", "N/A"), "Data center MW"),
        ("SSD Access/Day", snap.get("ssd_demand_fmt", "N/A"), "RAG + Vector DB"),
        ("GPU CapEx", snap.get("capex_gpu_fmt", "N/A"), "Hardware procurement"),
    ]
    for label, value, sub in metrics:
        metrics_html += f"""
        <div class="metric-card">
          <div class="label">{label}</div>
          <div class="value">{value}</div>
          <div class="sub">{sub}</div>
        </div>"""
    metrics_html += "</div>"

    # Service breakdown table
    table_html = f"""
    <div class="card" style="margin-top: 16px;">
      <table>
        <thead>
          <tr>
            <th>AI Service</th>
            <th>Tokens/Day (2024)</th>
            <th>H100 Equivalent</th>
            <th>B200 Equivalent</th>
            <th>HBM (GB)</th>
            <th>Power (MW)</th>
          </tr>
        </thead>
        <tbody>"""

    for service, data in service_breakdown.items():
        table_html += f"""
          <tr>
            <td style="font-weight: 600;">{service}</td>
            <td>{data.get('tokens_fmt', 'N/A')}</td>
            <td style="color: {COLORS['accent_cyan']};">{data.get('h100_fmt', 'N/A')}</td>
            <td style="color: {COLORS['accent_green']};">{data.get('b200_fmt', 'N/A')}</td>
            <td>{data.get('hbm_gb', 0)/1e6:.1f}M</td>
            <td>{data.get('power_mw', 0):.1f}</td>
          </tr>"""

    table_html += """
        </tbody>
      </table>
    </div>"""

    # Model explanation
    model_html = f"""
    <div class="card" style="margin-top: 12px; background: {COLORS['bg_card2']};">
      <div style="font-size: 13px; color: {COLORS['text_secondary']}; font-style: italic;">
        <strong style="color: {COLORS['accent_cyan']};">Model:</strong>
        GPU_needed = Tokens_per_day / (86400 × GPU_tokens_per_sec × utilization) |
        H100: 2,000 tok/s | B200: 4,500 tok/s | Utilization: 85% |
        HBM: 80GB/H100 | Power: TDP×1.3×PUE/1e6 (MW)
      </div>
    </div>"""

    return metrics_html + table_html + model_html


def _build_bottleneck_section(bottleneck_results):
    """Build bottleneck analysis section."""
    if not bottleneck_results:
        return "<p>Bottleneck data unavailable</p>"

    primary = bottleneck_results.get("current_bottleneck", "HBM")
    primary_util = bottleneck_results.get("current_utilization", 0)
    primary_sev = bottleneck_results.get("current_severity", "critical")
    next_b = bottleneck_results.get("next_bottleneck", "CoWoS")
    after_b = bottleneck_results.get("after_bottleneck", "Power_DC")
    resolution = bottleneck_results.get("resolution_timeline_months", 18)
    inv_window = bottleneck_results.get("investment_window", "H2 2025")

    primary_color = SEVERITY_COLORS.get(primary_sev, "#f44336")

    # Primary bottleneck card
    primary_html = f"""
    <div class="card-grid">
      <div class="metric-card" style="border-color: {primary_color};">
        <div class="label">Current Bottleneck</div>
        <div class="value" style="color: {primary_color}; font-size: 32px;">{primary}</div>
        <div style="margin: 8px 0;">
          <div class="progress-bar-wrap">
            <div class="progress-bar-bg">
              <div class="progress-bar-fill" style="width: {primary_util*100:.0f}%; background: {primary_color};"></div>
            </div>
            <span class="progress-label" style="color: {primary_color};">{primary_util:.0%}</span>
          </div>
        </div>
        <div class="sub">Resolution: ~{resolution} months | Window: {inv_window}</div>
      </div>
      <div class="metric-card">
        <div class="label">Next Bottleneck</div>
        <div class="value" style="color: {COLORS['high']}; font-size: 28px;">{next_b}</div>
        <div class="sub">As {primary} capacity is added, pressure shifts here</div>
      </div>
      <div class="metric-card">
        <div class="label">After Next</div>
        <div class="value" style="color: {COLORS['medium']}; font-size: 28px;">{after_b}</div>
        <div class="sub">Long-term structural constraint</div>
      </div>
    </div>"""

    # All components utilization bars
    all_scores = bottleneck_results.get("all_scores", {})
    sorted_scores = sorted(all_scores.items(), key=lambda x: x[1]["utilization"], reverse=True)

    util_html = f'<div class="card" style="margin-top: 16px;">'
    util_html += f'<table><thead><tr><th>Component</th><th>Utilization</th><th>Severity</th><th>Resolution</th><th>Investment Window</th></tr></thead><tbody>'

    for comp, data in sorted_scores:
        sev = data["severity"]
        util = data["utilization"]
        color = SEVERITY_COLORS.get(sev, COLORS["low"])
        bar = make_ascii_bar(util, 1.0, 20)
        util_html += f"""
        <tr>
          <td style="font-weight: 600;">{comp}</td>
          <td>
            <div class="progress-bar-wrap">
              <div class="progress-bar-bg" style="width: 120px;">
                <div class="progress-bar-fill" style="width: {util*100:.0f}%; background: {color};"></div>
              </div>
              <span style="color: {color}; font-size: 12px; font-weight: 600;">{util:.0%}</span>
            </div>
          </td>
          <td><span class="badge badge-{sev}">{sev}</span></td>
          <td style="color: {COLORS['text_secondary']};">{data.get('resolution_months', '?')}mo</td>
          <td style="color: {COLORS['text_secondary']}; font-size: 12px;">{data.get('investment_window', 'N/A')}</td>
        </tr>"""

    util_html += "</tbody></table></div>"

    # Cascade risks
    cascade_risks = bottleneck_results.get("cascade_risks", [])
    risk_html = ""
    if cascade_risks:
        risk_html = f'<div style="margin-top: 16px;"><div style="font-size: 14px; font-weight: 600; color: {COLORS["accent_cyan"]}; margin-bottom: 8px;">Cascade Risk Chains</div>'
        for risk in cascade_risks:
            sev_color = SEVERITY_COLORS.get(risk.get("severity", "medium"), COLORS["medium"])
            risk_html += f"""
            <div class="risk-card">
              <span class="badge badge-{risk.get('severity','medium')}" style="flex-shrink: 0;">{risk.get('severity','').upper()}</span>
              <div>
                <div style="font-weight: 600; color: {sev_color};">{risk.get('chain', '')}</div>
                <div style="font-size: 12px; color: {COLORS['text_secondary']};">{risk.get('description', '')}</div>
              </div>
            </div>"""
        risk_html += "</div>"

    return primary_html + util_html + risk_html


def _build_capex_section(data_results):
    """Build hyperscaler CapEx section with ASCII bar charts."""
    capex_data = {}
    if data_results:
        capex_data = data_results.get("hyperscaler_capex", {})
    if not capex_data:
        capex_data = config.HYPERSCALER_CAPEX

    years = ["2022", "2023", "2024", "2025", "2026_est"]
    max_val = 0
    for company, year_data in capex_data.items():
        for y in years:
            val = year_data.get(y, 0)
            if val > max_val:
                max_val = val

    html = f'<div class="card">'
    html += f'<table><thead><tr><th>Company</th>'
    for y in years:
        label = y.replace("_est", "E")
        html += f'<th style="text-align: center;">{label}</th>'
    html += '<th>Bar (2025)</th></tr></thead><tbody>'

    company_order = ["Amazon", "Microsoft", "Google", "Meta", "xAI"]
    company_colors = {
        "Amazon": "#ff9900",
        "Microsoft": "#00a4ef",
        "Google": "#4285f4",
        "Meta": "#0668e1",
        "xAI": "#ffffff",
    }

    total_by_year = {y: 0 for y in years}

    for company in company_order:
        year_data = capex_data.get(company, {})
        color = company_colors.get(company, COLORS["accent_cyan"])
        capex_2025 = year_data.get("2025", 0)
        bar = make_ascii_bar(capex_2025, max_val, 20)

        html += f'<tr><td style="font-weight: 600; color: {color};">{company}</td>'
        for y in years:
            val = year_data.get(y, "-")
            fmt_val = f"${val}B" if isinstance(val, (int, float)) else val
            html += f'<td style="text-align: center;">{fmt_val}</td>'
            if isinstance(val, (int, float)):
                total_by_year[y] += val

        html += f'<td style="font-family: monospace; color: {color}; font-size: 12px;">{bar}</td></tr>'

    # Total row
    html += f'<tr style="font-weight: 700; border-top: 2px solid {COLORS["accent_blue"]};">'
    html += f'<td style="color: {COLORS["accent_cyan"]};">TOTAL</td>'
    for y in years:
        total = total_by_year[y]
        html += f'<td style="text-align: center; color: {COLORS["accent_cyan"]};">${total:.0f}B</td>'
    total_2025 = total_by_year.get("2025", 0)
    bar = make_ascii_bar(total_2025, max_val * 5, 20)
    html += f'<td style="font-family: monospace; color: {COLORS["accent_cyan"]}; font-size: 12px;">{bar}</td></tr>'

    html += '</tbody></table></div>'

    # CAGR note
    html += f"""
    <div class="card" style="margin-top: 12px; background: {COLORS['bg_card2']};">
      <div style="font-size: 13px; color: {COLORS['text_secondary']};">
        <strong style="color: {COLORS['accent_cyan']};">Key Insight:</strong>
        Big 5 combined CapEx 2022: ~$154B → 2025: ~$340B → 2026E: ~$392B.
        2-year CAGR: ~30%. AI infrastructure spending has become the largest single CapEx category in tech history.
        Amazon leads at $105B (2025), driven by AWS + Alexa AI infrastructure.
      </div>
    </div>"""

    return html


def _build_signals_section(strategy_results):
    """Build investment signals table."""
    if not strategy_results:
        return "<p>Strategy data unavailable</p>"

    signals = strategy_results.get("investment_signals", [])
    phase = strategy_results.get("phase_position", "Phase 2")
    themes = strategy_results.get("key_themes", [])
    signal_summary = strategy_results.get("signal_summary", {})

    # Phase position header
    header_html = f"""
    <div class="card-grid">
      <div class="metric-card">
        <div class="label">Investment Phase</div>
        <div class="value" style="font-size: 20px;">{phase}</div>
        <div class="sub">Current cycle position</div>
      </div>"""

    for signal_type, count in signal_summary.items():
        color = SIGNAL_COLORS.get(signal_type, COLORS["text_secondary"])
        header_html += f"""
      <div class="metric-card">
        <div class="label">{signal_type}</div>
        <div class="value" style="color: {color};">{count}</div>
        <div class="sub">companies</div>
      </div>"""

    header_html += "</div>"

    # Themes
    themes_html = f'<div class="card" style="margin-top: 16px;"><div style="font-size: 13px; font-weight: 600; color: {COLORS["accent_cyan"]}; margin-bottom: 8px;">Key Investment Themes</div>'
    themes_html += '<div style="display: flex; flex-wrap: wrap; gap: 8px;">'
    for theme in themes:
        themes_html += f'<span style="background: rgba(30,136,229,0.15); border: 1px solid {COLORS["accent_blue"]}; border-radius: 16px; padding: 4px 12px; font-size: 12px;">{theme}</span>'
    themes_html += '</div></div>'

    # Signals table
    table_html = f'<div class="card" style="margin-top: 12px;"><table>'
    table_html += '<thead><tr><th>Company</th><th>Layer</th><th>Signal</th><th>Thesis</th><th>Catalyst</th><th>Risk</th><th>Timeframe</th></tr></thead><tbody>'

    # Sort: Strong Buy first, then Buy, Hold, Watch
    signal_order = {"Strong Buy": 0, "Buy": 1, "Hold": 2, "Watch": 3}
    sorted_signals = sorted(signals, key=lambda x: signal_order.get(x.get("signal", "Watch"), 99))

    for sig in sorted_signals[:12]:  # top 12
        signal = sig.get("signal", "Watch")
        signal_color = SIGNAL_COLORS.get(signal, COLORS["text_secondary"])
        badge_class = {
            "Strong Buy": "badge-sb",
            "Buy": "badge-buy",
            "Hold": "badge-hold",
            "Watch": "badge-watch",
        }.get(signal, "badge-watch")

        thesis = sig.get("thesis", "")[:100] + ("..." if len(sig.get("thesis", "")) > 100 else "")
        catalyst = sig.get("catalyst", "")[:80] + ("..." if len(sig.get("catalyst", "")) > 80 else "")
        risk = sig.get("risk", "")[:80] + ("..." if len(sig.get("risk", "")) > 80 else "")

        table_html += f"""
        <tr>
          <td style="font-weight: 700; color: {signal_color};">{sig.get('company', '')}</td>
          <td><span style="color: {COLORS['accent_cyan']}; font-size: 11px;">{sig.get('layer', '')}</span></td>
          <td><span class="badge {badge_class}">{signal}</span></td>
          <td style="font-size: 12px; max-width: 200px;">{thesis}</td>
          <td style="font-size: 12px; color: {COLORS['accent_green']}; max-width: 150px;">{catalyst}</td>
          <td style="font-size: 12px; color: {COLORS['high']}; max-width: 150px;">{risk}</td>
          <td style="font-size: 12px; color: {COLORS['text_secondary']};">{sig.get('timeframe', '')}</td>
        </tr>"""

    table_html += '</tbody></table></div>'

    # Phase roadmap
    phase_html = f'<div class="card" style="margin-top: 12px;"><div style="font-size: 14px; font-weight: 600; color: {COLORS["accent_cyan"]}; margin-bottom: 12px;">Investment Phase Roadmap</div>'
    phase_html += '<div class="phase-grid">'
    phase_roadmap = strategy_results.get("phase_roadmap", [])
    for p in phase_roadmap:
        is_current = p.get("is_current", False)
        phase_color = COLORS["accent_cyan"] if is_current else COLORS["text_secondary"]
        status_style = "border-color: " + COLORS["accent_cyan"] + ";" if is_current else ""
        status_badge = f'<span class="badge" style="background: {COLORS["accent_cyan"]}; color: #000;">CURRENT</span>' if is_current else f'<span style="font-size: 11px; color: {COLORS["text_secondary"]};">{p.get("status","")}</span>'
        phase_html += f"""
        <div class="phase-card {'active' if is_current else ''}" style="{status_style}">
          <div class="phase-number" style="color: {phase_color};">Phase {p['phase']}</div>
          <div style="font-weight: 700; margin-bottom: 4px;">{p.get('focus', '')}</div>
          <div style="font-size: 11px; color: {COLORS['text_secondary']}; margin-bottom: 8px;">{p.get('timeframe', '')}</div>
          {status_badge}
          <div style="font-size: 11px; color: {COLORS['text_secondary']}; margin-top: 8px;">{', '.join(p.get('key_players', []))}</div>
        </div>"""
    phase_html += '</div></div>'

    return header_html + themes_html + table_html + phase_html


def _build_forecast_section(modeling_results):
    """Build 2025-2027 scenario forecast section."""
    if not modeling_results:
        return "<p>Modeling data unavailable</p>"

    scenario_table = modeling_results.get("scenario_table", {})

    scenario_meta = {
        "bear": {"label": "Bear", "color": COLORS["critical"], "class": "scenario-bear"},
        "base": {"label": "Base", "color": COLORS["accent_blue"], "class": "scenario-base"},
        "bull": {"label": "Bull", "color": COLORS["low"], "class": "scenario-bull"},
    }

    metrics = [
        ("total_tokens_per_day_fmt", "Tokens/Day"),
        ("gpu_count_fmt", "GPU Required"),
        ("hbm_demand_fmt", "HBM Demand"),
        ("power_demand_fmt", "Power (MW)"),
        ("ssd_demand_fmt", "SSD Access/Day"),
    ]

    years = ["2024", "2025", "2026", "2027"]

    html = f'<div class="card"><table>'
    html += '<thead><tr><th>Scenario</th><th>Metric</th>'
    for y in years:
        html += f'<th style="text-align: center;">{y}</th>'
    html += '</tr></thead><tbody>'

    for scenario, meta in scenario_meta.items():
        color = meta["color"]
        scenario_data = scenario_table.get(scenario, {})
        label = meta["label"]

        for i, (metric_key, metric_label) in enumerate(metrics):
            row_style = f"border-top: 2px solid {color};" if i == 0 else ""
            html += f'<tr style="{row_style}">'
            if i == 0:
                html += f'<td rowspan="{len(metrics)}" style="font-weight: 700; color: {color}; vertical-align: middle; text-align: center; font-size: 18px;">{label}</td>'
            html += f'<td style="color: {COLORS["text_secondary"]}; font-size: 12px;">{metric_label}</td>'
            for y in years:
                year_data = scenario_data.get(y, {})
                val = year_data.get(metric_key, "-")
                html += f'<td style="text-align: center; color: {color};">{val}</td>'
            html += '</tr>'

    html += '</tbody></table></div>'

    # Key insight callout
    base_2027 = scenario_table.get("base", {}).get("2027", {})
    bull_2027 = scenario_table.get("bull", {}).get("2027", {})
    bear_2027 = scenario_table.get("bear", {}).get("2027", {})

    html += f"""
    <div class="card" style="margin-top: 12px; background: {COLORS['bg_card2']};">
      <div style="font-size: 14px; font-weight: 700; color: {COLORS['accent_cyan']}; margin-bottom: 10px;">2027 Forecast Summary</div>
      <div style="display: grid; grid-template-columns: repeat(3, 1fr); gap: 16px;">
        <div style="text-align: center;">
          <div style="color: {COLORS['critical']}; font-size: 12px; font-weight: 700; text-transform: uppercase;">Bear Case</div>
          <div style="font-size: 16px; font-weight: 700; margin-top: 4px;">{bear_2027.get('total_tokens_per_day_fmt','N/A')}/day</div>
          <div style="font-size: 12px; color: {COLORS['text_secondary']};">GPU: {bear_2027.get('gpu_count_fmt','N/A')}</div>
        </div>
        <div style="text-align: center; border-left: 1px solid {COLORS['border']}; border-right: 1px solid {COLORS['border']};">
          <div style="color: {COLORS['accent_blue']}; font-size: 12px; font-weight: 700; text-transform: uppercase;">Base Case</div>
          <div style="font-size: 16px; font-weight: 700; margin-top: 4px;">{base_2027.get('total_tokens_per_day_fmt','N/A')}/day</div>
          <div style="font-size: 12px; color: {COLORS['text_secondary']};">GPU: {base_2027.get('gpu_count_fmt','N/A')}</div>
        </div>
        <div style="text-align: center;">
          <div style="color: {COLORS['low']}; font-size: 12px; font-weight: 700; text-transform: uppercase;">Bull Case</div>
          <div style="font-size: 16px; font-weight: 700; margin-top: 4px;">{bull_2027.get('total_tokens_per_day_fmt','N/A')}/day</div>
          <div style="font-size: 12px; color: {COLORS['text_secondary']};">GPU: {bull_2027.get('gpu_count_fmt','N/A')}</div>
        </div>
      </div>
    </div>"""

    return html


def _build_network_graph_section(mapping_results):
    """D3.js force-directed 네트워크 그래프 HTML 섹션 생성"""
    import json as _json

    if not mapping_results or "network_graph" not in mapping_results:
        return "<p style='color:#8fa3c0;'>네트워크 그래프 데이터 없음</p>"

    network_graph = mapping_results["network_graph"]
    nodes_json = _json.dumps(network_graph.get("nodes", []), ensure_ascii=False)
    edges_json = _json.dumps(network_graph.get("edges", []), ensure_ascii=False)
    legend_node = network_graph.get("legend", {}).get("node_colors", {})
    legend_edge = network_graph.get("legend", {}).get("edge_colors", {})

    layer_ko = {
        "gpu_server":"GPU 서버","hbm":"HBM 메모리","cloud_dc":"클라우드/DC",
        "end_customer":"최종 사용자","ai_service":"AI 서비스",
        "packaging":"패키징","foundry":"파운드리","networking":"네트워킹",
        "power":"전력 인프라","dram":"DRAM","ssd_storage":"SSD",
        "cpu":"CPU","asic":"ASIC","edge_ai":"엣지 AI","sovereign_ai":"Sovereign AI",
    }
    edge_ko = {
        "investment":"투자","hardware_supply":"하드웨어 공급",
        "service":"서비스 계약","vc":"VC 투자","partnership":"파트너십",
    }

    node_legend_html = "".join(
        f'<span style="display:inline-flex;align-items:center;margin:4px 8px;">'
        f'<span style="width:12px;height:12px;border-radius:50%;background:{color};display:inline-block;margin-right:5px;"></span>'
        f'{layer_ko.get(layer, layer)}</span>'
        for layer, color in list(legend_node.items())[:8]
    )
    edge_legend_html = "".join(
        f'<span style="display:inline-flex;align-items:center;margin:4px 8px;">'
        f'<span style="width:24px;height:3px;background:{color};display:inline-block;margin-right:5px;"></span>'
        f'{edge_ko.get(rtype, rtype)}</span>'
        for rtype, color in legend_edge.items()
    )

    return f"""
<div style="background:#111827;border-radius:16px;padding:32px;margin:24px 0;">
  <h2 style="color:#f9fafb;font-size:1.4rem;margin-bottom:8px;">AI 공급망 투자 관계 네트워크</h2>
  <p style="color:#9ca3af;font-size:0.85rem;margin-bottom:16px;">
    노드 크기 = 시가총액 | 출처: Bloomberg, 각사 IR, Reuters
  </p>
  <div style="display:flex;flex-wrap:wrap;gap:4px;margin-bottom:8px;">
    <span style="color:#d1d5db;font-size:0.8rem;margin-right:8px;">노드:</span>
    {node_legend_html}
  </div>
  <div style="display:flex;flex-wrap:wrap;gap:4px;margin-bottom:16px;">
    <span style="color:#d1d5db;font-size:0.8rem;margin-right:8px;">관계:</span>
    {edge_legend_html}
  </div>
  <svg id="ai-network-svg" style="width:100%;height:700px;background:#0f1923;border-radius:12px;display:block;"></svg>
  <div id="network-tooltip" style="position:fixed;background:rgba(15,25,36,0.95);color:#f9fafb;
    padding:10px 14px;border-radius:8px;font-size:0.8rem;pointer-events:none;
    border:1px solid #374151;max-width:280px;display:none;z-index:9999;line-height:1.6;"></div>
</div>
<script src="https://d3js.org/d3.v7.min.js"></script>
<script>
(function() {{
  const graphNodes = {nodes_json};
  const graphEdges = {edges_json};

  const svg = d3.select("#ai-network-svg");
  const tooltip = document.getElementById("network-tooltip");
  const rect = svg.node().getBoundingClientRect();
  const W = rect.width || 900, H = 700;

  const defs = svg.append("defs");
  const edgeColors = {{}};
  graphEdges.forEach(e => edgeColors[e.type] = e.color);
  Object.entries(edgeColors).forEach(([type, color]) => {{
    defs.append("marker")
      .attr("id", "arrow-" + type)
      .attr("viewBox", "0 -5 10 10")
      .attr("refX", 22).attr("refY", 0)
      .attr("markerWidth", 6).attr("markerHeight", 6)
      .attr("orient", "auto")
      .append("path").attr("d", "M0,-5L10,0L0,5").attr("fill", color);
  }});

  const g = svg.append("g");

  svg.call(d3.zoom().scaleExtent([0.3, 4])
    .on("zoom", e => g.attr("transform", e.transform)));

  const link = g.selectAll(".link")
    .data(graphEdges).enter().append("line")
    .attr("class", "link")
    .attr("stroke", d => d.color)
    .attr("stroke-width", 1.5)
    .attr("stroke-opacity", 0.7)
    .attr("marker-end", d => "url(#arrow-" + d.type + ")")
    .on("mouseover", (event, d) => {{
      tooltip.style.display = "block";
      tooltip.innerHTML = "<b>" + (d.source.id || d.source) + " → " + (d.target.id || d.target) + "</b><br>"
        + "유형: " + d.type + "<br>" + d.description + "<br>"
        + "규모: " + (d.value || "미공개")
        + (d.source_url ? "<br><a href='" + d.source_url + "' style='color:#60a5fa' target='_blank'>출처 ↗</a>" : "");
      tooltip.style.left = (event.clientX + 12) + "px";
      tooltip.style.top  = (event.clientY - 28) + "px";
    }})
    .on("mouseout", () => tooltip.style.display = "none");

  const node = g.selectAll(".node")
    .data(graphNodes).enter().append("g")
    .attr("class", "node")
    .call(d3.drag()
      .on("start", (event, d) => {{ if (!event.active) sim.alphaTarget(0.3).restart(); d.fx=d.x; d.fy=d.y; }})
      .on("drag",  (event, d) => {{ d.fx=event.x; d.fy=event.y; }})
      .on("end",   (event, d) => {{ if (!event.active) sim.alphaTarget(0); d.fx=null; d.fy=null; }}));

  node.append("circle")
    .attr("r", d => d.size / 2)
    .attr("fill", d => d.color + "cc")
    .attr("stroke", d => d.color)
    .attr("stroke-width", 2)
    .on("mouseover", (event, d) => {{
      link.attr("stroke-opacity", l =>
        (l.source.id === d.id || l.target.id === d.id) ? 1 : 0.15);
      tooltip.style.display = "block";
      tooltip.innerHTML = "<b>" + d.id + "</b><br>레이어: " + d.layer + "<br>시가총액: $" + d.market_cap_b + "B";
      tooltip.style.left = (event.clientX + 12) + "px";
      tooltip.style.top  = (event.clientY - 28) + "px";
    }})
    .on("mouseout", () => {{
      link.attr("stroke-opacity", 0.7);
      tooltip.style.display = "none";
    }});

  node.append("text")
    .attr("dy", d => d.size / 2 + 12)
    .attr("text-anchor", "middle")
    .attr("fill", "#e5e7eb")
    .attr("font-size", d => Math.max(9, Math.min(13, d.size / 3)))
    .text(d => d.id);

  const sim = d3.forceSimulation(graphNodes)
    .force("link", d3.forceLink(graphEdges)
      .id(d => d.id).distance(100).strength(0.3))
    .force("charge", d3.forceManyBody().strength(-300))
    .force("center", d3.forceCenter(W/2, H/2))
    .force("collision", d3.forceCollide().radius(d => d.size/2 + 8))
    .on("tick", () => {{
      link
        .attr("x1", d => d.source.x).attr("y1", d => d.source.y)
        .attr("x2", d => d.target.x).attr("y2", d => d.target.y);
      node.attr("transform", d => "translate(" + d.x + "," + d.y + ")");
    }});
}})();
</script>
"""


# ============================================================
# PPT GENERATION
# ============================================================

def generate_pptx(data_results, mapping_results, modeling_results,
                  bottleneck_results, strategy_results, date_str, output_path):
    """Generate PowerPoint presentation."""
    if not PPTX_AVAILABLE:
        print("  [ReportAgent] Skipping PPT: python-pptx not available")
        return None

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Color scheme for PPT
    BG_COLOR = RGBColor(0x0a, 0x16, 0x28)
    ACCENT = RGBColor(0x00, 0xbc, 0xd4)
    WHITE = RGBColor(0xff, 0xff, 0xff)
    GRAY = RGBColor(0x8f, 0xa3, 0xc0)
    RED = RGBColor(0xf4, 0x43, 0x36)
    ORANGE = RGBColor(0xff, 0x98, 0x00)
    GREEN = RGBColor(0x00, 0xe6, 0x76)
    BLUE = RGBColor(0x1e, 0x88, 0xe5)

    def add_slide(layout_idx=6):
        slide_layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(slide_layout)
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        return slide

    def add_text_box(slide, text, left, top, width, height,
                     font_size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color if color else WHITE
        return txBox

    def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
        from pptx.util import Inches
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_color
        else:
            shape.fill.background()
        if line_color:
            shape.line.color.rgb = line_color
        else:
            shape.line.fill.background()
        return shape

    # ---- Slide 1: Cover ----
    slide1 = add_slide()
    add_rect(slide1, 0, 0, 13.33, 7.5, RGBColor(0x0a, 0x16, 0x28))
    add_rect(slide1, 0, 0, 13.33, 0.08, ACCENT)
    add_rect(slide1, 0, 7.42, 13.33, 0.08, ACCENT)
    add_text_box(slide1, "AI Supply Chain Intelligence Report",
                 0.8, 2.0, 11.73, 1.2, font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide1, "Token Demand → Hardware Demand | Bottleneck Detection | Investment Signals",
                 0.8, 3.4, 11.73, 0.6, font_size=16, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text_box(slide1, f"Generated: {date_str}",
                 0.8, 4.5, 11.73, 0.4, font_size=12, color=GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide1, "AI Supply Chain Intelligence System v1.0",
                 0.8, 5.5, 11.73, 0.4, font_size=11, color=GRAY, align=PP_ALIGN.CENTER)

    # ---- Slide 2: Value Chain Overview ----
    slide2 = add_slide()
    add_rect(slide2, 0, 0, 13.33, 0.6, RGBColor(0x0d, 0x47, 0xa1))
    add_text_box(slide2, "AI Supply Chain Value Chain Overview",
                 0.3, 0.1, 12.73, 0.5, font_size=20, bold=True, color=WHITE)

    chain_items = [
        ("End Customer", "critical", "OpenAI, Anthropic, Meta AI"),
        ("AI Service", "high", "API platforms, inference services"),
        ("Cloud / DC", "high", "Azure, AWS, GCP, CoreWeave"),
        ("GPU Server", "critical", "NVIDIA H100/B200, AMD MI300X"),
        ("HBM", "critical", "SK Hynix (50%), Samsung (40%), Micron (10%)"),
        ("Packaging", "critical", "TSMC CoWoS - KEY CONSTRAINT"),
        ("Power", "high", "Vertiv, Schneider, GE Vernova"),
        ("Networking", "high", "Broadcom, Marvell, Arista"),
        ("Foundry", "high", "TSMC (N3/N2), Samsung Foundry"),
    ]

    sev_colors_pptx = {
        "critical": RED,
        "high": ORANGE,
        "medium": RGBColor(0xff, 0xeb, 0x3b),
        "low": GREEN,
    }

    cols = 3
    items_per_col = (len(chain_items) + cols - 1) // cols
    col_w = 4.1
    for i, (name, sev, companies) in enumerate(chain_items):
        col = i // items_per_col
        row = i % items_per_col
        x = 0.3 + col * (col_w + 0.2)
        y = 0.9 + row * 0.72
        sev_color = sev_colors_pptx.get(sev, GREEN)
        add_rect(slide2, x, y, col_w, 0.62, RGBColor(0x0f, 0x1f, 0x3d),
                 line_color=sev_color)
        add_text_box(slide2, name, x + 0.1, y + 0.02, 2.5, 0.3,
                     font_size=11, bold=True, color=sev_color)
        add_text_box(slide2, companies, x + 0.1, y + 0.28, col_w - 0.2, 0.3,
                     font_size=9, color=GRAY)

    # Legend
    for j, (label, color) in enumerate([
        ("Critical (>85%)", RED), ("High (>70%)", ORANGE), ("Medium (>50%)", RGBColor(0xff, 0xeb, 0x3b)), ("Low (<50%)", GREEN)
    ]):
        add_rect(slide2, 0.3 + j * 3.3, 7.0, 0.2, 0.2, color)
        add_text_box(slide2, label, 0.6 + j * 3.3, 6.98, 2.5, 0.3, font_size=9, color=GRAY)

    # ---- Slide 3: Token -> Hardware Demand ----
    slide3 = add_slide()
    add_rect(slide3, 0, 0, 13.33, 0.6, RGBColor(0x0d, 0x47, 0xa1))
    add_text_box(slide3, "Token Demand → Hardware Demand Model",
                 0.3, 0.1, 12.73, 0.5, font_size=20, bold=True, color=WHITE)

    if modeling_results:
        snap = modeling_results.get("current_snapshot_2025", {})
        metrics_pptx = [
            ("Tokens/Day (2025 Base)", snap.get("total_tokens_per_day_fmt", "N/A")),
            ("GPU Required (H100 eq.)", snap.get("gpu_count_fmt", "N/A")),
            ("HBM Demand", snap.get("hbm_demand_fmt", "N/A")),
            ("Power Required", snap.get("power_demand_fmt", "N/A")),
        ]
        for i, (label, value) in enumerate(metrics_pptx):
            col = i % 2
            row = i // 2
            x = 0.3 + col * 6.5
            y = 0.9 + row * 1.4
            add_rect(slide3, x, y, 6.0, 1.2, RGBColor(0x0f, 0x1f, 0x3d),
                     line_color=RGBColor(0x1e, 0x3a, 0x5f))
            add_text_box(slide3, label, x + 0.2, y + 0.1, 5.6, 0.4,
                         font_size=11, color=GRAY)
            add_text_box(slide3, value, x + 0.2, y + 0.45, 5.6, 0.65,
                         font_size=26, bold=True, color=ACCENT)

    # Model formulas
    add_text_box(slide3, "Core Models:",
                 0.3, 3.9, 12.73, 0.4, font_size=12, bold=True, color=ACCENT)
    formulas = [
        "GPU = Tokens/day / (86400 × tok/sec × utilization)   |   H100: 2,000 tok/s | B200: 4,500 tok/s",
        "HBM (GB) = GPU_count × HBM_per_GPU   |   H100: 80GB | B200: 192GB | Utilization: 85%",
        "Power (MW) = GPU × TDP × 1.3 × PUE / 1e6   |   PUE: 1.4 | Server overhead: 1.3x",
    ]
    for i, f in enumerate(formulas):
        add_text_box(slide3, f, 0.3, 4.3 + i * 0.55, 12.73, 0.5,
                     font_size=10, color=RGBColor(0xb0, 0xd0, 0xff))

    # Service breakdown table
    if modeling_results:
        services = modeling_results.get("service_breakdown_2024", {})
        headers = ["Service", "Tokens/Day", "H100 eq.", "Power MW"]
        col_widths = [2.0, 2.5, 2.5, 2.0]
        x_start = 0.5
        y_start = 5.9

        add_rect(slide3, x_start - 0.1, y_start - 0.05,
                 sum(col_widths) + 0.2, 0.38, RGBColor(0x0d, 0x47, 0xa1))
        x_cur = x_start
        for i, h in enumerate(headers):
            add_text_box(slide3, h, x_cur, y_start, col_widths[i], 0.3,
                         font_size=10, bold=True, color=WHITE)
            x_cur += col_widths[i]

        for j, (svc, sdata) in enumerate(list(services.items())[:4]):
            y = y_start + 0.35 + j * 0.28
            x_cur = x_start
            row_vals = [
                svc,
                sdata.get("tokens_fmt", ""),
                sdata.get("h100_fmt", ""),
                f"{sdata.get('power_mw', 0):.1f}",
            ]
            for i, val in enumerate(row_vals):
                color = ACCENT if i > 0 else WHITE
                add_text_box(slide3, val, x_cur, y, col_widths[i], 0.28,
                             font_size=9, color=color)
                x_cur += col_widths[i]

    # ---- Slide 4: Bottleneck Analysis ----
    slide4 = add_slide()
    add_rect(slide4, 0, 0, 13.33, 0.6, RGBColor(0x0d, 0x47, 0xa1))
    add_text_box(slide4, "Bottleneck Analysis & Cascade Prediction",
                 0.3, 0.1, 12.73, 0.5, font_size=20, bold=True, color=WHITE)

    if bottleneck_results:
        primary = bottleneck_results.get("current_bottleneck", "HBM")
        primary_util = bottleneck_results.get("current_utilization", 0.95)
        primary_sev = bottleneck_results.get("current_severity", "critical")
        next_b = bottleneck_results.get("next_bottleneck", "CoWoS")
        resolution = bottleneck_results.get("resolution_timeline_months", 18)

        # Primary bottleneck box
        add_rect(slide4, 0.3, 0.8, 3.8, 2.0, RGBColor(0x2d, 0x06, 0x06),
                 line_color=RED)
        add_text_box(slide4, "PRIMARY BOTTLENECK", 0.4, 0.85, 3.6, 0.35,
                     font_size=9, color=GRAY)
        add_text_box(slide4, primary, 0.4, 1.15, 3.6, 0.7,
                     font_size=28, bold=True, color=RED)
        add_text_box(slide4, f"Utilization: {primary_util:.0%}", 0.4, 1.85, 3.6, 0.4,
                     font_size=14, color=RED)
        add_text_box(slide4, f"Resolution: ~{resolution} months", 0.4, 2.25, 3.6, 0.4,
                     font_size=11, color=GRAY)

        # Next bottleneck box
        add_rect(slide4, 4.3, 0.8, 3.8, 2.0, RGBColor(0x2d, 0x17, 0x06),
                 line_color=ORANGE)
        add_text_box(slide4, "NEXT BOTTLENECK", 4.4, 0.85, 3.6, 0.35,
                     font_size=9, color=GRAY)
        add_text_box(slide4, next_b, 4.4, 1.15, 3.6, 0.7,
                     font_size=28, bold=True, color=ORANGE)
        add_text_box(slide4, "Emerging as HBM relief builds", 4.4, 1.85, 3.6, 0.4,
                     font_size=11, color=GRAY)

        # Cascade arrow text
        add_text_box(slide4, "→", 4.05, 1.45, 0.4, 0.5,
                     font_size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text_box(slide4, "→", 8.05, 1.45, 0.4, 0.5,
                     font_size=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

        after_b = bottleneck_results.get("after_bottleneck", "Power_DC")
        add_rect(slide4, 8.3, 0.8, 3.8, 2.0, RGBColor(0x1a, 0x1a, 0x06),
                 line_color=RGBColor(0xff, 0xeb, 0x3b))
        add_text_box(slide4, "AFTER NEXT", 8.4, 0.85, 3.6, 0.35,
                     font_size=9, color=GRAY)
        add_text_box(slide4, after_b, 8.4, 1.15, 3.6, 0.7,
                     font_size=24, bold=True, color=RGBColor(0xff, 0xeb, 0x3b))

        # Utilization bars
        all_scores = bottleneck_results.get("all_scores", {})
        sorted_scores = sorted(all_scores.items(), key=lambda x: x[1]["utilization"], reverse=True)

        add_text_box(slide4, "Component Utilization (2025)",
                     0.3, 3.1, 12.73, 0.35, font_size=12, bold=True, color=ACCENT)

        bar_colors = {
            "critical": RED, "high": ORANGE,
            "medium": RGBColor(0xff, 0xeb, 0x3b), "low": GREEN
        }

        items_per_row = 4
        for i, (comp, data) in enumerate(sorted_scores[:8]):
            col = i % items_per_row
            row = i // items_per_row
            x = 0.3 + col * 3.2
            y = 3.55 + row * 0.85
            util = data["utilization"]
            sev = data["severity"]
            bar_color = bar_colors.get(sev, GREEN)

            add_text_box(slide4, comp, x, y, 3.0, 0.25, font_size=10, bold=True, color=WHITE)
            # Background bar
            add_rect(slide4, x, y + 0.28, 2.7, 0.2, RGBColor(0x1e, 0x3a, 0x5f))
            # Fill bar
            add_rect(slide4, x, y + 0.28, 2.7 * util, 0.2, bar_color)
            add_text_box(slide4, f"{util:.0%}", x + 2.75, y + 0.24, 0.5, 0.3,
                         font_size=10, bold=True, color=bar_color)

    # ---- Slide 5: Hyperscaler CapEx ----
    slide5 = add_slide()
    add_rect(slide5, 0, 0, 13.33, 0.6, RGBColor(0x0d, 0x47, 0xa1))
    add_text_box(slide5, "Hyperscaler CapEx Trends (USD Billion)",
                 0.3, 0.1, 12.73, 0.5, font_size=20, bold=True, color=WHITE)

    capex_data = config.HYPERSCALER_CAPEX
    years_pptx = ["2022", "2023", "2024", "2025", "2026_est"]
    headers_pptx = ["Company"] + [y.replace("_est", "E") for y in years_pptx] + ["CAGR 22-25"]
    col_w_pptx = [2.0, 1.7, 1.7, 1.7, 1.7, 1.7, 1.5]
    x_start_pptx = 0.3

    # Header row
    add_rect(slide5, x_start_pptx - 0.1, 0.75,
             sum(col_w_pptx) + 0.2, 0.4, RGBColor(0x0d, 0x47, 0xa1))
    x_cur = x_start_pptx
    for i, h in enumerate(headers_pptx):
        add_text_box(slide5, h, x_cur, 0.8, col_w_pptx[i], 0.3,
                     font_size=10, bold=True, color=WHITE)
        x_cur += col_w_pptx[i]

    company_colors_pptx = {
        "Amazon": RGBColor(0xff, 0x99, 0x00),
        "Microsoft": RGBColor(0x00, 0xa4, 0xef),
        "Google": RGBColor(0x42, 0x85, 0xf4),
        "Meta": RGBColor(0x06, 0x68, 0xe1),
        "xAI": WHITE,
    }

    for j, company in enumerate(["Amazon", "Microsoft", "Google", "Meta", "xAI"]):
        y_row = 1.25 + j * 0.5
        year_data = capex_data.get(company, {})
        color = company_colors_pptx.get(company, WHITE)
        x_cur = x_start_pptx
        add_text_box(slide5, company, x_cur, y_row, col_w_pptx[0], 0.42,
                     font_size=11, bold=True, color=color)
        x_cur += col_w_pptx[0]

        v_2022 = year_data.get("2022", 0)
        v_2025 = year_data.get("2025", 0)
        cagr = ((v_2025 / v_2022) ** (1/3) - 1) * 100 if v_2022 > 0 else 0

        for i, y_key in enumerate(years_pptx):
            val = year_data.get(y_key, "-")
            fmt = f"${val}B" if isinstance(val, (int, float)) and val > 0 else "-"
            add_text_box(slide5, fmt, x_cur, y_row, col_w_pptx[i + 1], 0.42,
                         font_size=11, color=color)
            x_cur += col_w_pptx[i + 1]

        add_text_box(slide5, f"+{cagr:.0f}%", x_cur, y_row, col_w_pptx[-1], 0.42,
                     font_size=11, bold=True, color=GREEN if cagr > 20 else ACCENT)

    # Total row
    y_total = 1.25 + 5 * 0.5
    add_rect(slide5, x_start_pptx - 0.1, y_total - 0.05,
             sum(col_w_pptx) + 0.2, 0.5, RGBColor(0x0d, 0x47, 0xa1))
    add_text_box(slide5, "TOTAL", x_start_pptx, y_total,
                 col_w_pptx[0], 0.4, font_size=11, bold=True, color=ACCENT)
    totals = {y: sum(capex_data.get(c, {}).get(y, 0) for c in capex_data) for y in years_pptx}
    x_cur = x_start_pptx + col_w_pptx[0]
    for i, y_key in enumerate(years_pptx):
        add_text_box(slide5, f"${totals[y_key]:.0f}B", x_cur, y_total, col_w_pptx[i+1], 0.4,
                     font_size=11, bold=True, color=ACCENT)
        x_cur += col_w_pptx[i+1]

    add_text_box(slide5,
                 "Big 5 combined CapEx: $154B (2022) → $340B (2025) → $392B (2026E)  |  "
                 "AI infrastructure is now the largest single CapEx category in tech history",
                 0.3, 4.3, 12.73, 0.5, font_size=10, color=GRAY)

    # ---- Slide 6: Investment Signals ----
    slide6 = add_slide()
    add_rect(slide6, 0, 0, 13.33, 0.6, RGBColor(0x0d, 0x47, 0xa1))
    add_text_box(slide6, "Investment Signals & Positioning",
                 0.3, 0.1, 12.73, 0.5, font_size=20, bold=True, color=WHITE)

    if strategy_results:
        phase = strategy_results.get("phase_position", "Phase 2 (HBM)")
        signals = strategy_results.get("investment_signals", [])
        signal_order = {"Strong Buy": 0, "Buy": 1, "Hold": 2, "Watch": 3}
        sorted_sigs = sorted(signals, key=lambda x: signal_order.get(x.get("signal", "Watch"), 99))

        add_text_box(slide6, f"Current Phase: {phase}",
                     0.3, 0.75, 8.0, 0.4, font_size=14, bold=True, color=ACCENT)

        signal_colors_pptx = {
            "Strong Buy": GREEN,
            "Buy": RGBColor(0x4c, 0xaf, 0x50),
            "Hold": RGBColor(0xff, 0xeb, 0x3b),
            "Watch": ORANGE,
        }

        headers_s = ["Company", "Layer", "Signal", "Thesis (Brief)", "Timeframe"]
        col_w_s = [2.0, 1.8, 1.5, 6.0, 1.5]
        x_start_s = 0.3

        add_rect(slide6, x_start_s - 0.1, 1.2, sum(col_w_s) + 0.2, 0.35,
                 RGBColor(0x0d, 0x47, 0xa1))
        x_cur = x_start_s
        for i, h in enumerate(headers_s):
            add_text_box(slide6, h, x_cur, 1.22, col_w_s[i], 0.3,
                         font_size=9, bold=True, color=WHITE)
            x_cur += col_w_s[i]

        for j, sig in enumerate(sorted_sigs[:8]):
            y_row = 1.62 + j * 0.52
            signal = sig.get("signal", "Watch")
            sig_color = signal_colors_pptx.get(signal, WHITE)
            x_cur = x_start_s

            vals = [
                sig.get("company", ""),
                sig.get("layer", ""),
                signal,
                (sig.get("thesis", "")[:90] + "...") if len(sig.get("thesis", "")) > 90 else sig.get("thesis", ""),
                sig.get("timeframe", ""),
            ]
            for i, val in enumerate(vals):
                color = sig_color if i == 2 else (WHITE if i == 0 else GRAY)
                bold = i in [0, 2]
                add_text_box(slide6, val, x_cur, y_row, col_w_s[i], 0.5,
                             font_size=9, bold=bold, color=color)
                x_cur += col_w_s[i]

    # ---- Slide 7: Risk Factors ----
    slide7 = add_slide()
    add_rect(slide7, 0, 0, 13.33, 0.6, RGBColor(0x0d, 0x47, 0xa1))
    add_text_box(slide7, "Key Risk Factors",
                 0.3, 0.1, 12.73, 0.5, font_size=20, bold=True, color=WHITE)

    if strategy_results:
        risks = strategy_results.get("risk_factors", [])
        sev_colors_risk = {
            "critical": RED, "high": ORANGE,
            "medium": RGBColor(0xff, 0xeb, 0x3b), "low": GREEN
        }

        for i, risk in enumerate(risks[:6]):
            col = i % 2
            row = i // 2
            x = 0.3 + col * 6.5
            y = 0.8 + row * 2.1
            sev = risk.get("severity", "medium")
            r_color = sev_colors_risk.get(sev, ORANGE)

            add_rect(slide7, x, y, 6.2, 1.85, RGBColor(0x0f, 0x1f, 0x3d),
                     line_color=r_color)
            add_text_box(slide7, risk.get("risk", "")[:60], x + 0.15, y + 0.1,
                         5.9, 0.4, font_size=11, bold=True, color=r_color)
            add_text_box(slide7, risk.get("description", "")[:120], x + 0.15, y + 0.5,
                         5.9, 0.6, font_size=9, color=GRAY)
            add_text_box(slide7, f"Mitigation: {risk.get('mitigation', '')[:100]}",
                         x + 0.15, y + 1.2, 5.9, 0.55, font_size=9,
                         color=RGBColor(0x66, 0xbb, 0x6a))

    # ---- Slide 8: 2025-2027 Forecast ----
    slide8 = add_slide()
    add_rect(slide8, 0, 0, 13.33, 0.6, RGBColor(0x0d, 0x47, 0xa1))
    add_text_box(slide8, "2025-2027 AI Demand Forecast Scenarios",
                 0.3, 0.1, 12.73, 0.5, font_size=20, bold=True, color=WHITE)

    if modeling_results:
        scenario_table = modeling_results.get("scenario_table", {})
        years_f = ["2024", "2025", "2026", "2027"]
        metrics_f = [
            ("total_tokens_per_day_fmt", "Tokens/Day"),
            ("gpu_count_fmt", "GPU Required"),
            ("hbm_demand_fmt", "HBM Demand"),
            ("power_demand_fmt", "Power (MW)"),
        ]
        scenario_meta_pptx = {
            "bear": (RED, "Bear Case (0.7x)"),
            "base": (BLUE, "Base Case (1.0x)"),
            "bull": (GREEN, "Bull Case (1.5x)"),
        }

        headers_f = ["Scenario", "Metric"] + years_f
        col_w_f = [2.2, 1.8, 2.1, 2.1, 2.1, 2.1]
        x_start_f = 0.2

        add_rect(slide8, x_start_f, 0.72, sum(col_w_f), 0.35,
                 RGBColor(0x0d, 0x47, 0xa1))
        x_cur = x_start_f
        for i, h in enumerate(headers_f):
            add_text_box(slide8, h, x_cur + 0.05, 0.75, col_w_f[i] - 0.1, 0.3,
                         font_size=10, bold=True, color=WHITE)
            x_cur += col_w_f[i]

        row_idx = 0
        for scenario, (s_color, s_label) in scenario_meta_pptx.items():
            sdata = scenario_table.get(scenario, {})
            for m_idx, (metric_key, metric_label) in enumerate(metrics_f):
                y_row = 1.12 + row_idx * 0.38
                x_cur = x_start_f

                if m_idx == 0:
                    add_text_box(slide8, s_label, x_cur + 0.05, y_row,
                                 col_w_f[0] - 0.1, 0.36, font_size=9, bold=True, color=s_color)
                x_cur += col_w_f[0]

                add_text_box(slide8, metric_label, x_cur + 0.05, y_row,
                             col_w_f[1] - 0.1, 0.36, font_size=9, color=GRAY)
                x_cur += col_w_f[1]

                for y_idx, yr in enumerate(years_f):
                    val = sdata.get(yr, {}).get(metric_key, "-")
                    add_text_box(slide8, str(val), x_cur + 0.05, y_row,
                                 col_w_f[y_idx + 2] - 0.1, 0.36, font_size=9, color=s_color)
                    x_cur += col_w_f[y_idx + 2]

                row_idx += 1

    prs.save(output_path)
    print(f"  [ReportAgent] PPT saved: {output_path}")
    return output_path


# ============================================================
# MAIN RUN
# ============================================================

def run(data_results=None, mapping_results=None, modeling_results=None,
        bottleneck_results=None, strategy_results=None):
    """Run the report agent and generate HTML + PPT."""
    print("[ReportAgent] Generating reports...")

    date_str = datetime.date.today().strftime("%Y-%m-%d")
    date_compact = datetime.date.today().strftime("%Y%m%d")

    os.makedirs(config.REPORTS_DIR, exist_ok=True)
    os.makedirs(config.PPTX_DIR, exist_ok=True)

    # Generate HTML dashboard
    html_filename = f"ai_scm_dashboard_{date_compact}.html"
    html_path = os.path.join(config.REPORTS_DIR, html_filename)
    latest_path = os.path.join(config.REPORTS_DIR, "latest.html")

    html_content = generate_html_report(
        data_results, mapping_results, modeling_results,
        bottleneck_results, strategy_results, date_str
    )

    with open(html_path, "w", encoding="utf-8") as f:
        f.write(html_content)
    with open(latest_path, "w", encoding="utf-8") as f:
        f.write(html_content)

    print(f"  [ReportAgent] HTML saved: {html_path}")
    print(f"  [ReportAgent] Latest HTML: {latest_path}")

    # Generate PPT
    pptx_filename = f"ai_scm_{date_compact}.pptx"
    pptx_path = os.path.join(config.PPTX_DIR, pptx_filename)

    pptx_result = generate_pptx(
        data_results, mapping_results, modeling_results,
        bottleneck_results, strategy_results, date_str, pptx_path
    )

    result = {
        "html_path": html_path,
        "html_latest": latest_path,
        "pptx_path": pptx_result,
        "date": date_str,
        "status": "success",
    }

    return result


if __name__ == "__main__":
    # Run standalone with no prior results (will use config defaults)
    result = run()
    print(f"\nHTML: {result['html_path']}")
    print(f"PPT: {result.get('pptx_path', 'Not generated')}")
