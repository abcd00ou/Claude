"""
마케팅 커뮤니케이션 에이전트 (MarCom Agent)
역할: 브랜드 캠페인 기획, 광고 아이템, Awareness 전략, 채널별 콘텐츠 플랜
산출물: PowerPoint — marcom_plan.pptx
"""
import sys
sys.path.insert(0, str(__import__('pathlib').Path(__file__).parent.parent))

import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from config import PPTX_DIR

C = {
    "navy":   "1B2A4A", "red":    "D4001E",
    "blue":   "4472C4", "orange": "ED7D31",
    "green":  "70AD47", "yellow": "FFC000",
    "purple": "7030A0", "teal":   "00B0F0",
    "gray":   "595959", "dark":   "0D1821",
}

def _rgb(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def _rect(slide, left, top, w, h, fill, text="", fs=11, bold=False, tc="FFFFFF",
          align=PP_ALIGN.CENTER, wrap=True):
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(fill)
    if text:
        tf = shape.text_frame; tf.word_wrap = wrap
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = text
        run.font.size = Pt(fs); run.font.bold = bold
        run.font.color.rgb = _rgb(tc)
    return shape

def _tbox(slide, text, left, top, w, h, fs=11, bold=False, color="000000",
          align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(fs); run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return tb

def _multi_para(slide, lines, left, top, w, h, base_fs=10, base_color="CCCCCC"):
    """여러 줄 텍스트박스 (리스트 형태)"""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (line, fs, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run(); run.text = line
        run.font.size = Pt(fs); run.font.bold = bold
        run.font.color.rgb = _rgb(color)
    return tb


def build_pptx():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    def ns(): return prs.slides.add_slide(blank)

    # ══════════════════════════════════════════════════════════
    # Slide 1: Title
    # ══════════════════════════════════════════════════════════
    sl = ns()
    _rect(sl, 0, 0, 13.33, 7.5, C["dark"])
    _rect(sl, 0, 3.2, 13.33, 0.08, C["red"])  # accent line
    _rect(sl, 0, 6.0, 13.33, 1.5, "0D1018")

    _tbox(sl, "SanDisk", 1.0, 0.5, 5, 1.0, fs=48, bold=True, color=C["red"])
    _tbox(sl, "B2C Storage", 1.0, 1.5, 7, 0.9, fs=36, bold=False, color="FFFFFF")
    _tbox(sl, "마케팅 커뮤니케이션 플랜  2025", 1.0, 2.5, 9, 0.7, fs=22, color="AAAAAA")
    _tbox(sl, "MarCom Agent  |  Brand & Campaign Strategy", 1.0, 3.35, 9, 0.55, fs=14, color="888888")

    campaigns = [
        ("Speed Untamed", C["red"]),
        ("Born to Create", C["blue"]),
        ("Go Extreme", C["orange"]),
        ("Always Ready", C["green"]),
    ]
    for i, (name, col) in enumerate(campaigns):
        _rect(sl, 1.0 + i*3.0, 4.4, 2.5, 0.6, col, f'"{name}"', fs=12, bold=True)

    _tbox(sl, f"마케팅 커뮤니케이션 에이전트  |  {datetime.date.today()}  |  CONFIDENTIAL",
          1, 6.2, 11, 0.4, fs=10, color="555555", align=PP_ALIGN.CENTER)
    _tbox(sl, "External SSD  ·  Internal SSD  ·  microSD",
          1, 6.7, 11, 0.4, fs=11, color="777777", align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    # Slide 2: 브랜드 포지셔닝
    # ══════════════════════════════════════════════════════════
    sl = ns()
    _rect(sl, 0, 0, 13.33, 0.9, C["navy"])
    _tbox(sl, "브랜드 포지셔닝 — SanDisk Brand Architecture", 0.3, 0.05, 12, 0.75,
          fs=22, bold=True, color="FFFFFF")

    # Brand house 구조
    _rect(sl, 0.3, 1.1, 12.7, 0.6, C["red"],
          "Brand Purpose: \"Empowering creators, gamers, and professionals to capture and access their world — fearlessly.\"",
          fs=11, bold=True)

    pillars = [
        ("Performance\n& Speed", C["blue"],
         "• Fastest in class\n• BiCS8 technology edge\n• PS5/gaming certified"),
        ("Reliability\n& Durability", C["orange"],
         "• IP68 water/dust proof\n• Drop proof 2m\n• Extreme temp range"),
        ("Accessibility\n& Value", C["green"],
         "• Best price/performance\n• Wide capacity range\n• Universal compatibility"),
        ("Creativity\n& Freedom", C["purple"],
         "• Creator-first design\n• 4K/8K workflow\n• Seamless ecosystem"),
    ]
    for i, (title, col, body) in enumerate(pillars):
        x = 0.3 + i * 3.2
        _rect(sl, x, 1.9, 3.0, 0.7, col, title, fs=13, bold=True)
        _rect(sl, x, 2.6, 3.0, 2.0, "1B2A3A", body, fs=9, align=PP_ALIGN.LEFT)

    # Target audiences
    _rect(sl, 0.3, 4.8, 12.7, 0.45, C["navy"], "TARGET AUDIENCES", fs=12, bold=True)
    audiences = [
        ("🎮 Gamers", "PlayStation / Xbox 커뮤니티\n내장 SSD 업그레이드\nPC 게이밍 빌드", C["blue"]),
        ("🎬 Creators", "유튜버 / 영상 크리에이터\n4K/8K 촬영·편집 워크플로우\n드론·액션캠", C["orange"]),
        ("📸 Photographers", "미러리스·DSLR 사용자\n고속 microSD\nExtreme Pro 라인업", C["green"]),
        ("🏃 Active Users", "아웃도어·스포츠 활동\n방수·방진 외장 SSD\n포터블 백업", C["purple"]),
        ("💼 Professionals", "재택근무·출장 전문직\nMy Passport 라인\n클라우드 백업 연동", C["teal"]),
    ]
    for i, (name, desc, col) in enumerate(audiences):
        x = 0.3 + i * 2.5
        _rect(sl, x, 5.35, 2.3, 0.45, col, name, fs=10, bold=True)
        _rect(sl, x, 5.8,  2.3, 1.4,  "1A2235", desc, fs=8, align=PP_ALIGN.LEFT)

    # ══════════════════════════════════════════════════════════
    # Slide 3: 2025 연간 캠페인 캘린더
    # ══════════════════════════════════════════════════════════
    sl = ns()
    _rect(sl, 0, 0, 13.33, 0.9, C["navy"])
    _tbox(sl, "2025 Annual Campaign Calendar — 캠페인 타임라인", 0.3, 0.05, 12, 0.75,
          fs=22, bold=True, color="FFFFFF")

    # 월별 헤더
    months = ["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]
    col_w = 0.95
    x_start = 2.0
    for mi, m in enumerate(months):
        _rect(sl, x_start + mi*col_w, 1.0, col_w-0.05, 0.4, C["navy"], m, fs=9, bold=True)

    campaigns_cal = [
        ("External SSD",   C["blue"],    [(1,3,"New Year Storage"), (4,6,"Creator Season"), (8,9,"Back to School"), (10,12,"Holiday Push")]),
        ("Internal SSD",   C["orange"],  [(1,2,"CES Launch"), (3,4,"Gaming Spring"), (6,8,"Summer Gaming"), (9,12,"Holiday Gaming")]),
        ("microSD",        C["green"],   [(3,5,"Spring Travel"), (6,7,"Action Season"), (8,9,"School/Gaming"), (10,12,"Holiday Bundle")]),
        ("Brand Awareness",C["red"],     [(1,12,"Always-on Digital")]),
        ("Influencer/PR",  C["purple"],  [(2,3,"CES Follow-up"), (6,7,"Summer Seeding"), (9,10,"Holiday Seeding")]),
    ]
    for ri, (track, color, events) in enumerate(campaigns_cal):
        y = 1.55 + ri * 1.0
        _rect(sl, 0.05, y, 1.88, 0.75, color, track, fs=9, bold=True)
        for (start_m, end_m, label) in events:
            x_ev = x_start + (start_m-1)*col_w
            w_ev = (end_m - start_m + 1)*col_w - 0.1
            alpha = "CC" if ri == 3 else ""
            _rect(sl, x_ev, y+0.1, w_ev, 0.55, color, label, fs=7.5, bold=False)

    _tbox(sl, "● 상시: Digital/Social  ● Q1: CES/신제품  ● Q3: Back-to-School  ● Q4: BF/CM/Holiday 집중",
          0.3, 7.1, 12, 0.35, fs=9, color="888888", align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    # Slide 4: 캠페인 상세 — "Go Extreme" (External SSD 플래그십)
    # ══════════════════════════════════════════════════════════
    sl = ns()
    _rect(sl, 0, 0, 13.33, 7.5, "0D1821")
    _rect(sl, 0, 0, 13.33, 0.9, C["red"])
    _tbox(sl, 'Campaign: "Go Extreme" — SanDisk Extreme/Extreme Pro External SSD', 0.3, 0.05,
          12, 0.75, fs=18, bold=True, color="FFFFFF")

    # Insight
    _rect(sl, 0.3, 1.0, 12.7, 0.5, "1A2235",
          "Consumer Insight: \"크리에이터들은 촬영 현장에서도 타협하지 않는다. 비가 와도, 모래밭이어도 내 데이터는 안전해야 한다.\"",
          fs=10, tc="CCDDFF")

    # Campaign elements
    elements = [
        ("BIG IDEA", C["red"],
         '"Go Extreme — Where Others Stop, You Start"\n'
         '극한의 환경에서도 멈추지 않는 크리에이터의 정신'),
        ("Key Visual", C["blue"],
         '• 액션: 폭포 앞, 사막, 설산에서 촬영 중인 크리에이터\n'
         '• 제품샷: IP68 물속 드롭 신 + 2000MB/s 속도 수치\n'
         '• Color tone: 강렬한 레드/블랙 + 자연 배경 대비'),
        ("Hero Message", C["orange"],
         '• "IP68. 2000MB/s. 4TB."\n'
         '• "Your adventure. Uninterrupted."\n'
         '• "Capture everything. Lose nothing."'),
        ("RTB (Reason to Believe)", C["green"],
         '• WD Lab 테스트 인증 (IP68, 2m 낙하, -20°C~60°C)\n'
         '• BiCS8 최신 NAND → 업계 최고 속도 클래스\n'
         '• 5년 보증 (Extreme Pro)'),
    ]
    for i, (title, col, body) in enumerate(elements):
        x = 0.3 + (i%2)*6.5
        y = 1.65 + (i//2)*2.4
        _rect(sl, x, y, 6.2, 0.5, col, title, fs=12, bold=True)
        _rect(sl, x, y+0.5, 6.2, 1.7, "1A2A3A", body, fs=9, align=PP_ALIGN.LEFT)

    # Channel mix
    _rect(sl, 0.3, 6.6, 12.7, 0.45, C["navy"],
          "채널: YouTube Pre-roll (30s) · Instagram Reels · TikTok · Amazon Sponsored · Tech YouTuber 협찬 · CES 미디어 브리핑",
          fs=9, bold=False)

    # ══════════════════════════════════════════════════════════
    # Slide 5: 캠페인 상세 — "Born to Create" (microSD)
    # ══════════════════════════════════════════════════════════
    sl = ns()
    _rect(sl, 0, 0, 13.33, 7.5, "0D1821")
    _rect(sl, 0, 0, 13.33, 0.9, C["blue"])
    _tbox(sl, 'Campaign: "Born to Create" — SanDisk Extreme/Pro Plus microSD', 0.3, 0.05,
          12, 0.75, fs=18, bold=True, color="FFFFFF")

    _rect(sl, 0.3, 1.0, 12.7, 0.5, "1A2235",
          "Consumer Insight: \"찰나의 순간을 놓치고 싶지 않다. 카메라가 버퍼링되는 순간, 그 장면은 영원히 사라진다.\"",
          fs=10, tc="CCDDFF")

    micro_elements = [
        ("BIG IDEA", C["blue"],
         '"Born to Create — Never Miss a Moment"\n'
         '모든 순간을 포착하는 크리에이터의 무기'),
        ("Target Segments", C["purple"],
         '① Action Camera (GoPro, DJI) — Extreme microSD A2 V30\n'
         '② Gaming Handhelds (Switch, Steam Deck) — Extreme/Ultra\n'
         '③ Drone / 4K Filming — Pro Plus 256GB~1TB\n'
         '④ Budget Android Phones (APAC) — Ultra 128/256GB'),
        ("Content Strategy", C["orange"],
         '• YouTube: 실제 크리에이터 협찬 영상 (언박싱 + 실사용)\n'
         '• Reddit/Discord: 게이밍 커뮤니티 리뷰 씨딩\n'
         '• Instagram: Before/After 포토 (버퍼링 vs 무버퍼링)\n'
         '• Amazon A+ Content: 속도 비교 인포그래픽'),
        ("Bundle Strategy", C["green"],
         '• Nintendo Switch 번들: Ultra 256GB + 케이스 할인\n'
         '• Steam Deck 번들: Extreme 512GB 최적화 마킹\n'
         '• GoPro 연계 프로모션: Extreme 공식 추천 카드\n'
         '• DJI Osmo 동봉 번들 (특정 리테일 전용)'),
    ]
    for i, (title, col, body) in enumerate(micro_elements):
        x = 0.3 + (i%2)*6.5
        y = 1.65 + (i//2)*2.4
        _rect(sl, x, y, 6.2, 0.5, col, title, fs=12, bold=True)
        _rect(sl, x, y+0.5, 6.2, 1.7, "1A2A3A", body, fs=9, align=PP_ALIGN.LEFT)

    _rect(sl, 0.3, 6.6, 12.7, 0.45, C["blue"],
          "KPI: Amazon BSR Top5 유지 · YouTube 누적 시청 500M · 리뷰 평점 4.6+ · 512GB/1TB 비중 35%↑",
          fs=9)

    # ══════════════════════════════════════════════════════════
    # Slide 6: 캠페인 상세 — Gaming / Internal SSD
    # ══════════════════════════════════════════════════════════
    sl = ns()
    _rect(sl, 0, 0, 13.33, 7.5, "080C14")
    _rect(sl, 0, 0, 13.33, 0.9, C["orange"])
    _tbox(sl, 'Campaign: "Speed Untamed" — WD_BLACK SN850X Gaming NVMe', 0.3, 0.05,
          12, 0.75, fs=18, bold=True, color="FFFFFF")

    _rect(sl, 0.3, 1.0, 12.7, 0.5, "1A2235",
          "Consumer Insight: \"게이머에게 로딩 시간은 적(敵)이다. 0.1초가 승패를 가른다. 내 PC/PS5는 최속이어야 한다.\"",
          fs=10, tc="FFDDAA")

    gaming_content = [
        ("Platform\nActivation", C["orange"],
         "PS5 공식 호환 인증 배지 활용\n"
         "Xbox Series X 최적화 마킹\n"
         "Steam Deck verified 캠페인\n"
         "→ 콘솔 커뮤니티 직접 공략"),
        ("Gaming\nInfluencer", C["blue"],
         "Tier-1 게이밍 유튜버 스폰서십\n"
         "(구독자 1M+ 3채널, 500K+ 5채널)\n"
         "실측 로딩 속도 비교 영상\n"
         "→ 경쟁사 대비 우위 데이터 시각화"),
        ("Esports\nSponsorship", C["red"],
         "LCK / LCS 공식 스폰서\n"
         "Pro 게이머 세팅 공개 협업\n"
         "대회장 배너 + 스트림 노출\n"
         "→ 게이밍 브랜드 이미지 강화"),
        ("Digital\nActivation", C["green"],
         "Reddit r/buildapc 네이티브 광고\n"
         "PCPartPicker 스폰서 리스팅\n"
         "YouTube Gaming Pre-roll\n"
         "→ 빌드 시즌(Q1·Q4) 집중 집행"),
    ]
    for i, (title, col, body) in enumerate(gaming_content):
        x = 0.3 + i*3.2
        _rect(sl, x, 1.7, 3.0, 0.75, col, title, fs=12, bold=True)
        _rect(sl, x, 2.45, 3.0, 2.3, "121C2E", body, fs=9, align=PP_ALIGN.LEFT)

    # KPI
    kpis = [
        ("Amazon BSR", "Gaming SSD Top3"),
        ("NPS Score",  "목표 72점"),
        ("Brand Recall","Awareness +8%p"),
        ("Revenue",    "SN850X $380M/yr"),
    ]
    for i, (k, v) in enumerate(kpis):
        _rect(sl, 0.3+i*3.2, 5.0, 3.0, 0.45, C["navy"], k, fs=9, bold=True)
        _rect(sl, 0.3+i*3.2, 5.45, 3.0, 0.6,  "D4001E", v, fs=14, bold=True)

    _rect(sl, 0.3, 6.2, 12.7, 0.45, C["orange"],
          "핵심 시즌: Q1(CES·빌드시즌), Q4(BF/CM·홀리데이 기프팅) 예산 집중 — 연간 총 $42M 집행",
          fs=9)

    # ══════════════════════════════════════════════════════════
    # Slide 7: 채널별 미디어 믹스 & 예산
    # ══════════════════════════════════════════════════════════
    sl = ns()
    _rect(sl, 0, 0, 13.33, 0.9, C["navy"])
    _tbox(sl, "미디어 믹스 & 마케팅 예산 배분 (2025년 기준)", 0.3, 0.05,
          12, 0.75, fs=22, bold=True, color="FFFFFF")

    channels = [
        ("Digital Paid\n(Search/Display)", 28, C["blue"],    "Amazon Sponsored, Google, Meta 광고"),
        ("Content/Influencer",             22, C["purple"],  "YouTube, TikTok, Instagram 인플루언서"),
        ("Retail/Trade",                   18, C["orange"],  "Best Buy, Amazon 프로모, In-store 디스플레이"),
        ("Brand/PR",                       12, C["red"],     "CES, Computex, Tech PR, Media Seeding"),
        ("Esports/Gaming",                 10, C["green"],   "LCK/LCS 스폰서, 게이밍 이벤트"),
        ("Social Media",                    7, C["teal"],    "Facebook, Instagram, Reddit 오가닉"),
        ("Research/Analytics",              3, C["gray"],    "시장 조사, 데이터 분석 툴"),
    ]

    total_budget = 458  # $M (전체 마케팅 예산)
    for i, (ch, pct, col, desc) in enumerate(channels):
        y = 1.1 + i * 0.85
        budget = total_budget * pct / 100
        bar_w = pct / 30 * 7.0
        _rect(sl, 0.3, y, 2.5, 0.7, col, ch, fs=9, bold=True)
        _rect(sl, 2.9, y+0.15, bar_w, 0.42, col)
        _tbox(sl, f"{pct}%  /  ${budget:.0f}M", 2.9+bar_w+0.1, y+0.15, 1.8, 0.42,
              fs=10, bold=True, color=col)
        _tbox(sl, desc, 5.5, y+0.2, 7.5, 0.42, fs=9, color="AAAAAA")

    _rect(sl, 0.3, 7.1, 12.7, 0.35, C["red"],
          f"총 마케팅 예산: ${total_budget}M  |  Revenue 대비 {total_budget/3527*100:.1f}%  |  "
          f"vs 업계 평균 12~15% 수준",
          fs=9)

    # ══════════════════════════════════════════════════════════
    # Slide 8: KPI 대시보드
    # ══════════════════════════════════════════════════════════
    sl = ns()
    _rect(sl, 0, 0, 13.33, 0.9, C["navy"])
    _tbox(sl, "MarCom KPI 대시보드 — 2025 목표", 0.3, 0.05,
          12, 0.75, fs=22, bold=True, color="FFFFFF")

    kpi_grid = [
        # (이름, 목표값, 단위, 색, 설명)
        ("Brand Awareness",    "65%",  "aided recall",    C["blue"],   "주요 3개국(NA/EMEA/JP) 소비자 대상"),
        ("YouTube Views",      "800M", "연간 누적",        C["red"],    "브랜드 채널 + 인플루언서 합산"),
        ("Amazon BSR Avg",     "Top 5", "카테고리 내",     C["orange"], "External·Internal·microSD 각 카테고리"),
        ("Product Review Avg", "4.7★",  "/5.0",           C["green"],  "Amazon 전체 SKU 평균 평점"),
        ("Influencer Reach",   "2.5B",  "impressions/yr", C["purple"], "Tier1+Tier2 인플루언서 합산"),
        ("Media Earned Value", "$120M", "EMV",            C["teal"],   "PR/세딩 via 미디어 노출 가치"),
        ("NPS Score",          "71",    "점",              C["blue"],   "Net Promoter Score — 제품 만족도"),
        ("SOM Growth",         "+2%p",  "점유율 증가",     C["red"],    "microSD/External SSD 합산 점유율"),
        ("ROAS",               "4.2x",  "평균",           C["orange"], "Paid digital 광고 투자대비 매출"),
    ]
    for i, (name, val, unit, col, desc) in enumerate(kpi_grid):
        x = 0.3 + (i%3)*4.3
        y = 1.1 + (i//3)*2.0
        _rect(sl, x, y, 4.0, 0.45, col, name, fs=11, bold=True)
        _rect(sl, x, y+0.45, 4.0, 0.8, C["navy"], f"{val} {unit}", fs=20, bold=True)
        _rect(sl, x, y+1.25, 4.0, 0.5, "1A2235", desc, fs=8, tc="AAAAAA")

    _rect(sl, 0.3, 7.1, 12.7, 0.35, C["navy"],
          "KPI 리뷰 주기: 주간(ROAS/BSR) · 월간(Awareness/NPS) · 분기(SOM/EMV) | 담당: MarCom 에이전트",
          fs=9)

    # ══════════════════════════════════════════════════════════
    # Slide 9: Awareness 아이템 아이디어
    # ══════════════════════════════════════════════════════════
    sl = ns()
    _rect(sl, 0, 0, 13.33, 0.9, C["navy"])
    _tbox(sl, "Awareness 아이템 & 크리에이티브 아이디어", 0.3, 0.05,
          12, 0.75, fs=22, bold=True, color="FFFFFF")

    ideas = [
        ("🎬 Hero Film 30s", C["red"],
         "캠페인: Go Extreme\n"
         "컨셉: 폭포·화산·설원에서 촬영하는 크리에이터\n"
         "채널: YouTube Pre-roll, TV (CES 주간)\n"
         "예산: $3.5M 제작 + $8M 집행"),
        ("📱 Short-form Series", C["blue"],
         "캠페인: Born to Create\n"
         "컨셉: 15초 TikTok/Reels — '버퍼링 없는 세상'\n"
         "채널: TikTok, Instagram Reels, YouTube Shorts\n"
         "예산: $800K 제작, 시리즈 12편"),
        ("🏆 Extreme Challenge", C["orange"],
         "캠페인: Go Extreme (UGC 이벤트)\n"
         "컨셉: 가장 극한의 환경 촬영 영상 공모전\n"
         "채널: Instagram #GoExtreme 해시태그\n"
         "상금: $50K + 제품 1년치 / 1만명 참여 목표"),
        ("🤝 Creator Partnership", C["purple"],
         "캠페인: Born to Create\n"
         "컨셉: MrBeast, MKBHD, Marques 등 협업\n"
         "형식: 24시간 제한 영상 제작 챌린지\n"
         "예산: $12M / 예상 조회수 300M+"),
        ("🎮 Gaming Integration", C["green"],
         "캠페인: Speed Untamed\n"
         "컨셉: 인기 게임 내 SSD 로딩 비교 인터랙션\n"
         "채널: Twitch 인터랙티브 스트림\n"
         "파트너: Ninja, Shroud 스폰서십"),
        ("📊 Benchmark Report", C["teal"],
         "캠페인: Always Ready (B2B PR용)\n"
         "컨셉: 3rd Party 독립 벤치마크 공식 발표\n"
         "채널: Tom's Hardware, AnandTech, Notebookcheck\n"
         "예산: $200K / 업계 신뢰도 강화"),
    ]
    for i, (title, col, body) in enumerate(ideas):
        x = 0.3 + (i%3)*4.3
        y = 1.1 + (i//3)*2.9
        _rect(sl, x, y, 4.0, 0.55, col, title, fs=11, bold=True)
        _rect(sl, x, y+0.55, 4.0, 2.1, "1A2235", body, fs=9, align=PP_ALIGN.LEFT)

    # ══════════════════════════════════════════════════════════
    # Slide 10: 분기별 실행 플랜
    # ══════════════════════════════════════════════════════════
    sl = ns()
    _rect(sl, 0, 0, 13.33, 0.9, C["navy"])
    _tbox(sl, "2025 분기별 MarCom 실행 플랜", 0.3, 0.05,
          12, 0.75, fs=22, bold=True, color="FFFFFF")

    quarters_plan = [
        ("Q1 2025\n(Jan~Mar)", C["blue"], [
            "CES 2025 미디어 브리핑 (Jan 7-10)",
            "SN850X 신규 4TB 런칭 캠페인",
            "YouTube Hero Film 집행 개시",
            "Amazon Sponsored 예산 +20%",
            "KPI: BSR Top3 내장SSD, Awareness +3%p",
        ]),
        ("Q2 2025\n(Apr~Jun)", C["orange"], [
            "Creator Season — microSD 캠페인 론칭",
            "Computex 제품 프리뷰 PR",
            "Born to Create TikTok 시리즈 시작",
            "GoPro/DJI 번들 파트너십 집행",
            "KPI: microSD 점유율 +1%p, EMV $25M",
        ]),
        ("Q3 2025\n(Jul~Sep)", C["green"], [
            "Back-to-School 패키지 프로모",
            "Gaming 캠페인 — LCK 스폰서십 집행",
            "Steam Deck 512GB Extreme 번들",
            "Extreme Challenge UGC 이벤트",
            "KPI: External SSD Amazon BSR Top5",
        ]),
        ("Q4 2025\n(Oct~Dec)", C["red"], [
            "Holiday Campaign 총력전 (Nov BF/CM)",
            "Go Extreme 히어로 필름 TV 집행",
            "MKBHD/MrBeast 콜라보 영상 집행",
            "Amazon 독점 한정판 패키지 출시",
            "KPI: Revenue $1,039M, NPS 71점 달성",
        ]),
    ]
    for i, (qname, col, actions) in enumerate(quarters_plan):
        x = 0.3 + i*3.25
        _rect(sl, x, 1.1, 3.1, 0.65, col, qname, fs=12, bold=True)
        for j, action in enumerate(actions):
            bg = "1A2235" if j % 2 == 0 else "15202E"
            _rect(sl, x, 1.75+j*1.0, 3.1, 0.9, bg, f"• {action}", fs=8.5, align=PP_ALIGN.LEFT)

    _rect(sl, 0.3, 7.1, 12.7, 0.35, C["red"],
          "총 연간 마케팅 예산: $458M  |  예산 집중: Q4 35% > Q1 27% > Q3 22% > Q2 16%",
          fs=9)

    out_path = PPTX_DIR / "marcom_plan.pptx"
    prs.save(out_path)
    print(f"[마케팅 커뮤니케이션 에이전트] ✅ PPTX 저장: {out_path}")
    return out_path


if __name__ == "__main__":
    build_pptx()
