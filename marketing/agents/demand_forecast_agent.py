"""
수요예측 에이전트 (Demand Forecast Agent)
역할: TAM/SAM/SOM 분석, 단기·중장기 분기별 수요 예측
산출물: Excel — demand_forecast.xlsx  +  PowerPoint — demand_forecast.pptx
"""
import sys
sys.path.insert(0, str(__import__('pathlib').Path(__file__).parent.parent))

import json
import datetime
import openpyxl
from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
from openpyxl.chart import BarChart, LineChart, Reference
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from config import QUARTERS, EXCEL_DIR, PPTX_DIR, DATA_DIR

# ── 시장 데이터 로드 ─────────────────────────────────────────
with open(DATA_DIR / "market_data.json") as f:
    MKT = json.load(f)

# ── 색상 ─────────────────────────────────────────────────────
C = {
    "navy":      "1B2A4A",
    "red":       "D4001E",
    "blue":      "4472C4",
    "orange":    "ED7D31",
    "green":     "70AD47",
    "yellow":    "FFC000",
    "gray_bg":   "F2F2F2",
    "white":     "FFFFFF",
    "green_ok":  "C6EFCE",
    "yellow_w":  "FFEB9C",
    "red_w":     "FFC7CE",
}

def _fill(h): return PatternFill("solid", fgColor=h)
def _font(bold=False, color="000000", size=10):
    return Font(bold=bold, color=color, size=size, name="Calibri")
def _border():
    s = Side(style="thin", color="BFBFBF")
    return Border(left=s, right=s, top=s, bottom=s)
def _center(): return Alignment(horizontal="center", vertical="center", wrap_text=True)

# ── TAM/SAM/SOM 모델 ─────────────────────────────────────────
CATEGORIES = {
    "external_ssd": {
        "tam_2024": MKT["tam_usd_b"]["external_ssd_2024"],
        "tam_2025": MKT["tam_usd_b"]["external_ssd_2025"],
        "cagr":     MKT["tam_usd_b"]["cagr_external_ssd"],
        "sam_ratio": 0.65,   # B2C retail/e-commerce only
        "som_share": 0.22,   # SanDisk/WD current share
        "som_target_2026": 0.24,
    },
    "internal_ssd": {
        "tam_2024": MKT["tam_usd_b"]["internal_ssd_consumer_2024"],
        "tam_2025": MKT["tam_usd_b"]["internal_ssd_consumer_2025"],
        "cagr":     MKT["tam_usd_b"]["cagr_internal_ssd_consumer"],
        "sam_ratio": 0.60,
        "som_share": 0.17,
        "som_target_2026": 0.19,
    },
    "microsd": {
        "tam_2024": MKT["tam_usd_b"]["microsd_2024"],
        "tam_2025": MKT["tam_usd_b"]["microsd_2025"],
        "cagr":     MKT["tam_usd_b"]["cagr_microsd"],
        "sam_ratio": 0.75,   # retail-accessible addressable
        "som_share": 0.32,   # SanDisk #1 in microSD
        "som_target_2026": 0.34,
    },
}

# 분기별 계절 지수
SEASONAL_IDX = {
    "Q1": 0.85, "Q2": 0.95, "Q3": 1.10, "Q4": 1.30
}

def _quarterly_forecast(cat_key, years=2):
    """분기별 수요 예측 (SOM, $M)"""
    cat = CATEGORIES[cat_key]
    result = {}
    for yr_offset in range(years):
        year = 2025 + yr_offset
        annual_tam = cat["tam_2025"] * (1 + cat["cagr"]) ** yr_offset
        annual_sam = annual_tam * cat["sam_ratio"]
        annual_som = annual_sam * (cat["som_share"] + yr_offset * (cat["som_target_2026"] - cat["som_share"]))
        for q in range(1, 5):
            q_key = f"{year}Q{q}"
            s_idx = SEASONAL_IDX[f"Q{q}"]
            result[q_key] = round(annual_som / 4 * s_idx * (1 + 0.005 * q), 1)
    return result


# ═══════════════════════════════════════════════════════════════
# EXCEL 빌드
# ═══════════════════════════════════════════════════════════════
def build_excel():
    wb = openpyxl.Workbook()

    # ── 시트 1: TAM/SAM/SOM Summary ──────────────────────────
    ws = wb.active
    ws.title = "📊 TAM-SAM-SOM"
    ws.column_dimensions["A"].width = 22
    for col in "BCDEFGHI":
        ws.column_dimensions[col].width = 16

    ws.merge_cells("A1:I1")
    c = ws["A1"]
    c.value = "SanDisk B2C Storage — TAM / SAM / SOM 분석 (단위: $B)"
    c.fill = _fill(C["navy"])
    c.font = Font(bold=True, color="FFFFFF", size=13, name="Calibri")
    c.alignment = _center()
    ws.row_dimensions[1].height = 32

    ws.merge_cells("A2:I2")
    ws["A2"].value = f"수요예측 에이전트 | 기준: 2025Q1 | 출처: TrendForce/IDC 추정 + 내부 시뮬레이션"
    ws["A2"].fill = _fill(C["red"])
    ws["A2"].font = Font(color="FFFFFF", size=9, name="Calibri")
    ws["A2"].alignment = _center()

    headers = ["카테고리", "지표", "2023A", "2024E", "2025E", "2026E", "2027E", "CAGR('24-27)", "비고"]
    ws.append(headers)
    rh = ws.max_row
    for ci, h in enumerate(headers, 1):
        c = ws.cell(row=rh, column=ci)
        c.fill = _fill(C["navy"])
        c.font = _font(True, "FFFFFF", 9)
        c.alignment = _center()
        c.border = _border()

    cat_label = {"external_ssd": "External SSD", "internal_ssd": "Internal SSD (Consumer)", "microsd": "microSD"}
    cat_colors = {"external_ssd": C["blue"], "internal_ssd": C["orange"], "microsd": C["green"]}

    for cat_key, cat in CATEGORIES.items():
        cagr = cat["cagr"]
        for metric, ratio in [("TAM", 1.0), ("SAM", cat["sam_ratio"]), ("SOM", cat["sam_ratio"] * cat["som_share"])]:
            t23 = cat["tam_2024"] / (1 + cagr) * ratio
            t24 = cat["tam_2024"] * ratio
            t25 = cat["tam_2025"] * ratio
            t26 = t25 * (1 + cagr)
            t27 = t26 * (1 + cagr)
            row = [
                cat_label[cat_key] if metric == "TAM" else "",
                metric,
                round(t23, 2), round(t24, 2), round(t25, 2), round(t26, 2), round(t27, 2),
                f"{cagr*100:.0f}%",
                "SAM=B2C retail" if metric == "SAM" else
                f"Share: {cat['som_share']*100:.0f}%" if metric == "SOM" else "",
            ]
            ws.append(row)
            r = ws.max_row
            for ci, val in enumerate(row, 1):
                c = ws.cell(row=r, column=ci)
                c.border = _border()
                c.alignment = _center()
                c.font = _font(size=9)
                if ci == 1 and metric == "TAM":
                    c.fill = _fill(cat_colors[cat_key])
                    c.font = _font(True, "FFFFFF", 9)
                elif ci == 2:
                    bg = {"TAM": C["gray_bg"], "SAM": "DDE8F0", "SOM": "D5E8D4"}.get(metric, C["white"])
                    c.fill = _fill(bg)
                    c.font = _font(True, size=9)
                elif ci in range(3, 8):
                    c.number_format = '"$"0.00 "B"'

    # ── 시트 2: 분기별 수요 예측 ─────────────────────────────
    ws2 = wb.create_sheet("📈 분기별 예측")
    ws2.column_dimensions["A"].width = 24
    ws2.column_dimensions["B"].width = 14
    for col in "CDEFGHI":
        ws2.column_dimensions[col].width = 14

    ws2.merge_cells("A1:I1")
    c = ws2["A1"]
    c.value = "분기별 SOM 수요 예측 (단위: $M) — 2025~2026"
    c.fill = _fill(C["navy"])
    c.font = Font(bold=True, color="FFFFFF", size=12, name="Calibri")
    c.alignment = _center()
    ws2.row_dimensions[1].height = 30

    q_keys = ["2025Q1","2025Q2","2025Q3","2025Q4","2026Q1","2026Q2","2026Q3","2026Q4"]
    h2 = ["카테고리", "연간합계"] + q_keys
    ws2.append(h2)
    rh = ws2.max_row
    for ci, h in enumerate(h2, 1):
        c = ws2.cell(row=rh, column=ci)
        c.fill = _fill(C["red"])
        c.font = _font(True, "FFFFFF", 9)
        c.alignment = _center()
        c.border = _border()

    chart_data_rows = {}
    for cat_key in CATEGORIES:
        fc = _quarterly_forecast(cat_key, years=2)
        annual = sum(v for k, v in fc.items() if "2025" in k)
        row_data = [cat_label[cat_key], round(annual, 1)] + [fc.get(q, 0) for q in q_keys]
        ws2.append(row_data)
        r = ws2.max_row
        chart_data_rows[cat_key] = r
        for ci, val in enumerate(row_data, 1):
            c = ws2.cell(row=r, column=ci)
            c.border = _border()
            c.alignment = _center()
            c.font = _font(size=9)
            if ci == 1:
                c.fill = _fill(cat_colors[cat_key])
                c.font = _font(True, "FFFFFF", 9)
            elif ci >= 2:
                c.fill = _fill(C["gray_bg"])
                c.number_format = '"$"#,##0.0 "M"'
                if ci >= 3:
                    q_label = q_keys[ci-3]
                    s = SEASONAL_IDX.get(f"Q{q_label[-1]}", 1.0)
                    if s >= 1.25:
                        c.fill = _fill(C["red_w"])
                    elif s >= 1.05:
                        c.fill = _fill(C["yellow_w"])
                    else:
                        c.fill = _fill(C["green_ok"])

    # 분기별 예측 라인 차트
    ws2.append([])
    ws2.append(["📝 범례: 빨강=Q4 성수기, 노랑=Q3, 초록=Q1~Q2 비수기"])

    from openpyxl.chart.series import SeriesLabel
    chart = LineChart()
    chart.title = "분기별 SOM 예측 ($M)"
    chart.style = 10
    chart.y_axis.title = "Revenue ($M)"
    chart.x_axis.title = "Quarter"
    chart.width = 22
    chart.height = 12

    for i, (cat_key, row_num) in enumerate(chart_data_rows.items()):
        data_ref = Reference(ws2, min_col=3, max_col=10, min_row=row_num, max_row=row_num)
        chart.add_data(data_ref)
        chart.series[i].title = SeriesLabel(v=cat_label[cat_key])

    cats_ref = Reference(ws2, min_col=3, max_col=10, min_row=2, max_row=2)
    chart.set_categories(cats_ref)
    ws2.add_chart(chart, "A12")

    # ── 시트 3: 경쟁사 비교 ──────────────────────────────────
    ws3 = wb.create_sheet("🏆 경쟁사 분석")
    ws3.column_dimensions["A"].width = 22
    ws3.column_dimensions["B"].width = 14
    ws3.column_dimensions["C"].width = 14
    ws3.column_dimensions["D"].width = 14
    ws3.column_dimensions["E"].width = 28

    ws3.merge_cells("A1:E1")
    c = ws3["A1"]
    c.value = "경쟁사 시장 점유율 현황 (2024E) — B2C Storage"
    c.fill = _fill(C["navy"])
    c.font = Font(bold=True, color="FFFFFF", size=12, name="Calibri")
    c.alignment = _center()
    ws3.row_dimensions[1].height = 30

    comp_headers = ["경쟁사", "External SSD", "Internal SSD(C)", "microSD", "전략적 포지션"]
    ws3.append(comp_headers)
    rh = ws3.max_row
    for ci, h in enumerate(comp_headers, 1):
        c = ws3.cell(row=rh, column=ci)
        c.fill = _fill(C["red"])
        c.font = _font(True, "FFFFFF", 9)
        c.alignment = _center()
        c.border = _border()

    comp_data = [
        ["SanDisk/WD ★",  22, 17, 32, "Vertically integrated NAND, #1 microSD, strong value brand"],
        ["Samsung",        33, 28, 23, "Premium leader, Galaxy ecosystem, T9/990Pro/EvoPlus"],
        ["Seagate",        13,  6,  0, "HDD heritage, creator/photographer external SSD"],
        ["SK Hynix",        0, 12,  0, "DRAM+NAND, Platinum P41 gaming internal"],
        ["Kingston",        3, 11,  9, "Value champion, DIY/channel focused"],
        ["Crucial (Micron)", 7, 11,  0, "Price/performance, mainstream consumer"],
        ["Lexar",           4,  4, 11, "Aggressive pricing, gaming/photography microSD"],
        ["Others",         18, 11, 25, "White-label, regional brands"],
    ]

    for row in comp_data:
        ws3.append(row)
        r = ws3.max_row
        for ci, val in enumerate(row, 1):
            c = ws3.cell(row=r, column=ci)
            c.border = _border()
            c.alignment = _center()
            c.font = _font(size=9)
            if ci == 1:
                if "SanDisk" in str(val):
                    c.fill = _fill(C["red"])
                    c.font = _font(True, "FFFFFF", 9)
                else:
                    c.fill = _fill(C["gray_bg"])
                    c.font = _font(True, size=9)
            elif 2 <= ci <= 4:
                share = val
                if share == 0:
                    c.fill = _fill("EEEEEE")
                elif share >= 25:
                    c.fill = _fill(C["red_w"])
                    c.font = _font(True, size=9)
                elif share >= 15:
                    c.fill = _fill(C["yellow_w"])
                else:
                    c.fill = _fill(C["green_ok"])
                if share > 0:
                    c.value = f"{share}%"
                    c.number_format = "@"
            elif ci == 5:
                c.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
        ws3.row_dimensions[r].height = 25

    # ── 시트 4: 성장 동인 / 리스크 ───────────────────────────
    ws4 = wb.create_sheet("📋 성장동인·리스크")
    ws4.column_dimensions["A"].width = 24
    ws4.column_dimensions["B"].width = 40
    ws4.column_dimensions["C"].width = 40

    ws4.merge_cells("A1:C1")
    c = ws4["A1"]
    c.value = "카테고리별 성장 동인 & 핵심 리스크"
    c.fill = _fill(C["navy"])
    c.font = Font(bold=True, color="FFFFFF", size=12, name="Calibri")
    c.alignment = _center()
    ws4.row_dimensions[1].height = 30

    ws4.append(["카테고리", "✅ 성장 동인", "⚠️ 핵심 리스크"])
    rh = ws4.max_row
    for ci in range(1, 4):
        c = ws4.cell(row=rh, column=ci)
        c.fill = _fill(C["red"])
        c.font = _font(True, "FFFFFF", 9)
        c.alignment = _center()
        c.border = _border()

    for cat_key in CATEGORIES:
        drivers = " / ".join(MKT["growth_drivers"][cat_key])
        risks   = " / ".join(MKT["key_risks"][cat_key])
        ws4.append([cat_label[cat_key], drivers, risks])
        r = ws4.max_row
        for ci, val in enumerate([cat_label[cat_key], drivers, risks], 1):
            c = ws4.cell(row=r, column=ci)
            c.border = _border()
            c.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
            c.font = _font(size=9)
            if ci == 1:
                c.fill = _fill(cat_colors[cat_key])
                c.font = _font(True, "FFFFFF", 9)
                c.alignment = _center()
        ws4.row_dimensions[r].height = 40

    out_path = EXCEL_DIR / "demand_forecast.xlsx"
    wb.save(out_path)
    print(f"[수요예측 에이전트] ✅ Excel 저장: {out_path}")
    return out_path


# ═══════════════════════════════════════════════════════════════
# POWERPOINT 빌드
# ═══════════════════════════════════════════════════════════════
def _rgb(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def _add_textbox(slide, text, left, top, width, height,
                 font_size=12, bold=False, color="000000",
                 bg_color=None, align=PP_ALIGN.LEFT, wrap=True):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    if bg_color:
        from pptx.oxml.ns import qn
        from lxml import etree
        sp = txBox._element
        spPr = sp.find(qn("p:spPr"))
        if spPr is None:
            spPr = etree.SubElement(sp, qn("p:spPr"))
        solidFill = etree.SubElement(spPr, qn("a:solidFill"))
        srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgb.set("val", bg_color)
    return txBox

def _add_rect(slide, left, top, width, height, fill_color, text="", font_size=11, bold=False, text_color="FFFFFF", align=PP_ALIGN.CENTER):
    from pptx.util import Inches, Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_color)
    shape.line.color.rgb = _rgb(fill_color)
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(text_color)
    return shape

def build_pptx():
    from pptx.util import Inches, Pt, Emu
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    def new_slide():
        return prs.slides.add_slide(blank_layout)

    # ── 슬라이드 1: Title ─────────────────────────────────────
    sl = new_slide()
    _add_rect(sl, 0, 0, 13.33, 7.5, C["navy"])
    _add_rect(sl, 0, 2.5, 13.33, 2.5, C["red"])
    _add_textbox(sl, "SanDisk B2C Storage", 1, 1.0, 11, 1.2,
                 font_size=32, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)
    _add_textbox(sl, "수요예측 보고서  |  Demand Forecast Report", 1, 2.7, 11, 0.8,
                 font_size=22, bold=False, color="FFFFFF", align=PP_ALIGN.CENTER)
    _add_textbox(sl, f"수요예측 에이전트  |  {datetime.date.today()}  |  CONFIDENTIAL", 1, 3.6, 11, 0.5,
                 font_size=12, color="FFCCCC", align=PP_ALIGN.CENTER)
    # Product category chips
    for i, (label, color) in enumerate([("External SSD", C["blue"]), ("Internal SSD", C["orange"]), ("microSD", C["green"])]):
        _add_rect(sl, 3.2 + i * 2.5, 4.5, 2.0, 0.5, color, label, font_size=13, bold=True)
    _add_textbox(sl, "NAND: BiCS5 | BiCS6 | BiCS8  ·  시장: NA / EMEA / APAC / Japan / China",
                 1, 5.5, 11, 0.5, font_size=11, color="AACCFF", align=PP_ALIGN.CENTER)

    # ── 슬라이드 2: TAM/SAM/SOM Overview ─────────────────────
    sl = new_slide()
    _add_rect(sl, 0, 0, 13.33, 0.9, C["navy"])
    _add_textbox(sl, "TAM / SAM / SOM — 2025 Overview", 0.2, 0.05, 10, 0.75,
                 font_size=22, bold=True, color="FFFFFF")

    tam_data = [
        ("External SSD", 4.0,  2.6,  0.57, C["blue"]),
        ("Internal SSD", 11.0, 6.6,  1.12, C["orange"]),
        ("microSD",       6.8, 5.1,  1.63, C["green"]),
    ]
    for i, (label, tam, sam, som, color) in enumerate(tam_data):
        x = 0.5 + i * 4.2
        _add_rect(sl, x, 1.1, 3.8, 0.55, color, label, font_size=14, bold=True)
        for j, (metric, val, bg) in enumerate([
            ("TAM (Total Addressable)", f"${tam}B", "2C3E5A"),
            ("SAM (Serviceable)",       f"${sam}B", "3A5068"),
            ("SOM (Our Revenue)",       f"${som}B", C["red"]),
        ]):
            _add_rect(sl, x, 1.75 + j*1.4, 3.8, 0.5, bg, metric, font_size=9, text_color="CCDDEE")
            _add_rect(sl, x, 2.25 + j*1.4, 3.8, 0.8, "1B2A4A" if j < 2 else C["red"],
                      val, font_size=24, bold=True, text_color="FFFFFF")

    _add_textbox(sl, "※ SAM = B2C retail/e-comm 채널 기준  |  SOM = SanDisk/WD 추정 매출  |  Source: TrendForce/IDC 추정",
                 0.3, 6.8, 12, 0.5, font_size=9, color="888888")

    # ── 슬라이드 3: 분기별 수요 예측 ─────────────────────────
    sl = new_slide()
    _add_rect(sl, 0, 0, 13.33, 0.9, C["navy"])
    _add_textbox(sl, "분기별 SOM 수요 예측 ($M) — 2025~2026", 0.2, 0.05, 10, 0.75,
                 font_size=22, bold=True, color="FFFFFF")

    q_labels = ["25Q1","25Q2","25Q3","25Q4","26Q1","26Q2","26Q3","26Q4"]
    cat_colors_list = [C["blue"], C["orange"], C["green"]]
    q_vals = {k: _quarterly_forecast(k, 2) for k in CATEGORIES}

    max_val = max(v for fc in q_vals.values() for v in fc.values())

    bar_w = 0.35
    gap   = 0.15
    group_w = 3 * bar_w + 2 * 0.05 + gap
    chart_left = 0.8
    chart_bot  = 7.0
    chart_h    = 4.8

    # Y축
    _add_textbox(sl, "$M", 0.3, 1.0, 0.4, 0.4, font_size=9, color="888888")
    for tick in [0, 50, 100, 150, 200]:
        y_pos = chart_bot - (tick / max_val) * chart_h - 0.1
        _add_textbox(sl, str(tick), 0.3, y_pos, 0.45, 0.3, font_size=8, color="555555", align=PP_ALIGN.RIGHT)

    for qi, q_key_full in enumerate(["2025Q1","2025Q2","2025Q3","2025Q4","2026Q1","2026Q2","2026Q3","2026Q4"]):
        x_base = chart_left + qi * group_w
        q_num = q_key_full[-1]
        s_idx = SEASONAL_IDX[f"Q{q_num}"]
        bg_col = C["red_w"] if s_idx >= 1.25 else (C["yellow_w"] if s_idx >= 1.05 else C["green_ok"])
        # Q라벨
        _add_rect(sl, x_base, chart_bot + 0.05, group_w - gap, 0.3, C["navy"],
                  q_labels[qi], font_size=8, bold=False)

        for ci, (cat_key, bar_color) in enumerate(zip(CATEGORIES.keys(), cat_colors_list)):
            val = q_vals[cat_key].get(q_key_full, 0)
            bar_h = (val / max_val) * chart_h
            bar_x = x_base + ci * (bar_w + 0.05)
            bar_y = chart_bot - bar_h
            _add_rect(sl, bar_x, bar_y, bar_w, bar_h, bar_color)
            if val > 30:
                _add_textbox(sl, f"{val:.0f}", bar_x, bar_y - 0.25, bar_w, 0.25,
                             font_size=7, color="333333", align=PP_ALIGN.CENTER)

    # 범례
    for i, (label, color) in enumerate([("External SSD", C["blue"]), ("Internal SSD", C["orange"]), ("microSD", C["green"])]):
        _add_rect(sl, 10.2, 1.2 + i * 0.45, 0.3, 0.3, color)
        _add_textbox(sl, label, 10.6, 1.2 + i * 0.45, 2.0, 0.3, font_size=10, color="222222")
    _add_textbox(sl, "🔴 Q4 성수기 (idx 1.30)  🟡 Q3 (1.10)  🟢 Q1-Q2 비수기",
                 0.3, 7.15, 12, 0.3, font_size=9, color="555555")

    # ── 슬라이드 4: 경쟁사 점유율 ────────────────────────────
    sl = new_slide()
    _add_rect(sl, 0, 0, 13.33, 0.9, C["navy"])
    _add_textbox(sl, "경쟁사 시장 점유율 현황 (2024E)", 0.2, 0.05, 10, 0.75,
                 font_size=22, bold=True, color="FFFFFF")

    comp = MKT["market_share_pct"]
    cats_ppt = [("external_ssd", "External SSD", C["blue"]),
                ("internal_ssd_consumer", "Internal SSD", C["orange"]),
                ("microsd", "microSD", C["green"])]

    for ci, (cat_key, cat_name, header_color) in enumerate(cats_ppt):
        x = 0.4 + ci * 4.3
        _add_rect(sl, x, 1.1, 3.9, 0.5, header_color, cat_name, font_size=14, bold=True)
        sorted_comp = sorted(comp[cat_key].items(), key=lambda x: -x[1])
        for ri, (brand, share) in enumerate(sorted_comp[:7]):
            bar_w_c = share / 100 * 3.5
            is_sd = "SanDisk" in brand or "WD" in brand
            bar_col = C["red"] if is_sd else (C["navy"] if share >= 20 else "607080")
            _add_rect(sl, x, 1.75 + ri * 0.7, bar_w_c + 0.3, 0.5, bar_col,
                      f"{brand}  {share}%", font_size=9, bold=is_sd)

    _add_textbox(sl, "Source: IDC/TrendForce 추정 2024 | ★ SanDisk 강점: microSD #1 (32%), External SSD #2 (22%), 수직계열화 NAND 경쟁력",
                 0.3, 7.1, 12.5, 0.35, font_size=9, color="555555")

    # ── 슬라이드 5: 성장 동인 & 리스크 ──────────────────────
    sl = new_slide()
    _add_rect(sl, 0, 0, 13.33, 0.9, C["navy"])
    _add_textbox(sl, "성장 동인 & 핵심 리스크 — 카테고리별", 0.2, 0.05, 10, 0.75,
                 font_size=22, bold=True, color="FFFFFF")

    growth_cats = [
        ("External SSD", C["blue"],   MKT["growth_drivers"]["external_ssd"], MKT["key_risks"]["external_ssd"]),
        ("Internal SSD", C["orange"], MKT["growth_drivers"]["internal_ssd"],  MKT["key_risks"]["internal_ssd"]),
        ("microSD",      C["green"],  MKT["growth_drivers"]["microsd"],       MKT["key_risks"]["microsd"]),
    ]
    for i, (name, color, drivers, risks) in enumerate(growth_cats):
        y = 1.1 + i * 2.0
        _add_rect(sl, 0.3, y, 1.5, 1.7, color, name, font_size=13, bold=True)
        _add_rect(sl, 2.0, y, 5.3, 1.7, "1E4D2B",
                  "✅ " + "  /  ".join(drivers), font_size=10, bold=False, align=PP_ALIGN.LEFT)
        _add_rect(sl, 7.5, y, 5.5, 1.7, "4D1E1E",
                  "⚠️ " + "  /  ".join(risks), font_size=10, bold=False, align=PP_ALIGN.LEFT)

    _add_textbox(sl, "▶ 기회 (초록)                                             ▶ 위협 (빨강)",
                 2.0, 7.05, 11, 0.35, font_size=10, color="888888")

    # ── 슬라이드 6: 중장기 전망 ───────────────────────────────
    sl = new_slide()
    _add_rect(sl, 0, 0, 13.33, 0.9, C["navy"])
    _add_textbox(sl, "중장기 수요 전망 — 2025~2027 Revenue Roadmap", 0.2, 0.05, 12, 0.75,
                 font_size=22, bold=True, color="FFFFFF")

    roadmap = [
        ("2025",  0.57+1.12+1.63, "현재 포지션 유지\n- microSD #1 수성\n- External SSD 성장 가속", C["blue"]),
        ("2026",  0.62+1.32+1.79, "점유율 확대 목표\n- External SSD 24%→목표\n- BiCS8 신제품 출시", C["orange"]),
        ("2027",  0.68+1.56+1.97, "Premium 전환 완료\n- BiCS8 라인업 주력화\n- EMEA/APAC 성장 집중", C["green"]),
    ]
    prev_h = 0
    for i, (year, total_som, strategy, color) in enumerate(roadmap):
        x = 2.0 + i * 3.5
        bar_h = total_som / 4.5 * 4.0
        bar_y = 6.0 - bar_h
        _add_rect(sl, x, bar_y, 2.8, bar_h, color)
        _add_textbox(sl, f"${total_som:.2f}B", x, bar_y - 0.4, 2.8, 0.4,
                     font_size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
        _add_textbox(sl, f"{year}E", x, 6.1, 2.8, 0.4, font_size=14, bold=True,
                     color="FFFFFF", align=PP_ALIGN.CENTER)
        _add_textbox(sl, strategy, x, 6.55, 2.8, 0.9, font_size=9, color="AAAAAA", align=PP_ALIGN.CENTER)
        if i > 0:
            prev_som = roadmap[i-1][1]
            growth = (total_som - prev_som) / prev_som * 100
            _add_rect(sl, x - 0.5, 4.0, 0.8, 0.4, C["red"], f"+{growth:.0f}%", font_size=10, bold=True)

    _add_textbox(sl, "SOM Total (External + Internal + microSD)",
                 3.5, 1.1, 7, 0.4, font_size=11, color="AAAAAA", align=PP_ALIGN.CENTER)

    out_path = PPTX_DIR / "demand_forecast.pptx"
    prs.save(out_path)
    print(f"[수요예측 에이전트] ✅ PPTX 저장: {out_path}")
    return out_path


if __name__ == "__main__":
    build_excel()
    build_pptx()
