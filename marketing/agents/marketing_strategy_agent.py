"""
마케팅 전략 에이전트 (Marketing Strategy Agent)
역할: 제품 Mix 전략, 가격 전략, P&L 시뮬레이션, 포트폴리오 로드맵
산출물: Excel — marketing_strategy.xlsx  +  PowerPoint — marketing_strategy.pptx
"""
import sys
sys.path.insert(0, str(__import__('pathlib').Path(__file__).parent.parent))

import json
import datetime
import openpyxl
from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from config import PRODUCTS, NAND_COST_PER_GB, QUARTERS, EXCEL_DIR, PPTX_DIR, DATA_DIR

with open(DATA_DIR / "market_data.json") as f:
    MKT = json.load(f)

C = {
    "navy":   "1B2A4A", "red":    "D4001E",
    "blue":   "4472C4", "orange": "ED7D31",
    "green":  "70AD47", "yellow": "FFC000",
    "gray":   "F2F2F2", "white":  "FFFFFF",
    "g_ok":  "C6EFCE", "y_warn": "FFEB9C", "r_bad": "FFC7CE",
}

def _fill(h): return PatternFill("solid", fgColor=h)
def _font(bold=False, color="000000", size=10):
    return Font(bold=bold, color=color, size=size, name="Calibri")
def _border():
    s = Side(style="thin", color="BFBFBF")
    return Border(left=s, right=s, top=s, bottom=s)
def _center(): return Alignment(horizontal="center", vertical="center", wrap_text=True)
def _rgb(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

def _rect(slide, left, top, w, h, fill, text="", fs=11, bold=False, tc="FFFFFF", align=PP_ALIGN.CENTER):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(fill)
    if text:
        tf = shape.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = text
        run.font.size = Pt(fs); run.font.bold = bold
        run.font.color.rgb = _rgb(tc)
    return shape

def _tbox(slide, text, left, top, w, h, fs=11, bold=False, color="000000", align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(fs); run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    return tb

ASP = {
    "WD_BLACK_SN850X_1TB":80,  "WD_BLACK_SN850X_2TB":145, "WD_BLACK_SN850X_4TB":275,
    "WD_BLACK_SN770_1TB":65,   "WD_BLACK_SN770_2TB":115,
    "SD_E30_1TB":60,           "SD_E30_2TB":105,
    "SD_EXTREME_1TB":90,       "SD_EXTREME_2TB":160,      "SD_EXTREME_4TB":310,
    "SD_EXTREME_PRO_1TB":120,  "SD_EXTREME_PRO_2TB":210,
    "WD_MY_PASSPORT_1TB":75,   "WD_MY_PASSPORT_2TB":130,
    "SD_EXTREME_MICRO_128G":18,"SD_EXTREME_MICRO_256G":28,"SD_EXTREME_MICRO_512G":50,
    "SD_EXTREME_MICRO_1TB":95,
    "SD_ULTRA_MICRO_128G":12,  "SD_ULTRA_MICRO_256G":20, "SD_ULTRA_MICRO_512G":38,
    "SD_PRO_PLUS_MICRO_256G":35,"SD_PRO_PLUS_MICRO_512G":65,
}

UNITS_Q_K = {
    "WD_BLACK_SN850X_1TB":950, "WD_BLACK_SN850X_2TB":700, "WD_BLACK_SN850X_4TB":200,
    "WD_BLACK_SN770_1TB":600,  "WD_BLACK_SN770_2TB":400,
    "SD_E30_1TB":450,          "SD_E30_2TB":280,
    "SD_EXTREME_1TB":1200,     "SD_EXTREME_2TB":800,      "SD_EXTREME_4TB":250,
    "SD_EXTREME_PRO_1TB":500,  "SD_EXTREME_PRO_2TB":350,
    "WD_MY_PASSPORT_1TB":900,  "WD_MY_PASSPORT_2TB":550,
    "SD_EXTREME_MICRO_128G":3000,"SD_EXTREME_MICRO_256G":4500,"SD_EXTREME_MICRO_512G":2800,
    "SD_EXTREME_MICRO_1TB":1200,
    "SD_ULTRA_MICRO_128G":5000,"SD_ULTRA_MICRO_256G":4000,"SD_ULTRA_MICRO_512G":2200,
    "SD_PRO_PLUS_MICRO_256G":1500,"SD_PRO_PLUS_MICRO_512G":1000,
}

def _bom_cost(sku):
    info = PRODUCTS[sku]
    nand = NAND_COST_PER_GB[info["nand"]] * info["capacity_gb"] * 1.12
    housing = 9.5 if info["cat"] == "external_ssd" else 4.0 if info["cat"] == "internal_ssd" else 1.5
    return round(nand + housing, 2)


def build_excel():
    wb = openpyxl.Workbook()

    # ── Sheet 1: Product Mix P&L ─────────────────────────────
    ws = wb.active
    ws.title = "💰 Product Mix P&L"
    ws.column_dimensions["A"].width = 26
    ws.column_dimensions["B"].width = 12
    for col in "CDEFGHIJ":
        ws.column_dimensions[col].width = 13

    ws.merge_cells("A1:J1")
    c = ws["A1"]
    c.value = "SanDisk B2C — 제품 Mix P&L 시뮬레이션 (2025Q1 기준)"
    c.fill = _fill(C["navy"]); c.font = Font(bold=True, color="FFFFFF", size=13, name="Calibri")
    c.alignment = _center(); ws.row_dimensions[1].height = 32

    headers = ["SKU","카테고리","NAND","BOM($)","ASP($)","GM($)","GM%",
               "Q1 판매(K)","분기매출($M)","분기이익($M)"]
    ws.append(headers)
    rh = ws.max_row
    for ci, h in enumerate(headers, 1):
        c = ws.cell(row=rh, column=ci)
        c.fill = _fill(C["red"]); c.font = _font(True, "FFFFFF", 9)
        c.alignment = _center(); c.border = _border()

    cat_label = {"internal_ssd":"Internal SSD","external_ssd":"External SSD","microsd":"microSD"}
    cat_colors = {"internal_ssd":C["orange"],"external_ssd":C["blue"],"microsd":C["green"]}
    cat_order  = ["internal_ssd","external_ssd","microsd"]
    last_cat = None
    totals = {"rev":0,"gp":0,"units":0}

    for sku, info in sorted(PRODUCTS.items(), key=lambda x: cat_order.index(x[1]["cat"])):
        if info["cat"] != last_cat:
            ws.append([f"▶ {cat_label[info['cat']]}"])
            r = ws.max_row
            ws.merge_cells(f"A{r}:J{r}")
            c = ws.cell(row=r, column=1)
            c.fill = _fill(C["navy"]); c.font = _font(True, "FFFFFF", 10)
            c.alignment = Alignment(horizontal="left", vertical="center")
            last_cat = info["cat"]

        bom   = _bom_cost(sku)
        asp   = ASP.get(sku, 50)
        gm    = asp - bom
        gm_p  = gm / asp * 100
        units = UNITS_Q_K.get(sku, 0)
        rev   = units * asp / 1000
        gp    = units * gm  / 1000
        totals["rev"] += rev; totals["gp"] += gp; totals["units"] += units

        row_data = [sku.replace("_"," "), cat_label[info["cat"]], info["nand"],
                    bom, asp, round(gm,2), round(gm_p,1), units, round(rev,2), round(gp,2)]
        ws.append(row_data)
        r = ws.max_row
        for ci, val in enumerate(row_data, 1):
            c = ws.cell(row=r, column=ci)
            c.border = _border(); c.alignment = _center(); c.font = _font(size=9)
            if ci == 1:
                c.fill = _fill(C["gray"])
            elif ci == 3:
                gc = {"BiCS5":C["blue"],"BiCS6":C["orange"],"BiCS8":C["green"]}.get(info["nand"],"888888")
                c.fill = _fill(gc); c.font = _font(True, "FFFFFF", 9)
            elif ci == 7:
                c.fill = _fill(C["g_ok"] if gm_p >= 40 else C["y_warn"] if gm_p >= 25 else C["r_bad"])
                c.number_format = '0.0"%"'
            elif ci in [4,5,6]:
                c.number_format = '"$"#,##0.00'
            elif ci in [9,10]:
                c.fill = _fill(C["gray"]); c.number_format = '"$"#,##0.0 "M"'

    ws.append(["TOTAL","","","","","","",totals["units"],round(totals["rev"],2),round(totals["gp"],2)])
    r = ws.max_row
    for ci in range(1, 11):
        c = ws.cell(row=r, column=ci)
        c.fill = _fill(C["navy"]); c.font = _font(True, "FFFFFF", 10)
        c.border = _border(); c.alignment = _center()
        if ci == 8: c.number_format = '#,##0 "K"'
        if ci in [9,10]: c.number_format = '"$"#,##0.0 "M"'

    # ── Sheet 2: 가격 전략 ────────────────────────────────────
    ws2 = wb.create_sheet("📊 가격 전략")
    ws2.column_dimensions["A"].width = 30
    for col in "BCDEF":
        ws2.column_dimensions[col].width = 16
    ws2.column_dimensions["F"].width = 35

    ws2.merge_cells("A1:F1")
    c = ws2["A1"]
    c.value = "가격 포지셔닝 전략 — 경쟁사 대비 ASP 분석"
    c.fill = _fill(C["navy"]); c.font = Font(bold=True, color="FFFFFF", size=12, name="Calibri")
    c.alignment = _center(); ws2.row_dimensions[1].height = 30

    ws2.append(["세그먼트","SanDisk/WD ASP","Samsung ASP","Kingston ASP","Gap vs Samsung","전략 방향"])
    rh = ws2.max_row
    for ci in range(1, 7):
        c = ws2.cell(row=rh, column=ci)
        c.fill = _fill(C["red"]); c.font = _font(True, "FFFFFF", 9)
        c.alignment = _center(); c.border = _border()

    pricing_data = [
        ["External SSD 1TB (Mainstream)",    88, 95, 70,  -7.4, "Samsung 대비 5~10% Discount 유지 → 가성비 포지션"],
        ["External SSD 2TB (Mainstream)",   145,160,115,  -9.4, "2TB 성장 집중 — Extreme 2TB 프로모 강화"],
        ["External SSD 1TB (Premium/Pro)",  120,130, 85,  -7.7, "Extreme Pro: rugged 차별화, IP68 강조"],
        ["Internal SSD 1TB (Gaming)",        80, 90, 65, -11.1, "SN850X: PS5 공식 인증 마케팅 레버리지"],
        ["Internal SSD 2TB (Gaming)",       145,160,105,  -9.4, "SN850X 4TB 프리미엄 포지션 확대"],
        ["Internal SSD 1TB (Value)",         65, 72, 55,  -9.7, "SN770: 가격 경쟁력 + Laptop OEM 타겟"],
        ["microSD 256GB (A2/V30)",           28, 30, 22,  -6.7, "Extreme: Action cam/Gaming 타겟 마케팅"],
        ["microSD 512GB (A2/V30)",           50, 55, 40,  -9.1, "512GB 성장세 집중 — 스위치/스팀덱 번들"],
        ["microSD 1TB (Premium)",            95,105, 80,  -9.5, "1TB: 드론/4K카메라 프리미엄 포지션"],
    ]
    for row in pricing_data:
        ws2.append(row)
        r = ws2.max_row
        gap = row[4]
        for ci, val in enumerate(row, 1):
            c = ws2.cell(row=r, column=ci)
            c.border = _border(); c.alignment = _center(); c.font = _font(size=9)
            if ci == 1: c.fill = _fill(C["gray"]); c.font = _font(True, size=9)
            elif ci in [2,3,4]: c.number_format = '"$"#,##0'
            elif ci == 5:
                c.value = f"{gap:.1f}%"
                c.fill = _fill(C["y_warn"] if gap > -10 else C["r_bad"])
                c.font = _font(True, size=9)
            elif ci == 6:
                c.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
        ws2.row_dimensions[r].height = 28

    # ── Sheet 3: Revenue 로드맵 ───────────────────────────────
    ws3 = wb.create_sheet("📈 Revenue 로드맵")
    ws3.column_dimensions["A"].width = 20
    for col in "BCDEFGHI":
        ws3.column_dimensions[col].width = 14

    ws3.merge_cells("A1:I1")
    c = ws3["A1"]
    c.value = "연간 Revenue 로드맵 ($M) — 카테고리별 목표"
    c.fill = _fill(C["navy"]); c.font = Font(bold=True, color="FFFFFF", size=12, name="Calibri")
    c.alignment = _center(); ws3.row_dimensions[1].height = 30

    ws3.append(["카테고리"] + QUARTERS)
    rh = ws3.max_row
    for ci in range(1, 10):
        c = ws3.cell(row=rh, column=ci)
        c.fill = _fill(C["red"]); c.font = _font(True, "FFFFFF", 9)
        c.alignment = _center(); c.border = _border()

    rev_plan = {
        "External SSD":  [228,240,264,306,252,268,294,342],
        "Internal SSD":  [172,181,199,230,190,200,220,255],
        "microSD":       [375,395,434,503,413,435,478,553],
        "Total":         [775,816,897,1039,855,903,992,1150],
    }
    row_colors = {"External SSD":C["blue"],"Internal SSD":C["orange"],"microSD":C["green"],"Total":C["navy"]}
    for cat_name, vals in rev_plan.items():
        ws3.append([cat_name]+vals)
        r = ws3.max_row
        for ci, val in enumerate([cat_name]+vals, 1):
            c = ws3.cell(row=r, column=ci)
            c.border = _border(); c.alignment = _center()
            if ci == 1:
                c.fill = _fill(row_colors[cat_name])
                c.font = _font(True, "FFFFFF", 10)
            else:
                c.font = _font(bold=(cat_name=="Total"), size=9)
                c.number_format = '"$"#,##0 "M"'
                if cat_name == "Total": c.fill = _fill(C["gray"])

    out_path = EXCEL_DIR / "marketing_strategy.xlsx"
    wb.save(out_path)
    print(f"[마케팅 전략 에이전트] ✅ Excel 저장: {out_path}")
    return out_path


def build_pptx():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    def ns(): return prs.slides.add_slide(blank)

    # Slide 1: Title
    sl = ns()
    _rect(sl, 0, 0, 13.33, 7.5, C["navy"])
    _rect(sl, 0, 2.8, 13.33, 2.0, C["red"])
    _tbox(sl, "SanDisk B2C Storage", 1, 0.8, 11, 1.2, fs=34, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)
    _tbox(sl, "마케팅 전략 보고서  |  Marketing Strategy", 1, 2.9, 11, 0.8, fs=22, color="FFFFFF", align=PP_ALIGN.CENTER)
    _tbox(sl, f"마케팅 전략 에이전트  |  {datetime.date.today()}", 1, 3.8, 11, 0.5, fs=12, color="FFCCCC", align=PP_ALIGN.CENTER)
    for i, (label, col) in enumerate([("제품 Mix 전략","4472C4"),("가격 전략","ED7D31"),("P&L 시뮬레이션","70AD47"),("Revenue 로드맵","D4001E")]):
        _rect(sl, 1.8+i*2.5, 4.8, 2.2, 0.55, col, label, fs=11, bold=True)

    # Slide 2: 포트폴리오 포지셔닝
    sl = ns()
    _rect(sl, 0, 0, 13.33, 0.9, C["navy"])
    _tbox(sl, "제품 포트폴리오 포지셔닝 전략", 0.3, 0.05, 12, 0.75, fs=22, bold=True, color="FFFFFF")

    _rect(sl, 0.5, 1.0, 5.8, 2.7, "1E3A5F", "High Performance (Premium)", fs=11, tc="AACCEE")
    _rect(sl, 6.5, 1.0, 5.8, 2.7, "1E3A2A", "High Performance + Best Value", fs=11, tc="AAEEBB")
    _rect(sl, 0.5, 4.0, 5.8, 2.7, "3A1E1E", "Entry / Value", fs=11, tc="EECCA0")
    _rect(sl, 6.5, 4.0, 5.8, 2.7, "2A1E3A", "Mid-Range", fs=11, tc="CCAAEE")

    products_matrix = [
        ("SN850X (Gaming NVMe)", 1.2, 1.5, C["orange"], "BiCS6 · PCIe 4.0 · PS5 공식"),
        ("Extreme Pro",          2.5, 2.0, C["blue"],   "BiCS8 · IP68 · 2000MB/s"),
        ("Extreme (Ext.SSD)",    7.2, 1.5, C["blue"],   "BiCS6 · 1050MB/s · Rugged"),
        ("SN770 (NVMe)",         7.8, 2.5, C["orange"], "BiCS5 · PCIe 4.0 · Value"),
        ("SD_E30 (SATA)",        1.5, 4.5, C["orange"], "BiCS5 · SATA · Budget"),
        ("Ultra microSD",        1.8, 5.2, C["green"],  "BiCS5 · A1 · Everyday"),
        ("Extreme microSD",      7.5, 4.3, C["green"],  "BiCS6 · A2 V30 · Action"),
        ("Pro Plus microSD",     8.5, 5.0, C["green"],  "BiCS6 · A2 V30 · Pro"),
        ("My Passport SSD",      7.0, 5.5, C["blue"],   "BiCS5 · USB3.2 · Portable"),
    ]
    for name, x, y, col, desc in products_matrix:
        _rect(sl, x, y, 2.0, 0.42, col, name, fs=9, bold=True)
        _tbox(sl, desc, x, y+0.42, 2.2, 0.3, fs=7, color="AAAAAA")
    _tbox(sl, "← Value ──────────────── Premium →", 4.0, 7.1, 5.5, 0.35, fs=10, color="888888", align=PP_ALIGN.CENTER)

    # Slide 3: 가격 전략
    sl = ns()
    _rect(sl, 0, 0, 13.33, 0.9, C["navy"])
    _tbox(sl, "가격 포지셔닝 — Samsung 대비 7~10% Discount 전략", 0.3, 0.05, 12, 0.75, fs=20, bold=True, color="FFFFFF")

    segments = [
        ("External SSD 1TB",     88, 95,  "Mainstream"),
        ("External SSD Pro 1TB",120,130,  "Premium"),
        ("Internal SSD 1TB (Gaming)",80,90,"Gaming"),
        ("Internal SSD 1TB (Value)",65,72, "Value"),
        ("microSD 256GB",        28, 30,  "A2/V30"),
        ("microSD 512GB",        50, 55,  "A2/V30"),
    ]
    max_price = 130
    for i, (seg, sd_asp, sam_asp, tier) in enumerate(segments):
        y = 1.1 + i * 0.96
        _rect(sl, 0.3, y, 2.7, 0.75, C["navy"], seg, fs=9, bold=True)
        sam_w = sam_asp / max_price * 7.0
        _rect(sl, 3.2, y+0.08, sam_w, 0.3, "888888", f"Samsung ${sam_asp}", fs=8)
        sd_w = sd_asp / max_price * 7.0
        _rect(sl, 3.2, y+0.42, sd_w, 0.3, C["red"], f"SanDisk ${sd_asp}", fs=8)
        gap = (sd_asp - sam_asp) / sam_asp * 100
        _tbox(sl, f"{gap:.1f}%", 3.2+max(sd_w,sam_w)+0.1, y+0.2, 0.9, 0.4, fs=9, bold=True, color=C["orange"])
    _rect(sl, 0.3, 7.1, 1.2, 0.25, "888888", "Samsung")
    _rect(sl, 2.0, 7.1, 1.2, 0.25, C["red"], "SanDisk")
    _tbox(sl, "전략: 가성비 포지셔닝 유지 (5~10% gap), 프리미엄 라인은 5% gap으로 프리미엄 이미지 훼손 방지",
          3.4, 7.05, 9.5, 0.4, fs=9, color="555555")

    # Slide 4: P&L Waterfall
    sl = ns()
    _rect(sl, 0, 0, 13.33, 0.9, C["navy"])
    _tbox(sl, "P&L 시뮬레이션 — 2025 연간 기준 ($M)", 0.3, 0.05, 12, 0.75, fs=22, bold=True, color="FFFFFF")

    waterfall = [
        ("Revenue",         3527, 0,    C["blue"]),
        ("COGS (NAND+BOM)", -1940, 3527, "FFC7CE"),
        ("Gross Profit",    1587, 0,    C["green"]),
        ("Marketing Opex",  -282, 1587, "FFD0CE"),
        ("Channel/Promo",   -176, 1305, "FFD9CC"),
        ("Op. Profit",      1129, 0,    C["orange"]),
    ]
    max_h   = 3527
    chart_h = 4.5
    chart_bot = 7.0
    bar_w = 1.5

    for i, (label, val, base, color) in enumerate(waterfall):
        x = 0.7 + i * 2.0
        if base == 0:
            bar_h_px = abs(val) / max_h * chart_h
            y = chart_bot - bar_h_px
        else:
            bar_h_px = abs(val) / max_h * chart_h
            y = chart_bot - (base / max_h * chart_h) if val < 0 else chart_bot - (base / max_h * chart_h) - bar_h_px
        _rect(sl, x, y, bar_w, bar_h_px, color)
        sign = "+" if val > 0 else ""
        _tbox(sl, f"{sign}${abs(val):,}M", x, y-0.38, bar_w, 0.35, fs=10, bold=True,
              color=color if val > 0 else C["red"], align=PP_ALIGN.CENTER)
        _rect(sl, x, chart_bot+0.05, bar_w, 0.42, C["navy"], label, fs=8, bold=True)

    _tbox(sl, f"GM%: {1587/3527*100:.1f}%  |  영업이익률: {1129/3527*100:.1f}%  |  마케팅 투자율: {(282+176)/3527*100:.1f}%",
          0.5, 7.1, 12, 0.35, fs=10, color="555555", align=PP_ALIGN.CENTER)

    # Slide 5: 제품 Mix 최적화
    sl = ns()
    _rect(sl, 0, 0, 13.33, 0.9, C["navy"])
    _tbox(sl, "제품 Mix 최적화 전략 — BiCS 세대 전환 로드맵", 0.3, 0.05, 12, 0.75, fs=20, bold=True, color="FFFFFF")

    items = [
        ("BiCS8 프리미엄화", C["green"],
         "Extreme Pro (Ext.SSD) → BiCS8 단독 적용\n"
         "• 218L TLC: 원가 $0.038/GB (BiCS6 대비 -10%)\n"
         "• 성능: 2000MB/s read → 프리미엄 ASP 유지\n"
         "• 목표: BiCS8 비중 2025→20%, 2026→35%"),
        ("BiCS6 주력화", C["orange"],
         "Extreme, SN850X, Pro Plus microSD → BiCS6\n"
         "• 162L TLC: 원가 $0.042/GB — 최적 가성비\n"
         "• 현재 주력 공정, Capa 확대 중 (+8% QoQ)\n"
         "• 2025년 매출의 60% BiCS6 기반 목표"),
        ("BiCS5 단계적 축소", C["blue"],
         "Ultra microSD, SN770, SD_E30 → 재고 소진 전략\n"
         "• 112L: 원가 $0.055/GB — 상대적 고원가\n"
         "• Q4 2025까지 신규 BiCS5 투입 중단\n"
         "• 기존 재고는 프로모션으로 적극 소화"),
        ("1TB→2TB 믹스 상향", C["red"],
         "전 카테고리 2TB 비중 확대 전략\n"
         "• 소비자 니즈: 4K영상/게임 용량 증가\n"
         "• GM% 우위: 2TB > 1TB (규모의 경제)\n"
         "• 2TB 번들 캠페인 + Amazon 프로모 집중"),
    ]
    for i, (title, color, body) in enumerate(items):
        x = 0.3 + (i%2)*6.5
        y = 1.1 + (i//2)*2.9
        _rect(sl, x, y, 6.0, 0.55, color, title, fs=13, bold=True)
        _rect(sl, x, y+0.55, 6.0, 2.2, "1E2A3A", body, fs=10, align=PP_ALIGN.LEFT)

    # Slide 6: Revenue 로드맵 누적 바차트
    sl = ns()
    _rect(sl, 0, 0, 13.33, 0.9, C["navy"])
    _tbox(sl, "Revenue 로드맵 — 2025~2026 분기별 목표 ($M)", 0.3, 0.05, 12, 0.75, fs=20, bold=True, color="FFFFFF")

    q_short = ["25Q1","25Q2","25Q3","25Q4","26Q1","26Q2","26Q3","26Q4"]
    rev_data = {
        "External SSD": [228,240,264,306,252,268,294,342],
        "Internal SSD": [172,181,199,230,190,200,220,255],
        "microSD":      [375,395,434,503,413,435,478,553],
    }
    cat_cols_ppt = {"External SSD":C["blue"],"Internal SSD":C["orange"],"microSD":C["green"]}
    totals_q = [sum(rev_data[k][i] for k in rev_data) for i in range(8)]
    max_rev = max(totals_q)
    bw = 1.1; xs = 0.5; cb = 7.0; ch = 4.8

    for qi in range(8):
        x = xs + qi*1.55
        cumul = 0
        for cat in ["External SSD","Internal SSD","microSD"]:
            val = rev_data[cat][qi]
            sh = val/max_rev*ch
            sy = cb - (cumul+val)/max_rev*ch
            _rect(sl, x, sy, bw, sh, cat_cols_ppt[cat])
            cumul += val
        total = totals_q[qi]
        _tbox(sl, f"${total:,}", x, cb-total/max_rev*ch-0.35, bw, 0.33, fs=8, bold=True, color="FFFFFF", align=PP_ALIGN.CENTER)
        _rect(sl, x, cb+0.05, bw, 0.35, C["navy"], q_short[qi], fs=9)

    for i, (cat, col) in enumerate(cat_cols_ppt.items()):
        _rect(sl, 0.4, 1.2+i*0.45, 0.3, 0.3, col)
        _tbox(sl, cat, 0.8, 1.2+i*0.45, 2.0, 0.3, fs=10, color="222222")
    _tbox(sl, "2026 목표: $3,900M  (+10.6% vs 2025 $3,527M)  |  Q4 성수기 집중 달성 전략",
          0.3, 7.1, 12, 0.35, fs=10, color="555555", align=PP_ALIGN.CENTER)

    out_path = PPTX_DIR / "marketing_strategy.pptx"
    prs.save(out_path)
    print(f"[마케팅 전략 에이전트] ✅ PPTX 저장: {out_path}")
    return out_path


if __name__ == "__main__":
    build_excel()
    build_pptx()
