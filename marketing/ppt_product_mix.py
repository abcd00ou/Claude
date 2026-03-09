"""
제품 믹스 전략 VP 보고서 — Product Mix Rationale
SanDisk B2C Storage Marketing Team

슬라이드 구성:
  1. 타이틀 (제품 이미지 콜라주)
  2. Executive Summary — 3대 믹스 원칙
  3. 포트폴리오 BCG 매트릭스 (Revenue × GM%)
  4. Revenue × Margin 최적화 수식
  5. NAND 아키텍처 배분 의사결정
  6. 고객 세그먼트 커버리지 매트릭스
  7. 가격 계단 (Good-Better-Best)
  8. 실시간 가격 인텔리전스 (DB 크롤링 데이터)
  9. 2025 믹스 전환 목표 vs 현재
 10. 실행 계획 & 승인 요청
"""
import sys
sys.path.insert(0, str(__import__('pathlib').Path(__file__).parent))

import json, datetime
from pathlib import Path
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

import ppt_design as D
from ppt_design import C, Layout, rect, tbox, tbox_multi, card, data_table
from ppt_design import slide_header, slide_footer, bar_chart, waterfall_chart
from ppt_design import exec_summary_block
from config import PPTX_DIR, DATA_DIR
from utils.image_fetcher import get_product_image, prefetch_images
from utils.db_price import get_latest_prices, get_run_stats

with open(DATA_DIR / "market_data.json")    as f: MKT  = json.load(f)
with open(DATA_DIR / "internal_sales.json") as f: SALE = json.load(f)

TODAY      = datetime.date.today()
PAGE_PREFIX = "Product Mix"

# ── RGB helper ─────────────────────────────────────────────
def _rgb(hex6: str):
    h = hex6.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

def oval(slide, x: float, y: float, w: float, h: float, fill: str,
         text: str = "", text_size: int = 9, text_color: str = "FFFFFF"):
    """타원(원) 추가 — BCG 버블용."""
    from pptx.util import Inches, Pt
    sp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    sp.fill.solid()
    sp.fill.fore_color.rgb = _rgb(fill)
    sp.line.color.rgb = _rgb("FFFFFF")
    sp.line.width = Pt(1)
    if text:
        tf = sp.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(text_size)
        run.font.bold = True
        run.font.color.rgb = _rgb(text_color)
        run.font.name = "Calibri"
    return sp


def _add_image(slide, sku_id: str, x: float, y: float,
               w: float, h: float, label: str = "") -> bool:
    """이미지 삽입, 없으면 styled placeholder."""
    path = get_product_image(sku_id)
    if path:
        try:
            slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
            return True
        except Exception:
            pass
    # placeholder
    rect(slide, x, y, w, h, C.NAVY_MID)
    if label:
        tbox(slide, label, x + 0.05, y + h/2 - 0.18, w - 0.1, 0.35,
             size=7.5, color=C.GRAY_400, align=PP_ALIGN.CENTER)
    return False


def build():
    # ── 이미지 사전 다운로드 (병렬) ───────────────────────────
    print("  [Product Mix] 제품 이미지 다운로드 중...")
    image_skus = [
        "sd_extreme_1tb", "sd_extreme_2tb",
        "sd_extreme_pro_1tb", "sd_extreme_pro_2tb",
        "wd_sn850x_1tb", "wd_sn850x_2tb",
        "sd_extreme_micro_256g", "sd_extreme_micro_512g",
        "sd_pro_plus_micro_512g",
        "sam_t9_1tb", "sam_t9_2tb", "sam_990pro_2tb",
    ]
    imgs = prefetch_images(image_skus, max_workers=4)
    got = sum(1 for v in imgs.values() if v)
    print(f"  [Product Mix] 이미지 {got}/{len(image_skus)}개 확보")

    # ── DB 가격 로드 ──────────────────────────────────────────
    print("  [Product Mix] DB 가격 데이터 로드 중...")
    db_prices = get_latest_prices("US")
    run_stats  = get_run_stats()
    last_run   = run_stats["runs"][0] if run_stats["runs"] else {}
    db_date    = last_run.get("started_at", "")[:10] if last_run else str(TODAY)

    prs = D.new_prs()
    pg  = 1
    SKU = SALE["sku_revenue_usd_m_2024"]

    # ══════════════════════════════════════════════════════════
    # SLIDE 1: Title (제품 이미지 콜라주)
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, Layout.W, Layout.H, C.NAVY)
    rect(sl, 0, 0, Layout.W, 0.06, C.RED)

    # 우측 제품 이미지 그리드 (4개)
    img_positions = [
        ("sd_extreme_1tb",         8.7, 0.1, 4.5, 3.65, "SD Extreme"),
        ("wd_sn850x_1tb",          8.7, 3.8, 4.5, 3.65, "WD Black SN850X"),
        ("sd_extreme_micro_256g",  8.7, 0.1, 4.5, 3.65, "SD Extreme Micro"),
        ("sam_t9_1tb",             8.7, 3.8, 4.5, 3.65, "Samsung T9"),
    ]
    # 실제 이미지 배치 (2x2 그리드)
    grid = [
        ("sd_extreme_1tb",        8.7, 0.1,  2.2, 3.65),
        ("wd_sn850x_1tb",        10.95, 0.1,  2.35, 3.65),
        ("sd_extreme_micro_256g",  8.7, 3.8,  2.2, 3.65),
        ("sam_t9_1tb",           10.95, 3.8,  2.35, 3.65),
    ]
    for sku_id, gx, gy, gw, gh in grid:
        label = sku_id.replace("_", " ").replace("sd ", "SD ").replace("wd ", "WD ").title()
        _add_image(sl, sku_id, gx, gy, gw, gh, label)

    # 반투명 오버레이 (좌측 텍스트 가독성)
    rect(sl, 0, 0, 8.55, Layout.H, C.NAVY_MID)
    rect(sl, 8.55, 0, 0.15, Layout.H, C.RED)

    tbox(sl, "Product Mix\nStrategy",
         0.55, 0.5, 7.7, 2.2,
         size=42, bold=True, color=C.WHITE, align=PP_ALIGN.LEFT)
    tbox(sl, "왜 이 제품 라인업인가?  포트폴리오 최적화 로직·NAND 배분·고객 세그먼트 커버리지",
         0.55, 2.85, 7.7, 0.65,
         size=12, color="A0AEC0", align=PP_ALIGN.LEFT)
    tbox(sl, f"SanDisk B2C Marketing  |  CONFIDENTIAL  |  {TODAY}",
         0.55, 3.6, 7.7, 0.45, size=10, color="6B7280")
    tbox(sl, "SanDisk Extreme · WD Black SN850X · SanDisk Extreme microSD · Samsung T9 (경쟁사 참조)",
         0.55, 4.1, 7.7, 0.4, size=8.5, color="4B5563", italic=True)

    # ══════════════════════════════════════════════════════════
    # SLIDE 2: Executive Summary — 3대 믹스 원칙
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Executive Summary — 제품 믹스 3대 원칙", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: Internal Sales Data FY2024, WD Annual Report 2024, TrendForce Market Share", pg)

    kpis = [
        ("$3.23B", "2024 B2C 총 매출",     "목표 2025E $3.53B"),
        ("47.6%",  "2024 Blended GM%",    "목표 2025E 49.8%"),
        ("+2.2pp", "GM 개선 목표",          "믹스 시프트 효과 +$110M"),
        ("19 SKU", "핵심 관리 SKU 수",      "Long-tail 합리화 진행"),
    ]
    for i, (val, label, sub) in enumerate(kpis):
        x  = Layout.MARGIN_L + i * 3.13
        lc = [C.RED, C.ORANGE, C.GREEN, C.BLUE][i]
        card(sl, x, 1.05, 2.95, 1.35, fill=C.NAVY_MID,
             title=label, title_size=9, title_color=C.GRAY_400,
             value=val, value_size=22, value_color=C.WHITE,
             sub=sub, sub_size=8, sub_color=lc, line_color=lc)

    exec_summary_block(sl, [
        ("📐", "원칙 1: GM% 최적화 — 마진 높은 SKU 집중",
         "microSD Pro Plus(64.8%), Extreme Micro 1TB(65.2%) 등 고마진 SKU 비중 확대."
         " BiCS8 원가 혁신으로 2025E Internal SSD GM 41.8% 달성. 믹스 시프트만으로 +$110M 추가 GM 창출."),
        ("🗺️", "원칙 2: 가격 계단 완성 — 전 가격대 커버리지",
         "External SSD $30~$303(ASP), Internal SSD $62~$266, microSD $12~$91."
         " 3단계(Good-Better-Best) 계층 유지로 업셀 경로 확보."
         " 공백 구간(External SSD $200~$280) Samsung T9 독점 → Extreme Pro 4TB 진입 시급."),
        ("🎯", "원칙 3: 경쟁 대응 — Samsung 라인별 Counter",
         "External SSD: T9 1TB $130 대응 → Extreme Pro 1TB $119(더 저렴·더 빠름)."
         " Internal SSD: 990Pro vs SN850X (PCIe 4.0 동급 경쟁, 원가 우위 활용)."
         " microSD: Samsung EVO Select vs SD Extreme(성능·내구성 우위 포지셔닝)."),
    ], Layout.MARGIN_L, 2.55, Layout.CONTENT_W, 4.1)

    # ══════════════════════════════════════════════════════════
    # SLIDE 3: BCG 매트릭스 (Revenue × GM%)
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Portfolio BCG Matrix — Revenue vs. Gross Margin%", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: Internal Sales Data FY2024 (Estimated), WD Flash Segment 10-K", pg)

    # BCG 그리드 배경 (X=Revenue, Y=GM%)
    # 차트 영역: x 0.4~9.2, y 1.1~6.8 (8.8 × 5.7 인치)
    CX, CY, CW, CH = 0.5, 1.1, 8.8, 5.7
    mid_x = CX + CW / 2  # Revenue 중간
    mid_y = CY + CH / 2  # GM 중간

    # 사분면 배경
    rect(sl, CX,       CY,       CW/2, CH/2, "1A2942")   # Q2: Stars (좌상)
    rect(sl, CX+CW/2, CY,       CW/2, CH/2, "0F2034")   # Q1: Stars (우상)
    rect(sl, CX,       CY+CH/2, CW/2, CH/2, "141E2E")   # Q3: Dogs (좌하)
    rect(sl, CX+CW/2, CY+CH/2, CW/2, CH/2, "1B2A4A")   # Q4: Cash Cows (우하)

    # 사분면 레이블
    for (label, sub, lx, ly, col) in [
        ("STARS ★",       "고매출·고마진",   CX+0.12,        CY+0.12,       "5CB85C"),
        ("QUESTION MARKS?","저매출·고마진",  CX+CW/2+0.12,  CY+0.12,       "E8A317"),
        ("DOGS ✗",         "저매출·저마진",  CX+0.12,        CY+CH/2+0.12,  "CC4444"),
        ("CASH COWS 💰",   "고매출·저마진",  CX+CW/2+0.12,  CY+CH/2+0.12,  "5B9BD5"),
    ]:
        tbox(sl, label, lx, ly, 3.5, 0.3, size=8.5, bold=True, color=col)
        tbox(sl, sub,   lx, ly+0.3, 3.5, 0.25, size=7.5, color="4B5563", italic=True)

    # 축 레이블
    tbox(sl, "← Revenue ($M) →", CX + CW/2 - 1.5, CY + CH + 0.05, 3.0, 0.28,
         size=8, color=C.GRAY_400, align=PP_ALIGN.CENTER)
    tbox(sl, "←\nG\nM\n%\n→", CX - 0.38, CY + CH/2 - 0.8, 0.35, 1.6,
         size=7.5, color=C.GRAY_400, align=PP_ALIGN.CENTER)

    # 제품 버블 배치 (Revenue vs GM%)
    # x_pos = CX + (rev - min_rev) / rev_range * CW
    # y_pos = CY + CH - (gm - min_gm) / gm_range * CH  (y는 위가 높음)
    rev_min, rev_max = 40, 400
    gm_min,  gm_max  = 30, 70
    def bx(rev): return CX + (rev - rev_min) / (rev_max - rev_min) * CW
    def by(gm):  return CY + CH - (gm - gm_min) / (gm_max - gm_min) * CH

    bubbles = [
        # (label, rev, gm%, color, bubble_r)
        ("SD\nEx. Micro\n256G",  378, 58.4, C.GREEN,    0.6),
        ("SD\nEx. Micro\n512G",  312, 62.1, C.GREEN,    0.55),
        ("WD Black\nSN850X 2TB", 194, 45.2, C.BLUE,     0.48),
        ("SD\nExtreme\n1TB",     215, 45.2, C.RED,      0.5),
        ("SD\nExtreme\n2TB",     198, 46.1, C.RED,      0.48),
        ("SD Ultra\nMicro 256G", 248, 52.3, "5CB85C",   0.52),
        ("SD Pro+\nMicro 256G",  195, 63.2, "16A34A",   0.48),
        ("WD Black\nSN850X 1TB", 162, 42.8, C.BLUE,     0.44),
        ("SD Ex.\nPro 2TB",      137, 54.7, C.ORANGE,   0.42),
        ("SD Ex.\nPro 1TB",      148, 52.3, C.ORANGE,   0.44),
        ("WD Black\nSN770 1TB",  138, 34.5, "6B7280",   0.42),
        ("WD My\nPassport 1TB",  106, 38.2, "6B7280",   0.38),
        ("SD Ex.\nMicro 1TB",     39, 65.2, "16A34A",   0.28),
        ("SD Pro+\nMicro 512G",  165, 64.8, "16A34A",   0.44),
    ]

    for label, rev, gm, color, r in bubbles:
        bx_ = bx(rev) - r / 2
        by_ = by(gm)  - r / 2
        # 경계 클리핑
        bx_ = max(CX + 0.05, min(bx_, CX + CW - r - 0.05))
        by_ = max(CY + 0.05, min(by_, CY + CH - r - 0.05))
        oval(sl, bx_, by_, r, r, color.lstrip("#") if "#" not in color else color, label, 6)

    # 중앙 기준선 (0.02인치 폭 rect로 대체)
    rect(sl, mid_x - 0.01, CY,   0.02, CH,   "4B5563")
    rect(sl, CX,   mid_y - 0.01, CW,   0.02, "4B5563")

    # 우측 범례
    rect(sl, 9.5, 1.1, 3.73, 5.7, "0D1926")
    tbox(sl, "범례 (카테고리)", 9.6, 1.15, 3.53, 0.3, size=9, bold=True, color=C.WHITE)
    legend_items = [
        (C.RED,    "External SSD (Extreme)"),
        (C.ORANGE, "External SSD (Extreme Pro)"),
        (C.BLUE,   "Internal SSD (WD Black)"),
        (C.GREEN,  "microSD (Extreme 계열)"),
        ("16A34A", "microSD (Pro Plus 계열)"),
        ("6B7280", "합리화 대상 SKU"),
    ]
    for i, (col, name) in enumerate(legend_items):
        y = 1.55 + i * 0.55
        rect(sl, 9.65, y + 0.08, 0.22, 0.22, col)
        tbox(sl, name, 9.95, y, 3.2, 0.4, size=8.5, color=C.GRAY_400)

    tbox(sl, "※ 버블 크기 = 상대 매출 규모", 9.6, 4.7, 3.53, 0.3,
         size=7.5, color="4B5563", italic=True)
    tbox(sl, "전략적 시사점:", 9.6, 5.0, 3.53, 0.3, size=9, bold=True, color=C.WHITE)
    insights_bcg = [
        "• 상단 우측(Stars) SKU 확대 집중",
        "• SD Ex. Micro + Pro Plus → 최우선 투자",
        "• SN770/My Passport → 합리화 검토",
        "• Ex. Pro + 4TB → Q.Mark→Star 전환 목표",
    ]
    for i, ins in enumerate(insights_bcg):
        tbox(sl, ins, 9.6, 5.3 + i * 0.37, 3.53, 0.34, size=8, color=C.GRAY_400)

    # ══════════════════════════════════════════════════════════
    # SLIDE 4: Revenue × Margin 최적화 수식
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Portfolio Optimization Math — Revenue × GM% Analysis", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: Internal Sales Data FY2024, Internal Strategy Planning FY2025", pg)

    # 헤더
    rect(sl, Layout.MARGIN_L, 1.1, Layout.CONTENT_W, 0.38, C.NAVY_MID)
    tbox(sl, "믹스 최적화 공식:  Total GM = Σ (Revenue_i × GM%_i)  →  SKU 믹스 시프트만으로 +$110M 달성 가능",
         Layout.MARGIN_L+0.15, 1.13, Layout.CONTENT_W-0.3, 0.32,
         size=9.5, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    # 카테고리별 현재 vs 목표 분석
    cat_analysis = [
        ("External SSD", "$950M", "44.8%", "$425M", "$1,036M", "47.2%", "$489M", "+$64M",
         "Extreme Pro 4TB 비중 확대 + ASP 방어"),
        ("Internal SSD", "$760M", "38.5%", "$293M", "$882M",   "41.8%", "$369M", "+$76M",
         "BiCS8 전환 → 원가 절감 + SN850X 집중"),
        ("microSD",     "$1,520M","57.3%","$871M","$1,610M",  "58.9%","$949M",  "+$78M",
         "Pro Plus 512G+ 비중 확대 (GM 64.8%)"),
        ("Total",      "$3,230M", "47.6%","$1,589M","$3,528M","49.8%","$1,807M","+$218M",
         "믹스 시프트 + 볼륨 성장 복합 효과"),
    ]
    headers = ["카테고리", "2024 Rev.", "GM%", "2024 GM$",
               "2025E Rev.", "GM%", "2025E GM$", "GM$ 증가", "드라이버"]
    rows = []
    for cat, rev24, gm24, gmv24, rev25, gm25, gmv25, delta, driver in cat_analysis:
        is_total = cat == "Total"
        rows.append([
            cat, rev24, gm24, gmv24,
            (rev25, C.GREEN), (gm25, C.GREEN), (gmv25, C.GREEN),
            (delta, C.RED if is_total else C.GREEN),
            driver,
        ])
    data_table(sl, Layout.MARGIN_L, 1.58, Layout.CONTENT_W, headers, rows,
               col_widths=[1.4, 0.95, 0.7, 1.05, 1.0, 0.7, 1.05, 0.88, 4.7],
               header_fill=C.RED, font_size=8.5)

    # 믹스 시프트 효과 워터폴
    rect(sl, Layout.MARGIN_L, 4.15, Layout.CONTENT_W, 0.28, C.NAVY_MID)
    tbox(sl, "GM$ 브릿지 — 볼륨 성장 vs 믹스 시프트 기여 분해",
         Layout.MARGIN_L+0.1, 4.17, Layout.CONTENT_W-0.2, 0.22,
         size=9, bold=True, color=C.WHITE)

    wf_items = [
        ("2024\nGM",          1589, "start"),
        ("볼륨\n성장",           108, "add"),
        ("External SSD\n믹스",   64, "add"),
        ("Internal SSD\n믹스",   76, "add"),
        ("microSD\n믹스",        78, "add"),
        ("가격 인하\n압력",       -25, "sub"),
        ("NAND\n비용 절감",       -83, "add"),
        ("2025E\nGM",             0, "total"),
    ]
    waterfall_chart(sl, wf_items, Layout.MARGIN_L, 4.5, Layout.CONTENT_W, 2.65,
                    pos_color=C.GREEN, neg_color=C.RED,
                    total_color=C.BLUE, bg_color="0D1926")

    # ══════════════════════════════════════════════════════════
    # SLIDE 5: NAND 아키텍처 배분 의사결정
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "NAND Architecture Decision — BiCS 세대별 SKU 배분 원칙", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: WD Technology Roadmap 2024, Internal NAND Planning, TrendForce NAND Cost Q4 2025", pg)

    # 의사결정 원칙 헤더
    rect(sl, Layout.MARGIN_L, 1.1, Layout.CONTENT_W, 0.38, C.NAVY_MID)
    tbox(sl, "원칙: 가장 새로운(저원가) NAND를 가장 마진 높은 제품에 배분 — BiCS8 = 프리미엄 SKU 독점 배분",
         Layout.MARGIN_L+0.15, 1.13, Layout.CONTENT_W-0.3, 0.32,
         size=9.5, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    # NAND 세대별 흐름도
    bics_data = [
        ("BiCS5\n112L", "$0.055/GB", "단종 예정", C.GRAY_400, [
            "WD My Passport 1TB/2TB",
            "SD Ultra microSD 128G/256G",
            "WD Blue SN580 (보급형)",
            "→ 2025년 내 전량 합리화",
        ], ["$38.2%", "$40.1%", "$36%", ""]),
        ("BiCS6\n162L", "$0.042/GB", "주력 전환 중", C.BLUE, [
            "SD Extreme 1TB/2TB",
            "WD Black SN770 1TB/2TB",
            "SD Extreme microSD 128~512G",
            "WD My Passport 4TB (유지)",
        ], ["45.2%", "34.5%~36.2%", "58.4%~62.1%", "~40%"]),
        ("BiCS8\n218L", "$0.038/GB", "프리미엄 전용", C.GREEN, [
            "SD Extreme Pro 1TB/2TB/4TB",
            "WD Black SN850X 1TB/2TB/4TB",
            "SD Pro Plus microSD 256~1TB",
            "SD Extreme microSD 1TB",
        ], ["52.3%~54.7%", "42.8%~47.1%", "63.2%~64.8%", "65.2%"]),
    ]

    for i, (bics, cost, status, color, skus, gms) in enumerate(bics_data):
        x = Layout.MARGIN_L + i * 4.18
        # 헤더 박스
        rect(sl, x, 1.6, 4.0, 0.85, color)
        tbox(sl, bics, x + 0.1, 1.62, 3.8, 0.45, size=14, bold=True, color=C.WHITE)
        tbox(sl, f"원가: {cost}  |  {status}", x + 0.1, 2.08, 3.8, 0.3,
             size=8.5, color=C.WHITE)
        # 본문
        rect(sl, x, 2.45, 4.0, 4.3, "0D1926" if i==2 else "131E2E")
        tbox(sl, "배분 SKU", x + 0.12, 2.5, 3.76, 0.28, size=8.5, bold=True, color=color)
        for j, (sku, gm) in enumerate(zip(skus, gms)):
            y = 2.85 + j * 0.92
            rect(sl, x+0.08, y, 3.84, 0.82, "1A2942")
            rect(sl, x+0.08, y, 0.06, 0.82, color)
            tbox(sl, sku, x+0.2, y+0.04, 2.5, 0.38, size=8.5, color=C.WHITE)
            if gm:
                tbox(sl, f"GM {gm}", x+0.2, y+0.44, 2.5, 0.3, size=8, color=color, bold=True)

    # 원가 절감 효과 요약
    rect(sl, Layout.MARGIN_L, 6.85, Layout.CONTENT_W, 0.4, C.NAVY)
    tbox(sl, "BiCS5→BiCS8 원가 31% 절감($0.055→$0.038)  ·  Pro Plus 512G 기준 GM 48%→65% (+17pp)  ·  "
             "2025E blended GM 49.8% 핵심 드라이버",
         Layout.MARGIN_L+0.15, 6.87, Layout.CONTENT_W-0.3, 0.35,
         size=9, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    # SLIDE 6: 고객 세그먼트 커버리지 매트릭스
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Customer Segment Coverage Matrix", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: GfK Consumer Insights 2024, IDC Storage End-user Survey, Internal Market Research", pg)

    # 세그먼트 × 제품 매트릭스
    segments = [
        "🎮 게이머\n(Hardcore)",
        "📹 Content\nCreator",
        "📸 사진작가\n(프로)",
        "🏢 전문직\n파워유저",
        "📱 일반\n소비자",
        "🎓 학생\n예산형",
    ]
    products = [
        ("WD SN850X\n1/2/4TB", C.BLUE),
        ("WD SN770\n1/2TB",    C.BLUE),
        ("SD Ext.\n1/2TB",     C.RED),
        ("SD Ext.Pro\n1/2/4TB",C.RED),
        ("SD Ex.Micro\n256-1TB",C.GREEN),
        ("SD Pro+\nMicro",     C.GREEN),
        ("WD My\nPassport",    C.GRAY_400),
    ]

    # 커버리지 매트릭스 데이터 (✅=주력 / ○=보조 / -=해당없음)
    coverage = [
        # SN850X, SN770, Extreme, Ext.Pro, Ex.Micro, Pro+Micro, MyPassport
        ["✅", "○",  "-",  "-",  "-",  "-",  "-" ],  # 게이머
        ["-",  "-",  "✅", "✅", "✅", "✅", "-" ],  # Creator
        ["-",  "-",  "○",  "✅", "✅", "✅", "-" ],  # 사진작가
        ["✅", "✅", "✅", "✅", "○",  "○",  "-" ],  # 전문직
        ["-",  "-",  "✅", "○",  "✅", "○",  "✅"],  # 일반 소비자
        ["-",  "○",  "✅", "-",  "✅", "-",  "✅"],  # 학생
    ]
    # Samsung 커버리지 비교
    sam_coverage = [
        ["✅", "-",  "-",  "-",  "-",  "-",  "-" ],
        ["-",  "-",  "✅", "-",  "✅", "-",  "-" ],
        ["-",  "-",  "✅", "✅", "✅", "-",  "-" ],
        ["✅", "-",  "✅", "✅", "-",  "-",  "-" ],
        ["-",  "-",  "✅", "-",  "✅", "✅", "-" ],
        ["-",  "-",  "○",  "-",  "✅", "-",  "✅"],
    ]

    # 헤더
    tbox(sl, "SanDisk/WD Portfolio", 0.4, 1.12, 6.5, 0.3,
         size=9, bold=True, color=C.WHITE)
    tbox(sl, "Samsung 비교 (경쟁사)", 7.0, 1.12, 6.13, 0.3,
         size=9, bold=True, color=C.ORANGE)

    cell_w = 0.85
    cell_h = 0.75
    start_x_sd = 0.4
    start_x_sam = 7.0
    start_y = 1.45

    # 제품 헤더 (SanDisk 측)
    for j, (prod, color) in enumerate(products):
        rect(sl, start_x_sd + j * cell_w, start_y, cell_w - 0.04, 0.45, color)
        tbox(sl, prod, start_x_sd + j * cell_w + 0.02, start_y, cell_w - 0.08, 0.42,
             size=6.5, color=C.WHITE, align=PP_ALIGN.CENTER)

    # Samsung 제품 헤더 (약식)
    sam_prods = ["990Pro\n1/2TB", "9100Pro\n(PCIe5)", "T9\n1/2/4TB", "EVO\nSelect", "PRO\nPlus"]
    for j, prod in enumerate(sam_prods):
        rect(sl, start_x_sam + j * cell_w, start_y, cell_w - 0.04, 0.45, C.ORANGE)
        tbox(sl, prod, start_x_sam + j * cell_w + 0.02, start_y, cell_w - 0.08, 0.42,
             size=6.5, color=C.WHITE, align=PP_ALIGN.CENTER)

    # 매트릭스 셀
    for i, (seg, cov_row) in enumerate(zip(segments, coverage)):
        y = start_y + 0.5 + i * cell_h
        # 세그먼트 레이블
        rect(sl, 0.0, y, start_x_sd - 0.05, cell_h - 0.04, "0D1926")
        tbox(sl, seg, 0.02, y + 0.12, start_x_sd - 0.08, cell_h - 0.2,
             size=7.5, color=C.WHITE, align=PP_ALIGN.CENTER)
        # SanDisk 커버리지
        for j, cell in enumerate(cov_row):
            bg = "1A3A22" if cell == "✅" else ("1A2942" if cell == "○" else "0D1926")
            rect(sl, start_x_sd + j * cell_w, y, cell_w - 0.04, cell_h - 0.04, bg)
            tbox(sl, cell, start_x_sd + j * cell_w, y + 0.18, cell_w - 0.04, 0.35,
                 size=14 if cell == "✅" else 11, color=C.GREEN if cell == "✅" else
                 (C.ORANGE if cell == "○" else "2A3A4A"), align=PP_ALIGN.CENTER)
        # Samsung 비교
        for j, cell in enumerate(sam_coverage[i][:5]):
            bg = "3A1A00" if cell == "✅" else "0D1926"
            rect(sl, start_x_sam + j * cell_w, y, cell_w - 0.04, cell_h - 0.04, bg)
            tbox(sl, cell, start_x_sam + j * cell_w, y + 0.18, cell_w - 0.04, 0.35,
                 size=12, color=C.ORANGE if cell == "✅" else "2A3A4A",
                 align=PP_ALIGN.CENTER)

    # 범례
    for sym, col, label in [("✅", C.GREEN, "주력 타겟 세그먼트"), ("○", C.ORANGE, "보조 커버리지")]:
        pass  # inline
    tbox(sl, "✅ = 주력 타겟  ○ = 보조  - = 비해당  |  공백: SanDisk 진입 기회",
         0.0, 6.92, Layout.W, 0.28, size=7.5, color=C.GRAY_400, italic=True,
         align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    # SLIDE 7: 가격 계단 Good-Better-Best
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Price Tier Ladder — Good-Better-Best 가격 계단 전략", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: Internal SKU ASP Data FY2024, Amazon Price Crawling (2026-03-09)", pg)

    tbox(sl, "가격 계단 설계 원칙:  각 계층은 명확한 성능 차별점으로 업셀 경로 형성  →  "
             "소비자가 '가성비 최대화' 지점으로 자연스럽게 상위 SKU 선택",
         Layout.MARGIN_L, 1.1, Layout.CONTENT_W, 0.38, size=9, color=C.GRAY_400,
         align=PP_ALIGN.CENTER)

    categories_ladder = [
        ("External SSD", [
            ("GOOD",   "WD My Passport",   "$72",  "540MB/s  USB-A  2년 보증",       C.GRAY_400),
            ("BETTER", "SD Extreme 1TB",   "$87",  "1050MB/s  USB-C  방수  IP65",    C.RED),
            ("BEST",   "SD Extreme Pro 2TB","$202", "2000MB/s  NVMe  forCreators",    C.ORANGE),
            ("ULTRA",  "SD Extreme Pro 4TB","$303", "2000MB/s  4TB 초대용량  NEW",    C.RED),
        ]),
        ("Internal SSD", [
            ("GOOD",   "WD Black SN770 1TB","$62",  "PCIe 4.0  5150MB/s  가성비",    C.GRAY_400),
            ("BETTER", "WD Black SN850X 1TB","$77", "PCIe 4.0  7300MB/s  게이밍",    C.BLUE),
            ("BEST",   "WD Black SN850X 2TB","$141","PCIe 4.0  7300MB/s  With HS",   C.BLUE),
            ("ULTRA",  "WD Black SN850X 4TB","$266","PCIe 4.0  7300MB/s  최고용량",  C.BLUE),
        ]),
        ("microSD", [
            ("GOOD",   "SD Ultra 128G",     "$12",  "120MB/s  A1  일상 사용",         C.GRAY_400),
            ("BETTER", "SD Extreme 256G",   "$26",  "190MB/s  A2  4K 촬영 지원",      C.RED),
            ("BEST",   "SD Extreme 512G",   "$43",  "190MB/s  4K  방수  -25°C",       C.RED),
            ("ULTRA",  "SD Pro Plus 512G",  "$58",  "200MB/s  V30  Creator급",        C.GREEN),
        ]),
    ]

    for i, (cat_name, tiers) in enumerate(categories_ladder):
        x = Layout.MARGIN_L + i * 4.18
        rect(sl, x, 1.55, 4.0, 0.35, C.NAVY_MID)
        tbox(sl, cat_name, x+0.1, 1.57, 3.8, 0.3, size=10, bold=True, color=C.WHITE)

        for j, (tier, prod, price, desc, color) in enumerate(tiers):
            tier_h = [0.72, 0.88, 1.0, 1.12][j]
            y = 1.95 + sum([0.72, 0.88, 1.0, 1.12][:j])
            # 계단식 너비 (낮을수록 좁음)
            tier_w = [2.8, 3.2, 3.6, 4.0][j]
            tier_x = x + (4.0 - tier_w)  # 우측 정렬

            rect(sl, tier_x, y, tier_w, tier_h - 0.06, color)
            tbox(sl, tier, tier_x + 0.08, y + 0.04, 0.7, 0.28,
                 size=8, bold=True, color=C.WHITE)
            tbox(sl, prod, tier_x + 0.08, y + 0.3, tier_w - 0.7, 0.3,
                 size=9, bold=True, color=C.WHITE)
            tbox(sl, desc, tier_x + 0.08, y + 0.58, tier_w - 1.0, 0.3,
                 size=7.5, color="DDDDDD")
            tbox(sl, price, tier_x + tier_w - 0.85, y + 0.12, 0.75, 0.35,
                 size=13, bold=True, color=C.WHITE, align=PP_ALIGN.RIGHT)

    # 업셀 전략 박스
    rect(sl, Layout.MARGIN_L, 6.82, Layout.CONTENT_W, 0.43, C.NAVY)
    tbox(sl, "업셀 전략:  GOOD→BETTER 전환 시 수익 +21% / BETTER→BEST +76% / BEST→ULTRA +40%  "
             "|  Bundle(External+microSD): ASP $111→$169(+52%) 달성",
         Layout.MARGIN_L+0.15, 6.84, Layout.CONTENT_W-0.3, 0.38,
         size=9, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    # SLIDE 8: 실시간 가격 인텔리전스 (DB 크롤링)
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, f"Real-Time Price Intelligence — DB 크롤링 데이터 ({db_date})", PAGE_PREFIX, pg)
    slide_footer(sl, f"Source: Price Intelligence DB (price_intel), Crawling Agent 자동 수집  |  Last Run: {db_date}", pg)

    # DB 상태 배너
    db_sku_count = len(db_prices)
    run_success  = last_run.get("success_count", "N/A")
    run_total    = last_run.get("total_targets", "N/A")

    rect(sl, Layout.MARGIN_L, 1.1, Layout.CONTENT_W, 0.38, C.NAVY_MID)
    tbox(sl, f"DB 수집 현황: {db_sku_count}개 SKU 가격 확보  ·  "
             f"최근 실행 성공률 {run_success}/{run_total}  ·  "
             f"수집 주기 1일 2회 (09:00 / 21:00 UTC)  ·  마지막 업데이트 {db_date}",
         Layout.MARGIN_L+0.15, 1.13, Layout.CONTENT_W-0.3, 0.32,
         size=9, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    # 가격 비교 테이블 (DB vs MSRP)
    price_table_headers = ["SKU", "MSRP", "DB 실측가", "프리미엄 %", "판매자", "수집일", "믹스 시사점"]

    # SALE의 amazon_real_prices에서 데이터 가져오기 (DB 폴백 포함)
    real = SALE["amazon_real_prices"]
    price_table_data = [
        ("SD Extreme 1TB",       99.99,  real.get("SD_EXTREME_1TB_USD",    {}).get("crawled", "-"),  "External SSD"),
        ("SD Extreme 2TB",      159.99,  real.get("SD_EXTREME_2TB_USD",    {}).get("crawled", "-"),  "External SSD"),
        ("SD SN850X 2TB",       149.99,  real.get("SD_SN850X_2TB_USD",     {}).get("crawled", "-"),  "Internal SSD"),
        ("SD Ultra 128G",        18.99,  real.get("SD_ULTRA_MICRO_128G_USD",{}).get("crawled", "-"),  "microSD"),
        ("SD Ex. Micro 256G",    28.99,  real.get("SD_EXTREME_MICRO_256G_USD",{}).get("crawled","-"),  "microSD"),
        ("Samsung T9 1TB",      129.99,  real.get("SAM_T9_1TB_USD",        {}).get("crawled", "-"),  "경쟁사"),
        ("Samsung T9 2TB",      229.99,  real.get("SAM_T9_2TB_USD",        {}).get("crawled", "-"),  "경쟁사"),
        ("Samsung 990Pro 2TB",  149.99,  real.get("SAM_990PRO_2TB_USD",    {}).get("crawled", "-"),  "경쟁사"),
    ]

    price_rows = []
    for prod, msrp, crawled, cat in price_table_data:
        if isinstance(crawled, (int, float)):
            premium = (crawled / msrp - 1) * 100
            premium_str = f"+{premium:.0f}%"
            prem_col = C.RED if premium > 50 else (C.ORANGE if premium > 20 else C.GREEN)
            crawled_str = f"${crawled:.2f}"
            seller_str = "Amazon 3P"
        else:
            premium_str = "N/A"
            prem_col = C.GRAY_400
            crawled_str = "데이터 없음"
            seller_str = "-"

        implication = {
            "External SSD": "3P 재고 부족 시사 → 공식 채널 강화",
            "Internal SSD": "공급 타이트 → 가격 인상 기회",
            "microSD":      "보급형 단종 가능성 → Pro 전환 촉진",
            "경쟁사":        "Samsung도 동반 프리미엄 → 업계 공급 부족",
        }.get(cat, "")

        price_rows.append([prod, f"${msrp}", crawled_str,
                           (premium_str, prem_col), seller_str, db_date, implication])

    data_table(sl, Layout.MARGIN_L, 1.58, Layout.CONTENT_W, price_table_headers, price_rows,
               col_widths=[2.0, 0.8, 1.0, 1.0, 1.0, 0.95, 5.58],
               header_fill=C.NAVY, font_size=8.5)

    # 인사이트 박스
    rect(sl, Layout.MARGIN_L, 6.0, Layout.CONTENT_W, 1.25, "0D1926")
    tbox(sl, "가격 인텔리전스 → 제품 믹스 시사점",
         Layout.MARGIN_L+0.15, 6.05, Layout.CONTENT_W-0.3, 0.3,
         size=9.5, bold=True, color=C.WHITE)
    implications = [
        "① 3P 프리미엄 급등 = 공식 채널 재고 부족 신호 → Extreme 1TB/2TB 우선 재고 확보 + Vendor Central 강화",
        "② SD Ultra 128G +337% 이상 → 단종/단계적 철수 신호 → SD Extreme 256G로 대체 마케팅 시급",
        "③ SN850X 2TB +195% 이상 → 게이밍 수요 폭발 + 공급 부족 → BiCS8 전환 가속으로 생산 확대",
        "④ Samsung 동반 프리미엄 → 업계 전체 공급 타이트 → 전략적 가격 인상(+5%) 및 ASP 방어 기회",
    ]
    for i, impl in enumerate(implications):
        tbox(sl, impl, Layout.MARGIN_L+0.15, 6.38 + i*0.22, Layout.CONTENT_W-0.3, 0.2,
             size=8, color=C.GRAY_400)

    # ══════════════════════════════════════════════════════════
    # SLIDE 9: 2025 믹스 목표 vs 현재
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "2025 Product Mix Target — 현재 vs 목표 비교", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: Internal Sales Data FY2024, Internal Strategy Planning FY2025", pg)

    # 믹스 시프트 테이블
    rect(sl, Layout.MARGIN_L, 1.1, Layout.CONTENT_W, 0.35, C.NAVY_MID)
    tbox(sl, "SKU별 Revenue 비중 목표 — 고마진 SKU 확대 · 저마진 SKU 합리화",
         Layout.MARGIN_L+0.1, 1.12, Layout.CONTENT_W-0.2, 0.3,
         size=10, bold=True, color=C.WHITE)

    mix_headers = ["SKU", "2024 Rev($M)", "2024 비중", "2025E Rev($M)", "2025E 비중", "변화", "전략"]
    mix_data = [
        # External SSD
        ("SD Extreme 1TB",      215, 6.7, 235, 6.7, "→ 유지",  "주력 볼륨 SKU 유지",         C.GRAY_400),
        ("SD Extreme 2TB",      198, 6.1, 218, 6.2, "↑ 확대",  "2TB 수요 성장 대응",          C.GREEN),
        ("SD Extreme Pro 1TB",  148, 4.6, 175, 5.0, "↑↑ 확대", "BiCS8 전환→마진 개선",       C.GREEN),
        ("SD Extreme Pro 2TB",  137, 4.2, 175, 5.0, "↑↑ 확대", "Creator 수요 집중 공략",      C.GREEN),
        ("SD Extreme Pro 4TB*",   0, 0.0,  88, 2.5, "🆕 신규",  "시장 공백 공략",              C.RED),
        ("WD My Passport 1TB",  106, 3.3,  85, 2.4, "↓ 합리화","보급형 정리, Extreme 전환",   C.ORANGE),
        # Internal SSD
        ("WD Black SN850X 1TB", 162, 5.0, 195, 5.5, "↑ 확대",  "게이밍 핵심 SKU 강화",        C.GREEN),
        ("WD Black SN850X 2TB", 194, 6.0, 248, 7.0, "↑↑ 확대", "고용량·고마진 전략",          C.GREEN),
        ("WD Black SN770 1TB",  138, 4.3, 110, 3.1, "↓ 합리화","SN850X 전환 유도",            C.ORANGE),
        # microSD
        ("SD Extreme Micro 256G",378,11.7, 415,11.8, "→ 유지", "최대 볼륨 SKU 유지",          C.GRAY_400),
        ("SD Extreme Micro 512G",312, 9.7, 375,10.6, "↑ 확대", "4K 영상 수요 성장",           C.GREEN),
        ("SD Pro+ Micro 256G",  195, 6.0, 240, 6.8, "↑ 확대",  "고마진 Pro 라인 강화",         C.GREEN),
        ("SD Pro+ Micro 512G",  165, 5.1, 225, 6.4, "↑↑ 확대", "Creator급 최고마진(64.8%)",   C.GREEN),
    ]

    rows = []
    total_24 = sum(d[1] for d in mix_data)
    total_25 = sum(d[3] for d in mix_data)
    for sku, rev24, pct24, rev25, pct25, change, strategy, color in mix_data:
        change_col = C.GREEN if "↑" in change or "🆕" in change else \
                     (C.ORANGE if "↓" in change else C.GRAY_400)
        rows.append([sku, f"${rev24}M", f"{pct24:.1f}%",
                     (f"${rev25}M", C.GREEN if rev25 > rev24 else C.ORANGE),
                     (f"{pct25:.1f}%", change_col),
                     (change, change_col), strategy])
    data_table(sl, Layout.MARGIN_L, 1.52, Layout.CONTENT_W, mix_headers, rows,
               col_widths=[2.4, 1.1, 0.85, 1.2, 0.9, 0.88, 5.0],
               header_fill=C.RED, font_size=8)

    # 요약 합계 바
    rect(sl, Layout.MARGIN_L, 6.82, Layout.CONTENT_W, 0.43, C.NAVY_MID)
    tbox(sl, f"전체 매출:  2024 ${total_24}M → 2025E ${total_25}M  ·  "
             f"주요 SKU Mix:  확대 7개 / 유지 3개 / 합리화 2개 / 신규 1개  ·  "
             f"가중평균 GM%: 47.6% → 49.8%(E)",
         Layout.MARGIN_L+0.15, 6.84, Layout.CONTENT_W-0.3, 0.38,
         size=9, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    # SLIDE 10: 실행 계획 & 승인 요청
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Action Plan — 제품 믹스 전환 실행 로드맵", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: Internal Strategy Planning, Supply Chain Roadmap FY2025", pg)

    rect(sl, Layout.MARGIN_L, 1.1, Layout.CONTENT_W, 0.38, C.RED)
    tbox(sl, "VP 승인 요청:  2025 믹스 전환 계획  ·  BiCS8 가속 투자  ·  합리화 SKU 단종 일정 확정",
         Layout.MARGIN_L+0.15, 1.13, Layout.CONTENT_W-0.3, 0.32,
         size=9.5, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    actions = [
        ("P0\n즉시\n실행", C.RED, "Extreme Pro 4TB 긴급 출시 가속",
         "현황: External SSD $200-$280 가격대 Samsung T9 독점 → 즉각 대응 필요\n"
         "조치: BiCS8 Extreme Pro 4TB 재고 확보 + Amazon Vendor Central 우선 등록\n"
         "효과: 신규 매출 $88M/년, GM% 49.8% (최고 마진 구간)\n"
         "오너: Product Marketing + Supply Chain  |  기한: 2025 Q1 완료"),

        ("P0\n즉시\n실행", C.RED, "SD Ultra 128G 단종 → SD Extreme 256G 전환",
         "현황: SD Ultra 128G Amazon 3P 가격 +337% → 공급 중단/단종 신호\n"
         "조치: 단종 공식화 + SD Extreme 256G 가격 $26 동결 + 전환 마케팅\n"
         "효과: microSD GM 52.3%→58.4% 믹스 개선 (매출 중립)\n"
         "오너: Category Manager microSD  |  기한: 2025 Q1"),

        ("P1\n2025Q2\n완료", C.ORANGE, "WD Black SN770 합리화 → SN850X 전환",
         "현황: SN770 GM 34.5-36.2% (SN850X 대비 -8pp) → 마진 유출 중\n"
         "전략: SN770 가격 인상 또는 단종 + SN850X 1TB 가격 $77 방어\n"
         "효과: Internal SSD GM 38.5%→41.8% 핵심 드라이버 (+$40M 기여)\n"
         "오너: Product Line Manager  |  기한: 2025 Q2"),

        ("P1\n2025Q2\n완료", C.ORANGE, "microSD Pro Plus 512G→1TB 확대",
         "현황: Pro Plus 512G GM 64.8%, 1TB GM 65.2% → 최고 마진 구간\n"
         "전략: Pro Plus 1TB 주력 SKU 전환 + $90 ASP 방어 + Creator 마케팅 집중\n"
         "효과: microSD 믹스 개선 +$78M GM 기여 (1TB 비중 2%→6%)\n"
         "오너: Category Manager microSD  |  기한: 2025 Q2"),

        ("P2\n2025Q3\n목표", C.BLUE, "WD My Passport 합리화 완료",
         "현황: WD My Passport GM 38.2-40.1% → 포트폴리오 최저 마진\n"
         "전략: 2025 Q3까지 단종 또는 대폭 가격 인상(+20%) 선택\n"
         "효과: External SSD 믹스 GM 44.8%→47.2% 달성에 기여\n"
         "오너: Product Line Manager  |  기한: 2025 Q3"),

        ("P2\n2025Q3\n목표", C.BLUE, "3P 가격 이상 Brand Protection 강화",
         "현황: Amazon 3P가 MSRP 대비 70~337% 프리미엄 → 브랜드 가치 훼손\n"
         "전략: MAP(최저 광고 가격) 정책 강화 + Vendor Central 재고 확대\n"
         "효과: ASP 방어 + 소비자 신뢰도 유지 + 공식 채널 매출 +$45M\n"
         "오너: VP of Marketing + Sales  |  기한: 2025 Q2 착수"),
    ]

    for i, (priority, color, title, detail) in enumerate(actions):
        col = i % 2
        row = i // 2
        x = Layout.MARGIN_L + col * 6.33
        y = 1.62 + row * 1.72
        rect(sl, x, y, 0.7, 1.62, color)
        tbox(sl, priority, x, y + 0.25, 0.7, 1.1,
             size=8, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
        rect(sl, x + 0.7, y, 5.5, 1.62, "0D1926")
        tbox(sl, title, x + 0.8, y + 0.06, 5.3, 0.35, size=10, bold=True, color=color)
        tbox(sl, detail, x + 0.8, y + 0.42, 5.3, 1.18, size=8, color=C.GRAY_600)

    # 최종 요약
    rect(sl, Layout.MARGIN_L, 6.82, Layout.CONTENT_W, 0.43, C.NAVY_MID)
    tbox(sl, "믹스 전환 목표:  GM% 47.6%→49.8% (+2.2pp)  ·  GM$ +$218M  ·  "
             "확대 7SKU / 합리화 3SKU / 신규 1SKU  ·  BiCS8 전환 2025 H2 완료",
         Layout.MARGIN_L+0.15, 6.84, Layout.CONTENT_W-0.3, 0.38,
         size=9, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    out = PPTX_DIR / "product_mix_vp.pptx"
    prs.save(str(out))
    print(f"[제품 믹스] ✅ VP급 PPTX: {out}")
    return str(out)


if __name__ == "__main__":
    build()
