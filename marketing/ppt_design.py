"""
VP-급 PowerPoint 디자인 시스템
SanDisk B2C Storage Marketing Team
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import datetime

# ══════════════════════════════════════════════════════════════
# SanDisk 브랜드 컬러 시스템
# ══════════════════════════════════════════════════════════════
class C:
    # Brand
    RED       = "D4001E"   # SanDisk 시그니처 레드
    DARK_RED  = "9B0016"   # 강조
    NAVY      = "0F1E35"   # 다크 배경
    NAVY_MID  = "1B2A4A"   # 중간 다크
    NAVY_LIGHT= "243554"   # 라이트 다크

    # Neutral
    WHITE     = "FFFFFF"
    LIGHT_BG  = "F8F9FA"
    GRAY_100  = "F3F4F6"
    GRAY_200  = "E5E7EB"
    GRAY_400  = "9CA3AF"
    GRAY_600  = "4B5563"
    GRAY_800  = "1F2937"

    # Data viz
    BLUE      = "2563EB"   # Primary data
    BLUE_LIGHT= "3B82F6"
    ORANGE    = "EA580C"
    ORANGE_LIGHT="F97316"
    GREEN     = "16A34A"
    GREEN_LIGHT="22C55E"
    PURPLE    = "7C3AED"
    TEAL      = "0891B2"
    YELLOW    = "D97706"

    # Status
    G_OK      = "DCFCE7"   # 초록 배경
    G_TEXT    = "166534"
    Y_WARN    = "FEF9C3"
    Y_TEXT    = "854D0E"
    R_BAD     = "FEE2E2"
    R_TEXT    = "991B1B"

    # Chart palette (순서대로 사용)
    CHART = ["D4001E", "2563EB", "16A34A", "EA580C", "7C3AED", "0891B2", "D97706"]


def rgb(h):
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))


# ══════════════════════════════════════════════════════════════
# 레이아웃 상수 (13.33" × 7.5")
# ══════════════════════════════════════════════════════════════
class Layout:
    W = 13.33   # 슬라이드 폭
    H = 7.5     # 슬라이드 높이

    # 영역 구분
    HEADER_H   = 0.95   # 헤더 높이
    FOOTER_Y   = 7.05   # 푸터 시작 Y
    FOOTER_H   = 0.45   # 푸터 높이
    CONTENT_Y  = 1.05   # 콘텐츠 시작 Y
    CONTENT_H  = 5.95   # 콘텐츠 높이

    # 여백
    MARGIN_L   = 0.4
    MARGIN_R   = 0.4
    CONTENT_W  = W - MARGIN_L - MARGIN_R  # 12.53"

    # 컬럼 레이아웃
    COL2_W     = (CONTENT_W - 0.2) / 2    # 2컬럼 각 너비
    COL3_W     = (CONTENT_W - 0.4) / 3    # 3컬럼 각 너비
    COL2_GAP   = 0.2
    COL3_GAP   = 0.2


# ══════════════════════════════════════════════════════════════
# 기본 프레젠테이션 생성
# ══════════════════════════════════════════════════════════════
def new_prs(theme="dark"):
    """새 프레젠테이션 생성 (16:9 와이드스크린)"""
    prs = Presentation()
    prs.slide_width  = Inches(Layout.W)
    prs.slide_height = Inches(Layout.H)
    return prs

def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


# ══════════════════════════════════════════════════════════════
# 도형 빌더
# ══════════════════════════════════════════════════════════════
def rect(slide, l, t, w, h, fill=C.NAVY, line=None, radius=0):
    """사각형 추가. l/t/w/h = inches"""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def line_shape(slide, l, t, w, h, color=C.RED, width_pt=1.5):
    """선 추가"""
    from pptx.util import Pt
    connector = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    connector.fill.background()
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width_pt)
    return connector


def tbox(slide, text, l, t, w, h,
         size=11, bold=False, color=C.WHITE,
         align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """텍스트박스 추가"""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    run.font.name = "Calibri"
    return tb


def tbox_multi(slide, lines, l, t, w, h, wrap=True):
    """
    여러 줄 텍스트박스.
    lines = [(text, size, bold, color, align, spacing_before), ...]
    """
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    for i, line_def in enumerate(lines):
        text, size, bold, color = line_def[:4]
        align = line_def[4] if len(line_def) > 4 else PP_ALIGN.LEFT
        spc   = line_def[5] if len(line_def) > 5 else 0

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if spc:
            p.space_before = Pt(spc)

        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        run.font.name = "Calibri"
    return tb


def shape_text(shape, text, size=11, bold=False, color=C.WHITE,
               align=PP_ALIGN.CENTER, wrap=True):
    """기존 shape에 텍스트 설정"""
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    # 기존 runs 제거
    for run in p.runs:
        run.text = ""
    if not p.runs:
        run = p.add_run()
    else:
        run = p.runs[0]
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    run.font.name = "Calibri"


def card(slide, l, t, w, h, fill=C.NAVY_MID,
         title="", title_size=11, title_color=C.GRAY_400,
         value="", value_size=28, value_color=C.WHITE,
         sub="", sub_size=9, sub_color=C.GRAY_400,
         line_color=None):
    """KPI 카드"""
    bg = rect(slide, l, t, w, h, fill=fill)
    if line_color:
        accent = rect(slide, l, t, 0.04, h, fill=line_color)

    title_l = l + (0.1 if line_color else 0.15)
    if title:
        tbox(slide, title, title_l, t+0.1, w-0.2, 0.25,
             size=title_size, color=title_color, align=PP_ALIGN.LEFT)
    if value:
        tbox(slide, value, title_l, t+0.38, w-0.2, h*0.45,
             size=value_size, bold=True, color=value_color, align=PP_ALIGN.LEFT)
    if sub:
        tbox(slide, sub, title_l, t+h-0.32, w-0.2, 0.28,
             size=sub_size, color=sub_color, align=PP_ALIGN.LEFT)
    return bg


def data_table(slide, l, t, w, headers, rows,
               col_widths=None, header_fill=C.RED,
               alt_fill=C.NAVY_MID, base_fill=C.NAVY_LIGHT,
               font_size=9, header_size=9):
    """
    데이터 테이블 (pptx native table)
    headers: [str, ...]
    rows: [[val, ...], ...]
    col_widths: [inches, ...] (합이 w여야 함)
    """
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header
    row_h = 0.32

    if not col_widths:
        col_widths = [w / n_cols] * n_cols

    table = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(l), Inches(t),
        Inches(w), Inches(row_h * n_rows)
    ).table

    # 컬럼 너비 설정
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = Inches(cw)

    # 행 높이
    for ri in range(n_rows):
        table.rows[ri].height = Inches(row_h)

    # 헤더 행
    for ci, h_text in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(header_fill)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(h_text)
        run.font.size = Pt(header_size)
        run.font.bold = True
        run.font.color.rgb = rgb(C.WHITE)
        run.font.name = "Calibri"
        # border
        _set_cell_border(cell)

    # 데이터 행
    for ri, row in enumerate(rows):
        fill_color = alt_fill if ri % 2 == 0 else base_fill
        for ci, val in enumerate(row):
            cell = table.cell(ri+1, ci)

            # 색상 오버라이드 (tuple인 경우)
            if isinstance(val, tuple):
                display_val, cell_fill = val[0], val[1]
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(cell_fill)
            else:
                display_val = val
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(fill_color)

            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(display_val)
            run.font.size = Pt(font_size)
            run.font.color.rgb = rgb(C.WHITE)
            run.font.name = "Calibri"
            _set_cell_border(cell)

    return table


def _set_cell_border(cell, color="2D3748", width=Pt(0.5)):
    """테이블 셀 테두리 설정"""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for tag in ["a:lnL","a:lnR","a:lnT","a:lnB"]:
        existing = tcPr.find(qn(tag))
        if existing is not None:
            tcPr.remove(existing)
        ln = etree.SubElement(tcPr, qn(tag))
        ln.set("w", str(int(width)))
        solidFill = etree.SubElement(ln, qn("a:solidFill"))
        srgb = etree.SubElement(solidFill, qn("a:srgbClr"))
        srgb.set("val", color)


# ══════════════════════════════════════════════════════════════
# 슬라이드 레이아웃 템플릿
# ══════════════════════════════════════════════════════════════
def slide_background(slide, style="dark"):
    """슬라이드 배경 설정"""
    if style == "dark":
        rect(slide, 0, 0, Layout.W, Layout.H, fill=C.NAVY)
    elif style == "dark2":
        rect(slide, 0, 0, Layout.W, Layout.H, fill="0A1628")
    elif style == "light":
        rect(slide, 0, 0, Layout.W, Layout.H, fill=C.LIGHT_BG)
    elif style == "white":
        rect(slide, 0, 0, Layout.W, Layout.H, fill=C.WHITE)


def slide_header(slide, title, subtitle="", style="dark"):
    """표준 슬라이드 헤더"""
    # 배경
    slide_background(slide, style)

    # 헤더 바
    rect(slide, 0, 0, Layout.W, Layout.HEADER_H, fill=C.NAVY_MID)

    # 레드 액센트 라인
    rect(slide, 0, Layout.HEADER_H - 0.04, Layout.W, 0.04, fill=C.RED)

    # 제목
    tbox(slide, title,
         Layout.MARGIN_L, 0.12,
         Layout.CONTENT_W - 2.5, 0.65,
         size=20, bold=True, color=C.WHITE)

    # 부제목 (우측 정렬)
    if subtitle:
        tbox(slide, subtitle,
             10.0, 0.35, 3.0, 0.45,
             size=9, color=C.GRAY_400, align=PP_ALIGN.RIGHT)


def slide_footer(slide, source="", page_num="", confidential=True):
    """표준 슬라이드 푸터"""
    rect(slide, 0, Layout.FOOTER_Y, Layout.W, Layout.FOOTER_H, fill="0A1628")

    # 구분선
    rect(slide, 0, Layout.FOOTER_Y, Layout.W, 0.02, fill=C.NAVY_LIGHT)

    # 출처
    if source:
        tbox(slide, f"Source: {source}",
             Layout.MARGIN_L, Layout.FOOTER_Y + 0.08,
             8.0, 0.3,
             size=7.5, color=C.GRAY_400, italic=True)

    # Confidential
    if confidential:
        tbox(slide, "CONFIDENTIAL — Internal Use Only",
             5.5, Layout.FOOTER_Y + 0.08, 4.5, 0.3,
             size=7.5, color=C.GRAY_600, align=PP_ALIGN.CENTER)

    # 날짜 + 페이지
    date_str = datetime.date.today().strftime("%Y.%m.%d")
    tbox(slide, f"{date_str}  |  {page_num}",
         Layout.W - 2.5, Layout.FOOTER_Y + 0.08, 2.2, 0.3,
         size=7.5, color=C.GRAY_400, align=PP_ALIGN.RIGHT)

    # SanDisk 워터마크
    tbox(slide, "SanDisk",
         Layout.W - 2.5, Layout.FOOTER_Y + 0.24, 2.2, 0.2,
         size=7, bold=True, color=C.RED, align=PP_ALIGN.RIGHT)


def divider(slide, y, color=C.NAVY_LIGHT, width=Layout.CONTENT_W):
    """수평 구분선"""
    rect(slide, Layout.MARGIN_L, y, width, 0.015, fill=color)


# ══════════════════════════════════════════════════════════════
# 차트 빌더 (matplotlib 없이 shape 기반)
# ══════════════════════════════════════════════════════════════
def bar_chart(slide, data, l, t, w, h,
              max_val=None, show_labels=True, show_grid=True,
              bar_color=C.RED, bg_color=C.NAVY_MID,
              label_color=C.GRAY_400, value_color=C.WHITE):
    """
    data = [(label, value), ...]
    세로 막대 차트 (shape 기반)
    """
    n = len(data)
    if not max_val:
        max_val = max(v for _, v in data) * 1.15 if data else 1

    chart_bg = rect(slide, l, t, w, h, fill=bg_color)

    total_bar_w = w * 0.85
    bar_w = total_bar_w / n * 0.6
    gap_w = total_bar_w / n * 0.4
    x_start = l + w * 0.075

    # 그리드 라인
    if show_grid:
        for pct in [0.25, 0.5, 0.75, 1.0]:
            gy = t + h * 0.9 - (pct * h * 0.82)
            rect(slide, l + 0.05, gy, w - 0.1, 0.008, fill="1E3048")

    for i, (label, value) in enumerate(data):
        bh = (value / max_val) * h * 0.82 if max_val > 0 else 0
        bx = x_start + i * (bar_w + gap_w)
        by = t + h * 0.9 - bh

        col = bar_color if not isinstance(bar_color, list) else bar_color[i % len(bar_color)]
        rect(slide, bx, by, bar_w, bh, fill=col)

        if show_labels:
            # 값 레이블 (바 위)
            if value > 0:
                tbox(slide, _fmt_val(value),
                     bx - 0.1, by - 0.3, bar_w + 0.2, 0.28,
                     size=8, bold=True, color=value_color, align=PP_ALIGN.CENTER)
            # X축 레이블
            tbox(slide, str(label),
                 bx - 0.1, t + h * 0.92, bar_w + 0.2, 0.28,
                 size=7.5, color=label_color, align=PP_ALIGN.CENTER)

    return chart_bg


def stacked_bar(slide, data, labels, l, t, w, h,
                colors=None, max_val=None,
                show_total=True, show_legend=True,
                bg_color=C.NAVY_MID):
    """
    누적 막대 차트
    data = {series_name: [val1, val2, ...], ...}
    labels = [x1_label, x2_label, ...]
    """
    n = len(labels)
    series_names = list(data.keys())
    if not colors:
        colors = C.CHART
    if not max_val:
        totals = [sum(data[s][i] for s in series_names) for i in range(n)]
        max_val = max(totals) * 1.15 if totals else 1

    rect(slide, l, t, w, h, fill=bg_color)

    total_bar_w = w * 0.82
    bar_w = total_bar_w / n * 0.55
    gap_w = total_bar_w / n * 0.45
    x_start = l + w * 0.09

    for i in range(n):
        cumul = 0
        bx = x_start + i * (bar_w + gap_w)
        total = sum(data[s][i] for s in series_names)

        for si, sname in enumerate(series_names):
            val = data[sname][i]
            if val == 0:
                cumul += val
                continue
            seg_h = (val / max_val) * h * 0.80
            seg_y = t + h * 0.88 - (cumul + val) / max_val * h * 0.80
            rect(slide, bx, seg_y, bar_w, seg_h,
                 fill=colors[si % len(colors)])
            cumul += val

        # Total 레이블
        if show_total:
            total_h = (total / max_val) * h * 0.80
            ty = t + h * 0.88 - total_h - 0.28
            tbox(slide, _fmt_val(total),
                 bx - 0.05, ty, bar_w + 0.1, 0.25,
                 size=7.5, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

        # X 레이블
        tbox(slide, str(labels[i]),
             bx - 0.1, t + h * 0.90, bar_w + 0.2, 0.28,
             size=7.5, color=C.GRAY_400, align=PP_ALIGN.CENTER)

    # 범례
    if show_legend:
        leg_y = t + 0.05
        for si, sname in enumerate(series_names):
            lx = l + 0.15 + si * (w / len(series_names))
            rect(slide, lx, leg_y, 0.18, 0.18, fill=colors[si % len(colors)])
            tbox(slide, sname, lx + 0.22, leg_y, 1.5, 0.2,
                 size=7.5, color=C.GRAY_400)


def waterfall_chart(slide, items, l, t, w, h,
                    pos_color=C.BLUE, neg_color=C.RED,
                    total_color=C.ORANGE, bg_color=C.NAVY_MID):
    """
    워터폴 차트
    items = [(label, value, type)] type: "start"|"add"|"sub"|"total"
    """
    rect(slide, l, t, w, h, fill=bg_color)
    n = len(items)
    bar_w = w * 0.7 / n
    gap_w = w * 0.3 / n
    x_start = l + w * 0.05

    max_val = sum(abs(v) for _, v, _ in items) * 1.2
    chart_bot = t + h * 0.88
    chart_h_px = h * 0.80

    cumul = 0
    for i, (label, val, typ) in enumerate(items):
        bx = x_start + i * (bar_w + gap_w)

        if typ == "total":
            bh = abs(cumul) / max_val * chart_h_px
            by = chart_bot - bh
            rect(slide, bx, by, bar_w, bh, fill=total_color)
            tbox(slide, _fmt_val(abs(cumul)),
                 bx-0.05, by-0.28, bar_w+0.1, 0.25,
                 size=8, bold=True, color=total_color, align=PP_ALIGN.CENTER)
        elif typ in ("start", "add"):
            bh = abs(val) / max_val * chart_h_px
            by = chart_bot - (cumul + val) / max_val * chart_h_px if typ == "add" else chart_bot - bh
            if typ == "add": by = chart_bot - (cumul + val) / max_val * chart_h_px
            rect(slide, bx, by, bar_w, bh, fill=pos_color if val > 0 else neg_color)
            tbox(slide, ("+" if val > 0 else "") + _fmt_val(val),
                 bx-0.05, by-0.28, bar_w+0.1, 0.25,
                 size=8, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
            cumul += val
        elif typ == "sub":
            bh = abs(val) / max_val * chart_h_px
            by = chart_bot - cumul / max_val * chart_h_px
            rect(slide, bx, by, bar_w, bh, fill=neg_color)
            tbox(slide, _fmt_val(val),
                 bx-0.05, by-0.28, bar_w+0.1, 0.25,
                 size=8, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
            cumul += val

        # X 레이블
        tbox(slide, label,
             bx - 0.05, t + h * 0.91, bar_w + 0.1, 0.3,
             size=7.5, color=C.GRAY_400, align=PP_ALIGN.CENTER)


def horizontal_bar(slide, data, l, t, w, h,
                   max_val=None, bar_color=C.RED, bg_color=C.NAVY_MID,
                   show_pct=False):
    """
    수평 막대 차트 (시장점유율 등)
    data = [(label, value, [color]), ...]
    """
    n = len(data)
    if not max_val:
        max_val = max(v for item in data for v in [item[1]]) * 1.1

    rect(slide, l, t, w, h, fill=bg_color)
    row_h = (h - 0.1) / n
    bar_max_w = w * 0.55

    for i, item in enumerate(data):
        label = item[0]; val = item[1]
        col = item[2] if len(item) > 2 else bar_color

        y = t + 0.05 + i * row_h
        bw = val / max_val * bar_max_w if max_val > 0 else 0
        label_x = l + 0.1

        # 레이블
        tbox(slide, str(label), label_x, y + 0.04, 2.8, row_h - 0.08,
             size=8.5, color=C.WHITE, align=PP_ALIGN.LEFT)

        # 바
        bar_x = l + 3.1
        if bw > 0:
            rect(slide, bar_x, y + 0.06, bw, row_h * 0.65, fill=col)

        # 값
        val_str = f"{val:.0f}%" if show_pct else _fmt_val(val)
        tbox(slide, val_str,
             bar_x + bw + 0.05, y + 0.04, 0.8, row_h - 0.08,
             size=8.5, bold=True, color=col, align=PP_ALIGN.LEFT)


def _fmt_val(val):
    """숫자 포맷팅"""
    if isinstance(val, str): return val
    if abs(val) >= 1000:
        return f"${val/1000:.1f}B" if abs(val) >= 1000 else f"${val:.0f}M"
    if abs(val) >= 100:
        return f"${val:.0f}M"
    return f"${val:.1f}M"


# ══════════════════════════════════════════════════════════════
# 타이틀 슬라이드
# ══════════════════════════════════════════════════════════════
def title_slide(prs, title, subtitle, author="SanDisk Marketing Team",
                date=None, deck_type=""):
    """전문적인 타이틀 슬라이드"""
    sl = new_slide(prs)

    # 배경: 그라데이션 효과를 레이어로
    rect(sl, 0, 0, Layout.W, Layout.H, fill="0A1628")
    rect(sl, 0, 0, Layout.W/2, Layout.H, fill="0F1E35")  # 왼쪽 패널
    rect(sl, 0, 0, 0.06, Layout.H, fill=C.RED)            # 레드 사이드바

    # 슬로건 영역
    rect(sl, 0.06, 3.5, Layout.W - 0.06, 0.04, fill=C.RED)

    # 제목
    tbox(sl, title,
         0.6, 1.0, 8.0, 2.2,
         size=40, bold=True, color=C.WHITE)

    # 부제목
    tbox(sl, subtitle,
         0.6, 3.2, 9.0, 0.7,
         size=18, color=C.GRAY_400)

    # 메타 정보
    date_str = (date or datetime.date.today()).strftime("%B %Y")
    tbox(sl, f"{author}  |  {date_str}",
         0.6, 3.65, 8.0, 0.4,
         size=11, color=C.GRAY_400)

    # Deck type 뱃지
    if deck_type:
        rect(sl, 0.6, 4.2, 2.5, 0.45, fill=C.RED)
        tbox(sl, deck_type, 0.6, 4.2, 2.5, 0.45,
             size=12, bold=True, align=PP_ALIGN.CENTER)

    # 기밀 표시
    rect(sl, 0, Layout.H - 0.55, Layout.W, 0.55, fill="060E1A")
    tbox(sl, "CONFIDENTIAL — Internal Use Only  |  Not for Distribution",
         0.4, Layout.H - 0.45, Layout.W - 0.8, 0.35,
         size=9, color=C.GRAY_600, align=PP_ALIGN.CENTER)
    tbox(sl, "SanDisk Corporation",
         Layout.W - 2.8, Layout.H - 0.45, 2.5, 0.35,
         size=9, bold=True, color=C.RED, align=PP_ALIGN.RIGHT)

    return sl


# ══════════════════════════════════════════════════════════════
# 특수 슬라이드 컴포넌트
# ══════════════════════════════════════════════════════════════
def exec_summary_block(slide, findings, l, t, w, h):
    """
    Executive Summary 박스
    findings = [(icon, title, text), ...]
    """
    n = len(findings)
    block_h = (h - (n-1)*0.1) / n

    for i, (icon, title, text) in enumerate(findings):
        by = t + i * (block_h + 0.1)
        rect(slide, l, by, 0.04, block_h, fill=C.RED)
        rect(slide, l+0.04, by, w-0.04, block_h, fill=C.NAVY_MID)

        tbox_multi(slide, [
            (f"{icon}  {title}", 10, True, C.WHITE, PP_ALIGN.LEFT),
            (text, 9, False, C.GRAY_400, PP_ALIGN.LEFT, 2),
        ], l+0.2, by+0.06, w-0.3, block_h-0.1)


def insight_tag(slide, text, l, t, color=C.RED):
    """인사이트 태그"""
    rect(slide, l, t, 0.15, 0.25, fill=color)
    tbox(slide, text, l+0.2, t, 3.0, 0.28,
         size=8.5, bold=True, color=color)


def source_note(slide, text, page=""):
    """슬라이드 하단 출처 노트"""
    slide_footer(slide, source=text, page_num=page)
