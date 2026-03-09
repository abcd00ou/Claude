"""
마케팅 커뮤니케이션 VP 보고서 — MarCom Plan
SanDisk B2C Storage Marketing Team
"""
import sys
sys.path.insert(0, str(__import__('pathlib').Path(__file__).parent))

import json, datetime
from pathlib import Path
from pptx.enum.text import PP_ALIGN

import ppt_design as D
from ppt_design import C, Layout, rect, tbox, tbox_multi, card, data_table
from ppt_design import slide_header, slide_footer, bar_chart, stacked_bar
from ppt_design import horizontal_bar, exec_summary_block, waterfall_chart
from config import PPTX_DIR, DATA_DIR

with open(DATA_DIR / "market_data.json")    as f: MKT  = json.load(f)
with open(DATA_DIR / "internal_sales.json") as f: SALE = json.load(f)

TODAY = datetime.date.today()
PAGE_PREFIX = "MarCom Plan"


def build():
    prs = D.new_prs()
    pg = 1

    # ══════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, Layout.W, Layout.H, C.NAVY)
    rect(sl, 0, 0, Layout.W, 3.8, C.NAVY_MID)
    rect(sl, 0, 3.8, Layout.W, 0.06, C.RED)
    rect(sl, 0, 0, 0.45, Layout.H, C.RED)
    tbox(sl, "SanDisk B2C Storage\nMarketing Communications Plan",
         1.0, 0.9, 11.5, 2.0,
         size=32, bold=True, color=C.WHITE, align=PP_ALIGN.LEFT)
    tbox(sl, "FY2025 브랜드 전략  ·  4대 캠페인  ·  미디어 믹스  ·  KPI 프레임워크",
         1.0, 3.0, 11.5, 0.65,
         size=13, color="A0AEC0", align=PP_ALIGN.LEFT)
    tbox(sl, "SanDisk B2C Marketing Communications  |  CONFIDENTIAL — Internal Use Only",
         1.0, 4.15, 11.5, 0.5, size=11, color="6B7280")
    tbox(sl, f"Prepared: {TODAY}  |  Total MarCom Budget: $458M FY2025E",
         1.0, 4.65, 11.5, 0.5, size=10, color="6B7280")

    # ══════════════════════════════════════════════════════════
    # SLIDE 2: Executive Summary
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Executive Summary — MarCom FY2025", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: Internal MarCom Budget Plan FY2025, WD Marketing Annual Report 2024", pg)

    kpi_data = [
        ("$458M",  "FY2025 MarCom Budget", "+12.3% vs 2024"),
        ("3.1×",   "Blended Campaign ROI",  "목표: 3.0× 이상"),
        ("4대",    "전략 캠페인",            "Premium·Creator·Back-to-School·Holiday"),
        ("34%",    "Aided Brand Awareness",  "+4pp 목표 (2024: 30%)"),
    ]
    for i, (val, label, sub) in enumerate(kpi_data):
        x = Layout.MARGIN_L + i * 3.13
        lc = [C.RED, C.ORANGE, C.BLUE, C.GREEN][i]
        card(sl, x, 1.05, 2.95, 1.35, fill=C.NAVY_MID,
             title=label, title_size=9, title_color=C.GRAY_400,
             value=val, value_size=22, value_color=C.WHITE,
             sub=sub, sub_size=8, sub_color=lc, line_color=lc)

    exec_summary_block(sl, [
        ("🎯", "FY2025 마케팅 핵심 목표",
         "External SSD 시장점유율 22%→24% 달성. microSD #1 수성(32%→34%). "
         "Internal SSD 신규 AI PC 연계 마케팅으로 +16.1% 성장 지원."),
        ("📢", "4대 캠페인 집중 전략",
         "① Premium Launch(Q1): Extreme Pro 4TB · ② Creator Summer(Q2): microSD+External SSD · "
         "③ Back-to-School(Q3): $52M ROI 2.89× · ④ Holiday Peak(Q4): $45M → $138M 매출 (ROI 3.07×)"),
        ("📊", "디지털 퍼스트 미디어 믹스",
         "Digital 62% / TV·OTT 18% / Retail 14% / PR 6%. "
         "Amazon DSP + Google PMax + Meta Advantage+ 통합 운영. "
         "실적: 2024 Holiday ROI 3.07×, BTS ROI 2.89×"),
    ], Layout.MARGIN_L, 2.55, Layout.CONTENT_W, 4.1)

    # ══════════════════════════════════════════════════════════
    # SLIDE 3: Brand Architecture
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Brand Architecture — SanDisk Portfolio", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: WD Brand Guidelines 2024, Internal Brand Tracking Study Q4 2024", pg)

    # 브랜드 피라미드 시각화
    rect(sl, Layout.MARGIN_L, 1.1, 7.5, 0.35, C.NAVY_MID)
    tbox(sl, "SanDisk 브랜드 포트폴리오 아키텍처",
         Layout.MARGIN_L+0.1, 1.1, 7.3, 0.35,
         size=10, bold=True, color=C.WHITE)

    brands = [
        ("SanDisk Extreme Pro",  C.RED,      "프리미엄 퍼포먼스",
         "타겟: 프로 크리에이터·포토그래퍼\nExternal SSD 4TB + microSD Pro 512G~1TB\nGM: 52-65% (최고 마진 구간)\nASP: $120~303"),
        ("SanDisk Extreme",      "E65C00",   "고성능 메인스트림",
         "타겟: 액티브 크리에이터·게이머\nExternal SSD 1-2TB + microSD 128G-512G\nGM: 45-62% (핵심 볼륨 드라이버)\nASP: $26~154"),
        ("WD Black",             C.NAVY_MID, "PC 고성능 내장",
         "타겟: 게이머·파워유저\nNVMe SN850X 1-4TB + SN770 1-2TB\nGM: 42-47% (Internal SSD 주력)\nASP: $62~266"),
        ("WD My Passport / Ultra","4B5563",  "일반 소비자 보급형",
         "타겟: 일상 스토리지 사용자\nExternal HDD+SSD 1-2TB\nGM: 38-40% (볼륨 기여)\nASP: $72~123"),
    ]

    for i, (name, color, pos, desc) in enumerate(brands):
        x = Layout.MARGIN_L + i * 1.9
        w = 1.85 - i * 0.05
        rect(sl, x, 1.55, w, 0.38, color)
        tbox(sl, name, x+0.05, 1.55, w-0.1, 0.38,
             size=8.5, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
        tbox(sl, pos, x+0.05, 1.96, w-0.1, 0.28,
             size=8, color=color, align=PP_ALIGN.CENTER, bold=True)
        rect(sl, x, 2.28, w, 3.0, "F8F9FA")
        tbox(sl, desc, x+0.08, 2.35, w-0.16, 2.85,
             size=8, color=C.GRAY_600)

    # 우측: 브랜드 KPI
    rect(sl, 8.15, 1.1, 4.98, 0.35, C.NAVY_MID)
    tbox(sl, "브랜드 Health Tracker (2024 기준)",
         8.25, 1.1, 4.78, 0.35, size=10, bold=True, color=C.WHITE)

    brand_kpis = [
        ("Aided Brand Awareness",    "30%",   "+4pp 목표 (2025E: 34%)",    C.BLUE),
        ("Unaided Brand Recall",     "18%",   "Samsung 22% 추격 중",       C.ORANGE),
        ("Brand Preference (ExSSD)", "22%",   "시장점유율 동일 수준 확인",  C.GREEN),
        ("Net Promoter Score",       "42",    "+5pt 목표 (B2C 스토리지)",   C.BLUE),
        ("Digital Share of Voice",   "28%",   "카테고리 내 검색 점유",      C.RED),
        ("Amazon Search Rank #1",    "8SKU",  "External SSD 카테고리 내",   C.GREEN),
    ]

    for i, (metric, val, note, color) in enumerate(brand_kpis):
        y = 1.55 + i * 0.82
        rect(sl, 8.15, y, 4.98, 0.75, "F8F9FA" if i % 2 == 0 else "FFFFFF")
        rect(sl, 8.15, y, 0.08, 0.75, color)
        tbox(sl, metric, 8.3, y+0.05, 2.5, 0.3, size=8.5, bold=True, color=C.NAVY)
        tbox(sl, val, 10.85, y+0.05, 1.0, 0.35, size=14, bold=True, color=color, align=PP_ALIGN.RIGHT)
        tbox(sl, note, 8.3, y+0.38, 3.75, 0.3, size=7.5, color=C.GRAY_600)

    # ══════════════════════════════════════════════════════════
    # SLIDE 4: 4대 전략 캠페인 Overview
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "FY2025 Campaign Calendar — 4대 전략 캠페인", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: Internal MarCom Plan FY2025, 2024 Campaign ROI Data (Holiday ROI 3.07×, BTS ROI 2.89×)", pg)

    # 타임라인 헤더
    months = ["Jan", "Feb", "Mar", "Apr", "May", "Jun",
              "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
    timeline_x = Layout.MARGIN_L
    timeline_w = Layout.CONTENT_W / 12
    rect(sl, timeline_x, 1.1, Layout.CONTENT_W, 0.35, C.NAVY_MID)
    for i, m in enumerate(months):
        tbox(sl, m, timeline_x + i * timeline_w, 1.12,
             timeline_w, 0.3, size=8.5, color=C.WHITE, align=PP_ALIGN.CENTER)
        if i > 0:
            rect(sl, timeline_x + i * timeline_w, 1.1, 0.01, 0.35, "4B5563")

    # Q 구분선
    for qi, (ql, qx) in enumerate([("Q1", 0), ("Q2", 3), ("Q3", 6), ("Q4", 9)]):
        rect(sl, timeline_x + qx * timeline_w, 1.45, 3 * timeline_w, 0.22, C.NAVY)
        tbox(sl, ql, timeline_x + qx * timeline_w, 1.45, 3 * timeline_w, 0.22,
             size=8, bold=True, color=C.GRAY_400, align=PP_ALIGN.CENTER)

    # 캠페인 바
    campaigns = [
        ("① Premium Launch: Extreme Pro 4TB", 0, 3, C.RED,
         "예산 $95M | 타겟: 프로 크리에이터·파워유저 | 채널: YouTube·Instagram·Amazon DSP"),
        ("② Creator Summer: microSD + External SSD", 3, 3, C.ORANGE,
         "예산 $82M | 타겟: 여름 여행·Creator | 채널: TikTok·YouTube·Influencer"),
        ("③ Back-to-School: Student+Gamer Pack", 6, 3, C.BLUE,
         "예산 $78M (2024 ROI 2.89×) | 타겟: 대학생·게이머 | 채널: Amazon·Reddit·Discord"),
        ("④ Holiday Peak: All Category Push", 9, 3, "1A472A",
         "예산 $203M (2024 ROI 3.07×) | 타겟: 전체 소비자 | 채널: TV·Amazon·Retail Display"),
    ]

    for i, (name, start, dur, color, desc) in enumerate(campaigns):
        y_bar = 1.75 + i * 1.3
        bx = timeline_x + start * timeline_w
        bw = dur * timeline_w
        rect(sl, bx, y_bar, bw, 0.42, color)
        tbox(sl, name, bx+0.08, y_bar+0.06, bw-0.16, 0.32,
             size=8.5, bold=True, color=C.WHITE)
        tbox(sl, desc, bx+0.08, y_bar+0.5, Layout.CONTENT_W - start*timeline_w - 0.1, 0.32,
             size=7.5, color=C.GRAY_600)

    # 총 예산 요약
    rect(sl, Layout.MARGIN_L, 7.0, Layout.CONTENT_W, 0.25, C.NAVY_MID)
    tbox(sl, "FY2025 Total MarCom Budget:  $458M  ·  Digital 62% / TV+OTT 18% / Retail 14% / PR 6%",
         Layout.MARGIN_L+0.15, 7.02, Layout.CONTENT_W-0.3, 0.22,
         size=9, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    # SLIDE 5: 캠페인 ① Premium Launch Detail
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Campaign ①: Premium Launch — Extreme Pro 4TB", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: Internal Campaign Brief Q1 2025, GfK Consumer Insights 2024", pg)

    rect(sl, Layout.MARGIN_L, 1.1, Layout.CONTENT_W, 0.38, C.RED)
    tbox(sl, "Budget $95M  ·  기간 Jan~Mar 2025  ·  목표 SKU: Extreme Pro 4TB External SSD + Pro Plus 512G microSD",
         Layout.MARGIN_L+0.15, 1.13, Layout.CONTENT_W-0.3, 0.32,
         size=9.5, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    sections = [
        ("타겟 오디언스", C.RED, [
            "Primary: 프로 사진작가·비디오그래퍼 (25-45세)",
            "Secondary: Content Creator·유튜버 (20-35세)",
            "3rd: 고성능 워크플로우 파워유저",
            "지역: NA 45% / EMEA 35% / APAC 20%",
        ]),
        ("핵심 메시지", C.ORANGE, [
            "\"Ultra Speed for Ultra Creativity\"",
            "4TB 초대용량 + 2000MB/s 전송속도",
            "ProGrade 인증 / IP65 방수방진",
            "Adobe Creative Cloud 공식 파트너",
        ]),
        ("채널 믹스 ($95M)", C.BLUE, [
            "YouTube Pre-roll: $22M (23%)",
            "Instagram/Meta: $18M (19%)",
            "Amazon DSP: $25M (26%)",
            "Influencer (500K+): $15M (16%)",
            "PR + Seeding: $8M (8%)",
            "Retail Display: $7M (7%)",
        ]),
        ("KPI", C.GREEN, [
            "매출 기여: +$190M (ROI 2.0×)",
            "Extreme Pro 4TB 인지도: 45%→65%",
            "Amazon BSR #1 목표 (External SSD 4TB)",
            "YouTube View-through Rate: 30%+",
            "Creator 파트너십: 2,500개 콘텐츠",
            "신규 고객 획득: +85K",
        ]),
    ]

    for i, (title, color, items) in enumerate(sections):
        col = i % 2
        row = i // 2
        x = Layout.MARGIN_L + col * 6.33
        y = 1.6 + row * 2.65
        rect(sl, x, y, 6.15, 0.35, color)
        tbox(sl, title, x+0.1, y, 5.95, 0.35,
             size=10, bold=True, color=C.WHITE)
        rect(sl, x, y+0.35, 6.15, 2.2, "F8F9FA")
        for j, item in enumerate(items):
            tbox(sl, f"• {item}", x+0.12, y+0.45 + j*0.33, 5.9, 0.31,
                 size=9, color=C.GRAY_600)

    # ══════════════════════════════════════════════════════════
    # SLIDE 6: 캠페인 ③ Back-to-School (실적 기반)
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Campaign ③: Back-to-School 2025 — 2024 ROI 2.89× 기반 플랜", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: 2024 BTS Campaign ROI Report (Internal), Amazon Gaming Category 2024", pg)

    # 2024 실적 요약
    rect(sl, Layout.MARGIN_L, 1.1, Layout.CONTENT_W, 0.38, C.BLUE)
    tbox(sl, "2024 BTS 실적: 투자 $18M → 매출 기여 $52M (ROI 2.89×)  →  2025: 예산 $78M으로 확대",
         Layout.MARGIN_L+0.15, 1.13, Layout.CONTENT_W-0.3, 0.32,
         size=9.5, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    # 좌측: 2025 실행 계획
    rect(sl, Layout.MARGIN_L, 1.58, 6.0, 0.3, C.NAVY_MID)
    tbox(sl, "2025 BTS 실행 계획 ($78M)", Layout.MARGIN_L+0.1, 1.6, 5.8, 0.25,
         size=9.5, bold=True, color=C.WHITE)

    bts_plan = [
        ("Phase 1: Awareness (Jul)", "$18M",
         "YouTube Campus 타겟 Pre-roll + TikTok Back-to-School Challenge"),
        ("Phase 2: Consideration (Aug)", "$28M",
         "Amazon Storefront 학생 전용 페이지 + 인플루언서 Unboxing 200개"),
        ("Phase 3: Conversion (Sep)", "$22M",
         "Amazon Prime Day 2025 연계 Flash Deal + Retail 번들 프로모션"),
        ("Phase 4: Retention", "$10M",
         "구매 후 리뷰 캠페인 + Discord/Reddit 게이밍 커뮤니티 유지"),
    ]
    for i, (phase, budget, desc) in enumerate(bts_plan):
        y = 1.95 + i * 1.1
        rect(sl, Layout.MARGIN_L, y, 0.5, 0.95, C.BLUE)
        tbox(sl, budget, Layout.MARGIN_L, y+0.2, 0.5, 0.5,
             size=10, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
        rect(sl, Layout.MARGIN_L+0.5, y, 5.5, 0.95, "F0F4FF" if i%2==0 else "F8F9FA")
        tbox(sl, phase, Layout.MARGIN_L+0.6, y+0.05, 5.25, 0.32,
             size=9.5, bold=True, color=C.BLUE)
        tbox(sl, desc, Layout.MARGIN_L+0.6, y+0.38, 5.25, 0.52,
             size=8.5, color=C.GRAY_600)

    # 우측: 타겟 SKU
    rect(sl, 6.6, 1.58, 6.33, 0.3, C.NAVY_MID)
    tbox(sl, "BTS 주요 추진 SKU", 6.7, 1.6, 6.13, 0.25,
         size=9.5, bold=True, color=C.WHITE)

    bts_skus = [
        ("WD Black SN850X 1TB",  "$77", "GM 42.8%", "게이머 #1 픽",         C.RED),
        ("WD Black SN850X 2TB",  "$141","GM 45.2%", "콘텐츠 크리에이터",     C.RED),
        ("WD Black SN770 1TB",   "$62", "GM 34.5%", "가성비 업그레이드",     C.BLUE),
        ("SD Extreme Micro 256G","$26", "GM 58.4%", "Switch·핸드헬드 게임",  C.ORANGE),
        ("SD Extreme Micro 512G","$43", "GM 62.1%", "4K 영상 제작 학생",     C.ORANGE),
        ("SD Pro+ Micro 512G",   "$58", "GM 64.8%", "Creator 최고 마진",     C.GREEN),
    ]

    sku_headers = ["SKU", "ASP", "GM%", "타겟", "라인"]
    sku_rows = []
    for name, asp, gm, tgt, color in bts_skus:
        sku_rows.append([name, asp, (gm, color), tgt, ""])
    data_table(sl, 6.6, 1.95, 6.33, sku_headers, sku_rows,
               col_widths=[2.2, 0.55, 0.7, 1.8, 0.68],
               header_fill=C.BLUE, font_size=8.5)

    # 2024 vs 2025 비교
    rect(sl, 6.6, 5.05, 6.33, 0.28, C.NAVY_MID)
    tbox(sl, "2024 실적 vs 2025 목표", 6.7, 5.08, 6.13, 0.22,
         size=8.5, bold=True, color=C.WHITE)
    comp_data = [
        ("예산",       "$18M",  "$78M",  "+333%"),
        ("예상 매출기여", "$52M",  "$220M", "+323%"),
        ("ROI",        "2.89×", "2.82×", "유지"),
        ("신규 고객",   "12K",   "44K",   "+267%"),
    ]
    comp_headers = ["항목", "2024 실적", "2025 목표", "변화"]
    comp_rows = [[m, a, (b, C.GREEN), c] for m, a, b, c in comp_data]
    data_table(sl, 6.6, 5.38, 6.33, comp_headers, comp_rows,
               col_widths=[1.8, 1.4, 1.5, 1.33], header_fill=C.NAVY, font_size=9)

    # ══════════════════════════════════════════════════════════
    # SLIDE 7: 캠페인 ④ Holiday Peak (최대 예산)
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Campaign ④: Holiday Peak — $203M 최대 투자 캠페인", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: 2024 Holiday Campaign ROI Report (Internal), Amazon Holiday 2024 Data", pg)

    rect(sl, Layout.MARGIN_L, 1.1, Layout.CONTENT_W, 0.38, "1A472A")
    tbox(sl, "2024 Holiday 실적: 투자 $45M → 매출 기여 $138M (ROI 3.07×) — 역대 최고 성과",
         Layout.MARGIN_L+0.15, 1.13, Layout.CONTENT_W-0.3, 0.32,
         size=9.5, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    # 3단 구성
    holiday_sections = [
        ("Black Friday (Nov)", "1A472A", [
            "예산: $68M",
            "Amazon Lightning Deal: All SKU",
            "TV: CTV/OTT $20M",
            "Google Shopping Priority 입찰",
            "Retail: Best Buy / Costco Feature",
            "목표 매출: $210M (+ROI 3.1×)",
        ]),
        ("Cyber Monday (Dec 1)", C.BLUE, [
            "예산: $45M",
            "Amazon Deal of the Day 확보",
            "Email/Push 기존 고객 재타겟",
            "Bundle 전략: External+microSD",
            "Social Flash Sale (12시간)",
            "목표 매출: $138M (+ROI 3.1×)",
        ]),
        ("Christmas Season (Dec)", "8B1A1A", [
            "예산: $90M",
            "Gift Guide: Top 선물 추천 등록",
            "YouTube 감성 브랜드 영상",
            "Retail 박스 패키징 업그레이드",
            "스토어 내 선물 포장 캠페인",
            "목표 매출: $280M (+ROI 3.1×)",
        ]),
    ]

    for i, (title, color, items) in enumerate(holiday_sections):
        x = Layout.MARGIN_L + i * 4.18
        rect(sl, x, 1.6, 4.0, 0.38, color)
        tbox(sl, title, x+0.1, 1.62, 3.8, 0.32,
             size=10, bold=True, color=C.WHITE)
        rect(sl, x, 1.98, 4.0, 3.2, "F8F9FA")
        for j, item in enumerate(items):
            tbox(sl, f"• {item}", x+0.12, 2.05+j*0.5, 3.76, 0.47,
                 size=9, color=C.GRAY_600)

    # 채널별 Holiday 예산 배분
    rect(sl, Layout.MARGIN_L, 5.28, Layout.CONTENT_W, 0.28, C.NAVY_MID)
    tbox(sl, "Holiday 채널별 예산 배분 ($203M)",
         Layout.MARGIN_L+0.1, 5.3, Layout.CONTENT_W-0.2, 0.22,
         size=9, bold=True, color=C.WHITE)

    ch_alloc = [
        ("TV/CTV/OTT", 44, C.NAVY_MID),
        ("Amazon DSP", 52, C.ORANGE),
        ("Google/Meta", 38, C.BLUE),
        ("Retail",      36, "1A472A"),
        ("Influencer",  22, C.PURPLE if hasattr(C, "PURPLE") else "7C3AED"),
        ("PR/Content",  11, C.GRAY_600),
    ]
    total_ch = sum(v for _, v, _ in ch_alloc)
    x_cursor = Layout.MARGIN_L
    for ch_name, ch_val, ch_color in ch_alloc:
        ch_w = (ch_val / total_ch) * Layout.CONTENT_W
        rect(sl, x_cursor, 5.6, ch_w, 0.55, ch_color)
        if ch_w > 0.8:
            tbox(sl, f"${ch_val}M", x_cursor+0.05, 5.65, ch_w-0.1, 0.25,
                 size=8.5, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
            tbox(sl, ch_name, x_cursor+0.05, 5.88, ch_w-0.1, 0.22,
                 size=7.5, color=C.WHITE, align=PP_ALIGN.CENTER)
        x_cursor += ch_w

    # ══════════════════════════════════════════════════════════
    # SLIDE 8: 미디어 믹스 전략
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Media Mix Strategy — Digital First FY2025", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: Internal MarCom Budget Plan FY2025, Nielsen Digital Ad Spend 2024", pg)

    # 미디어 믹스 파이 (bar로 시각화)
    rect(sl, Layout.MARGIN_L, 1.1, 6.5, 0.35, C.NAVY_MID)
    tbox(sl, "FY2025 미디어 채널 배분 ($458M)",
         Layout.MARGIN_L+0.1, 1.1, 6.3, 0.35,
         size=10, bold=True, color=C.WHITE)

    media_mix = [
        ("Digital",          62, 284, C.RED,     "Amazon DSP + Google PMax + Meta Advantage+ + YouTube"),
        ("TV / OTT",         18, 82,  C.BLUE,    "Connected TV 중심 전환 (Linear 30%→20%)"),
        ("Retail Media",     14, 64,  C.GREEN,   "Best Buy / Costco / Target 인스토어 디스플레이"),
        ("PR / Content",      6, 28,  C.ORANGE,  "리뷰 미디어 + Influencer Seeding + 언론 홍보"),
    ]

    for i, (ch, pct, budget, color, desc) in enumerate(media_mix):
        y = 1.55 + i * 1.3
        # 바 시각화
        bar_w = pct * 0.085
        rect(sl, Layout.MARGIN_L, y+0.02, bar_w, 0.55, color)
        tbox(sl, f"{pct}%", Layout.MARGIN_L + bar_w + 0.1, y+0.1, 0.55, 0.35,
             size=13, bold=True, color=color)
        tbox(sl, f"${budget}M", Layout.MARGIN_L + bar_w + 0.7, y+0.1, 0.8, 0.35,
             size=11, bold=True, color=C.GRAY_600)
        tbox(sl, ch, Layout.MARGIN_L, y-0.22, 2.0, 0.2,
             size=9, bold=True, color=C.NAVY)
        tbox(sl, desc, Layout.MARGIN_L + bar_w + 1.55, y+0.1, 4.5, 0.35,
             size=8.5, color=C.GRAY_600, italic=True)

    # 디지털 채널 세부 분류
    rect(sl, 7.1, 1.1, 6.03, 0.35, C.NAVY_MID)
    tbox(sl, "Digital $284M 세부 채널 배분",
         7.2, 1.1, 5.83, 0.35, size=10, bold=True, color=C.WHITE)

    digital_detail = [
        ("Amazon DSP",       90, "매출 직결 최우선",   C.ORANGE),
        ("Google PMax",      62, "검색+쇼핑 통합",     C.BLUE),
        ("Meta Advantage+",  48, "리타겟팅 강화",      C.NAVY_MID),
        ("YouTube",          42, "브랜드 인지도",      C.RED),
        ("TikTok",           24, "Gen-Z Creator",      "7C3AED"),
        ("기타 Digital",     18, "Reddit·Discord·etc", C.GRAY_600),
    ]
    dd_data = [(ch, v) for ch, v, _, _ in digital_detail]
    bar_chart(sl, dd_data, 7.1, 1.55, 6.03, 5.1,
              max_val=100, bar_color=C.RED, show_labels=True)

    # ══════════════════════════════════════════════════════════
    # SLIDE 9: KPI 대시보드
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "KPI Dashboard — FY2025 성과 측정 프레임워크", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: Internal KPI Framework 2025, Google Analytics, Amazon Vendor Central", pg)

    rect(sl, Layout.MARGIN_L, 1.1, Layout.CONTENT_W, 0.35, C.NAVY_MID)
    tbox(sl, "측정 프레임워크:  인지 → 고려 → 전환 → 충성도 단계별 KPI",
         Layout.MARGIN_L+0.15, 1.12, Layout.CONTENT_W-0.3, 0.3,
         size=10, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    funnel_stages = [
        ("인지\n(Awareness)", C.BLUE, [
            ("Aided Brand Awareness",   "30%→34%",   "Brand Tracking 분기"),
            ("Digital Share of Voice",  "28%→33%",   "Kantar SOV Monthly"),
            ("YouTube View Reach",      "180M/year",  "YouTube Analytics"),
            ("PR 미디어 노출",           "2,500건",   "Cision 모니터링"),
        ]),
        ("고려\n(Consideration)", C.ORANGE, [
            ("Amazon 검색 CTR",         "8.5%→10%",  "Vendor Central"),
            ("Product Page Conversion", "12%→15%",   "A+ 콘텐츠 A/B"),
            ("Review Rating",           "4.5★ 유지",  "모든 주요 SKU"),
            ("Influencer 콘텐츠 수",    "5,000건",   "Creator 파트너십"),
        ]),
        ("전환\n(Conversion)", C.RED, [
            ("Total MarCom ROI",        "3.0× 목표",  "Channel Attribution"),
            ("e-comm Revenue",          "+15% YoY",   "Amazon+Direct"),
            ("CPA (Cost per Acq.)",     "$38 이하",   "GA4 + AMC"),
            ("Cart Abandonment Rate",   "72%→65%",   "Retargeting 강화"),
        ]),
        ("충성도\n(Loyalty)", C.GREEN, [
            ("NPS (Net Promoter)",      "42→47",      "분기별 패널 조사"),
            ("Repeat Purchase Rate",    "22%→27%",   "CRM 분석"),
            ("Amazon 5★ Review 비율",  "78%→82%",   "Review Management"),
            ("Brand Community Size",    "+50K 팔로워", "YouTube+Discord"),
        ]),
    ]

    for i, (stage, color, kpis) in enumerate(funnel_stages):
        x = Layout.MARGIN_L + i * 3.15
        rect(sl, x, 1.55, 3.05, 0.45, color)
        tbox(sl, stage, x, 1.55, 3.05, 0.45,
             size=10, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
        for j, (metric, target, source) in enumerate(kpis):
            y = 2.1 + j * 1.2
            rect(sl, x, y, 3.05, 1.1, "F8F9FA" if j%2==0 else "FFFFFF")
            rect(sl, x, y, 0.06, 1.1, color)
            tbox(sl, metric, x+0.12, y+0.06, 2.82, 0.32,
                 size=8.5, bold=True, color=C.NAVY)
            tbox(sl, target, x+0.12, y+0.38, 2.82, 0.35,
                 size=11, bold=True, color=color)
            tbox(sl, f"측정: {source}", x+0.12, y+0.75, 2.82, 0.28,
                 size=7.5, color=C.GRAY_600, italic=True)

    # ══════════════════════════════════════════════════════════
    # SLIDE 10: 실행 계획 & Next Steps
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Execution Plan & Approval Request", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: Internal MarCom Plan FY2025, Finance Budget Approval Process", pg)

    rect(sl, Layout.MARGIN_L, 1.1, Layout.CONTENT_W, 0.38, C.RED)
    tbox(sl, "VP 승인 요청:  FY2025 MarCom Budget $458M (+12.3% YoY)  |  캠페인 4건  |  KPI 16개 지표",
         Layout.MARGIN_L+0.15, 1.13, Layout.CONTENT_W-0.3, 0.32,
         size=9.5, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    # 분기별 실행 로드맵
    quarters = [
        ("Q1 2025\nJan-Mar", C.RED, [
            "① Premium Launch 캠페인 실행",
            "Extreme Pro 4TB 런칭 이벤트",
            "Amazon A+ 콘텐츠 전면 업그레이드",
            "Brand Protection: MAP 정책 강화",
            "예산: $95M",
        ]),
        ("Q2 2025\nApr-Jun", C.ORANGE, [
            "② Creator Summer 캠페인",
            "TikTok Creator 파트너십 500개",
            "microSD 여행+Creator 패키지",
            "APAC 신흥시장 microSD 특공대",
            "예산: $82M",
        ]),
        ("Q3 2025\nJul-Sep", C.BLUE, [
            "③ Back-to-School 캠페인",
            "Amazon Prime Day 2025 연계",
            "Gaming 커뮤니티 (Discord/Reddit)",
            "학생 번들 프로모션 (SN850X+microSD)",
            "예산: $78M",
        ]),
        ("Q4 2025\nOct-Dec", "1A472A", [
            "④ Holiday Peak 캠페인 (최대 규모)",
            "Black Friday: $68M 집중 투자",
            "Cyber Monday: Deal of the Day 확보",
            "크리스마스: Gift Guide 상위 노출",
            "예산: $203M",
        ]),
    ]

    for i, (qtr, color, items) in enumerate(quarters):
        x = Layout.MARGIN_L + i * 3.15
        rect(sl, x, 1.6, 3.05, 0.45, color)
        tbox(sl, qtr, x, 1.6, 3.05, 0.45,
             size=10, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
        rect(sl, x, 2.05, 3.05, 3.5, "F8F9FA")
        for j, item in enumerate(items):
            tbox(sl, f"• {item}", x+0.1, 2.12 + j*0.65, 2.85, 0.6,
                 size=8.5, color=C.GRAY_600)

    # 승인 필요 항목
    rect(sl, Layout.MARGIN_L, 5.65, Layout.CONTENT_W, 0.28, C.NAVY_MID)
    tbox(sl, "승인 요청 항목 (VP 결재 필요)",
         Layout.MARGIN_L+0.1, 5.67, Layout.CONTENT_W-0.2, 0.22,
         size=9.5, bold=True, color=C.WHITE)

    approvals = [
        ("예산 승인",         "$458M FY2025 MarCom 총 예산 (Finance VP 승인 필요)",          C.RED),
        ("Holiday 캠페인",    "Holiday Peak $203M — 역대 최대 단일 캠페인 예산 승인",          C.ORANGE),
        ("MAP 정책 강화",     "Amazon 3P 가격 이상 대응: Brand Protection 전담팀 예산 $5M",   C.BLUE),
        ("Agency 계약",       "Creative Agency 재계약: $22M (BBH/TBWA 최종 협상 중)",        C.GREEN),
    ]

    for i, (topic, detail, color) in enumerate(approvals):
        col = i % 2
        row = i // 2
        x = Layout.MARGIN_L + col * 6.33
        y = 6.0 + row * 0.65
        rect(sl, x, y, 0.4, 0.55, color)
        tbox(sl, "▶", x+0.05, y+0.1, 0.3, 0.35,
             size=11, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
        tbox(sl, topic, x+0.45, y+0.04, 1.5, 0.25,
             size=9, bold=True, color=color)
        tbox(sl, detail, x+0.45, y+0.28, 5.75, 0.25,
             size=8, color=C.GRAY_600)

    # ══════════════════════════════════════════════════════════
    out = PPTX_DIR / "marcom_plan_vp.pptx"
    prs.save(str(out))
    print(f"[마케팅 커뮤니케이션] ✅ VP급 PPTX: {out}")
    return str(out)


if __name__ == "__main__":
    build()
