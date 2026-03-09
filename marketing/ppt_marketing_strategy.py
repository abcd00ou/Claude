"""
마케팅 전략 VP 보고서 — Marketing Strategy
SanDisk B2C Storage Marketing Team
"""
import sys
sys.path.insert(0, str(__import__('pathlib').Path(__file__).parent))

import json, datetime
from pathlib import Path
from pptx.enum.text import PP_ALIGN

import ppt_design as D
from ppt_design import C, Layout, rect, tbox, tbox_multi, card, data_table
from ppt_design import slide_header, slide_footer, bar_chart, stacked_bar
from ppt_design import horizontal_bar, exec_summary_block, waterfall_chart
from config import PPTX_DIR, DATA_DIR

with open(DATA_DIR / "market_data.json")    as f: MKT  = json.load(f)
with open(DATA_DIR / "internal_sales.json") as f: SALE = json.load(f)

TODAY = datetime.date.today()
PAGE_PREFIX = "Marketing Strategy"


def build():
    prs = D.new_prs()
    pg = 1

    # ══════════════════════════════════════════════════════════
    # SLIDE 1: Title
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, Layout.W, Layout.H, C.NAVY)
    rect(sl, 0, 0, Layout.W, 3.8, C.NAVY_MID)
    rect(sl, 0, 3.8, Layout.W, 0.06, C.RED)
    tbox(sl, "SanDisk B2C Storage\nMarketing Strategy FY2025",
         0.7, 0.95, 11.9, 1.9,
         size=34, bold=True, color=C.WHITE, align=PP_ALIGN.LEFT)
    tbox(sl, "제품 포트폴리오 최적화  ·  가격 전략  ·  채널 성장 로드맵",
         0.7, 2.95, 11.9, 0.7,
         size=14, color="A0AEC0", align=PP_ALIGN.LEFT)
    tbox(sl, "SanDisk Marketing Team  |  CONFIDENTIAL — Internal Use Only",
         0.7, 4.15, 11.9, 0.55, size=11, color="6B7280")
    tbox(sl, f"Prepared: {TODAY}",
         0.7, 4.65, 11.9, 0.5, size=10, color="6B7280")

    # ══════════════════════════════════════════════════════════
    # SLIDE 2: Executive Summary
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Executive Summary", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: WD FY2024 10-K, TrendForce 2025, Internal Sales Data", pg)

    # 4 KPI 카드
    kpi_data = [
        ("$3.53B", "2025E B2C Revenue", "+9.2% YoY"),
        ("49.8%",  "Blended Gross Margin", "+2.2pp vs 2024"),
        ("24%",    "Ext. SSD Market Share", "2025 Target"),
        ("$882M",  "Internal SSD Rev.", "+16.1% YoY"),
    ]
    for i, (val, label, sub) in enumerate(kpi_data):
        x = Layout.MARGIN_L + i * 3.13
        lc = C.RED if i == 0 else C.BLUE
        card(sl, x, 1.05, 2.95, 1.35, fill=C.NAVY_MID,
             title=label, title_size=9, title_color=C.GRAY_400,
             value=val, value_size=22, value_color=C.WHITE,
             sub=sub, sub_size=8, sub_color=lc,
             line_color=lc)

    exec_summary_block(sl, [
        ("⚙️", "BiCS8 전환 가속 → 원가 혁신",
         "BiCS8(218L) 전환으로 GB당 원가 $0.038 달성. Internal SSD GM +3.3pp → FY2025 $882M 목표 실현 가속."),
        ("📦", "External SSD: 프리미엄 방어 + 4TB 성장",
         "Extreme Pro 4TB 출시 확대 + Amazon 3P 가격 anomaly 대응 (Extreme 1TB: $198.59 vs MSRP $99.99, +99%)."
         " 공식 채널 강화로 ASP 방어."),
        ("📱", "microSD: 고용량·고마진 믹스 개선",
         "512GB+ 비중 확대 (GM 62-65% 구간). Pro Plus / Extreme 라인 집중 투자."
         " 2025E microSD Rev $1.61B, 시장점유율 34% 목표."),
    ], Layout.MARGIN_L, 2.55, Layout.CONTENT_W, 4.1)

    # ══════════════════════════════════════════════════════════
    # SLIDE 3: Strategic Framework
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "FY2025 Strategic Framework", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: Internal Strategy Planning, WD Annual Report 2024", pg)

    rect(sl, Layout.MARGIN_L, 1.1, Layout.CONTENT_W, 0.4, C.NAVY_MID)
    tbox(sl, "3대 전략 축:  Ⅰ. 원가 혁신  ·  Ⅱ. 포트폴리오 최적화  ·  Ⅲ. 채널·지역 믹스 개선",
         Layout.MARGIN_L+0.12, 1.12, Layout.CONTENT_W-0.24, 0.38,
         size=11, bold=True, align=PP_ALIGN.CENTER, color=C.WHITE)

    pillars = [
        ("Ⅰ", "원가 혁신", C.RED, [
            "BiCS8 전환 완료: FY2025 H2",
            "GB단가: $0.080 → $0.038 (BiCS8)",
            "Internal SSD GM: 38.5% → 41.8%",
            "External SSD GM: 44.8% → 47.2%",
            "공정 수율 개선: 3.2pp YoY 목표",
        ]),
        ("Ⅱ", "포트폴리오 최적화", C.NAVY_MID, [
            "4TB SKU 집중 확대 (최고 GM: 49-65%)",
            "Ext SSD Extreme Pro 라인 강화",
            "microSD Pro Plus 512G→1TB 상향",
            "보급형 WD Blue/Ultra 합리화",
            "신규: PCIe 5.0 NVMe 준비",
        ]),
        ("Ⅲ", "채널·지역 믹스", C.BLUE, [
            "E-commerce: 52%→55% (External SSD)",
            "NA/EMEA 시장 성장 집중",
            "APAC 신흥시장 microSD 확대",
            "직영몰 ASP 프리미엄 방어 전략",
            "B2B 채널 15% 유지 (Internal SSD)",
        ]),
    ]
    for i, (num, title, color, items) in enumerate(pillars):
        x = Layout.MARGIN_L + i * 4.18
        rect(sl, x, 1.65, 4.0, 0.65, color)
        tbox(sl, f"  {num}  {title}", x+0.05, 1.65, 3.9, 0.65,
             size=14, bold=True, color=C.WHITE)
        rect(sl, x, 2.3, 4.0, 4.6, "F8F9FA")
        for j, item in enumerate(items):
            tbox(sl, f"• {item}", x+0.15, 2.4 + j*0.82, 3.7, 0.72,
                 size=10, color=C.GRAY_600)

    # ══════════════════════════════════════════════════════════
    # SLIDE 4: Product P&L Analysis
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Product Portfolio P&L Analysis — SKU-Level", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: Internal Sales Data FY2024 (Estimated), WD Flash Segment 10-K", pg)

    rect(sl, Layout.MARGIN_L, 1.1, 7.6, 0.35, C.NAVY_MID)
    tbox(sl, "Top 10 SKU — Revenue vs. Gross Margin (FY2024)",
         Layout.MARGIN_L+0.1, 1.1, 7.4, 0.35,
         size=10, bold=True, color=C.WHITE)

    top_skus = [
        ("SD Extreme Micro 256G",  378, 58.4, "microSD"),
        ("SD Extreme Micro 512G",  312, 62.1, "microSD"),
        ("WD Black SN850X 2TB",    194, 45.2, "Internal SSD"),
        ("SD Extreme 1TB",         215, 45.2, "External SSD"),
        ("SD Extreme 2TB",         198, 46.1, "External SSD"),
        ("SD Ultra Micro 256G",    248, 52.3, "microSD"),
        ("SD Pro+ Micro 256G",     195, 63.2, "microSD"),
        ("WD Black SN850X 1TB",    162, 42.8, "Internal SSD"),
        ("SD Extreme Pro 2TB",     137, 54.7, "External SSD"),
        ("SD Extreme Pro 1TB",     148, 52.3, "External SSD"),
    ]

    headers = ["SKU", "Rev ($M)", "GM%", "카테고리", "GM 등급"]
    rows = []
    for name, rev, gm, cat in top_skus:
        if gm >= 60:
            grade, gc = ("★★★ Premium", C.GREEN)
        elif gm >= 50:
            grade, gc = ("★★  High",    C.BLUE)
        elif gm >= 40:
            grade, gc = ("★   Standard", C.NAVY_MID)
        else:
            grade, gc = ("◦   Low",      C.GRAY_400)
        rows.append([name, f"${rev}M", f"{gm:.1f}%", cat, (grade, gc)])

    data_table(sl, Layout.MARGIN_L, 1.5, 7.6, headers, rows,
               col_widths=[2.2, 0.9, 0.7, 1.3, 2.5],
               header_fill=C.RED, font_size=9)

    # GM 추이 차트 (우측)
    rect(sl, 8.2, 1.1, 4.73, 0.35, C.NAVY_MID)
    tbox(sl, "카테고리별 GM% 추이",
         8.3, 1.1, 4.53, 0.35, size=10, bold=True, color=C.WHITE)

    gm_data = SALE["gross_margin_pct"]
    cats_gm = [
        ("External SSD", [gm_data["2022"]["external_ssd"],
                          gm_data["2023"]["external_ssd"],
                          gm_data["2024"]["external_ssd"],
                          gm_data["2025E"]["external_ssd"]]),
        ("Internal SSD", [gm_data["2022"]["internal_ssd"],
                          gm_data["2023"]["internal_ssd"],
                          gm_data["2024"]["internal_ssd"],
                          gm_data["2025E"]["internal_ssd"]]),
        ("microSD",      [gm_data["2022"]["microsd"],
                          gm_data["2023"]["microsd"],
                          gm_data["2024"]["microsd"],
                          gm_data["2025E"]["microsd"]]),
    ]
    years = ["2022", "2023", "2024", "2025E"]
    for i, (cat_name, vals) in enumerate(cats_gm):
        y_start = 1.6 + i * 1.65
        color = [C.RED, C.BLUE, C.GREEN][i]
        tbox(sl, cat_name, 8.2, y_start, 4.73, 0.3,
             size=9, bold=True, color=color)
        gm_pairs = list(zip(years, vals))
        bar_chart(sl, gm_pairs, 8.2, y_start + 0.32, 4.73, 1.15,
                  max_val=65, bar_color=color, show_labels=True)

    # ══════════════════════════════════════════════════════════
    # SLIDE 5: 가격 전략 — Market Price Intelligence
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Price Strategy — Market Price Intelligence", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: Amazon Crawling Agent (2026-03-09), WD/Samsung 공식 사이트", pg)

    rect(sl, Layout.MARGIN_L, 1.1, Layout.CONTENT_W, 0.35, "FEF3C7")
    tbox(sl, "⚠  3P Seller Premium Alert: SanDisk/Samsung 주요 SKU, Amazon 3P 가격이 MSRP 대비 70~337% 프리미엄 형성 중",
         Layout.MARGIN_L+0.1, 1.12, Layout.CONTENT_W-0.2, 0.32,
         size=10, bold=True, align=PP_ALIGN.LEFT, color="7C2D12")

    price_headers = ["제품", "MSRP", "Amazon 3P", "프리미엄", "리스크"]
    price_rows = [
        ["SD Extreme 1TB",      "$99.99",  "$198.59", ("+99%",  C.RED),    ("HIGH", C.RED)],
        ["SD Extreme 2TB",      "$159.99", "$284.88", ("+78%",  C.RED),    ("HIGH", C.RED)],
        ["SD SN850X 2TB",       "$149.99", "$442.99", ("+195%", C.RED),    ("CRIT", C.RED)],
        ["SD Ultra Micro 128G", "$18.99",  "$82.99",  ("+337%", C.RED),    ("CRIT", C.RED)],
        ["SD Ex. Micro 256G",   "$28.99",  "$35.23",  ("+21%",  C.ORANGE), ("MED",  C.ORANGE)],
        ["Samsung T9 1TB",      "$129.99", "$224.99", ("+73%",  C.RED),    ("HIGH", C.RED)],
        ["Samsung T9 2TB",      "$229.99", "$376.77", ("+64%",  C.RED),    ("HIGH", C.RED)],
        ["Samsung 990Pro 2TB",  "$149.99", "$384.97", ("+157%", C.RED),    ("CRIT", C.RED)],
    ]
    data_table(sl, Layout.MARGIN_L, 1.55, 12.53, price_headers, price_rows,
               col_widths=[2.8, 1.2, 1.3, 1.2, 1.1],
               header_fill=C.NAVY, font_size=9.5)

    rect(sl, Layout.MARGIN_L, 5.65, 12.53, 1.6, "FFF7ED")
    tbox(sl, "전략적 시사점 (Price Intelligence)",
         Layout.MARGIN_L+0.15, 5.7, 12.2, 0.3,
         size=10, bold=True, color=C.NAVY)
    insights = [
        "① 공급 부족 신호: 3P 프리미엄 급등은 공식 채널 재고 부족 가능성 — 즉각적인 재고 점검 필요",
        "② 공식 Direct-to-Consumer 채널 강화: Amazon Vendor Central / 자사몰 재고 확대로 ASP 방어",
        "③ 경쟁사(Samsung) 동반 프리미엄: 시장 전반적 공급 타이트 → 가격 인상 기회 검토",
        "④ 보급형(Ultra 128G +337%) 이상 급등: 공급 중단 또는 단종 신호 — 후속 SKU 전환 계획 필요",
    ]
    for i, ins in enumerate(insights):
        tbox(sl, ins, Layout.MARGIN_L+0.15, 6.02 + i*0.29, 12.2, 0.27,
             size=8.5, color=C.GRAY_600)

    # ══════════════════════════════════════════════════════════
    # SLIDE 6: BiCS NAND Roadmap
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "BiCS NAND Technology Roadmap", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: WD Technology Roadmap 2024, TrendForce NAND Cost Analysis Q4 2025", pg)

    bics_info = [
        ("BiCS5",  "112L",  "$0.055/GB", "2022-2023", "성숙기·합리화 대상",  C.GRAY_400, [
            "외부 SSD 보급형",
            "WD My Passport",
            "SD Ultra microSD",
        ]),
        ("BiCS6",  "162L",  "$0.042/GB", "2023-2024", "주력 세대·전환 중",    C.BLUE, [
            "SD Extreme 1TB/2TB",
            "WD Black SN770",
            "SD Extreme Micro",
        ]),
        ("BiCS8",  "218L",  "$0.038/GB", "2025-2026", "차세대 주력·확대 중",  C.GREEN, [
            "SD Extreme Pro (전 용량)",
            "WD Black SN850X",
            "SD Pro Plus Micro",
        ]),
    ]
    for i, (gen, layers, cost, era, status, color, skus) in enumerate(bics_info):
        x = Layout.MARGIN_L + i * 4.18
        rect(sl, x, 1.1, 4.0, 5.6, "F8F9FA")
        rect(sl, x, 1.1, 4.0, 0.55, color)
        tbox(sl, f"{gen}  ·  {layers}  ·  {cost}",
             x+0.1, 1.1, 3.8, 0.55, size=12, bold=True, color=C.WHITE)
        tbox(sl, f"⏱ {era}",
             x+0.1, 1.73, 3.8, 0.32, size=9.5, color=C.GRAY_600)
        tbox(sl, status, x+0.1, 2.05, 3.8, 0.35,
             size=10, bold=True, color=color)
        tbox(sl, "대표 적용 SKU", x+0.1, 2.45, 3.8, 0.3,
             size=9, bold=True, color=C.NAVY)
        for j, sku in enumerate(skus):
            tbox(sl, f"• {sku}", x+0.2, 2.8 + j*0.5, 3.6, 0.45,
                 size=9.5, color=C.GRAY_600)
        gm_impact = ["GM ~38-45%", "GM ~43-54%", "GM ~46-65%"][i]
        rect(sl, x+0.1, 5.55, 3.8, 0.4, color)
        tbox(sl, gm_impact, x+0.15, 5.55, 3.7, 0.4,
             size=10, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    tbox(sl, "→", 4.5, 3.0, 0.4, 0.4,
         size=18, bold=True, color=C.RED, align=PP_ALIGN.CENTER)
    tbox(sl, "→", 8.68, 3.0, 0.4, 0.4,
         size=18, bold=True, color=C.RED, align=PP_ALIGN.CENTER)
    rect(sl, Layout.MARGIN_L, 6.8, Layout.CONTENT_W, 0.45, C.NAVY)
    tbox(sl, "BiCS5→BiCS8 전환 시 GB당 원가 31% 절감($0.055→$0.038)  ·  2025E blended GM 49.8% 달성 핵심 드라이버",
         Layout.MARGIN_L+0.15, 6.82, Layout.CONTENT_W-0.3, 0.4,
         size=9.5, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    # SLIDE 7: Revenue Waterfall
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Revenue Bridge — FY2024 to FY2025E", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: Internal Sales Data, WD FY2024 10-K, Internal Strategy Estimates", pg)

    wf_items = [
        ("2024\nActual",     3230, "start"),
        ("Volume\nGrowth",    198, "add"),
        ("Price/Mix\nEffect",  87, "add"),
        ("Channel\nMix",       34, "add"),
        ("NAND Cost\nPass",   -21, "sub"),
        ("FX\nEffect",          0, "add"),
        ("2025E\nForecast",     0, "total"),
    ]
    waterfall_chart(sl, wf_items, Layout.MARGIN_L, 1.15, 8.5, 5.5,
                    pos_color=C.GREEN, neg_color=C.RED,
                    total_color=C.BLUE, bg_color=C.NAVY_MID)

    # 우측 분석 패널
    rect(sl, 9.1, 1.15, 3.83, 5.5, "F8F9FA")
    tbox(sl, "주요 성장 드라이버",
         9.2, 1.22, 3.6, 0.35, size=10, bold=True, color=C.NAVY)

    drivers = [
        ("+$198M", "Volume Growth",     "Internal SSD +16.1%\nExternal SSD +9.1%", C.GREEN),
        ("+$87M",  "Price/Mix Upgrade", "4TB 확대, Pro 라인\n비중 증가",             C.BLUE),
        ("+$34M",  "Channel Mix",       "E-comm 비중 확대\nDTC 채널 강화",           C.NAVY_MID),
        ("-$21M",  "NAND Passthrough",  "BiCS8 전환 기간\n일부 가격 조정",           C.RED),
    ]
    for i, (amt, name, detail, color) in enumerate(drivers):
        y = 1.65 + i * 1.2
        rect(sl, 9.15, y, 0.7, 1.05, color)
        tbox(sl, amt, 9.15, y+0.12, 0.7, 0.55,
             size=10, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
        tbox(sl, name, 9.9, y, 2.95, 0.38,
             size=9.5, bold=True, color=C.NAVY)
        tbox(sl, detail, 9.9, y+0.38, 2.95, 0.65,
             size=8.5, color=C.GRAY_600)

    # ══════════════════════════════════════════════════════════
    # SLIDE 8: Market Share Strategy
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Market Share Strategy — Competitive Positioning", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: TrendForce 2024 Market Share Report, IDC Storage Tracker Q4 2024", pg)

    rect(sl, Layout.MARGIN_L, 1.1, 6.0, 0.35, C.NAVY_MID)
    tbox(sl, "카테고리별 시장 점유율 현황 및 목표",
         Layout.MARGIN_L+0.1, 1.1, 5.8, 0.35,
         size=10, bold=True, color=C.WHITE)

    share_data = [
        ("External SSD",           21, 22, 24, "Samsung 25%\nSeagate 18%"),
        ("Internal SSD (Consumer)", 16, 17, 19, "Samsung 27%\nCrucial 14%"),
        ("microSD",                 31, 32, 34, "Samsung 33%\nLexar 12%"),
    ]
    sh_headers = ["카테고리", "2023", "2024", "2025T", "주요 경쟁사"]
    sh_rows = []
    for cat, y23, y24, y25, comp in share_data:
        delta = y25 - y24
        sh_rows.append([cat, f"{y23}%", f"{y24}%",
                         (f"{y25}% (+{delta}pp)", C.GREEN), comp])
    data_table(sl, Layout.MARGIN_L, 1.5, 6.0, sh_headers, sh_rows,
               col_widths=[1.8, 0.8, 0.8, 1.3, 1.3], header_fill=C.RED, font_size=9.5)

    # 점유율 시각화 바
    for i, (cat, y23, y24, y25, _) in enumerate(share_data):
        y_pos = 3.5 + i * 1.1
        cat_label = cat.split("(")[0].strip()
        tbox(sl, cat_label, Layout.MARGIN_L, y_pos, 1.6, 0.4,
             size=9, bold=True, color=C.NAVY)
        for j, (val, color) in enumerate([(y23, C.GRAY_400), (y24, C.BLUE), (y25, C.GREEN)]):
            bar_w = val * 0.12
            rect(sl, 2.1 + j * 1.5, y_pos+0.05, bar_w, 0.28, color)
            tbox(sl, f"{val}%", 2.1 + j * 1.5 + bar_w + 0.05, y_pos+0.05, 0.5, 0.28,
                 size=8.5, color=color)

    for j, (label, color) in enumerate([("2023", C.GRAY_400), ("2024", C.BLUE), ("2025T", C.GREEN)]):
        rect(sl, 2.1 + j * 1.5, 6.5, 0.2, 0.15, color)
        tbox(sl, label, 2.35 + j * 1.5, 6.5, 0.8, 0.2, size=8, color=C.GRAY_600)

    # SWOT 분석 (우측)
    rect(sl, 6.6, 1.1, 6.33, 0.35, C.NAVY_MID)
    tbox(sl, "경쟁사 대비 SanDisk 강점/약점 분석",
         6.7, 1.1, 6.13, 0.35, size=10, bold=True, color=C.WHITE)

    sw_data = [
        ("강점", C.GREEN, [
            "Extreme 라인 브랜드 인지도 세계 1위",
            "BiCS NAND 내재화로 원가 경쟁력",
            "microSD 시장 32% 점유율 (2024)",
        ]),
        ("약점", C.RED, [
            "Premium Portable: Samsung T9에 밀림",
            "Consumer Internal: 990Pro 대비 ASP 열위",
            "3P 가격 통제력 부족 → 브랜드 훼손",
        ]),
        ("기회", C.BLUE, [
            "AI PC 확산 → PCIe 5.0 NVMe 수요",
            "Creator 시장 고용량 microSD 수요",
            "4TB+ 외장 SSD 시장 초기 단계",
        ]),
        ("위협", C.ORANGE, [
            "Samsung 공격적 가격 인하 가능성",
            "Lexar/Crucial 가성비 라인 성장",
            "NAND 공급 과잉 시 ASP 하락 압력",
        ]),
    ]
    for i, (cat, color, items) in enumerate(sw_data):
        col = i % 2
        row = i // 2
        x = 6.6 + col * 3.2
        y = 1.55 + row * 2.5
        rect(sl, x, y, 3.1, 0.3, color)
        tbox(sl, cat, x+0.1, y, 2.9, 0.3,
             size=9.5, bold=True, color=C.WHITE)
        for j, item in enumerate(items):
            tbox(sl, f"• {item}", x+0.1, y+0.35+j*0.6, 2.9, 0.55,
                 size=8.5, color=C.GRAY_600)

    # ══════════════════════════════════════════════════════════
    # SLIDE 9: Channel Strategy
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Channel Strategy — E-commerce · Retail · B2B", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: Internal Channel Mix Data 2024, Amazon Vendor Central Analytics", pg)

    rect(sl, Layout.MARGIN_L, 1.1, 6.0, 0.35, C.NAVY_MID)
    tbox(sl, "채널 믹스 현황 (2024) 및 2025E 목표",
         Layout.MARGIN_L+0.1, 1.1, 5.8, 0.35,
         size=10, bold=True, color=C.WHITE)

    ch_headers = ["카테고리", "E-comm\n2024", "E-comm\n2025E", "Retail\n2024", "Retail\n2025E", "B2B"]
    ch_mix = SALE["channel_mix_pct"]
    ch_rows = [
        ["External SSD",
         f"{ch_mix['2024']['external_ssd']['ecommerce']}%",
         (f"{ch_mix['2025E']['external_ssd']['ecommerce']}% ↑", C.GREEN),
         f"{ch_mix['2024']['external_ssd']['retail']}%",
         f"{ch_mix['2025E']['external_ssd']['retail']}%",
         f"{ch_mix['2024']['external_ssd']['b2b']}%"],
        ["Internal SSD",
         f"{ch_mix['2024']['internal_ssd']['ecommerce']}%",
         (f"{ch_mix['2025E']['internal_ssd']['ecommerce']}% ↑", C.GREEN),
         f"{ch_mix['2024']['internal_ssd']['retail']}%",
         f"{ch_mix['2025E']['internal_ssd']['retail']}%",
         f"{ch_mix['2024']['internal_ssd']['b2b']}%"],
        ["microSD",
         f"{ch_mix['2024']['microsd']['ecommerce']}%",
         (f"{ch_mix['2025E']['microsd']['ecommerce']}% ↑", C.GREEN),
         f"{ch_mix['2024']['microsd']['retail']}%",
         f"{ch_mix['2025E']['microsd']['retail']}%",
         f"{ch_mix['2024']['microsd']['b2b']}%"],
    ]
    data_table(sl, Layout.MARGIN_L, 1.5, 6.0, ch_headers, ch_rows,
               col_widths=[1.5, 0.9, 1.0, 0.85, 0.85, 0.81],
               header_fill=C.RED, font_size=9)

    channel_strats = [
        ("E-Commerce", C.RED, [
            "Amazon: Vendor Central 강화 + A+ 콘텐츠 업그레이드",
            "Direct (sandisk.com): 번들 프로모션으로 ASP+15%",
            "3P 셀러 가격 모니터링 → MAP 정책 강화",
            "Prime Day / Black Friday 집중 투자 (ROI 3.0x)",
        ]),
        ("Retail", C.BLUE, [
            "Best Buy / Costco: 4TB 프리미엄 전용 디스플레이",
            "Shop-in-shop: Extreme/Pro 라인 전용 존",
            "Retail 판매 인센티브 재설계 (매출 기반)",
            "microSD: 계산대 근처 충동구매 디스플레이 최적화",
        ]),
        ("B2B / Enterprise", C.GREEN, [
            "Internal SSD: 15% B2B 채널 유지 + 확대 검토",
            "System Integrator 파트너십 강화",
            "교육/공공기관 microSD 대량 공급 계약",
            "기업용 WD Black 라인 별도 B2B 가격 정책",
        ]),
    ]
    for i, (ch_name, color, items) in enumerate(channel_strats):
        x = Layout.MARGIN_L + i * 4.18
        y_start = 3.3
        rect(sl, x, y_start, 4.0, 0.38, color)
        tbox(sl, ch_name, x+0.1, y_start, 3.8, 0.38,
             size=11, bold=True, color=C.WHITE)
        for j, item in enumerate(items):
            tbox(sl, f"• {item}", x+0.12, y_start+0.44+j*0.72, 3.76, 0.68,
                 size=9, color=C.GRAY_600)

    # ══════════════════════════════════════════════════════════
    # SLIDE 10: Strategic Recommendations
    # ══════════════════════════════════════════════════════════
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    pg += 1
    slide_header(sl, "Strategic Recommendations & Action Plan", PAGE_PREFIX, pg)
    slide_footer(sl, "Source: Internal Strategy Planning, Market Intelligence Synthesis", pg)

    rect(sl, Layout.MARGIN_L, 1.1, Layout.CONTENT_W, 0.38, C.NAVY)
    tbox(sl, "우선순위 기반 실행 계획 — FY2025 H1/H2 분리",
         Layout.MARGIN_L+0.15, 1.1, Layout.CONTENT_W-0.3, 0.38,
         size=11, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    recs = [
        ("P0\n즉시", C.RED, "BiCS8 전환 가속",
         "목표: 2025 H2까지 Internal SSD 100% BiCS8 전환\n"
         "효과: GM +3.3pp, FY2025 Internal SSD $882M 달성\n"
         "담당: 생산팀 + 공급망팀  |  오너: VP of Supply Chain"),

        ("P0\n즉시", C.RED, "3P 가격 이상 대응",
         "현황: Extreme 1TB $198.59 (+99%), SN850X 2TB $442.99 (+195%)\n"
         "조치: Brand Protection 전담팀 + 공식채널 재고 확대\n"
         "오너: VP of Marketing + VP of Sales"),

        ("P1\n2025Q2", C.ORANGE, "4TB SKU 확대 출시",
         "목표: Extreme 4TB + Extreme Pro 4TB + SN850X 4TB 주력 전환\n"
         "근거: 4TB SKU GM 49-65% (최고 마진 구간)\n"
         "오너: Product Marketing"),

        ("P1\n2025Q2", C.ORANGE, "microSD 고용량 믹스 개선",
         "목표: 512GB+ 비중 2024년 대비 +15pp 확대\n"
         "근거: Pro Plus 512G GM 64.8%, Extreme 1TB GM 65.2%\n"
         "오너: Category Manager microSD"),

        ("P2\n2025Q3", C.BLUE, "직접 채널(DTC) 성장",
         "목표: sandisk.com Direct 매출 +30% YoY\n"
         "전술: 번들 프로모션, 구독형 보증 서비스 → ASP +15%\n"
         "오너: E-Commerce Team"),

        ("P2\n2025Q3", C.BLUE, "APAC/신흥시장 microSD",
         "목표: APAC microSD 점유율 28%→31%\n"
         "기회: Creator 경제 성장 (인도네시아, 베트남)\n"
         "오너: Regional Marketing APAC"),
    ]

    for i, (priority, color, title, detail) in enumerate(recs):
        col = i % 2
        row = i // 2
        x = Layout.MARGIN_L + col * 6.33
        y = 1.6 + row * 1.7
        rect(sl, x, y, 0.65, 1.55, color)
        tbox(sl, priority, x, y+0.3, 0.65, 0.95,
             size=9, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
        rect(sl, x+0.65, y, 5.55, 1.55, "F8F9FA")
        tbox(sl, title, x+0.75, y+0.06, 5.35, 0.35,
             size=10.5, bold=True, color=color)
        tbox(sl, detail, x+0.75, y+0.42, 5.35, 1.1,
             size=8.5, color=C.GRAY_600)

    rect(sl, Layout.MARGIN_L, 6.75, Layout.CONTENT_W, 0.5, C.NAVY_MID)
    tbox(sl, "FY2025 Target:  Revenue $3.53B (+9.2%)  ·  GM 49.8% (+2.2pp)  ·  "
             "Ext SSD Share 24%  ·  Internal SSD $882M (+16.1%)  ·  microSD Share 34%",
         Layout.MARGIN_L+0.15, 6.77, Layout.CONTENT_W-0.3, 0.45,
         size=9.5, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════
    out = PPTX_DIR / "marketing_strategy_vp.pptx"
    prs.save(str(out))
    print(f"[마케팅 전략] ✅ VP급 PPTX: {out}")
    return str(out)


if __name__ == "__main__":
    build()
